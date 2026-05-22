"""Inject teacher notes into a Cranky Bear lesson PPTX.

Usage:
  python inject.py <lesson_number>

Reads notes from notes_lesson{N}.py (must define NOTES dict of slide_idx -> notes_text)
and writes notes into the matching deck, saving as new file in same folder.

Notes formatting strategy:
- Each slide's notes is a single plain-text string with sections separated by blank lines.
- Section headers (SAY:, DO:, TEACHER NOTES:, WATCH FOR:, etc.) appear on their own paragraph.
- Content lines under each section appear as separate paragraphs.
- We set explicit a:buNone on every paragraph to suppress any auto-bullet styling from the notes master.
- ASCII-safe: hyphens for bullets, straight quotes, no em dashes or smart quotes.
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os
import sys
import importlib.util

FOLDER = r"c:\Users\09560329\Downloads\0en07a-the-very-cranky-bear-persuasive"
LESSON_FILES = {
    1: "17. Lesson Slides Extended Daily lesson 1 0En07A The Very Cranky Bear - Persuasive.pptx",
    2: "18. Lesson Slides Extended Daily lesson 2 0En07A The Very Cranky Bear - Persuasive.pptx",
    3: "19. Lesson Slides Extended Daily lesson 3 0En07A The Very Cranky Bear - Persuasive.pptx",
    4: "20. Lesson Slides Extended Daily lesson 4 0En07A The Very Cranky Bear - Persuasive.pptx",
    5: "21. Lesson Slides Extended Daily lesson 5 0En07A The Very Cranky Bear - Persuasive.pptx",
    6: "22. Lesson Slides Extended Daily lesson 6 0En07A The Very Cranky Bear - Persuasive.pptx",
    7: "23. Lesson Slides Extended Daily lesson 7 0En07A The Very Cranky Bear - Persuasive.pptx",
    8: "24. Lesson Slides Extended Daily lesson 8 0En07A The Very Cranky Bear - Persuasive.pptx",
}

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def make_paragraph(text, is_header=False):
    """Build an a:p paragraph element with explicit no-bullet and a single run."""
    p = etree.SubElement(etree.Element("root"), qn("a:p"))
    # paragraph properties: no bullet
    pPr = etree.SubElement(p, qn("a:pPr"))
    etree.SubElement(pPr, qn("a:buNone"))
    # run
    r = etree.SubElement(p, qn("a:r"))
    rPr = etree.SubElement(r, qn("a:rPr"))
    rPr.set("lang", "en-AU")
    rPr.set("dirty", "0")
    if is_header:
        rPr.set("b", "1")
    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    return p


def make_blank_paragraph():
    """Empty paragraph for spacing between sections."""
    p = etree.SubElement(etree.Element("root"), qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    etree.SubElement(pPr, qn("a:buNone"))
    endParaRPr = etree.SubElement(p, qn("a:endParaRPr"))
    endParaRPr.set("lang", "en-AU")
    return p


HEADER_PREFIXES = (
    "SOURCES:", "SAY:", "DO:", "CFU CHECKPOINT:", "TEACHER NOTES:",
    "ENABLING & EXTENDING:", "MISCONCEPTIONS:", "SENSITIVITY ADVISORY:",
    "WATCH FOR:",
)
SUBHEADER_PREFIXES = (
    "Technique:", "Script:", "PROCEED:", "PIVOT:",
    "ENABLING PROMPT:", "EXTENDING PROMPT:",
)


def write_notes(slide, notes_text):
    """Replace the notes text frame content with formatted paragraphs."""
    notes_slide = slide.notes_slide
    # Find the notes placeholder text frame (placeholder type 'body')
    tf = notes_slide.notes_text_frame
    txBody = tf._txBody
    # Remove existing a:p children (preserve bodyPr and lstStyle)
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    lines = notes_text.split("\n")
    for raw in lines:
        line = raw.rstrip()
        if not line.strip():
            txBody.append(make_blank_paragraph())
            continue
        stripped = line.strip()
        is_header = (
            stripped in HEADER_PREFIXES
            or any(stripped.startswith(h) for h in HEADER_PREFIXES if stripped.endswith(":"))
            or stripped in SUBHEADER_PREFIXES
            or any(stripped.startswith(h) for h in SUBHEADER_PREFIXES)
        )
        # More precise: header only if the line IS one of the known headers (whole line)
        is_header = stripped in HEADER_PREFIXES or stripped in SUBHEADER_PREFIXES or any(
            stripped.startswith(p) and stripped.endswith(":") for p in HEADER_PREFIXES + SUBHEADER_PREFIXES
        )
        txBody.append(make_paragraph(stripped, is_header=is_header))


def load_notes_module(lesson_n):
    spec = importlib.util.spec_from_file_location(
        f"notes_lesson{lesson_n}",
        os.path.join(os.path.dirname(__file__), f"notes_lesson{lesson_n}.py")
    )
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod.NOTES


def main(lesson_n):
    src = os.path.join(FOLDER, LESSON_FILES[lesson_n])
    stem, ext = os.path.splitext(LESSON_FILES[lesson_n])
    out_path = os.path.join(FOLDER, f"{stem} - with teacher notes{ext}")

    notes = load_notes_module(lesson_n)
    pres = Presentation(src)

    missing = []
    for idx, slide in enumerate(pres.slides, 1):
        if idx not in notes:
            missing.append(idx)
            continue
        write_notes(slide, notes[idx])

    if missing:
        print(f"WARNING: lesson {lesson_n} missing notes for slides: {missing}")

    pres.save(out_path)
    print(f"Saved: {out_path}")
    return len(pres.slides), len(missing)


if __name__ == "__main__":
    lesson_n = int(sys.argv[1])
    main(lesson_n)
