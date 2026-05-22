"""Extract slide content from all 8 Cranky Bear lesson PPTX files."""
from pptx import Presentation
import os
import sys

FOLDER = r"c:\Users\09560329\Downloads\0en07a-the-very-cranky-bear-persuasive"
OUT = r"c:\Users\09560329\Documents\Scripts\Claude Powerpoint\cranky_notes_workspace"

files = sorted([f for f in os.listdir(FOLDER) if f.endswith(".pptx") and not f.startswith("~")])

for fname in files:
    path = os.path.join(FOLDER, fname)
    pres = Presentation(path)
    # Lesson number from "17. ..." -> lesson_number = position - 16
    lesson_n = int(fname.split(".")[0]) - 16
    out_path = os.path.join(OUT, f"lesson{lesson_n}_extract.txt")
    with open(out_path, "w", encoding="utf-8") as out:
        out.write(f"FILE: {fname}\n")
        out.write(f"LESSON: {lesson_n}\n")
        out.write(f"SLIDE COUNT: {len(pres.slides)}\n")
        out.write("=" * 70 + "\n\n")
        for i, slide in enumerate(pres.slides, 1):
            out.write(f"=== SLIDE {i} ===\n")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if t:
                            out.write(f"  TEXT | {t}\n")
                elif shape.shape_type == 13:
                    out.write(f"  IMG  | {shape.name}\n")
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    out.write(f"  EXISTING_NOTES |\n")
                    for line in notes_text.split("\n"):
                        out.write(f"    {line}\n")
            out.write("\n")
    print(f"Wrote {out_path}")
