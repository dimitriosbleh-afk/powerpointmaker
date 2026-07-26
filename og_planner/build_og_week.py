#!/usr/bin/env python3
"""Build Orton-Gillingham session decks from the master template.

Usage:
    python3 og_planner/build_og_week.py og_planner/weeks/<week>.json [--only Monday] [--outdir output]

Reads the week spec JSON (schema documented in IMPORTANT/OG_MEGA_PROMPT.md),
clones slides from og_planner/OG_MASTER_TEMPLATE.pptx at raw-XML level so that
fonts, positions and click animations are preserved exactly. Card and sound-bank
fills are normalised to the school key (root green, prefix yellow, suffix red).
One PPTX per session is written into output/<unit_folder>/.

The master template must not be edited. All content decisions live in the spec.
"""

import copy
import json
import math
import re
import shutil
import sys
import zipfile
from io import BytesIO
from pathlib import Path
from xml.sax.saxutils import escape

from lxml import etree
from PIL import Image

NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}


def q(tag):
    pfx, local = tag.split(":")
    return "{%s}%s" % (NS[pfx], local)


def make_compatible_yellow_meter(green_png):
    """Recolour the known-good green gauge without the broken yellow asset.

    The master's original broad yellow PNG makes LibreOffice shift/clamp other
    slide elements during PDF export. Recolouring the compatible green gauge
    preserves the meter artwork and yellow difficulty cue without that defect.
    """
    image = Image.open(BytesIO(green_png)).convert("RGBA")
    pixels = image.load()
    for y in range(image.height):
        for x in range(image.width):
            red, green, blue, alpha = pixels[x, y]
            if alpha and green > red * 1.2 and green > blue * 1.2:
                pixels[x, y] = (255, 193, 7, alpha)
    output = BytesIO()
    image.save(output, format="PNG", dpi=(150, 150))
    return output.getvalue()


ROOT_DIR = Path(__file__).resolve().parent
MASTER = ROOT_DIR / "OG_MASTER_TEMPLATE.pptx"

# ---------------------------------------------------------------- template map
# 1-based slide indices in the master template.
T = {
    "title": 1,
    "overview": 2,
    "hdr_morph_review": 3,
    "card_review_root": 4,
    "card_review_prefix": 5,
    "card_review_suffix": 6,
    "hdr_wtr_review": 7,
    "wtr_review_table": 8,
    "hdr_sound_bank": 9,
    "sound_bank": 10,
    "hdr_wts_review": 11,
    "spell_word": 12,
    "hdr_new_morph": 13,
    "card_new_root": 14,
    "card_new_prefix": 15,
    "card_new_suffix": 16,
    "hdr_wtr_new": 17,
    "wtr_new_grid": 18,
    "hdr_wts_new": 19,
    "wts_new_grid": 20,
    "hdr_review_lw": 21,
    "learned_word": 22,
    "hdr_new_lw": 24,
    "learned_word_new": 25,
    "hdr_dictation": 26,
    "dictation": {"green": 27, "yellow": 28},
    "hdr_grammar": 29,
    "gradual_release": {"i_do": 30, "we_do": 31, "you_do": 32},
}

# Shape ids inside each template slide (stable - taken from the master).
SHAPE = {
    "title_heading": 6391,       # 'OG Planner'
    "title_sub": 6392,           # 'Enrichment / Term 3 Week XYZ'
    "overview_table": 6406,
    "card_review_word": {4: 6431, 5: 6439, 6: 6447},
    "wtr_review_tbl": 6465,
    "sound_boxes": [6482, 6483, 6484, 6486, 6485, 6487, 6488, 6489, 6490],
    # row-major order: (r1: x0.74,3.84,6.93) (r2: 0.74=6486, 3.84=6485, 6.93=6487) (r3)
    "spell_word_box": 6505,
    "hdr_new_morph_txt": 6512,
    "card_new_word": {14: 6523, 15: 6531, 16: 6539},
    "hdr_wtr_new_txt": 6544,
    "wtr_cols": [6556, 6555, 6554],  # left, middle, right
    "hdr_wts_new_txt": 6561,
    "wts_boxes": [6571, 6572],       # left (2 paras), right (2 paras)
    "lw_word": {22: 6598, 25: 6634},
    "dictation_title": {27: 6660, 28: 6671},
    "dictation_sentence": {27: 6663, 28: 6675},
    "dictation_meter": {27: 6664, 28: 6674},
    "gr_title": {30: 6687, 31: 6693, 32: 6699},
    "gr_body": {30: 6688, 31: 6694, 32: 6700},
}

MORPH_FILL = {"root": "ABEA90", "prefix": "FFF2CC", "suffix": "E06666"}
CARD_REVIEW = {"root": 4, "prefix": 5, "suffix": 6}
CARD_NEW = {"root": 14, "prefix": 15, "suffix": 16}

GRAMMAR_PURPLE = "5E58A6"   # sampled from the template's Grammar section header
GRAMMAR_TINT = "EFEEF8"
MORPH_MAGENTA = "A64D79"    # sampled from the New Morphology section header
MORPH_TINT = "F6EBF1"
CHARCOAL = "26262B"
PUNCT_RED = "C00000"
CAPITAL_GREEN = "008A45"
ANSWER_GREEN = "008A45"

HEADER_NOTES = {
    "morph_review": (
        "PROCEDURE:\n"
        "1. SHOW one card at a time in a mixed order.\n"
        "2. SAY: \"Read the morpheme, then give its meaning and keyword.\"\n"
        "3. WAIT for everyone to think before calling on one student.\n"
        "4. REPEAT any uncertain card at the end."
    ),
    "wtr_review": (
        "PROCEDURE:\n"
        "1. READ all rows together before asking the review questions.\n"
        "2. CLARIFY any unfamiliar word briefly.\n"
        "3. READ selected rows again for fluency."
    ),
    "sound_bank": (
        "PROCEDURE:\n"
        "1. SAY each sound and have students repeat it.\n"
        "2. WRITE the matching morpheme before revealing the next one.\n"
        "3. CHECK and fix each response immediately."
    ),
    "wts_review": "Students spell each review word, then use the reveal to check and fix.",
    "new_morph": (
        "PROCEDURE:\n"
        "1. SHOW the morpheme card and say the sound.\n"
        "2. TRACE three times while saying the sound.\n"
        "3. GIVE the letter names, keyword and card meaning.\n"
        "4. WRITE the morpheme three times and underline it in each new word."
    ),
    "wtr_new": "Read each word as it appears, then connect the morpheme to the whole-word meaning.",
    "wts_new": "Students spell the four selected new-family words, then complete the extension.",
    "morph_activity": "Students apply the new morpheme independently before the answer-check slide.",
    "learned_review": (
        "PROCEDURE:\n"
        "1. SAY the word and have students repeat it.\n"
        "2. NAME the unfair part and say the letter names.\n"
        "3. WRITE three times using far-point practice.\n"
        "4. COVER, write from memory, then check and fix."
    ),
    "learned_new": (
        "PROCEDURE:\n"
        "1. SAY the new word and have students repeat it.\n"
        "2. NAME the unfair part and explain why it is unexpected.\n"
        "3. WRITE three times using far-point practice.\n"
        "4. COVER, write from memory, then check and fix."
    ),
    "dictation": (
        "PROCEDURE:\n"
        "1. REVIEW the target morphemes, learned word and punctuation.\n"
        "2. SAY the whole sentence; students repeat it before writing.\n"
        "3. REPEAT by phrase only when the sentence is long.\n"
        "4. REVEAL after all writing is complete.\n"
        "5. TICK capitals, order, punctuation and spelling targets.\n"
        "6. FIX every error before moving on."
    ),
}

WARNINGS = []


def warn(msg):
    WARNINGS.append(msg)
    print(f"  WARN: {msg}")


_BANK = None
_MEANINGS = None
_MEANINGS_ANY = None
_PHOTO_CARDS = None


def bank():
    """morph -> legacy verified-card entry."""
    global _BANK
    if _BANK is None:
        data = json.loads((ROOT_DIR / "morpheme_bank.json").read_text())
        _BANK = {m["morph"]: m for m in data["morphemes"]}
    return _BANK


def canon_morph_key(morph):
    """Normalize display punctuation without collapsing genuine variant groups."""
    return str(morph).strip().strip("-").replace(" ", "")


def photo_card_catalogue():
    """(type, normalized morph) -> photographed Yoshimoto card entries.

    A list is retained for each key because the physical set contains genuine
    collisions, notably two different di- cards.  The photographed wording is
    authoritative; the legacy bank is only a compatibility fallback.
    """
    global _PHOTO_CARDS
    if _PHOTO_CARDS is None:
        files = {
            "suffix": "yoshimoto_cards_suffixes.json",
            "prefix": "yoshimoto_cards_prefixes.json",
            "root": "yoshimoto_cards_latin_roots.json",
        }
        catalogue = {}
        for morph_type, filename in files.items():
            data = json.loads((ROOT_DIR / filename).read_text(encoding="utf-8"))
            for card in data["cards"]:
                entry = {
                    "morph": card["morpheme"],
                    "type": morph_type,
                    "meaning": card["meaning"],
                    "keyword": card.get("keyword", ""),
                    "part_of_speech": card.get("part_of_speech", []),
                    "associated_words": card.get("associated_words", []),
                    "excluded_words": card.get("excluded_words", []),
                    "source": f"Photographed Yoshimoto card: {card['source_image']}",
                    "verified": True,
                }
                key = (morph_type, canon_morph_key(card["morpheme"]))
                catalogue.setdefault(key, []).append(entry)
        _PHOTO_CARDS = catalogue
    return _PHOTO_CARDS


def photo_card_entries(morph, morph_type=None):
    key = canon_morph_key(morph)
    catalogue = photo_card_catalogue()
    if morph_type:
        return catalogue.get((morph_type, key), [])
    matches = []
    for (entry_type, entry_key), entries in catalogue.items():
        if entry_key == key:
            matches.extend(entries)
    return matches


def meaning_catalogue():
    """(type, normalized morph) -> canonical meaning from the complete scope."""
    global _MEANINGS, _MEANINGS_ANY
    if _MEANINGS is None:
        data = json.loads((ROOT_DIR / "morpheme_meanings.json").read_text())
        section_types = {
            "suffixes_lower": "suffix",
            "continue_suffixes": "suffix",
            "suffix_card_additions": "suffix",
            "prefixes_first": "prefix",
            "continue_prefixes_chameleon": "prefix",
            "continue_prefixes_other": "prefix",
            "number_prefixes": "prefix",
            "lower_level_roots_start": "root",
            "lower_level_roots_continue": "root",
            "upper_level_roots": "root",
            "greek_combining_forms": "root",
        }
        typed = {}
        any_type = {}
        for section, morph_type in section_types.items():
            for morph, meaning in data[section].items():
                key = canon_morph_key(morph)
                typed[(morph_type, key)] = meaning
                any_type.setdefault(key, set()).add(meaning)
        _MEANINGS = typed
        _MEANINGS_ANY = {
            key: next(iter(values)) for key, values in any_type.items()
            if len(values) == 1
        }
    return _MEANINGS


def catalogue_meaning(morph, morph_type=None):
    typed = meaning_catalogue()
    key = canon_morph_key(morph)
    return typed.get((morph_type, key)) or _MEANINGS_ANY.get(key)


def verified_card_entry(morph, morph_type=None, requested_meaning=None):
    """Return a teacher/card-confirmed entry, tolerant of display hyphens."""
    photographed = photo_card_entries(morph, morph_type)
    if len(photographed) == 1:
        return photographed[0]
    if len(photographed) > 1:
        matching = [entry for entry in photographed
                    if requested_meaning and entry["meaning"] == requested_meaning]
        return matching[0] if len(matching) == 1 else None

    direct = bank().get(morph)
    if direct and (not morph_type or direct.get("type") == morph_type):
        return direct
    key = canon_morph_key(morph)
    for entry in bank().values():
        if (not morph_type or entry.get("type") == morph_type) and canon_morph_key(entry["morph"]) == key:
            return entry
    return None


def resolved_photo_card(morph, morph_type=None, requested_meaning=None):
    """Return the one intended photographed card, including collision handling."""
    entries = photo_card_entries(morph, morph_type)
    if len(entries) == 1:
        return entries[0]
    if len(entries) > 1 and requested_meaning:
        matches = [entry for entry in entries
                   if entry["meaning"] == requested_meaning]
        if len(matches) == 1:
            return matches[0]
    return None


def check_banked(day, where, morph, morph_type=None, requested_meaning=None):
    photographed = photo_card_entries(morph, morph_type)
    if len(photographed) > 1:
        matching = [entry for entry in photographed
                    if requested_meaning and entry["meaning"] == requested_meaning]
        if len(matching) == 1:
            return
        meanings = "; ".join(sorted({entry["meaning"] for entry in photographed}))
        warn(f"{day}: {where} morph {morph!r} matches multiple photographed cards "
             f"({meanings}); include the intended card meaning in the spec")
        return
    if verified_card_entry(morph, morph_type):
        return
    if catalogue_meaning(morph, morph_type):
        warn(f"{day}: {where} morph {morph!r} has a reference meaning, but its "
             "exact physical-card wording has not been confirmed")
    else:
        warn(f"{day}: {where} morph {morph!r} has no verified meaning in the "
             "photographed catalogue, morpheme_bank.json, or morpheme_meanings.json")


def resolve_card(day, where, card):
    """Resolve optional legacy keyword plus required canonical meaning."""
    kw = card.get("keyword", "")
    mean = card.get("meaning", "")
    photographed = photo_card_entries(card["morph"], card.get("type"))
    entry = verified_card_entry(card["morph"], card.get("type"), mean)
    if len(photographed) > 1 and not entry:
        meanings = "; ".join(sorted({item["meaning"] for item in photographed}))
        warn(f"{day}: {where} {card['morph']!r} is ambiguous in the photographed "
             f"card set ({meanings}); supplied keyword and meaning retained")
        return kw, mean
    if entry:
        if entry.get("keyword"):
            if kw and kw != entry["keyword"]:
                warn(f"{day}: {where} {card['morph']!r} spec keyword {kw!r} differs "
                     f"from canonical card value {entry['keyword']!r} - canonical "
                     "value used")
            kw = entry["keyword"]
    if entry:
        if mean and mean != entry["meaning"]:
            warn(f"{day}: {where} {card['morph']!r} spec meaning differs from "
                 "the confirmed physical-card meaning - confirmed value used")
        mean = entry["meaning"]
    else:
        reference_meaning = catalogue_meaning(card["morph"], card.get("type"))
        if reference_meaning:
            if mean and mean != reference_meaning:
                warn(f"{day}: {where} {card['morph']!r} spec meaning differs from "
                     "the unconfirmed reference catalogue - reference value used")
            mean = reference_meaning
    return kw, mean


def card_pos_note(card):
    """Teacher-note line for part of speech printed on the physical card."""
    parts = card_parts_of_speech(card)
    if not parts:
        return ""
    return f"\nPart of speech: {', '.join(parts)}"


def card_parts_of_speech(card):
    """Part-of-speech labels printed on the photographed physical card."""
    entry = resolved_photo_card(card["morph"], card.get("type"),
                                card.get("meaning"))
    return entry.get("part_of_speech", []) if entry else []


def card_meaning_note(meaning):
    """Present a meaning card as clear note lines, including sound-pattern cards."""
    text = str(meaning or "").strip()
    pattern = re.fullmatch(
        r"ending that says\s+([^()]+?)\s*\(the\s+(.+?)\s+part means\s+(.+?)\)",
        text,
        flags=re.IGNORECASE,
    )
    if pattern:
        sound, part, part_meaning = (value.strip() for value in pattern.groups())
        return f"Sound: {sound}\nMeaning of {part}: {part_meaning}"
    pattern = re.fullmatch(r"(?:ending that )?says\s+(.+)", text,
                           flags=re.IGNORECASE)
    if pattern:
        return f"Sound: {pattern.group(1).strip()}"
    return f"Meaning: {text}"


def extra_task_note(card):
    """Optional card retrieval prompt in clear teacher-facing language."""
    task = card.get("extra_task")
    if isinstance(task, dict):
        prompt = str(task.get("prompt", "")).strip()
        answers = task.get("answers", [])
        answer_text = ", ".join(str(answer) for answer in answers)
    else:
        legacy = str(card.get("derivative_ask", "")).strip()
        if not legacy:
            return ""
        warn("legacy derivative_ask is not allowed; use an extra_task object "
             "with prompt and answers")
        prompt, _, answer_text = legacy.partition("EXPECT:")
        prompt = prompt.strip()
        answer_text = answer_text.strip()
    if not prompt:
        return ""
    prompt = prompt[0].upper() + prompt[1:]
    if prompt[-1] not in ".?!":
        prompt += "."
    lines = ["Extra task:", f"SAY: \"{prompt}\""]
    if answer_text:
        lines.append(f"Possible answers: {answer_text}")
    return "\n" + "\n".join(lines)


def validate_session_card_words(day, session):
    """Enforce photographed-card word families for the new morphology block."""
    nm = session.get("new_morphology")
    if not nm:
        return
    entry = resolved_photo_card(nm["morph"], nm.get("type"), nm.get("meaning"))
    if not entry:
        return  # Uncaptured/ambiguous cards are reported by check_banked().

    exceptions = {}
    for item in session.get("associated_word_exceptions", []):
        if not isinstance(item, dict) or not item.get("word"):
            warn(f"{day}: malformed associated_word_exceptions entry {item!r}")
            continue
        word_key = str(item["word"]).strip().casefold()
        if not item.get("reason") or not item.get("source"):
            warn(f"{day}: outside derivative {item['word']!r} needs both reason and source")
        exceptions[word_key] = item

    allowed = {str(word).strip().casefold()
               for word in entry.get("associated_words", [])}
    excluded = {}
    for item in entry.get("excluded_words", []):
        if isinstance(item, dict):
            excluded[str(item.get("word", "")).strip().casefold()] = item.get("reason", "")
        else:
            excluded[str(item).strip().casefold()] = ""

    new_read = [item.get("word", "") if isinstance(item, dict) else str(item)
                for item in session.get("words_to_read_new", [])]
    selected = {str(word).strip().casefold() for word in new_read}
    for word in new_read:
        word_key = str(word).strip().casefold()
        if word_key in excluded:
            reason = excluded[word_key]
            suffix = f" ({reason})" if reason else ""
            warn(f"{day}: new-family word {word!r} is excluded from automatic lesson use{suffix}")
        elif word_key not in allowed and word_key not in exceptions:
            warn(f"{day}: new-family word {word!r} is not in the photographed card's "
                 "associated_words and has no verified associated_word_exceptions entry")

    for word_key, item in exceptions.items():
        if word_key not in selected:
            warn(f"{day}: associated-word exception {item['word']!r} is not used in "
                 "words_to_read_new")

    new_spell = [str(word).strip().casefold()
                 for word in session.get("words_to_spell_new", [])]
    for word in new_spell:
        if word not in selected:
            warn(f"{day}: new spelling word {word!r} is not in words_to_read_new")

    review_read = {str(word).strip().casefold()
                   for word in session.get("words_to_read_review", {}).get("words", [])}
    review_spell = {str(item.get("word", "")).strip().casefold()
                    for item in session.get("words_to_spell_review", [])}
    overlap = sorted(review_read & review_spell)
    if overlap:
        warn(f"{day}: same-day review reading/spelling overlap: {', '.join(overlap)}")


# ---------------------------------------------------------------- text helpers
def fit_font(text, width_in, base_pt, min_pt=20, char_em=0.62):
    """Largest size <= base_pt at which `text` fits on one line in width_in."""
    n = max(1, len(text))
    fit = int(width_in * 72 / (char_em * n))
    return max(min_pt, min(base_pt, fit))


def fit_font_block(text, width_in, height_in, base_pt, min_pt=20, char_em=0.55,
                   line_factor=1.32):
    """Largest size <= base_pt at which wrapped `text` fits in the box."""
    for sz in range(base_pt, min_pt - 1, -2):
        cpl = max(1, int(width_in * 72 / (char_em * sz)))
        lines = math.ceil(len(text) / cpl)
        if lines * sz * line_factor / 72 <= height_in:
            return sz
    return min_pt


def find_sp(root, spid):
    for sp in root.iter(q("p:sp")):
        nv = sp.find(q("p:nvSpPr") + "/" + q("p:cNvPr"))
        if nv is not None and nv.get("id") == str(spid):
            return sp
    raise KeyError(f"shape id {spid} not found")


def txbody(sp):
    tb = sp.find(q("p:txBody"))
    if tb is None:
        raise KeyError("shape has no txBody")
    return tb


def set_runs(sp, paragraphs, size_pt=None):
    """Replace a shape's paragraphs.

    paragraphs: list of paragraphs; each paragraph is a list of runs;
    each run is (text, color_hex_or_None, bold_or_None, underline_or_None).
    The first existing run's rPr is cloned as base formatting for every run;
    the first existing paragraph's pPr is kept for all paragraphs.
    """
    tb = txbody(sp)
    paras = tb.findall(q("a:p"))
    if not paras:
        raise KeyError("no paragraphs in shape")
    proto_p = paras[0]
    proto_ppr = proto_p.find(q("a:pPr"))
    proto_rpr = None
    for r_el in proto_p.findall(q("a:r")):
        proto_rpr = r_el.find(q("a:rPr"))
        if proto_rpr is not None:
            break
    if proto_rpr is None:  # look in later paragraphs
        for p_el in paras[1:]:
            for r_el in p_el.findall(q("a:r")):
                proto_rpr = r_el.find(q("a:rPr"))
                break
            if proto_rpr is not None:
                break
    for p_el in paras:
        tb.remove(p_el)
    for runs in paragraphs:
        p_el = etree.SubElement(tb, q("a:p"))
        if proto_ppr is not None:
            p_el.append(copy.deepcopy(proto_ppr))
        for text, color, bold, underline in runs:
            r_el = etree.SubElement(p_el, q("a:r"))
            if proto_rpr is not None:
                rpr = copy.deepcopy(proto_rpr)
            else:
                rpr = etree.Element(q("a:rPr"), lang="en-US")
            if size_pt is not None:
                rpr.set("sz", str(int(size_pt * 100)))
            if bold is not None:
                rpr.set("b", "1" if bold else "0")
            if underline:
                rpr.set("u", "sng")
            if color is not None:
                for fill in rpr.findall(q("a:solidFill")):
                    rpr.remove(fill)
                fill = etree.Element(q("a:solidFill"))
                clr = etree.SubElement(fill, q("a:srgbClr"))
                clr.set("val", color)
                fill.insert(0, clr)
                rpr.insert(0, fill)
            r_el.append(rpr)
            t_el = etree.SubElement(r_el, q("a:t"))
            t_el.text = text
        if not runs:
            if proto_rpr is not None:
                end = copy.deepcopy(proto_rpr)
                end.tag = q("a:endParaRPr")
                p_el.append(end)


def set_text(sp, text, size_pt=None, color=None):
    set_runs(sp, [[(text, color, None, None)]], size_pt=size_pt)


def replace_run_text(sp, old, new):
    """Replace `old` with `new` inside existing runs, keeping formatting."""
    hit = False
    for t_el in sp.iter(q("a:t")):
        if t_el.text and old in t_el.text:
            t_el.text = t_el.text.replace(old, new)
            hit = True
    return hit


def set_solid_fill(sp, hexval):
    sppr = sp.find(q("p:spPr"))
    for fill in sppr.findall(q("a:solidFill")):
        sppr.remove(fill)
    fill = etree.Element(q("a:solidFill"))
    clr = etree.SubElement(fill, q("a:srgbClr"))
    clr.set("val", hexval)
    ln = sppr.find(q("a:ln"))
    if ln is not None:
        ln.addprevious(fill)
    else:
        sppr.append(fill)


def set_slide_background(root, hexval):
    """Set a cloned slide's solid background without changing its type icon."""
    csld = root.find(q("p:cSld"))
    if csld is None:
        raise KeyError("slide common data missing")
    bg = csld.find(q("p:bg"))
    if bg is None:
        bg = etree.Element(q("p:bg"))
        csld.insert(0, bg)
    for child in list(bg):
        bg.remove(child)
    bgpr = etree.SubElement(bg, q("p:bgPr"))
    solid = etree.SubElement(bgpr, q("a:solidFill"))
    etree.SubElement(solid, q("a:srgbClr"), val=hexval)


def add_textbox(root, sid, x, y, w, h, lines, size_pt, color=CHARCOAL, bold=False,
                fill=None, rounded=False, align="ctr", anchor="ctr", italic=False):
    """Add a generated text box/card to a slide. `**bold**` markup supported."""
    a, p = NS["a"], NS["p"]
    geom = "roundRect" if rounded else "rect"
    fill_xml = (f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>'
                if fill else "<a:noFill/>")
    paras = []
    for ln in lines:
        runs = ""
        for i, seg in enumerate(re.split(r"\*\*(.*?)\*\*", ln)):
            if seg == "":
                continue
            b = "1" if (bold or i % 2 == 1) else "0"
            ital = ' i="1"' if italic else ""
            runs += (f'<a:r><a:rPr lang="en-AU" sz="{int(size_pt * 100)}" b="{b}"{ital}>'
                     f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                     f'<a:latin typeface="Lexend"/></a:rPr><a:t>{escape(seg)}</a:t></a:r>')
        if not runs:
            runs = f'<a:endParaRPr lang="en-AU" sz="{int(size_pt * 100)}"/>'
        paras.append(f'<a:p><a:pPr algn="{align}"/>{runs}</a:p>')
    xml = (
        f'<p:sp xmlns:p="{p}" xmlns:a="{a}">'
        f'<p:nvSpPr><p:cNvPr id="{sid}" name="OGGen{sid}"/><p:cNvSpPr txBox="1"/>'
        f'<p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{int(x * 914400)}" y="{int(y * 914400)}"/>'
        f'<a:ext cx="{int(w * 914400)}" cy="{int(h * 914400)}"/></a:xfrm>'
        f'<a:prstGeom prst="{geom}"><a:avLst/></a:prstGeom>{fill_xml}'
        f'<a:ln><a:noFill/></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr anchor="{anchor}" wrap="square" lIns="91440" '
        f'tIns="45720" rIns="91440" bIns="45720"/><a:lstStyle/>{"".join(paras)}'
        f'</p:txBody></p:sp>')
    sptree = root.find(q("p:cSld") + "/" + q("p:spTree"))
    sptree.append(etree.fromstring(xml))


def table_cells(root, spid):
    for gf in root.iter(q("p:graphicFrame")):
        nv = gf.find(q("p:nvGraphicFramePr") + "/" + q("p:cNvPr"))
        if nv is not None and nv.get("id") == str(spid):
            return gf.findall(".//" + q("a:tr"))
    raise KeyError(f"table {spid} not found")


def set_cell_paras(tc, lines, size_pt=None):
    """Replace a table cell's paragraphs with `lines` (list of strings)."""
    tx = tc.find(q("a:txBody"))
    paras = tx.findall(q("a:p"))
    proto_p = paras[0]
    proto_ppr = proto_p.find(q("a:pPr"))
    proto_rpr = None
    for p_el in paras:
        for r_el in p_el.findall(q("a:r")):
            proto_rpr = r_el.find(q("a:rPr"))
            break
        if proto_rpr is not None:
            break
    for p_el in paras:
        tx.remove(p_el)
    for line in lines:
        p_el = etree.SubElement(tx, q("a:p"))
        if proto_ppr is not None:
            p_el.append(copy.deepcopy(proto_ppr))
        r_el = etree.SubElement(p_el, q("a:r"))
        rpr = copy.deepcopy(proto_rpr) if proto_rpr is not None else etree.Element(q("a:rPr"), lang="en-US")
        if size_pt is not None:
            rpr.set("sz", str(int(size_pt * 100)))
        r_el.append(rpr)
        t_el = etree.SubElement(r_el, q("a:t"))
        t_el.text = line


# ---------------------------------------------------------------- timing regen
def regen_timing(root, targets):
    """Rebuild the click-to-reveal sequence.

    targets: list of (spid, para_index_or_None). One click per entry, in order.
    Uses the slide's existing timing block as the structural pattern.
    """
    timing = root.find(q("p:timing"))
    if timing is None:
        raise KeyError("slide has no timing block to use as pattern")
    seq_ctn = timing.find(".//" + q("p:seq") + "/" + q("p:cTn"))
    child_lst = seq_ctn.find(q("p:childTnLst"))
    pars = child_lst.findall(q("p:par"))
    pattern = copy.deepcopy(pars[0])
    for p_el in pars:
        child_lst.remove(p_el)
    ids = [2000]

    def renumber(el):
        for ctn in el.iter(q("p:cTn")):
            if ctn.get("id"):
                ids[0] += 1
                ctn.set("id", str(ids[0]))

    for spid, para in targets:
        node = copy.deepcopy(pattern)
        renumber(node)
        tgt = node.find(".//" + q("p:spTgt"))
        tgt.set("spid", str(spid))
        txel = tgt.find(q("p:txEl"))
        if para is None:
            if txel is not None:
                tgt.remove(txel)
        else:
            if txel is None:
                txel = etree.SubElement(tgt, q("p:txEl"))
                etree.SubElement(txel, q("p:pRg"))
            prg = txel.find(q("p:pRg"))
            if prg is None:
                for ch in list(txel):
                    txel.remove(ch)
                prg = etree.SubElement(txel, q("p:pRg"))
            prg.set("st", str(para))
            prg.set("end", str(para))
        child_lst.append(node)
    # bldLst: ensure one bldP per shape used, paragraph builds
    bld = timing.find(q("p:bldLst"))
    if bld is not None:
        for b in list(bld):
            bld.remove(b)
        for spid in dict.fromkeys(s for s, _ in targets):
            bp = etree.SubElement(bld, q("p:bldP"))
            bp.set("spid", str(spid))
            bp.set("grpId", "0")
            if any(p is not None for s, p in targets if s == spid):
                bp.set("build", "p")


def remove_timing(root):
    """Make an instructional follow-up slide visible immediately."""
    timing = root.find(q("p:timing"))
    if timing is not None:
        root.remove(timing)


# ---------------------------------------------------------------- notes
NOTE_ANCHOR_RE = re.compile(
    r"(?:\d+\.\s+)?(?:ANSWER|SAY|ASK|EXPECT|ACCEPT|SCAN|TRAP|FIX|STRETCH|HELP|"
    r"CARE|TIME|CIRCULATE|CHECK|PREP|POINT|SHOW|MODEL|DRAW|BUILD|COVER|REVEAL|COLLECT|"
    r"READ|ASSIGN|CHORAL|WAIT|REPEAT|CLARIFY|WRITE|TICK|PROCEDURE|CORRECTION|"
    r"Extra task|Possible answers|Root|Prefix|Suffix|Keyword|Meaning|Sound|"
    r"Meaning of [^:]+|Part of speech|Word to spell|Sentence|After checking|"
    r"Answer|Why learned|Say it|Link|C - Capitals(?:, shown green)?|U - Understanding|"
    r"P - Punctuation, shown red|S - Spelling targets, underlined|Score|Focus|"
    r"Targets|Memory hook|Purpose|SOURCES):"
)
NUMBERED_ACTION_RE = re.compile(
    r"^\d+\.\s+(?:POINT|SHOW|MODEL|DRAW|BUILD|COVER|REVEAL|TIME|COLLECT|"
    r"CIRCULATE|CHECK|SCAN|READ|WRITE|TICK|WAIT|REPEAT|CLARIFY)\b"
)


def prepare_notes_text(text):
    """Backstop old dense scripts while keeping well-authored line breaks intact."""
    text = sanitize(str(text or "")).replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\x0b", "\n")
    # Older specs sometimes placed every numbered beat and live action on one line.
    # (?<!\d) keeps the break out of the middle of a two-digit beat number ("10. ASK").
    text = re.sub(
        r"(?<!\n)(?<!^)(?<!\d)(?=\d+\.\s+(?:SAY|ASK|SCAN|MODEL|POINT|SHOW|READ|"
        r"REVEAL|CIRCULATE|TIME|WRITE|BUILD|TICK):)",
        "\n",
        text,
    )
    text = re.sub(
        r"(?<!\n)(?<!\d\. )(?=(?:TIME|CIRCULATE|CHECK|SCAN|TRAP|STRETCH|HELP|CARE):)",
        "\n",
        text,
    )
    text = re.sub(r"\s*---\s*", "\n---\n", text)
    return "\n".join(line.rstrip() for line in text.splitlines()).strip()


def validate_notes_source(template_idx, original_text, prepared_text):
    """Fail the build on note sources that will become an iPad text wall."""
    if "**" in original_text:
        warn(f"template slide {template_idx}: notes contain markdown ** markers; "
             "use plain labels and let the builder apply bold formatting")
    if re.search(r"[^\x09\x0A\x0D\x20-\x7E]", prepared_text):
        warn(f"template slide {template_idx}: notes contain non-ASCII characters; "
             "use straight quotes, hyphens and plain arrows")
    live_lines = []
    for line_no, line in enumerate(prepared_text.splitlines(), 1):
        if line == "---":
            break
        if line.strip():
            live_lines.append(line)
        if len(line) > 220:
            warn(f"template slide {template_idx}: notes line {line_no} is "
                 f"{len(line)} characters; split each new thought onto its own line")
        if re.fullmatch(r"\d+", line.strip()):
            warn(f"template slide {template_idx}: notes line {line_no} is a bare "
                 "number - a numbered beat has been split mid-line")
        numbered_beats = re.findall(
            r"\b\d+\.\s+(?:SAY|ASK|SCAN|MODEL|POINT|SHOW|READ|REVEAL|"
            r"CIRCULATE|TIME|WRITE|BUILD|TICK|CHECK|WAIT|REPEAT|CLARIFY)\b",
            line,
        )
        if len(numbered_beats) > 1:
            warn(f"template slide {template_idx}: notes line {line_no} contains "
                 "multiple numbered beats")
        if re.match(r"^\d+\.\s+", line) and (
                line.count("EXPECT:") > 1 or line.count("?") > 1):
            warn(f"template slide {template_idx}: notes line {line_no} contains "
                 "multiple questions or answers; use one retrieval prompt per line")
    if "---" in prepared_text.splitlines() and len(live_lines) > 8:
        warn(f"template slide {template_idx}: notes live zone has "
             f"{len(live_lines)} lines; Glance Format allows at most 8")


def note_runs(line, underline_terms=()):
    """Return (text, bold, underline) triples using real PowerPoint formatting."""
    plain = line.replace("**", "")
    bold_ranges = []
    underline_ranges = []
    # Preserve legacy bold intent for old specs while emitting no markdown.
    plain_pos = 0
    cursor = 0
    bold = False
    while cursor < len(line):
        if line.startswith("**", cursor):
            bold = not bold
            cursor += 2
            continue
        start = cursor
        while cursor < len(line) and not line.startswith("**", cursor):
            cursor += 1
        length = cursor - start
        if bold and length:
            bold_ranges.append((plain_pos, plain_pos + length))
        plain_pos += length
    for match in NOTE_ANCHOR_RE.finditer(plain):
        bold_ranges.append(match.span())
    action_match = NUMBERED_ACTION_RE.search(plain)
    if action_match:
        bold_ranges.append(action_match.span())
    for term in underline_terms:
        term = str(term).strip()
        if not term:
            continue
        pattern = rf"(?<![A-Za-z]){re.escape(term)}(?![A-Za-z])"
        underline_ranges.extend(
            match.span() for match in re.finditer(pattern, plain, re.IGNORECASE)
        )
    if not plain:
        return []
    boundaries = {0, len(plain)}
    for start, end in bold_ranges + underline_ranges:
        boundaries.update((start, end))
    points = sorted(boundaries)
    runs = []
    for start, end in zip(points, points[1:]):
        if start == end:
            continue
        is_bold = any(bs <= start and end <= be for bs, be in bold_ranges)
        is_underlined = any(
            us <= start and end <= ue for us, ue in underline_ranges
        )
        runs.append((plain[start:end], is_bold, is_underlined))
    return runs


def set_notes(notes_root, text):
    """Replace the body placeholder text of a notes slide.

    Source notes are plain text. The builder bolds recognised labels as real
    PowerPoint formatting and preserves one paragraph per source line.
    """
    body = None
    for sp in notes_root.iter(q("p:sp")):
        ph = sp.find(".//" + q("p:ph"))
        if ph is not None and ph.get("type") == "body":
            body = sp
            break
    if body is None:
        raise KeyError("notes body placeholder not found")
    original_text = str(text or "")
    text = prepare_notes_text(original_text)
    spelling_match = re.search(
        r"^Word to spell:\s*(\S.*)$",
        text,
        flags=re.MULTILINE,
    )
    spelling_word = spelling_match.group(1).strip() if spelling_match else ""
    tb = txbody(body)
    body_pr = tb.find(q("a:bodyPr"))
    if body_pr is not None:
        body_pr.set("anchor", "t")
    for p_el in tb.findall(q("a:p")):
        tb.remove(p_el)
    for line in text.split("\n"):
        line = re.sub(r" {2,}", " ", line.rstrip())
        p_el = etree.SubElement(tb, q("a:p"))
        # kill the notes master's default bullets - they read as dot-point
        # walls on the iPad speaker-notes view
        ppr = etree.SubElement(p_el, q("a:pPr"))
        ppr.set("algn", "l")
        etree.SubElement(ppr, q("a:buNone"))
        underline_terms = (
            (spelling_word,)
            if spelling_word and line.startswith("Sentence:")
            else ()
        )
        runs = note_runs(line, underline_terms=underline_terms)
        for seg, bold, underlined in runs:
            r_el = etree.SubElement(p_el, q("a:r"))
            rpr = etree.SubElement(r_el, q("a:rPr"), lang="en-AU", sz="1100")
            if bold:
                rpr.set("b", "1")
            if underlined:
                rpr.set("u", "sng")
            t_el = etree.SubElement(r_el, q("a:t"))
            t_el.text = seg
        if not runs:
            etree.SubElement(p_el, q("a:endParaRPr"), lang="en-AU", sz="1100")


def tidy_notes_whitespace(notes_root):
    """Normalise fixed template notes (double spaces read badly on the iPad)."""
    for t_el in notes_root.iter(q("a:t")):
        if t_el.text:
            t_el.text = re.sub(r" {2,}", " ", t_el.text)


# ---------------------------------------------------------------- package
class Package:
    def __init__(self, path):
        with zipfile.ZipFile(path) as z:
            self.parts = {n: z.read(n) for n in z.namelist()}
        if ("ppt/media/image12.png" in self.parts
                and "ppt/media/image13.png" in self.parts):
            self.parts["ppt/media/image12.png"] = make_compatible_yellow_meter(
                self.parts["ppt/media/image13.png"]
            )
        self.pres = etree.fromstring(self.parts["ppt/presentation.xml"])
        self.pres_rels = etree.fromstring(self.parts["ppt/_rels/presentation.xml.rels"])
        # ordered template slide part names
        rid_to_target = {
            rel.get("Id"): rel.get("Target")
            for rel in self.pres_rels.findall(q("rel:Relationship"))
        }
        self.slide_parts = []
        for sld in self.pres.findall(".//" + q("p:sldId")):
            rid = sld.get(q("r:id"))
            self.slide_parts.append("ppt/" + rid_to_target[rid].lstrip("/").replace("../", ""))

    def slide_xml(self, idx1):
        return self.parts[self.slide_parts[idx1 - 1]]

    def slide_rels(self, idx1):
        name = self.slide_parts[idx1 - 1]
        rn = name.rsplit("/", 1)
        return self.parts.get(f"{rn[0]}/_rels/{rn[1]}.rels")

    def notes_name(self, idx1):
        rels = self.slide_rels(idx1)
        if not rels:
            return None
        root = etree.fromstring(rels)
        for rel in root.findall(q("rel:Relationship")):
            if rel.get("Type").endswith("/notesSlide"):
                tgt = rel.get("Target").replace("../", "ppt/")
                return tgt
        return None


class DeckBuilder:
    """Assembles one output PPTX from template slide clones."""

    def __init__(self, pkg):
        self.pkg = pkg
        self.slides = []  # (slide_xml_root, rels_root, notes_root_or_None)

    def add(self, template_idx, fill=None, notes=None, fill_rels=None):
        """Clone template slide `template_idx`; apply fill(root) mutations."""
        root = etree.fromstring(self.pkg.slide_xml(template_idx))
        rels_bytes = self.pkg.slide_rels(template_idx)
        rels = etree.fromstring(rels_bytes) if rels_bytes else None
        notes_name = self.pkg.notes_name(template_idx)
        notes_root = None
        if notes is not None and notes_name is None:
            # borrow a notes slide from a template slide that has one
            notes_name = self.pkg.notes_name(T["spell_word"])
        if notes_name:
            notes_root = etree.fromstring(self.pkg.parts[notes_name])
            notes_rels = etree.fromstring(
                self.pkg.parts[notes_name.replace("notesSlides/", "notesSlides/_rels/") + ".rels"]
            )
        else:
            notes_rels = None
        if fill:
            fill(root)
        if fill_rels:
            fill_rels(rels)
        if notes is not None and notes_root is not None:
            prepared_notes = prepare_notes_text(str(notes))
            validate_notes_source(template_idx, str(notes), prepared_notes)
            set_notes(notes_root, notes)
        elif notes_root is not None:
            tidy_notes_whitespace(notes_root)
        self.slides.append((root, rels, notes_root, notes_rels))

    def write(self, out_path):
        parts = {}
        # copy every part except slides/notesSlides and their rels and ctypes/pres
        skip_prefixes = ("ppt/slides/", "ppt/notesSlides/")
        for name, data in self.pkg.parts.items():
            if name.startswith(skip_prefixes):
                continue
            if name in ("[Content_Types].xml", "ppt/presentation.xml",
                        "ppt/_rels/presentation.xml.rels"):
                continue
            parts[name] = data

        # slides + notes
        for i, (root, rels, notes_root, notes_rels) in enumerate(self.slides, 1):
            sname = f"ppt/slides/slide{i}.xml"
            parts[sname] = etree.tostring(root, xml_declaration=True,
                                          encoding="UTF-8", standalone=True)
            if rels is not None:
                rels2 = copy.deepcopy(rels)
                for rel in rels2.findall(q("rel:Relationship")):
                    if rel.get("Type").endswith("/notesSlide"):
                        if notes_root is not None:
                            rel.set("Target", f"../notesSlides/notesSlide{i}.xml")
                        else:
                            rels2.remove(rel)
                parts[f"ppt/slides/_rels/slide{i}.xml.rels"] = etree.tostring(
                    rels2, xml_declaration=True, encoding="UTF-8", standalone=True)
            if notes_root is not None:
                nname = f"ppt/notesSlides/notesSlide{i}.xml"
                parts[nname] = etree.tostring(notes_root, xml_declaration=True,
                                              encoding="UTF-8", standalone=True)
                nrels = copy.deepcopy(notes_rels)
                for rel in nrels.findall(q("rel:Relationship")):
                    if rel.get("Type").endswith("/slide"):
                        rel.set("Target", f"../slides/slide{i}.xml")
                parts[f"ppt/notesSlides/_rels/notesSlide{i}.xml.rels"] = etree.tostring(
                    nrels, xml_declaration=True, encoding="UTF-8", standalone=True)

        # presentation.xml
        pres = copy.deepcopy(self.pkg.pres)
        sldlst = pres.find(q("p:sldIdLst"))
        for sld in list(sldlst):
            sldlst.remove(sld)
        for i in range(1, len(self.slides) + 1):
            sld = etree.SubElement(sldlst, q("p:sldId"))
            sld.set("id", str(255 + i))
            sld.set(q("r:id"), f"rIdSl{i}")
        parts["ppt/presentation.xml"] = etree.tostring(
            pres, xml_declaration=True, encoding="UTF-8", standalone=True)

        # presentation rels
        prels = copy.deepcopy(self.pkg.pres_rels)
        for rel in prels.findall(q("rel:Relationship")):
            if rel.get("Type").endswith("/slide"):
                prels.remove(rel)
        for i in range(1, len(self.slides) + 1):
            rel = etree.SubElement(prels, q("rel:Relationship"))
            rel.set("Id", f"rIdSl{i}")
            rel.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
            rel.set("Target", f"slides/slide{i}.xml")
        parts["ppt/_rels/presentation.xml.rels"] = etree.tostring(
            prels, xml_declaration=True, encoding="UTF-8", standalone=True)

        # content types
        ct = etree.fromstring(self.pkg.parts["[Content_Types].xml"])
        for ov in ct.findall(q("ct:Override")):
            pn = ov.get("PartName")
            if pn.startswith("/ppt/slides/") or pn.startswith("/ppt/notesSlides/"):
                ct.remove(ov)
        for i in range(1, len(self.slides) + 1):
            ov = etree.SubElement(ct, q("ct:Override"))
            ov.set("PartName", f"/ppt/slides/slide{i}.xml")
            ov.set("ContentType",
                   "application/vnd.openxmlformats-officedocument.presentationml.slide+xml")
            if self.slides[i - 1][2] is not None:
                ov = etree.SubElement(ct, q("ct:Override"))
                ov.set("PartName", f"/ppt/notesSlides/notesSlide{i}.xml")
                ov.set("ContentType",
                       "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml")
        parts["[Content_Types].xml"] = etree.tostring(
            ct, xml_declaration=True, encoding="UTF-8", standalone=True)

        out_path.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as z:
            for name, data in sorted(parts.items()):
                z.writestr(name, data)


# ---------------------------------------------------------------- fills
def draw_structured_slide(root, blk, accent, tint, ctx):
    """Rule banner / hero card / support lines / routine chip / footer layout.

    Shared by grammar I-We-You do slides and the new-morphology You Do activity.
    """
    sid, y = 90100, 1.15
    if blk.get("rule"):
        rsz = fit_font_block(blk["rule"], 8.6, 0.58, 18, min_pt=15)
        add_textbox(root, sid, 0.5, y, 9.0, 0.64, [blk["rule"]], rsz,
                    color="FFFFFF", bold=True, fill=accent, rounded=True)
        sid += 1
        y += 0.76
    if blk.get("example"):
        ex_h = 1.65 if blk.get("items") else 2.45
        sz = fit_font_block(blk["example"], 7.7, ex_h - 0.25, 30, min_pt=20,
                            char_em=0.58)
        add_textbox(root, sid, 0.9, y, 8.2, ex_h, [blk["example"]], sz,
                    fill=tint, rounded=True)
        sid += 1
        y += ex_h + 0.12
    item_size = int(blk.get("item_size", 15))
    item_h = 0.54 if item_size >= 20 else 0.42
    item_step = item_h + 0.04
    for it in blk.get("items", [])[:4]:
        add_textbox(root, sid, 0.9, y, 8.2, item_h, [it], item_size,
                    color=blk.get("item_color", CHARCOAL),
                    bold=bool(blk.get("items_bold")), align="l")
        sid += 1
        y += item_step
    if y > 4.60:
        warn(f"{ctx}: content reaches y={y:.2f} - trim items")
    if blk.get("routine"):
        csz = fit_font(blk["routine"], 3.55, 15, min_pt=13, char_em=0.58)
        add_textbox(root, sid, 0.5, 4.66, 3.9, 0.44, [blk["routine"]], csz,
                    color="FFFFFF", bold=True, fill=accent, rounded=True)
        sid += 1
    if blk.get("footer"):
        add_textbox(root, sid, 4.55, 4.66, 4.95, 0.44, [blk["footer"]], 14,
                    align="r", italic=True, color="595959")


def has_fixed_answer_notes(blk):
    """True when a task's notes open with a definite ANSWER (not 'ANSWER: open')."""
    lines = prepare_notes_text(blk.get("notes", "")).splitlines()
    return bool(
        lines
        and lines[0].upper().startswith("ANSWER:")
        and not lines[0].lower().startswith("answer: open")
    )


def add_answer_check_slide(d, blk, tidx, accent, tint, ctx,
                           default_footer="Explain which card meaning fits each word"):
    """Emit the green 'Tick it or fix it' answer slide for a fixed-answer task.

    Shared by the new-morphology You Do and the grammar You Do. Returns True
    when the block supplied check_items and a slide was added.
    """
    check_items = blk.get("check_items", [])
    if not check_items:
        return False
    check_blk = {
        "rule": blk.get(
            "check_rule",
            "Tick each correct answer. Fix any answer that does not match.",
        ),
        "items": check_items,
        "item_color": ANSWER_GREEN,
        "items_bold": True,
        "item_size": blk.get("check_item_size", 22),
        "routine": blk.get("check_routine", "Tick - fix - explain"),
        "footer": blk.get("check_footer", default_footer),
    }

    def fill_check(root, check_blk=check_blk, blk=blk, tidx=tidx):
        tsp = find_sp(root, SHAPE["gr_title"][tidx])
        check_title = blk.get(
            "check_title",
            "Tick it or fix it - " + blk.get("title", ""),
        )
        set_text(
            tsp,
            check_title,
            size_pt=fit_font(check_title, 7.22, 28, min_pt=20, char_em=0.5),
        )
        bsp = find_sp(root, SHAPE["gr_body"][tidx])
        set_runs(bsp, [[]])
        draw_structured_slide(root, check_blk, accent, tint, ctx)

    d.add(tidx, fill_check, notes=blk.get("check_notes", ""))
    return True


def sanitize(s):
    if s is None:
        return s
    for a, b in [("–", "-"), ("—", "-"), ("‘", "'"), ("’", "'"),
                 ("“", '"'), ("”", '"'), ("…", "...")]:
        s = s.replace(a, b)
    return s


def deep_sanitize(obj):
    if isinstance(obj, str):
        return sanitize(obj)
    if isinstance(obj, list):
        return [deep_sanitize(x) for x in obj]
    if isinstance(obj, dict):
        return {k: deep_sanitize(v) for k, v in obj.items()}
    return obj


DICTATION_CUE_RE = re.compile(r"([A-Z]|[,.;:!?\"'-]+)")


def cue_runs(text, bold=None, underline=None):
    """Colour capital letters green and punctuation red for CUPS checking."""
    runs = []
    for seg in DICTATION_CUE_RE.split(text):
        if not seg:
            continue
        if re.fullmatch(r"[A-Z]", seg):
            color = CAPITAL_GREEN
        elif re.fullmatch(r"[,.;:!?\"'-]+", seg):
            color = PUNCT_RED
        else:
            color = None
        runs.append((seg, color, bold, underline))
    return runs


def dictation_runs(sentence, targets):
    """Reveal runs: capitals green, punctuation red, targets bold+underlined."""
    if not targets:
        return [cue_runs(sentence)]
    pattern = re.compile(
        r"\b(" + "|".join(re.escape(t) for t in sorted(targets, key=len, reverse=True)) + r")\b",
        re.IGNORECASE)
    runs = []
    pos = 0
    for m in pattern.finditer(sentence):
        if m.start() > pos:
            runs.extend(cue_runs(sentence[pos:m.start()]))
        runs.extend(cue_runs(m.group(0), bold=True, underline=True))
        pos = m.end()
    if pos < len(sentence):
        runs.extend(cue_runs(sentence[pos:]))
    return [runs]


def lw_runs(word, unfair):
    """White word with the unfair part in yellow."""
    i = word.lower().find(unfair.lower()) if unfair else -1
    if unfair and i < 0:
        warn(f"learned word '{word}': unfair part '{unfair}' not found; highlighting nothing")
    if i < 0:
        return [[(word, "FFFFFF", None, None)]]
    runs = []
    if i > 0:
        runs.append((word[:i], "FFFFFF", None, None))
    runs.append((word[i:i + len(unfair)], "FFFF00", None, None))
    if i + len(unfair) < len(word):
        runs.append((word[i + len(unfair):], "FFFFFF", None, None))
    return [runs]


def build_session(pkg, week, session, out_path):
    d = DeckBuilder(pkg)
    day = session["day"]
    stype = session.get("type", "new")  # new | review | week_review
    term, wknum = week["term"], week["week"]
    cohort = week.get("cohort", "Enrichment")

    # --- 1 title
    def fill_title(root):
        sp = find_sp(root, SHAPE["title_sub"])
        set_runs(sp, [[(cohort, None, None, None)],
                      [(f"Term {term} Week {wknum} - {day}", None, None, None)]])
        # widen the subtitle box (kept centred) so the day fits on one line
        xfrm = sp.find(q("p:spPr") + "/" + q("a:xfrm"))
        off, ext = xfrm.find(q("a:off")), xfrm.find(q("a:ext"))
        new_w = int(4.6 * 914400)
        cx = int(off.get("x")) + int(ext.get("cx")) // 2
        off.set("x", str(cx - new_w // 2))
        ext.set("cx", str(new_w))
    d.add(T["title"], fill_title, notes=f"{cohort} OG session deck. Term {term} Week {wknum}, {day}.")

    # --- 2 weekly overview
    ov = week["overview"]  # dict day -> {morphology, grammar, learned_words}
    days = week.get("days", ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday"])

    def fill_overview(root):
        rows = table_cells(root, SHAPE["overview_table"])
        # rows[0] headers; rows[1] morphology; rows[2] grammar; rows[3] learned words
        for ci, dy in enumerate(days, start=1):
            cells_m = rows[1].findall(q("a:tc"))
            cells_g = rows[2].findall(q("a:tc"))
            cells_l = rows[3].findall(q("a:tc"))
            if ci >= len(cells_m):
                break
            o = ov.get(dy, {})
            set_cell_paras(cells_m[ci], o.get("morphology", [""]) if isinstance(o.get("morphology"), list) else [o.get("morphology", "")])
            set_cell_paras(cells_g[ci], [o.get("grammar", "")])
            set_cell_paras(cells_l[ci], o.get("learned_words", [""]) if isinstance(o.get("learned_words"), list) else [o.get("learned_words", "")])
    d.add(T["overview"], fill_overview,
          notes="\n".join(f"{dy} - {ov.get(dy, {}).get('morphology','')} | Grammar: {ov.get(dy, {}).get('grammar','')} | LW: {ov.get(dy, {}).get('learned_words','')}" for dy in days))

    # --- 3 morphology review header (fixed drill notes stay)
    d.add(T["hdr_morph_review"], notes=HEADER_NOTES["morph_review"])

    # --- review cards x10
    for card in session["morphology_review"]:
        check_banked(day, "review card", card["morph"], card.get("type"),
                     card.get("meaning"))
    for item in session.get("sound_bank", []):
        check_banked(day, "sound bank", item["morph"], item.get("type"),
                     item.get("meaning"))
    if session.get("new_morphology"):
        new_morph = session["new_morphology"]
        check_banked(day, "new morphology", new_morph["morph"],
                     new_morph.get("type"), new_morph.get("meaning"))
    validate_session_card_words(day, session)
    for card in session["morphology_review"]:
        tidx = CARD_REVIEW[card["type"]]

        def fill_card(root, card=card, tidx=tidx):
            set_slide_background(root, MORPH_FILL[card["type"]])
            sp = find_sp(root, SHAPE["card_review_word"][tidx])
            set_text(sp, card["morph"],
                     size_pt=fit_font(card["morph"], 5.75, 56))
        kw, mean = resolve_card(day, "review card", card)
        card_note = (f"{card['type'].capitalize()}: {card['morph']}\n"
                     f"Keyword: {kw}\n"
                     f"{card_meaning_note(mean)}"
                     f"{card_pos_note(card)}"
                     f"{extra_task_note(card)}")
        d.add(tidx, fill_card, notes=card_note)

    # --- words to read review
    d.add(T["hdr_wtr_review"], notes=HEADER_NOTES["wtr_review"])
    wtr = session["words_to_read_review"]
    if len(wtr["words"]) != 15:
        warn(f"{day}: words_to_read_review has {len(wtr['words'])} words (template holds 15)")
    review_notes = prepare_notes_text(wtr.get("notes", ""))
    if len(re.findall(r"^\d+\. READ:", review_notes, flags=re.MULTILINE)) < 3:
        warn(f"{day}: Words to Read Review needs an all-class read, a fun "
             "group-row allocation, and an everyone row")
    if "Every student reads at least two rows" not in review_notes:
        warn(f"{day}: Words to Read Review must state that every student "
             "reads at least two rows in the group round")

    def fill_wtr_table(root):
        rows = table_cells(root, SHAPE["wtr_review_tbl"])
        n = 0
        for row in rows:
            for tc in row.findall(q("a:tc")):
                if n < len(wtr["words"]):
                    set_cell_paras(tc, [f"{n + 1}) {wtr['words'][n]}"])
                else:
                    set_cell_paras(tc, [""])
                n += 1
    d.add(T["wtr_review_table"], fill_wtr_table, notes=wtr["notes"])

    # --- sound bank
    d.add(T["hdr_sound_bank"], notes=HEADER_NOTES["sound_bank"])
    sb = session["sound_bank"]
    if len(sb) > 9:
        warn(f"{day}: sound bank has {len(sb)} entries; only 9 boxes - extras dropped")

    def fill_sound_bank(root):
        for i, spid in enumerate(SHAPE["sound_boxes"]):
            sp = find_sp(root, spid)
            if i < len(sb):
                item = sb[i]
                set_text(sp, item["morph"],
                         size_pt=fit_font(item["morph"], 2.33, 45, min_pt=24))
                set_solid_fill(sp, MORPH_FILL[item["type"]])
            else:
                set_text(sp, "")
                set_solid_fill(sp, "FFFFFF")
    d.add(T["sound_bank"], fill_sound_bank,
          notes="Students copy sounds into books before spelling words. Box colour = card type: green root, yellow prefix, red suffix.")

    # --- words to spell review (one reveal slide per word)
    d.add(T["hdr_wts_review"], notes=HEADER_NOTES["wts_review"])
    wsr = session["words_to_spell_review"]
    if not 8 <= len(wsr) <= 12:
        warn(f"{day}: words_to_spell_review has {len(wsr)} words (expected 10)")
    for item in wsr:
        def fill_spell(root, item=item):
            sp = find_sp(root, SHAPE["spell_word_box"])
            set_text(sp, item["word"], size_pt=fit_font(item["word"], 4.96, 60, min_pt=28))
        if not item.get("prompt"):
            warn(f"{day}: spell review word '{item['word']}' has no after-check prompt")
        elif not item.get("answer"):
            warn(f"{day}: spell review word '{item['word']}' prompt has no answer")
        note = (f"Word to spell: {item['word']}\n"
                f"Sentence: {item['sentence']}\n"
                "Click to reveal answer.")
        if item.get("prompt"):
            note += f"\nAfter checking: {item['prompt']}"
        if item.get("answer"):
            note += f"\nAnswer: {item['answer']}"
        d.add(T["spell_word"], fill_spell, notes=note)

    # --- new/review morphology section (skipped on week_review days)
    if stype != "week_review":
        section = "New Morphology" if stype == "new" else "Review Morphology"
        nm = session["new_morphology"]

        def fill_nm_hdr(root):
            sp = find_sp(root, SHAPE["hdr_new_morph_txt"])
            replace_run_text(sp, "New Morphology", section)
        d.add(T["hdr_new_morph"], fill_nm_hdr, notes=HEADER_NOTES["new_morph"])

        tidx = CARD_NEW[nm["type"]]

        def fill_nm_card(root, tidx=tidx, nm=nm):
            set_slide_background(root, MORPH_FILL[nm["type"]])
            sp = find_sp(root, SHAPE["card_new_word"][tidx])
            set_text(sp, nm["morph"], size_pt=fit_font(nm["morph"], 5.75, 56))
        nm_kw, nm_mean = resolve_card(day, "new morphology", nm)
        d.add(tidx, fill_nm_card,
              notes=(f"{nm['type'].capitalize()}: {nm['morph']}\n"
                     f"Keyword: {nm_kw}\n"
                     f"{card_meaning_note(nm_mean)}"
                     f"{card_pos_note(nm)}"))

        # Immediately follow the delivered morpheme card with a stable,
        # non-animated reference students can copy into their books.
        def fill_nm_copy(root, tidx=tidx, nm=nm, nm_kw=nm_kw,
                         nm_mean=nm_mean):
            set_slide_background(root, MORPH_FILL[nm["type"]])
            sp = find_sp(root, SHAPE["card_new_word"][tidx])
            meaning_lines = card_meaning_note(nm_mean).splitlines()
            fields = [(nm["type"].capitalize(), nm["morph"]),
                      ("Keyword", nm_kw)]
            for line in meaning_lines:
                label, _, value = line.partition(":")
                fields.append((label, value.strip()))
            parts = card_parts_of_speech(nm)
            if parts:
                fields.append(("Part of speech", ", ".join(parts)))
            display_text = "\n".join(
                f"{label}: {value}" for label, value in fields
            )
            sz = fit_font_block(display_text, 6.2, 3.0, 32,
                                min_pt=24, char_em=0.55)
            set_runs(
                sp,
                [[(f"{label}: ", None, True, None),
                  (value, None, False, None)]
                 for label, value in fields],
                size_pt=sz,
            )
            xfrm = sp.find(q("p:spPr") + "/" + q("a:xfrm"))
            xfrm.find(q("a:off")).set("x", str(int(1.7 * 914400)))
            xfrm.find(q("a:off")).set("y", str(int(1.35 * 914400)))
            xfrm.find(q("a:ext")).set("cx", str(int(6.6 * 914400)))
            xfrm.find(q("a:ext")).set("cy", str(int(3.1 * 914400)))
            bodypr = sp.find(q("p:txBody") + "/" + q("a:bodyPr"))
            if bodypr is not None:
                bodypr.set("anchor", "ctr")
            for p_el in sp.findall(q("p:txBody") + "/" + q("a:p")):
                ppr = p_el.find(q("a:pPr"))
                if ppr is not None:
                    ppr.set("algn", "l")
            remove_timing(root)

        copy_note_fields = "morpheme, keyword and meaning"
        if card_parts_of_speech(nm):
            copy_note_fields += ", including the part of speech"
        d.add(
            tidx,
            fill_nm_copy,
            notes=(
                "1. SAY: \"Copy this morpheme card into your book exactly as shown.\"\n"
                f"2. CHECK: Scan the {copy_note_fields} before moving on."
            ),
        )

        def fill_wtr_new_hdr(root):
            sp = find_sp(root, SHAPE["hdr_wtr_new_txt"])
            replace_run_text(sp, "New Morphology", section)
        d.add(T["hdr_wtr_new"], fill_wtr_new_hdr, notes=HEADER_NOTES["wtr_new"])

        words = session["words_to_read_new"]  # list of {word, meaning}
        if not 8 <= len(words) <= 12:
            warn(f"{day}: words_to_read_new has {len(words)} words (aim 9-12, min 8)")

        def fill_wtr_grid(root):
            # always the template's 3-column grid; with 10-12 words the columns
            # widen and stretch down the page so a 4th row still reads big
            n = len(words)
            col_spids = SHAPE["wtr_cols"]
            per = math.ceil(n / 3)
            cols = [words[0:per], words[per:2 * per], words[2 * per:]]
            box_y, box_h = (1.15, 3.95) if n > 9 else (1.48, 3.30)
            col_w = 3.15
            for spid, x in zip(col_spids, (0.25, 3.425, 6.60)):
                sp = find_sp(root, spid)
                xfrm = sp.find(q("p:spPr") + "/" + q("a:xfrm"))
                xfrm.find(q("a:off")).set("x", str(int(x * 914400)))
                xfrm.find(q("a:off")).set("y", str(int(box_y * 914400)))
                xfrm.find(q("a:ext")).set("cx", str(int(col_w * 914400)))
                xfrm.find(q("a:ext")).set("cy", str(int(box_h * 914400)))
                bodypr = sp.find(q("p:txBody") + "/" + q("a:bodyPr"))
                if bodypr is not None:
                    bodypr.set("anchor", "ctr")
                    bodypr.set("lIns", "0")
                    bodypr.set("rIns", "0")
            longest = max(len(w["word"]) for w in words)
            rows_max = max(len(c) for c in cols)
            # Classroom projection needs a readable floor. Use a conservative
            # Lexend estimate, but do not shrink the whole grid to tiny text.
            sz = fit_font("x" * longest, col_w, 50 if rows_max <= 3 else 44,
                          min_pt=26, char_em=0.67)
            sz = min(sz, int(box_h / rows_max * 72 / 1.5))
            if sz < 27:
                warn(f"{day}: new morphology word grid would render at {sz} pt; "
                     "replace/rebalance long words so the grid is at least 27 pt")
            targets = []
            for spid, col in zip(col_spids, cols):
                sp = find_sp(root, spid)
                set_runs(sp, [[(w["word"], None, None, None)] for w in col] or [[]],
                         size_pt=sz)
                for pi in range(len(col)):
                    targets.append((spid, pi))
            regen_timing(root, targets)
        d.add(T["wtr_new_grid"], fill_wtr_grid,
              notes=session.get(
                  "wtr_new_notes",
                  "\n".join(f"{w['word']} - {w['meaning']}" for w in words)))

        def fill_wts_new_hdr(root):
            sp = find_sp(root, SHAPE["hdr_wts_new_txt"])
            replace_run_text(sp, "New Morphology", section)
        d.add(T["hdr_wts_new"], fill_wts_new_hdr, notes=HEADER_NOTES["wts_new"])

        spell = session["words_to_spell_new"]
        if len(spell) != 4:
            warn(f"{day}: words_to_spell_new has {len(spell)} words (template holds 4)")

        def fill_wts_grid(root):
            pairs = [spell[0:2], spell[2:4]]
            longest = max((len(w) for w in spell), default=1)
            sz = fit_font("x" * longest, 4.8, 52, min_pt=28)
            targets = []
            for spid, pair in zip(SHAPE["wts_boxes"], pairs):
                sp = find_sp(root, spid)
                set_runs(sp, [[(w, None, None, None)] for w in pair] or [[]], size_pt=sz)
                for pi in range(len(pair)):
                    targets.append((spid, pi))
            regen_timing(root, targets)
        d.add(T["wts_new_grid"], fill_wts_grid,
              notes=session.get("extension",
                                "Once finished, students write a fun sentence with one or more of the words."))

        # --- Yoshimoto You Do activity (magenta header + designed task slide)
        act = session.get("new_morph_activity")
        if not act:
            warn(f"{day}: no new_morph_activity (Yoshimoto You Do) provided")
        else:
            def fill_act_hdr(root, section=section):
                sp = find_sp(root, SHAPE["hdr_wts_new_txt"])
                replace_run_text(sp, "Words to Spell- New Morphology",
                                 f"{section} - You Do")
                replace_run_text(sp, "2 mins", act.get("time", "3 mins"))
            d.add(T["hdr_wts_new"], fill_act_hdr,
                  notes=HEADER_NOTES["morph_activity"])

            act_tidx = T["gradual_release"]["you_do"]

            def fill_act(root, act=act, act_tidx=act_tidx):
                tsp = find_sp(root, SHAPE["gr_title"][act_tidx])
                set_text(tsp, "You do - " + act.get("title", ""))
                bsp = find_sp(root, SHAPE["gr_body"][act_tidx])
                set_runs(bsp, [[]])
                draw_structured_slide(root, act, MORPH_MAGENTA, MORPH_TINT,
                                      f"{day} new_morph_activity")
            d.add(act_tidx, fill_act, notes=act.get("notes", ""))

            added_check = add_answer_check_slide(
                d, act, act_tidx, MORPH_MAGENTA, MORPH_TINT,
                f"{day} new_morph_activity answer check")
            if has_fixed_answer_notes(act) and not added_check:
                warn(f"{day}: fixed-answer new_morph_activity needs a duplicate "
                     "check_items answer slide")

    # --- learned words
    lw = session["learned_words"]
    d.add(T["hdr_review_lw"], notes=HEADER_NOTES["learned_review"])
    for item in lw.get("review", []):
        def fill_lw(root, item=item):
            sp = find_sp(root, SHAPE["lw_word"][22])
            sz = fit_font(item["word"], 7.78, 67, min_pt=32)
            set_runs(sp, lw_runs(item["word"], item.get("unfair", "")), size_pt=sz)
        d.add(T["learned_word"], fill_lw, notes=item["notes"])
    if lw.get("new"):
        item = lw["new"]
        d.add(T["hdr_new_lw"], notes=HEADER_NOTES["learned_new"])

        def fill_lw_new(root, item=item):
            sp = find_sp(root, SHAPE["lw_word"][25])
            sz = fit_font(item["word"], 7.78, 67, min_pt=32)
            set_runs(sp, lw_runs(item["word"], item.get("unfair", "")), size_pt=sz)
        d.add(T["learned_word_new"], fill_lw_new, notes=item["notes"])

    # --- dictation
    d.add(T["hdr_dictation"], notes=HEADER_NOTES["dictation"])
    dictation_items = session["dictation"]
    if len(dictation_items) == 2:
        meters = [item.get("meter") for item in dictation_items]
        if meters != ["green", "yellow"]:
            warn(f"{day}: dictation meters must progress from green on the "
                 "first sentence to yellow on the trickier second sentence")
    for dictation_index, dt in enumerate(dictation_items):
        meter = dt.get("meter")
        if meter not in T["dictation"]:
            warn(f"{day}: dictation {dictation_index + 1} needs meter "
                 "'green' or 'yellow'")
            meter = "green" if dictation_index == 0 else "yellow"
        dictation_tidx = T["dictation"][meter]
        # The locked yellow slide has no native animation. Clone the green
        # animated slide for both meters, then restore a compatible yellow
        # meter image. This preserves a true appear effect without the
        # LibreOffice/Google Slides position shift caused by transplanting a
        # timing block into the yellow slide.
        clone_tidx = T["dictation"]["green"]
        sentence_spid = SHAPE["dictation_sentence"][clone_tidx]
        cups = dt.get("cups")
        targets = dt.get("targets", [])
        score = dt.get("score")
        if score is None and cups:
            score = (len(cups.get("capitals", [])) + len(cups.get("punctuation", []))
                     + len(targets) + 1)  # +1 = Understanding (dictated order)
        if not cups:
            warn(f"{day}: dictation '{dt['sentence'][:30]}...' has no cups marking block")

        def fill_dict(root, dt=dt, score=score,
                      sentence_spid=sentence_spid):
            sp = find_sp(root, sentence_spid)
            # LibreOffice and Google Slides use slightly taller Lexend metrics
            # than PowerPoint. Keep long 14-16 word reveals inside the locked
            # sentence box so they cannot expand upward over the title.
            sz = fit_font_block(
                dt["sentence"], 9.0, 2.45, 32, min_pt=28,
                char_em=0.62, line_factor=1.5,
            )
            set_runs(sp, dictation_runs(dt["sentence"], dt.get("targets", [])), size_pt=sz)
            bodypr = sp.find(q("p:txBody") + "/" + q("a:bodyPr"))
            if bodypr is not None:
                bodypr.set("anchor", "t")
            regen_timing(root, [(sentence_spid, None)])
            if score:
                add_textbox(root, 90001, 6.8, 4.55, 2.7, 0.45,
                            [f"Score: ___ /{score}"], 20, align="r", anchor="ctr")

        def fill_dict_rels(rels, meter=meter):
            if meter != "yellow":
                return
            yellow_rels = etree.fromstring(
                pkg.slide_rels(T["dictation"]["yellow"])
            )
            yellow_image_rel = next(
                rel for rel in yellow_rels.findall(q("rel:Relationship"))
                if rel.get("Id") == "rId3"
            )
            current_image_rel = next(
                rel for rel in rels.findall(q("rel:Relationship"))
                if rel.get("Id") == "rId3"
            )
            current_image_rel.set("Target", yellow_image_rel.get("Target"))
        lines = [dt["sentence"], ""]
        if cups:
            caps = cups.get("capitals", [])
            punct = cups.get("punctuation", [])
            lines += ["Mark with CUPS, tick each line:",
                      f"C - Capitals, shown green: {', '.join(caps)} ({len(caps)})",
                      "U - Understanding: words written in the dictated order (1)",
                      f"P - Punctuation, shown red: {', '.join(punct)} ({len(punct)})",
                      f"S - Spelling targets, underlined: {', '.join(targets)} ({len(targets)})",
                      f"Score: /{score}"]
        elif targets:
            lines.append(f"Targets: {', '.join(targets)}")
        if dt.get("focus"):
            lines.append(f"Focus: {dt['focus']}")
        d.add(clone_tidx, fill_dict, notes="\n".join(lines),
              fill_rels=fill_dict_rels)

    # --- grammar
    gr = session["grammar"]

    def fill_gr_hdr(root):
        pass
    d.add(T["hdr_grammar"], fill_gr_hdr,
          notes=gr.get("header_notes", ""))
    for key, label in [("i_do", "I do - "), ("we_do", "We do - "), ("you_do", "You do - ")]:
        tidx = T["gradual_release"][key]
        blk = gr[key]
        if "lines" in blk:
            warn(f"{day}: grammar {key} uses the legacy lines layout; use the "
                 "structured rule/example/items format so classroom text stays large")

        def fill_gr(root, tidx=tidx, blk=blk, label=label):
            tsp = find_sp(root, SHAPE["gr_title"][tidx])
            set_text(tsp, label + blk.get("title", ""))
            bsp = find_sp(root, SHAPE["gr_body"][tidx])
            if "lines" in blk:  # legacy plain layout
                set_runs(bsp, [[(ln, None, None, None)] for ln in blk["lines"]] or [[]])
                return
            set_runs(bsp, [[]])  # clear the body placeholder; we draw instead
            draw_structured_slide(root, blk, GRAMMAR_PURPLE, GRAMMAR_TINT,
                                  f"{day} grammar {key}")
        d.add(tidx, fill_gr, notes=blk.get("notes", ""))
        if key == "you_do":
            added_check = add_answer_check_slide(
                d, blk, tidx, GRAMMAR_PURPLE, GRAMMAR_TINT,
                f"{day} grammar you_do answer check",
                default_footer="Check your work against each green answer")
            if has_fixed_answer_notes(blk) and not added_check:
                warn(f"{day}: fixed-answer grammar you_do needs a duplicate "
                     "check_items answer slide")

    d.write(out_path)
    return len(d.slides)


# ---------------------------------------------------------------- QA gate
def qa_deck(path, session=None):
    """Reopen the deck and enforce OG readability/teacher-note invariants."""
    problems = []
    try:
        from pptx import Presentation
        from pptx.enum.dml import MSO_COLOR_TYPE

        pres = Presentation(path)

        with zipfile.ZipFile(path) as zf:
            slide_names = sorted(
                (name for name in zf.namelist()
                 if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
                key=lambda name: int(re.search(r"slide(\d+)", name).group(1)),
            )
            raw_slide_roots = [
                etree.fromstring(zf.read(name)) for name in slide_names
            ]
            raw_slide_rels = []
            for name in slide_names:
                slide_file = name.rsplit("/", 1)[1]
                rel_name = f"ppt/slides/_rels/{slide_file}.rels"
                raw_slide_rels.append(
                    etree.fromstring(zf.read(rel_name))
                    if rel_name in zf.namelist() else None
                )

        def slide_text(slide):
            return "\n".join(
                shape.text for shape in slide.shapes
                if shape.has_text_frame and shape.text
            )

        def rgb(run):
            try:
                if run.font.color.type == MSO_COLOR_TYPE.RGB:
                    return str(run.font.color.rgb)
            except (AttributeError, TypeError):
                pass
            return None

        def fill_rgb(fill):
            try:
                if fill.fore_color.type == MSO_COLOR_TYPE.RGB:
                    return str(fill.fore_color.rgb)
            except (AttributeError, TypeError):
                pass
            return None

        for i, slide in enumerate(pres.slides, 1):
            for sh in slide.shapes:
                if sh.has_text_frame and "XYZ" in sh.text_frame.text:
                    problems.append(f"slide {i}: leftover XYZ placeholder")
            if slide.has_notes_slide and "XYZ" in slide.notes_slide.notes_text_frame.text:
                problems.append(f"slide {i}: leftover XYZ in notes")

        # Notes must be actual PowerPoint paragraphs, not a single iPad text wall.
        with zipfile.ZipFile(path) as zf:
            note_names = sorted(
                (name for name in zf.namelist()
                 if re.fullmatch(r"ppt/notesSlides/notesSlide\d+\.xml", name)),
                key=lambda name: int(re.search(r"notesSlide(\d+)", name).group(1)),
            )
            for note_name in note_names:
                slide_no = int(re.search(r"notesSlide(\d+)", note_name).group(1))
                root = etree.fromstring(zf.read(note_name))
                body = None
                for sp in root.iter(q("p:sp")):
                    ph = sp.find(".//" + q("p:ph"))
                    if ph is not None and ph.get("type") == "body":
                        body = sp
                        break
                if body is None:
                    problems.append(f"slide {slide_no}: notes body placeholder missing")
                    continue
                tb = body.find(q("p:txBody"))
                lines = []
                spelling_word = ""
                spelling_sentence_checked = False
                for p_el in tb.findall(q("a:p")) if tb is not None else []:
                    run_texts = [t_el.text or "" for t_el in p_el.iter(q("a:t"))]
                    line = "".join(run_texts).strip()
                    if not line:
                        continue
                    lines.append(line)
                    if any("\n" in text or "\r" in text or "\x0b" in text
                           for text in run_texts):
                        problems.append(
                            f"slide {slide_no}: note run contains an embedded line break"
                        )
                    if "**" in line:
                        problems.append(f"slide {slide_no}: raw Markdown in notes")
                    if re.search(r"[^\x09\x20-\x7E]", line):
                        problems.append(f"slide {slide_no}: non-ASCII character in notes")
                    if len(line) > 220:
                        problems.append(
                            f"slide {slide_no}: {len(line)}-character note paragraph"
                        )
                    beats = re.findall(
                        r"\b\d+\.\s+(?:SAY|ASK|SCAN|MODEL|POINT|SHOW|READ|"
                        r"REVEAL|CIRCULATE|TIME|WRITE|BUILD|TICK|CHECK|WAIT|"
                        r"REPEAT|CLARIFY)\b",
                        line,
                    )
                    if len(beats) > 1:
                        problems.append(
                            f"slide {slide_no}: multiple numbered beats in one paragraph"
                        )
                    if re.match(r"^\d+\.\s+", line) and (
                            line.count("EXPECT:") > 1 or line.count("?") > 1):
                        problems.append(
                            f"slide {slide_no}: multiple retrieval questions in one paragraph"
                        )
                    ppr = p_el.find(q("a:pPr"))
                    if ppr is None or ppr.find(q("a:buNone")) is None:
                        problems.append(
                            f"slide {slide_no}: notes paragraph may inherit a bullet"
                        )
                    if ppr is None or ppr.get("algn") != "l":
                        problems.append(
                            f"slide {slide_no}: notes paragraph is not left-aligned"
                        )
                    needs_bold_label = bool(
                        NOTE_ANCHOR_RE.search(line) or NUMBERED_ACTION_RE.search(line)
                    )
                    has_bold_run = any(
                        (run.find(q("a:rPr")) is not None
                         and run.find(q("a:rPr")).get("b") == "1")
                        for run in p_el.findall(q("a:r"))
                    )
                    if needs_bold_label and not has_bold_run:
                        problems.append(
                            f"slide {slide_no}: teacher-note label is not PowerPoint bold"
                        )
                    if line.startswith("Word to spell:"):
                        spelling_word = line.split(":", 1)[1].strip()
                    if line.startswith("Sentence:") and spelling_word:
                        spelling_sentence_checked = True
                        has_underlined_word = any(
                            (run.find(q("a:rPr")) is not None
                             and run.find(q("a:rPr")).get("u") == "sng"
                             and "".join(
                                 t_el.text or "" for t_el in run.iter(q("a:t"))
                             ).strip().casefold() == spelling_word.casefold())
                            for run in p_el.findall(q("a:r"))
                        )
                        if not has_underlined_word:
                            problems.append(
                                f"slide {slide_no}: spelling target is not underlined "
                                "inside the Sentence line"
                            )
                if spelling_word and not spelling_sentence_checked:
                    problems.append(
                        f"slide {slide_no}: spelling notes have no Sentence line "
                        "to carry the underlined target"
                    )
                if "---" in lines and lines.index("---") > 8:
                    problems.append(
                        f"slide {slide_no}: notes live zone exceeds 8 Glance Format lines"
                    )
                joined = "\n".join(lines)
                if "Derivative ask:" in joined:
                    problems.append(
                        f"slide {slide_no}: unclear 'Derivative ask' label in notes"
                    )
                if "Meaning: ending that says" in joined:
                    problems.append(
                        f"slide {slide_no}: sound and component meaning are crammed together"
                    )

        if session:
            # The school colour key is global: root green, prefix yellow, suffix red.
            review_cards = session.get("morphology_review", [])
            for offset, card in enumerate(review_cards):
                slide_idx = 3 + offset  # output slides 4-13, zero-based here
                if slide_idx >= len(pres.slides):
                    problems.append("morphology review card slide is missing")
                    break
                slide = pres.slides[slide_idx]
                expected = MORPH_FILL[card["type"]]
                actual = fill_rgb(slide.background.fill)
                if actual != expected:
                    problems.append(
                        f"morphology review card '{card['morph']}' uses {actual} "
                        f"instead of {expected} for {card['type']}"
                    )

            sound_bank = session.get("sound_bank", [])
            if sound_bank and len(pres.slides) > 16:
                sound_shapes = {
                    shape.shape_id: shape for shape in pres.slides[16].shapes
                }
                for idx, item in enumerate(sound_bank[:9]):
                    shape = sound_shapes.get(SHAPE["sound_boxes"][idx])
                    expected = MORPH_FILL[item["type"]]
                    actual = fill_rgb(shape.fill) if shape is not None else None
                    if actual != expected:
                        problems.append(
                            f"sound-bank box '{item['morph']}' uses {actual} "
                            f"instead of {expected} for {item['type']}"
                        )

            new_morph = session.get("new_morphology")
            if new_morph:
                if len(pres.slides) <= 29:
                    problems.append("new morphology card slide is missing")
                else:
                    expected = MORPH_FILL[new_morph["type"]]
                    actual = fill_rgb(pres.slides[29].background.fill)
                    if actual != expected:
                        problems.append(
                            f"new morphology card '{new_morph['morph']}' uses "
                            f"{actual} instead of {expected} for "
                            f"{new_morph['type']}"
                        )

                if len(pres.slides) <= 30:
                    problems.append("new morphology copy-card slide is missing")
                else:
                    copy_slide = pres.slides[30]
                    copy_text = slide_text(copy_slide)
                    entry = verified_card_entry(
                        new_morph["morph"], new_morph.get("type"),
                        new_morph.get("meaning"),
                    )
                    expected_keyword = (
                        entry.get("keyword", "") if entry
                        else new_morph.get("keyword", "")
                    )
                    expected_meaning = (
                        entry.get("meaning", "") if entry
                        else catalogue_meaning(
                            new_morph["morph"], new_morph.get("type")
                        ) or new_morph.get("meaning", "")
                    )
                    expected_lines = [
                        f"{new_morph['type'].capitalize()}: {new_morph['morph']}",
                        f"Keyword: {expected_keyword}",
                        *card_meaning_note(expected_meaning).splitlines(),
                    ]
                    parts = card_parts_of_speech(new_morph)
                    if parts:
                        expected_lines.append(
                            f"Part of speech: {', '.join(parts)}"
                        )
                    for line in expected_lines:
                        if line not in copy_text:
                            problems.append(
                                f"new morphology copy-card slide is missing {line!r}"
                            )
                    actual = fill_rgb(copy_slide.background.fill)
                    expected = MORPH_FILL[new_morph["type"]]
                    if actual != expected:
                        problems.append(
                            "new morphology copy-card slide does not retain the "
                            f"{new_morph['type']} card colour"
                        )
                    if raw_slide_roots[30].find(q("p:timing")) is not None:
                        problems.append(
                            "new morphology copy-card slide should be visible "
                            "immediately, without a reveal animation"
                        )

            # New-morphology grids must stay large enough for classroom reading.
            new_words = [
                item["word"] if isinstance(item, dict) else str(item)
                for item in session.get("words_to_read_new", [])
            ]
            if new_words:
                grid_slide = next(
                    (slide for slide in pres.slides
                     if all(word in slide_text(slide) for word in new_words)),
                    None,
                )
                if grid_slide is None:
                    problems.append("new morphology word grid not found")
                else:
                    sizes = {}
                    for shape in grid_slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                word = run.text.strip()
                                if word in new_words and run.font.size is not None:
                                    sizes[word] = run.font.size.pt
                    missing = [word for word in new_words if word not in sizes]
                    if missing:
                        problems.append(
                            "new morphology words lack explicit font sizes: "
                            + ", ".join(missing)
                        )
                    elif min(sizes.values()) < 27:
                        problems.append(
                            f"new morphology grid falls to {min(sizes.values()):.0f} pt"
                        )

            # Fixed-answer activities require an immediate green check-and-fix slide.
            for activity in (session.get("new_morph_activity", {}),
                             session.get("grammar", {}).get("you_do", {})):
                check_items = activity.get("check_items", [])
                if not check_items:
                    continue
                answer_idx = next(
                    (idx for idx, slide in enumerate(pres.slides)
                     if all(item in slide_text(slide) for item in check_items)),
                    None,
                )
                if answer_idx is None:
                    problems.append("fixed-answer check slide not found")
                else:
                    if "Tick it or fix it -" not in slide_text(
                            pres.slides[answer_idx]):
                        problems.append(
                            "fixed-answer slide does not use the 'Tick it or fix it' title"
                        )
                    if answer_idx == 0 or "You do -" not in slide_text(
                            pres.slides[answer_idx - 1]):
                        problems.append(
                            "fixed-answer check slide does not immediately follow the task"
                        )
                    coloured_items = set()
                    for shape in pres.slides[answer_idx].shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if (run.text.strip() in check_items
                                        and rgb(run) == ANSWER_GREEN):
                                    coloured_items.add(run.text.strip())
                    if coloured_items != set(check_items):
                        problems.append(
                            "check-and-fix answers are not all shown in green"
                        )

            # Dictation slides use colour/format as CUPS information, so verify it.
            punctuation = set(",.;:!?\"'-")
            for item in session.get("dictation", []):
                sentence = item["sentence"]
                target_shape = None
                target_slide = None
                target_slide_index = None
                for slide_index, slide in enumerate(pres.slides):
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip() == sentence:
                            target_shape = shape
                            target_slide = slide
                            target_slide_index = slide_index
                            break
                    if target_shape is not None:
                        break
                if target_shape is None:
                    problems.append(f"dictation sentence not found: {sentence[:32]}")
                    continue
                expected_meter = item.get("meter")
                meter_shape_id = SHAPE["dictation_meter"][T["dictation"]["green"]]
                if (meter_shape_id is None
                        or meter_shape_id not in {
                            shape.shape_id for shape in target_slide.shapes
                        }):
                    problems.append(
                        f"dictation does not use the requested {expected_meter} meter"
                    )
                meter_pic = next(
                    (pic for pic in raw_slide_roots[target_slide_index].iter(q("p:pic"))
                     if pic.find(q("p:nvPicPr") + "/" + q("p:cNvPr")).get("id")
                     == str(meter_shape_id)),
                    None,
                )
                if meter_pic is not None:
                    blip = meter_pic.find(q("p:blipFill") + "/" + q("a:blip"))
                    rel_id = blip.get(q("r:embed")) if blip is not None else None
                    rels = raw_slide_rels[target_slide_index]
                    image_target = next(
                        (rel.get("Target") for rel in rels.findall(q("rel:Relationship"))
                         if rel.get("Id") == rel_id),
                        "",
                    ) if rels is not None else ""
                    expected_image = (
                        "image13.png" if expected_meter == "green"
                        else "image12.png"
                    )
                    if not image_target.endswith(expected_image):
                        problems.append(
                            f"dictation does not use the requested {expected_meter} meter image"
                        )
                timing = raw_slide_roots[target_slide_index].find(q("p:timing"))
                timed_shapes = {
                    target.get("spid")
                    for target in timing.findall(".//" + q("p:spTgt"))
                } if timing is not None else set()
                if str(target_shape.shape_id) not in timed_shapes:
                    problems.append(
                        f"{expected_meter} dictation sentence has no click-to-reveal "
                        "animation"
                    )
                chars = []
                for paragraph in target_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        colour = rgb(run)
                        for char in run.text:
                            chars.append((char, colour, bool(run.font.bold),
                                          bool(run.font.underline)))
                rendered = "".join(char for char, _, _, _ in chars)
                if rendered != sentence:
                    problems.append(
                        f"dictation run text changed: {sentence[:32]}"
                    )
                    continue
                for pos, (char, colour, _, _) in enumerate(chars):
                    if char.isupper() and colour != CAPITAL_GREEN:
                        problems.append(
                            f"dictation capital at character {pos + 1} is not green"
                        )
                    if char in punctuation and colour != PUNCT_RED:
                        problems.append(
                            f"dictation punctuation at character {pos + 1} is not red"
                        )
                lower_sentence = sentence.casefold()
                for target in item.get("targets", []):
                    start = lower_sentence.find(str(target).casefold())
                    if start < 0:
                        problems.append(
                            f"dictation target '{target}' is absent from its sentence"
                        )
                        continue
                    end = start + len(target)
                    if not all(chars[pos][2] and chars[pos][3]
                               for pos in range(start, end)):
                        problems.append(
                            f"dictation target '{target}' is not bold and underlined"
                        )

            # Main grammar examples must not regress to the old tiny template text.
            for key in ("i_do", "we_do", "you_do"):
                block = session.get("grammar", {}).get(key, {})
                example = str(block.get("example", "")).strip()
                if not example:
                    continue
                example_shape = next(
                    (shape for slide in pres.slides for shape in slide.shapes
                     if shape.has_text_frame and shape.text.strip() == example),
                    None,
                )
                if example_shape is None:
                    problems.append(f"grammar {key} example not found")
                    continue
                sizes = [
                    run.font.size.pt
                    for paragraph in example_shape.text_frame.paragraphs
                    for run in paragraph.runs if run.font.size is not None
                ]
                if not sizes or min(sizes) < 20:
                    problems.append(
                        f"grammar {key} example is below the 20 pt readability floor"
                    )
    except Exception as e:
        problems.append(f"failed to reopen: {e}")
    return problems


def main():
    WARNINGS.clear()
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    only = None
    outdir = Path("output")
    for i, a in enumerate(sys.argv[1:]):
        if a == "--only":
            only = sys.argv[1:][i + 1]
        if a == "--outdir":
            outdir = Path(sys.argv[1:][i + 1])
    if not args:
        sys.exit("usage: build_og_week.py <week_spec.json> [--only Day]")
    spec_path = Path(args[0])
    week = deep_sanitize(json.loads(spec_path.read_text()))
    pkg = Package(MASTER)
    unit = week.get("unit_folder", f"OG_Term{week['term']}_Week{week['week']}")
    built = []
    failed = False
    morph_idx = 0
    for session in week["sessions"]:
        # team filename convention: "1a. Monday (-eer).pptx", "1b. Tuesday
        # (-eer review).pptx", ..., "3. Friday (week review).pptx"
        stype = session.get("type", "new")
        if stype == "new":
            morph_idx += 1
            label = f"{morph_idx}a"
        elif stype == "review":
            label = f"{max(morph_idx, 1)}b"
        else:
            morph_idx += 1
            label = str(morph_idx)
        morph = session.get("new_morphology", {}).get("morph", "")
        morph = morph.replace(" / ", "-").replace("/", "-")
        if stype == "week_review" or not morph:
            paren = "week review"
        else:
            paren = morph + (" review" if stype == "review" else "")
        fname = session.get("file_name", f"{label}. {session['day']} ({paren}).pptx")
        if only and session["day"].lower() != only.lower():
            continue
        out_path = outdir / unit / fname
        print(f"Building {out_path} ...")
        n = build_session(pkg, week, session, out_path)
        problems = qa_deck(out_path, session)
        for p in problems:
            print(f"  GATE FAIL: {p}")
        if problems:
            failed = True
        print(f"  {n} slides written.")
        built.append(out_path)
    if WARNINGS:
        print(f"\n{len(WARNINGS)} warning(s) - review above.")
        failed = True
    if failed:
        sys.exit("QA gate failed.")
    print("\nDone:")
    for b in built:
        print(f"  {b}")


if __name__ == "__main__":
    main()
