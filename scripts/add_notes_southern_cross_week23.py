"""Add teacher notes (v1.0 specialist format) to Under the Southern Cross lessons 3-8.

Targets the files in output/literacy_discovery_week23/ and writes ` - with notes.pptx`
copies next to each original.
"""
from pptx import Presentation
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
FOLDER = ROOT / "output" / "literacy_discovery_week23"

FILES = {
    3: "Copy of 3. Lesson Slides Extended Daily lesson 3 2En01V Under the Southern Cross - Informative.pptx",
    4: "Copy of 4. Lesson Slides Extended Daily lesson 4 2En01V Under the Southern Cross - Informative.pptx",
    5: "Copy of 5. Lesson Slides Extended Daily lesson 5 2En01V Under the Southern Cross - Informative.pptx",
    6: "Copy of 6. Lesson Slides Extended Daily lesson 6 2En01V Under the Southern Cross - Informative.pptx",
    7: "Copy of 7. Lesson Slides Extended Daily lesson 7 2En01V Under the Southern Cross - Informative.pptx",
    8: "Copy of 8. Lesson Slides Extended Daily lesson 8 2En01V Under the Southern Cross - Informative.pptx",
}


# =====================================================================
# LESSON 3
# =====================================================================
LESSON_3_NOTES = {
    10: """SOURCES:
Teacher Handbook, Under the Southern Cross (2En01V) - Lesson 3 read-aloud notes.

SAY:
- Today we are going to keep exploring our book, Under the Southern Cross.
- I am going to read the cover and pages 1 to 28 out loud. Your job is to listen carefully and look at the pictures.
- Ask: How do the illustrations help us to understand more about each place and the Southern Cross? Expected: the pictures show what each place looks like at night, who lives there, and where the stars are.

DO:
- Hold the book so all students can see the illustrations.
- On each page, ask students to represent the orientation of the Southern Cross with their hands in a diamond shape.
- When Banjo is found on each page, ask: What do you notice about what Banjo is doing in the picture? Does he help show us something about the place?
- Page 6, pause after "outback Queensland." Ask: What do you see in this picture and how do the faces in the car match the text?
- Page 10, pause after "ping-pong balls." Ask: How does this picture help us understand the role of the man and woman helping the hatchlings?
- Page 20, pause after "mainland too." Ask: What can you tell about the people's feelings by looking at their faces in this picture?

TEACHER NOTES:
Pause points, the Southern Cross hand gesture and the Banjo-spotting routine are taken directly from the Teacher Handbook. Keep them intact and pace the read-aloud so each pause is a real pause, not a rushed aside.

WATCH FOR:
- Students drifting off - bring them back by pointing to an illustration and asking what they can see.
- Students calling out - remind the class we listen first, then share.
- Students forgetting the diamond hand shape - model it again and re-cue on the next page.""",

    12: """SOURCES:
Collins Primary School Dictionary, https://schools.collinsdictionary.com/dictionary/primary, accessed 19.10.25

SAY:
- In the text, the Min Min lights in the Australian outback are described as eerie.
- Something that is eerie is strange or frightening.
- Say the word with me... eerie!
- An old, empty house might look eerie. An owl might sound eerie when it hoots at night.

DO:
- Point to the image on the slide.
- Use a slightly spooky voice so students hear the feeling of the word.
- Have students repeat the word twice.

TEACHER NOTES:
This is the first meeting with the word. Keep it concrete and memorable - students will practise using it across the next few slides.

WATCH FOR:
- Students mixing up eerie with movie-scary - confirm eerie can feel strange or uneasy, not only frightening.""",

    13: """SAY:
- Listen: travellers have seen eerie lights in their travels near Boulia.
- Which answer fits? A - the lights are unusual and frightening. B - the lights are beautiful to look at.
- Show me A or B on your whiteboard.

DO:
- Read both options clearly.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Write A or B. 3, 2, 1, show me.
- Scan for: most boards show A.
PROCEED:
- If most students show A, move on to the next practice slide.
PIVOT:
- Most likely misconception: students pick B because lights can look pretty.
- Reteach: re-read the eerie definition - strange or frightening. Remind students that unusual and frightening is the match.
- Fresh re-check: Say the meaning of eerie with me.

WATCH FOR:
- Students guessing quickly - protect wait time before they show.""",

    14: """SAY:
- Look carefully at the pictures on the slide.
- Point to the picture that could be described as eerie.
- Remember - eerie means strange or frightening.

DO:
- Give 5 seconds wait time.
- Say: 1, 2, 3, point.
- Cold call one student to explain why their picture looks eerie.

TEACHER NOTES:
Active image choice helps students link the word to a real-looking scene. Ask why, not just which.

WATCH FOR:
- Students pointing to the brightest picture just because it stands out - ask them to match to the meaning.""",

    15: """SAY:
- Let's finish this sentence together.
- The forest looked eerie because... what could come next?
- Turn and tell your partner one idea.

DO:
- Give 20 seconds of partner talk.
- Take 2-3 ideas from the room.
- Reveal the example answer and compare with student ideas.

TEACHER NOTES:
Students are stretching the word into new sentences. Accept any reason that matches strange or frightening.

WATCH FOR:
- Answers that only describe the forest (e.g. the forest was big) without a strange or frightening reason - redirect with: strange or frightening?""",

    16: """SAY:
- I will say a feeling. Thumbs up if an eerie thing might make you feel that way, thumbs down if not.
- happy?... scared?... uncomfortable?... safe?

DO:
- Pause between each word so thumbs come up together.
- Scan the room on each one.

TEACHER NOTES:
This checks whether students connect eerie to a feeling, not a shape or size. scared and uncomfortable are thumbs up; happy and safe are thumbs down.

WATCH FOR:
- Students putting thumbs up for all four to avoid choosing - re-cue: only if eerie could make you feel that way.""",

    17: """SAY:
- Show me which word fits best with eerie.
- Hold up 1 for normal, 2 for strange, 3 for bright, 4 for interesting.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which word fits best with eerie? On three, show me 1, 2, 3 or 4.
- Scan for: fingers showing 2 (strange).
PROCEED:
- If the class is on 2, move on to the next vocabulary word.
PIVOT:
- Most likely misconception: students pick bright because some eerie things have unusual light (Min Min lights).
- Reteach: eerie is about the feeling, not whether there is light.
- Fresh re-check: Say the meaning of eerie with me... strange or frightening.

WATCH FOR:
- Students mirroring the fastest hand - keep wait time firm.""",

    18: """SOURCES:
Collins Primary School Dictionary, https://schools.collinsdictionary.com/dictionary/primary, accessed 29.10.25

SAY:
- In the text, aspiring astronomers are keen to use their telescopes to gaze at the Milky Way.
- If you are keen to do something, or for something to happen, you want very much to do it or for it to happen.
- Say the word with me... keen!
- You might be keen to go to the beach. The children were keen to play a game of soccer.

DO:
- Point to the image on the slide.
- Use an enthusiastic voice for keen so students hear the feeling.
- Have students repeat the word twice.

TEACHER NOTES:
Keen is a strong want-to word. Connect it to something the class genuinely wants to do so it sticks.

WATCH FOR:
- Students mixing up keen with good at - keen is about wanting to, not being skilled.""",

    19: """SAY:
- Aspiring astronomers are keen to use their telescopes to look at the Milky Way.
- A - they really want to use their telescopes to look at the Milky Way. B - they are scared to use their telescopes to look at the Milky Way.
- Show me A or B on your whiteboard.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Write A or B. 3, 2, 1, show me.
- Scan for: most boards show A.
PROCEED:
- If the class shows A, move on to the next practice slide.
PIVOT:
- Most likely misconception: students confuse keen with nervous.
- Reteach: keen means you want to do it very much.
- Fresh re-check: Are you keen for the weekend? Thumbs up or down.

WATCH FOR:
- Students guessing from the word aspiring - remind them the key word today is keen.""",

    20: """SAY:
- Which activity would you be keen to do?
- Jump on a trampoline, or clean your room?
- Point to your choice.

DO:
- Give wait time.
- Students point.
- Take 1-2 explanations - why are you keen to do that one?

TEACHER NOTES:
Accept either answer if students can explain keen in the reason. The goal is using the word, not agreeing with a favourite.

WATCH FOR:
- Students changing their choice when a friend picks differently - protect independent thinking with wait time first.""",

    21: """SAY:
- Look at the two children. Which child looks keen to go on the rollercoaster?
- Show me A or B.

DO:
- Wait time.
- Show Me Boards.
- Ask one student how they could tell - what were the clues?

TEACHER NOTES:
Students read facial and body clues to match the word to a feeling. Name the clues (big smile, eyes wide, hands up) when you reveal.

WATCH FOR:
- Students picking the scared child - rewind: keen means want to do it.""",

    22: """SAY:
- I will read four sentences. Thumbs up if the sentence shows the person is keen. Thumbs down if not.
- I don't want to get my hair cut.
- I can't wait to go to the beach!
- I am so excited to ride my bike!
- I am too tired to go to swimming lessons.

DO:
- Pause between sentences for thumbs.
- Scan each time.

TEACHER NOTES:
Students are matching the word to language clues (can't wait, so excited, don't want to, too tired). Name the clue words when you reveal each answer.

WATCH FOR:
- Students copying their neighbour - scan quickly on each one.""",

    23: """SAY:
- Show me which word fits best with keen.
- 1 for bored, 2 for scared, 3 for excited, 4 for surprised.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which word fits best with keen? On three, show me.
- Scan for: fingers showing 3 (excited).
PROCEED:
- If the class is on 3, move on to the sentence-level writing section.
PIVOT:
- Most likely misconception: students pick surprised because excited and surprised feel similar.
- Reteach: keen is wanting to do something before it happens.
- Fresh re-check: Are you keen for something after school? Thumbs up or down.

WATCH FOR:
- Students mirroring the fastest hand - keep wait time firm.""",

    25: """SAY:
- A sentence is a complete thought that makes sense.
- Every sentence has a who or what - the subject.
- Every sentence has a verb - a doing or did word.
- Every sentence starts with a capital letter and ends with a punctuation mark.

DO:
- Point to each icon as you say each rule.
- Have students echo each rule back.

TEACHER NOTES:
This is the anchor for the rest of the lesson. Keep it short and clear - students will practise soon.

WATCH FOR:
- Students getting stuck on the word subject - tie it to who or what every time.""",

    26: """SAY:
- Listen to this sentence: The penguins swim through the icy water to the shore.
- Who or what is doing the action? The penguins. That is the subject.
- What are they doing? They swim. That is the verb.
- Does it start with a capital letter? Yes. Does it end with a full stop? Yes.
- So this is a complete sentence.

DO:
- Point to The penguins and swim as you name each part.
- Read the explanation at the bottom aloud to confirm.

TEACHER NOTES:
Model the thinking aloud step by step so students can copy it on the next slides.

WATCH FOR:
- Students who name icy water as the subject - point back to who is doing the action.""",

    27: """SAY:
- If a group of words is missing a subject or a verb, we call it a fragment.
- A fragment is not a complete sentence.
- Listen: watch the dazzling fireworks. This has a verb, but no subject. Who is watching? We do not know. This is a fragment.
- Listen: In Sydney, the excited families. This has a subject, but no verb. What are they doing? We do not know. This is a fragment.

DO:
- Point to each example as you read it.
- Ask: Is this a complete thought? Answer: no.

TEACHER NOTES:
Keep hammering who/what + what doing. The missing piece is what makes it a fragment.

WATCH FOR:
- Students thinking any short group of words is a fragment - link it back to missing subject or missing verb.""",

    28: """SAY:
- We can turn a fragment into a sentence by adding the missing piece.
- If the fragment has a who or what, we add a what doing.
- If the fragment has a what doing, we add a who or what.
- Watch: In Sydney, the excited families plus watch the dazzling fireworks makes a full sentence.
- In Sydney, the excited families watch the dazzling fireworks.

DO:
- Point to the subject box, then the verb box, then the joined sentence.
- Read the full sentence aloud with expression.

TEACHER NOTES:
This is the core I Do for the lesson. Emphasise that we add only the missing part - we do not rewrite everything.

WATCH FOR:
- Students trying to replace the fragment instead of adding to it - redirect: what is missing?""",

    29: """SAY:
- I will show you a group of words. You decide - is it a sentence or a fragment?
- Hands on head for a sentence. Hands on shoulders for a fragment.
- Let's try each one: the excited supporters... the crocodiles' eyes glow in the torchlight... the mysterious lights shine at night... rises above the ocean... the crowd quietly listens to The Last Post... hatch from their eggs and scuttle towards the ocean.

DO:
- Read each group clearly.
- Give wait time before the signal.
- Scan for hands on head or shoulders.
- Reveal the answer and the corrected sentence each time.

CFU CHECKPOINT:
Technique: Movement (hands on head for sentence, hands on shoulders for fragment)
Script:
- Sentence - hands on head. Fragment - hands on shoulders. Choose now.
- Scan for: correct response for each one.
PROCEED:
- If 80% or more are correct on each, move to the I Do on subject vs verb fragments.
PIVOT:
- Most likely misconception: students think a capital letter or a full stop makes it a sentence.
- Reteach: say the rule again - a sentence needs a subject AND a verb.
- Fresh re-check: the happy dogs - sentence or fragment?

WATCH FOR:
- Students copying their neighbour's movement - scan quickly and change order next time.""",

    30: """SAY:
- The subject is the person or thing doing the action - the who or what.
- A subject fragment has a subject, but is missing a verb.
- Here is an example: In Brisbane, the magical Ferris wheel.
- We have a what - the Ferris wheel - but no doing word. So this is a subject fragment.

DO:
- Point to the magical Ferris wheel and name it as the subject.
- Ask: Is there a doing word? Shake head.

TEACHER NOTES:
Make the name stick - subject fragment has a subject, missing a verb.

WATCH FOR:
- Students reading extra detail words as verbs (e.g. magical) - confirm that verbs show action.""",

    31: """SAY:
- A verb is an action word. It tells what the subject is doing.
- A verb fragment has a verb, but is missing a subject.
- Here is an example: spins rapidly up to the stars.
- There is a doing word - spins - but no who or what. So this is a verb fragment.

DO:
- Point to spins as the verb.
- Ask: Who or what spins? Shake head - we don't know.

TEACHER NOTES:
Mirror the slide before. Say the matching names aloud - verb fragment has a verb, missing a subject.

WATCH FOR:
- Students naming the missing piece as an adjective or noun - steer back to who or what is doing the action.""",

    32: """SAY:
- Listen: The waddling penguins with short legs and large feet.
- Is this a sentence or a fragment? What is missing?
- It has a subject - the waddling penguins - but no verb. This is a subject fragment.
- Let's add a verb: are heading for the dunes.
- Now our sentence is complete: The waddling penguins with short legs and large feet are heading for the dunes.

DO:
- Read the fragment, then ask the question.
- Think aloud as you identify the subject and the missing verb.
- Read the whole sentence aloud when it is complete.

TEACHER NOTES:
This is the last teacher-led example before We Do. Students should see the full thinking before they try one.

WATCH FOR:
- Students who think the description (short legs, large feet) is the verb - point back to what doing.""",

    33: """SAY:
- Read with me: twinkle like stars in the night sky.
- Is this a subject fragment or a verb fragment? What is missing?
- Turn and tell your partner what is missing.
- Now let's add it together - The city lights twinkle like stars in the night sky.

DO:
- Give 20 seconds of partner talk.
- Take 1-2 responses.
- Build the full sentence on the board, or reveal the completed one.

CFU CHECKPOINT:
Technique: Turn and Tell, then cold call
Script:
- Is this a subject fragment or a verb fragment?
- Scan for: students saying verb fragment.
PROCEED:
- If most pairs say verb fragment, move to the next example.
PIVOT:
- Most likely misconception: students say subject fragment because it looks short.
- Reteach: name the verb (twinkle), then ask who or what twinkles.
- Fresh re-check: float on the sea - subject or verb fragment?

WATCH FOR:
- Students adding extra detail but not a subject - redirect: who or what twinkles?""",

    34: """SAY:
- Read with me: In the evening, the excited children with huge smiles on their faces.
- What is missing? Is this a subject fragment or a verb fragment?
- This is a subject fragment. It needs a verb.
- Let's add one together - go prawn fishing on the beach.
- Full sentence: In the evening, the excited children with huge smiles on their faces go prawn fishing on the beach.

DO:
- Read the fragment.
- Partner talk: what is missing?
- Take a response and write the verb on the board.
- Read the completed sentence with expression.

TEACHER NOTES:
Keep the focus on identifying the fragment type first, then adding the missing piece.

WATCH FOR:
- Students who try to replace the long description - remind them only the missing piece needs adding.""",

    35: """SAY:
- Read with me: Near the ocean, the tiny hatchlings.
- Is this a subject fragment or a verb fragment?
- This is a subject fragment. It needs a verb.
- Let's add one - break free from their eggs.
- Full sentence: Near the ocean, the tiny hatchlings break free from their eggs.

DO:
- Read the fragment.
- Partner talk - what doing word could we add?
- Take a response, write the verb, read the full sentence.

TEACHER NOTES:
Third We Do fragment - pace should pick up as confidence grows.

WATCH FOR:
- Students adding only are - push for a stronger action verb.""",

    36: """SAY:
- Read with me: cheer loudly for their favourite team.
- Is this a subject fragment or a verb fragment?
- This is a verb fragment. It needs a subject.
- Let's add one - At the Melbourne Cricket Ground, the excited fans.
- Full sentence: At the Melbourne Cricket Ground, the excited fans cheer loudly for their favourite team.

DO:
- Read the fragment.
- Partner talk - who or what cheers?
- Take a response and complete the sentence.

TEACHER NOTES:
Students should now notice verb fragments quickly. Celebrate correct identification before rewriting.

WATCH FOR:
- Students using a pronoun only (they cheer) - accept it, then stretch: can we add more detail about who?""",

    37: """SAY:
- Look at the starting words: In the evening, local families.
- Which example shows this fragment correctly turned into a complete sentence - A or B?
- A: In the evening, local families with big smiles and excited faces.
- B: In the evening, local families watch exciting movies at the outdoor cinema.
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Write A or B. 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move to the booklet task.
PIVOT:
- Most likely misconception: students pick A because it sounds longer and more detailed.
- Reteach: check for a verb. A has no verb. B has watch - that is the verb that completes the sentence.
- Fresh re-check: The tired koala sleeps in the tree - sentence or fragment?

WATCH FOR:
- Students who assume more words always means a sentence - pull back to subject plus verb.""",

    38: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 3 - Sentence level writing.
- Read each fragment carefully.
- Decide if it is a subject fragment or a verb fragment.
- Add the missing piece to make a complete sentence.

DO:
- Hand out booklets or direct students to the page.
- Circulate and check the first 2-3 answers for each student.
- Remind students of the rule - every sentence needs a subject and a verb.

TEACHER NOTES:
Use this time to scan for students who still confuse subject and verb fragments. Pull a small group if needed.

WATCH FOR:
- Students who write new sentences that still miss a verb - prompt: what is the doing word?
- Students finishing fast - ask them to read their sentences aloud and check for a capital letter and full stop.""",
}


# =====================================================================
# LESSON 4
# =====================================================================
LESSON_4_NOTES = {
    10: """SAY:
- Authors write texts for different reasons.
- Some texts are written to entertain - like a funny story or a picture book we read for fun.
- Some texts are written to inform - like a fact book about animals.
- Some texts are written to persuade - like a poster that tries to get you to buy something.

DO:
- Point to each purpose as you name it.
- Give one quick familiar example for each (e.g. a picture book we have read, a non-fiction animal book, a school fundraiser poster).

TEACHER NOTES:
This is the anchor for the next few slides. Keep the three purposes short and clear so students can sort texts by purpose.

WATCH FOR:
- Students who think all books are to entertain - show the non-fiction example again.""",

    11: """SAY:
- Why do you think Frane Lessac wrote Under the Southern Cross?
- Was it to entertain us, to inform us, or to persuade us?
- Turn and tell your partner what you think and why.

DO:
- Give 20 seconds of partner talk.
- Take 2-3 responses from the room.
- Affirm that the book is a mix - we learn facts about Australia AND enjoy the journey. We will unpack this next.

TEACHER NOTES:
Accept to entertain, to inform, or both, as long as students can give a reason. This sets up the idea of a hybrid text on the next slide.

WATCH FOR:
- Students who only say to entertain - prompt back to the factual parts of the book (the facts about each place).""",

    12: """SOURCES:
Image credit: Free SVG

SAY:
- The word hybrid means something that is a mix of two different things.
- For example, a hybrid car has two ways of getting power - electricity and petrol.
- Say the word with me... hybrid!

DO:
- Point to the hybrid car image.
- Use your fingers to show two parts coming together.

TEACHER NOTES:
Use the car as a concrete anchor before applying the word to a text. Students need the everyday meaning first.

WATCH FOR:
- Students who hear hybrid for the first time - have them repeat it twice.""",

    13: """SAY:
- A hybrid text combines features of two types of texts together into one text.
- Under the Southern Cross combines a narrative journey layer with an informative layer.
- The narrative journey layer tells a story as we travel around Australia at night.
- The informative layer gives us facts about each place.

DO:
- Point to the word narrative, then to the word informative.
- Hold up the book and flick through a page so students can see the two different print sizes.

TEACHER NOTES:
Name both layers clearly. Students will be asked to spot them on the next slides.

WATCH FOR:
- Students who think the two layers are two different books - re-show the single page with both types of text.""",

    14: """SOURCES:
Image credit: Free SVG

SAY:
- Open the book to pages 1 and 2.
- Can you see two layers on the page? One tells the story, one gives us facts.
- Why do you think the two layers are printed and positioned differently?
- What is different about the language features used in both layers?

DO:
- Make sure every student can see pages 1-2 (show the book, or direct them to their copies).
- Partner talk for 30 seconds.
- Take 2-3 responses and highlight the bigger print (narrative) versus the smaller factual text (informative).

TEACHER NOTES:
The layout itself signals the two layers. Name the signals: bigger print and story-style is narrative; smaller print with facts is informative.

WATCH FOR:
- Students who focus only on the illustrations - pull them back to the two text layers.""",

    15: """SAY:
- Let's look at pages 15 and 16.
- On these pages, we can see the language features that belong to each layer.
- Narrative layer uses adjectives, rhythmic phrases, repetition of Under the Southern Cross, and imagery.
- Informative layer uses factual information, technical language, proper nouns, geographical references, and explanatory paragraphs.

DO:
- Read a short example of narrative language from the page.
- Read a short example of informative language from the page.
- Point to each feature on the slide as you name it.

TEACHER NOTES:
This is a big list - do not teach every feature deeply. Name them and use the examples on the page to anchor the key difference: story-sounding words versus fact-sounding words.

WATCH FOR:
- Students overwhelmed by the vocabulary - focus on two contrasts: descriptive words for narrative, and proper nouns or facts for informative.""",

    16: """SAY:
- Listen to this sentence: Across Australia, the Southern Cross is a constellation that can be seen in the night sky year-round.
- Is this narrative or informative?
- I can hear a fact - year-round - and a technical word - constellation. This is the informative layer.

DO:
- Read the sentence aloud twice.
- Think aloud as you identify the clues (fact, technical word).
- Point to the word informative as you confirm.

TEACHER NOTES:
This is the I Do for sorting by layer. Model the thinking students will copy on the next slides.

WATCH FOR:
- Students who just vote without clues - name the clue word (constellation, year-round) each time.""",

    17: """SAY:
- Listen: In the enchanted rainforest, eyes gleam in the dark while frogs sing their midnight songs.
- Is this the narrative layer or the informative layer?
- Turn and tell your partner your answer and one clue word.

DO:
- Read the sentence clearly.
- Partner talk for 20 seconds.
- Cold call 1-2 pairs to share.
- Reveal and name the clues (enchanted, gleam, sing - descriptive story words).

TEACHER NOTES:
Expected answer: narrative layer. Students should point to descriptive or story-sounding words.

WATCH FOR:
- Students pointing to rainforest as a fact - remind them that facts sound like information, not a story scene.""",

    18: """SAY:
- Listen: In Broome, the Staircase to the Moon appears when the full moon rises over the ocean at low tide.
- Is this narrative or informative?
- Show me with your hand - thumbs up if informative, thumbs down if narrative.

DO:
- Read the sentence twice.
- Give 5 seconds wait time.
- Scan thumbs.
- Reveal and name the clues (Broome as proper noun, factual description of when it happens).

TEACHER NOTES:
Expected answer: informative. Even though it sounds lovely, the clues are proper noun and factual explanation.

WATCH FOR:
- Students who pick narrative because of the word moonlight-type images - re-cue on the fact that it explains when and where.""",

    19: """SAY:
- Read with me.
- A: In the outback, glowing balls called Min Min lights are sometimes seen at night.
- B: The desert glows softly as strange orbs of light float gently across the sand. A dingo watches quietly nearby.
- Which sentence is written with narrative features? Show me A or B.

DO:
- Read both options clearly.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which has narrative features? 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move on to the next CfU.
PIVOT:
- Most likely misconception: students pick A because Min Min lights sound interesting.
- Reteach: narrative uses descriptive story words - glows softly, float gently, quietly. A uses facts - are sometimes seen.
- Fresh re-check: Read a short narrative sentence and a short informative sentence. Which is which?

WATCH FOR:
- Students who choose on length - remind them the clue is in the type of words, not how long the sentence is.""",

    20: """SAY:
- Read with me.
- A: Under the Southern Cross, the waves glitter as sea turtles crawl towards the moonlit beach, while a young girl watches in silence from the dunes.
- B: The Aurora Australis is a natural light display in the sky, seen in the southern parts of the world.
- Which sentence is written with informative features? Show me A or B.

DO:
- Read both options clearly.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which has informative features? 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move to the booklet task.
PIVOT:
- Most likely misconception: students pick A because it mentions the Southern Cross and feels factual.
- Reteach: B explains what the Aurora Australis is and where it is seen - that is an informative explanation. A describes a scene with characters.
- Fresh re-check: Say one clue word for an informative sentence - expected: explain, fact, describes what something is.

WATCH FOR:
- Students who mix up the two layers - re-anchor: narrative = story, informative = facts.""",

    21: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 4 - Understanding, responding to and creating text.
- Read each sentence and decide which layer it belongs to.

DO:
- Hand out booklets or direct students to the page.
- Circulate and check the first 2-3 answers for each student.
- Remind students to look for the clue words each time.

TEACHER NOTES:
This is an independent practice moment. Pull a small group if many students are still unsure after the CfU checkpoints.

WATCH FOR:
- Students who pick randomly - ask them to underline one clue word in each sentence before they decide.""",

    23: """SAY:
- A noun is a naming word.
- It names a person, place, animal or thing.
- We also call this a common noun.
- A common noun does not need a capital letter.
- Say these examples with me... penguin... cinema... turtles.

DO:
- Point to each example on the slide.
- Echo read the word list together.

TEACHER NOTES:
Keep the noun definition ultra simple before introducing proper nouns on the next slide.

WATCH FOR:
- Students who try to add a capital to penguin - remind them common nouns are general names, no capital needed.""",

    24: """SOURCES:
The Grammar Book: Understanding and teaching primary grammar, Zoe Paramour and Timothy Paramour (2020).

SAY:
- A proper noun is special - it names a specific person, place, object, animal or title.
- A proper noun always begins with a capital letter.
- Say these examples with me... Southern Cross... Phillip Island... Milky Way.

DO:
- Point to each proper noun on the slide.
- Show the capital letter at the start of each one.

TEACHER NOTES:
This is the key rule of the lesson - proper nouns always start with a capital letter. Students will practise this next.

WATCH FOR:
- Students who think all nouns need capitals - remind them penguin is a common noun, Phillip Island is a proper noun.""",

    25: """SAY:
- Let's look at common and proper noun examples from our book.
- Common noun examples: penguin, cinema, lights, turtle, fireworks, bush - no capital letter.
- Proper noun examples: Southern Cross, Phillip Island, Queensland, Min Min Hotel, Australian Football League, Milky Way - capital letter at the start of each word.

DO:
- Read one column, then the other.
- Point to a capital letter on each proper noun.

TEACHER NOTES:
Connect the grammar content back to Under the Southern Cross so students see these nouns in context.

WATCH FOR:
- Students who think Min Min Hotel needs only one capital - point to each word that needs a capital.""",

    26: """SAY:
- Read the words on the screen.
- Which words are proper nouns? Which ones need a capital letter?
- Animals, eggs, November, Darwin, ocean, South Pacific.
- Turn and tell your partner which ones are proper nouns.

DO:
- Give 20 seconds of partner talk.
- Cold call 2-3 students to share.
- Reveal: November, Darwin, South Pacific.

TEACHER NOTES:
Students pick out the proper nouns from a mixed list. The clue is that proper nouns name specific people, places or things.

WATCH FOR:
- Students who pick ocean or animals - remind them these are general names (common nouns), not specific ones.""",

    27: """SAY:
- Read the words.
- Gallipoli, fire, South Pole, dawn, Perth, crocodiles.
- Which ones are proper nouns? Write them down with a capital letter at the start.

DO:
- Give 30 seconds of independent writing time.
- Circulate and scan 3-4 books.
- Reveal: Gallipoli, South Pole, Perth.

CFU CHECKPOINT:
Technique: Independent write, then cold call for answers
Script:
- Write the proper nouns on your whiteboard with a capital letter at the start.
- Scan for: Gallipoli, South Pole, Perth.
PROCEED:
- If most students have all three, move on to capital letter correction.
PIVOT:
- Most likely misconception: students capitalise crocodiles or fire because they sound important.
- Reteach: Is it a specific place, person, or title? Crocodiles names a type of animal - common noun.
- Fresh re-check: Is Sydney a proper noun? Thumbs up or down.

WATCH FOR:
- Students who miss Perth because it is a single word - remind them that cities are proper nouns.""",

    28: """SAY:
- Read this sentence: The southern cross can be seen all year round.
- What is the proper noun? Southern Cross - a specific name.
- So we need capital letters - S for Southern, C for Cross.
- The corrected sentence is: The Southern Cross can be seen all year round.

DO:
- Read the original sentence aloud.
- Point to southern cross and circle it with your finger.
- Read the corrected sentence with the capital letters emphasised.

TEACHER NOTES:
Model the thinking: find the proper noun, then check every word in the name for a capital letter.

WATCH FOR:
- Students who only capitalise Southern and leave cross lowercase - point out that the full name needs capitals on every word.""",

    29: """SAY:
- Read this sentence with me: There are 100 to 400 billion stars in the milky way.
- What is the proper noun that needs a capital? The Milky Way.
- Rewrite the sentence with a capital M for Milky and a capital W for Way.
- Turn and tell your partner the corrected version.

DO:
- Give 20 seconds of partner talk.
- Take one response.
- Reveal: There are 100 - 400 billion stars in the Milky Way.

TEACHER NOTES:
Keep the focus on the proper noun only. Other words in the sentence do not change.

WATCH FOR:
- Students who capitalise stars or billion - redirect to what is a specific name here.""",

    30: """SAY:
- Read this sentence with me: Loggerhead turtles hatch at mon repos near bundaberg.
- There are two proper nouns that need capitals - Mon Repos and Bundaberg.
- Rewrite the sentence with the capital letters in the right places.

DO:
- Partner talk for 30 seconds.
- Take a response.
- Reveal: Loggerhead turtles hatch at Mon Repos near Bundaberg.

TEACHER NOTES:
This one has two proper nouns, which is the stretch. Prompt students to find both names before writing.

WATCH FOR:
- Students who only fix one name - ask: is there another specific place here?""",

    31: """SAY:
- Look at the two sentences.
- The Deckchair Cinema is home to great movies and lots of animals.
- The wheel of brisbane looks over the Brisbane River.
- Are these sentences punctuated correctly? Show me thumbs up if yes, thumbs down if no.

DO:
- Read each sentence clearly.
- Give 10 seconds of thinking time.
- Scan thumbs.
- Reveal the correction: The Wheel of Brisbane looks over the Brisbane River.

CFU CHECKPOINT:
Technique: Thumbs up or thumbs down
Script:
- Is sentence 2 correct? Thumbs up or down.
- Scan for: thumbs down (it is not correct - wheel of brisbane needs capitals).
PROCEED:
- If most students show thumbs down, move on to sentence expansion.
PIVOT:
- Most likely misconception: students miss that wheel of brisbane is the name of a specific thing.
- Reteach: the Wheel of Brisbane is a specific landmark, just like the MCG - it is a title, so every main word needs a capital.
- Fresh re-check: the sydney opera house - correct or not?

WATCH FOR:
- Students who only capitalise Wheel - remind them of, and, the stay lowercase but Brisbane is a proper noun.""",

    32: """SAY:
- We can use what we know about common and proper nouns to expand sentences and add more detail.
- Listen: They are watching a movie at night.
- Who is doing it? We do not know. Let's add a who/what.
- The families are watching a movie at night.
- Now we can read pages 3 and 4 to get even more detail.

DO:
- Point to They, then to the families.
- Point to each part of the build as you name it.
- Read the final expanded sentence aloud with expression.

TEACHER NOTES:
This is the I Do for expanding sentences. The move students need to copy: swap or add a specific who/what to add detail.

WATCH FOR:
- Students who rewrite the whole sentence - remind them we only add or swap the who/what piece.""",

    33: """SAY:
- Read pages 7 and 8 with me.
- Here is a simple sentence: He plays football at the MCG.
- Who is he? Can we add more detail to the who/what?
- Let's try: the athlete plays football.
- Full expanded sentence: The athlete plays football at the MCG.

DO:
- Point to He, then swap it for the athlete.
- Partner talk - can we add more detail about who?
- Take a response and reveal the expanded sentence.

TEACHER NOTES:
The MCG is a great proper noun example - link this back to capital letters for proper nouns.

WATCH FOR:
- Students who add more detail after football - gently redirect: we are expanding the who/what piece today.""",

    34: """SAY:
- Read pages 15 and 16 with me.
- Here is a simple sentence: They exploded in bright colours over the sky.
- Who or what is they?
- Let's add detail: the fireworks exploded in bright colours.
- Full sentence: The fireworks exploded in bright colours over the sky.

DO:
- Point to They on the slide.
- Partner talk - who or what exploded?
- Take a response, then reveal.

TEACHER NOTES:
Students should be faster at this move now. Celebrate the swap - they becomes the fireworks.

WATCH FOR:
- Students who say the word they again - prompt: can we say exactly what kind of thing?""",

    35: """SAY:
- Read pages 21 and 22 with me.
- Simple sentence: They can be caught in the Swan River.
- Who or what is they?
- Let's add: prawns can be caught.
- Full sentence: Prawns can be caught in the Swan River.

DO:
- Point to They.
- Partner talk to identify the who/what.
- Take a response and reveal.

TEACHER NOTES:
Swan River is another proper noun in context - point to the capital S and capital R while reading.

WATCH FOR:
- Students who forget the capital letter on Prawns at the start of the sentence - remind them every sentence starts with a capital letter.""",

    36: """SAY:
- Read the sentence: It tells us about our military's history.
- We need a who/what to expand this sentence. Read pages 26 and 27 to help.
- Which detail fits best? A - The Australian War Memorial. B - The Last Post.
- Show me A or B.

DO:
- Read both options clearly.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which best expands the sentence? 3, 2, 1, show me.
- Scan for: most boards show A (The Australian War Memorial).
PROCEED:
- If most students show A, move to the booklet task.
PIVOT:
- Most likely misconception: students pick B because The Last Post appears in the book and is memorable.
- Reteach: The Last Post is a tune, not a place that tells us about history. The Australian War Memorial is the place that tells us about our military's history.
- Fresh re-check: point to a proper noun on the book cover - is it a thing or a tune?

WATCH FOR:
- Students who forget to look at the original sentence meaning - re-read: it tells us about our military's history. Which one does that?""",

    37: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 4 - Sentence level writing.
- Remember to use a capital letter for every proper noun.
- Add a who/what detail to expand each sentence.

DO:
- Hand out booklets or direct students to the page.
- Circulate and check the first 2-3 answers for each student.
- Remind students that common nouns do not need a capital, but proper nouns do.

TEACHER NOTES:
Scan for two errors: missing capital letters on proper nouns, and replacing instead of adding detail. Pull a small group if needed.

WATCH FOR:
- Students who leave proper nouns lowercase - point back to the capital letter rule.
- Students finishing fast - ask them to read their sentences aloud and check that every proper noun has a capital letter.""",
}


# =====================================================================
# LESSON 5
# =====================================================================
LESSON_5_NOTES = {
    10: """SOURCES:
Teacher Handbook, Under the Southern Cross (2En01V) - Lesson 5 read-aloud notes.

SAY:
- Last time we read, we noticed how the pictures help us understand more about each place.
- Today, we are reading the cover and pages 1 to 28 again. This time, listen for the two layers - the journey story and the facts.
- Ask: Why do you think the author created a hybrid text with both a journey and factual information?

DO:
- Hold the book so all students can see both text layers on each page.
- Page 2, pause after "rocking motion." Ask: How does the factual information help us understand why the penguins are waddling up the beach at night for the crowd to watch?
- Page 6, pause after "outback Queensland." Ask: Does the factual information help explain why the light is so mysterious?
- Page 12, pause after "year treat." Ask: Why might the author describe the Ferris wheel as giving a sweeping view?
- Page 18, pause after "mud here!" Ask: How does the factual layer help us understand what the staircase of gold actually is?
- Page 22, pause after "the wild." Ask: What information is surprising in the informative layer, compared to the description in the narrative layer?

TEACHER NOTES:
Pause points come straight from the Teacher Handbook. Use the read-aloud to set up today's bigger thinking - why a hybrid text works for this topic.

WATCH FOR:
- Students drifting on a long read - re-anchor with a quick gesture or page-point.
- Students hearing only the narrative layer - re-read the smaller factual text after each pause.""",

    12: """SOURCES:
Collins Primary School Dictionary, https://schools.collinsdictionary.com/dictionary/primary

SAY:
- In the text, the Min Min lights are described as a bizarre mystery in the outback.
- Something that is bizarre is very strange and weird.
- Say the word with me... bizarre!
- A purple cow would look bizarre. A fish riding a bike would be bizarre.

DO:
- Point to the image on the slide.
- Use a slightly surprised voice for bizarre so students hear the strangeness in the word.
- Have students repeat the word twice.

TEACHER NOTES:
Bizarre is funny-strange, not scary-strange. Use a playful tone to anchor the difference from eerie.

WATCH FOR:
- Students mixing up bizarre with bad - bizarre just means very strange or weird, not bad.""",

    13: """SAY:
- The Min Min lights are a bizarre mystery in the outback.
- Which answer fits? A - something funny that happens in the outback. B - something odd or weird that happens in the outback.
- Show me A or B on your whiteboard.

DO:
- Read both options clearly.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Write A or B. 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move on to the next practice slide.
PIVOT:
- Most likely misconception: students pick A because bizarre things can be funny too.
- Reteach: bizarre means very strange or weird. Funny is sometimes part of it, but the main meaning is strange.
- Fresh re-check: a singing pizza - bizarre or not bizarre?

WATCH FOR:
- Students guessing on funny because the word sounds light - bring them back to strange or weird.""",

    14: """SAY:
- Listen carefully. It would be bizarre if my teacher...
- gave me homework... wore a clown costume to school... read a book to the class... had a nap in the reading corner.
- Show me a thumb up if it would be bizarre. Thumbs down if it would not.

DO:
- Read each option clearly.
- Pause after each one for thumbs.
- Reveal the bizarre options - clown costume, nap in the reading corner.

TEACHER NOTES:
Two answers are bizarre - the clown costume and the nap. Use a playful voice to keep the energy high.

WATCH FOR:
- Students putting all four thumbs up - re-cue: only if it would be very strange or weird.""",

    15: """SAY:
- Look at the two pictures.
- Which one describes something bizarre - your friends wishing you happy birthday on your birthday, or a frog cooking you dinner?
- Show me A or B.

DO:
- Read both options.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Cold call one student to explain why their picture is bizarre.

TEACHER NOTES:
Expected answer is the frog cooking dinner. Birthday wishes are normal, not bizarre.

WATCH FOR:
- Students picking birthday wishes - remind them bizarre means very strange. A birthday is not strange.""",

    16: """SAY:
- Turn and tell your partner one bizarre thing it would be to see at school.
- Remember - bizarre means very strange or weird.

DO:
- Give 30 seconds of partner talk.
- Take 2-3 ideas from the room.
- Affirm strange or weird ideas - playful tone is fine.

TEACHER NOTES:
Students try the word in their own sentence. Accept any imaginative idea that fits very strange or weird.

WATCH FOR:
- Students sharing normal school events (like a maths lesson) - prompt: is that strange or weird? What would be more bizarre?""",

    17: """SAY:
- Show me which word fits best with bizarre.
- A for normal, B for strange, C for boring.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Me Boards (A, B or C)
Script:
- Which word fits best with bizarre? 3, 2, 1, show me.
- Scan for: most boards show B (strange).
PROCEED:
- If the class is on B, move on to the next vocabulary word.
PIVOT:
- Most likely misconception: students pick boring or normal because they cannot recall the meaning.
- Reteach: bizarre means very strange or weird. Boring is the opposite.
- Fresh re-check: Say the meaning of bizarre with me.

WATCH FOR:
- Students mirroring a friend - keep wait time firm before showing.""",

    18: """SOURCES:
Collins Primary School Dictionary, https://schools.collinsdictionary.com/dictionary/primary

SAY:
- In the text, gigantic lanterns are carried along the riverbank in Brisbane.
- Something that is gigantic is extremely large.
- Say the word with me... gigantic!
- A whale is gigantic. The MCG is gigantic.

DO:
- Point to the image on the slide.
- Stretch your arms wide as you say gigantic.
- Have students repeat the word twice with arms stretched out.

TEACHER NOTES:
Gigantic is bigger than big and bigger than huge. Use the body gesture so students feel the size in the word.

WATCH FOR:
- Students using gigantic for everyday objects - remind them it means extremely large, not just big.""",

    19: """SAY:
- We learnt that gigantic lanterns are carried along the riverbank in Brisbane.
- Which answer fits? A - the lanterns are very large. B - the lanterns are very colourful.
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which answer fits gigantic? 3, 2, 1, show me.
- Scan for: most boards show A.
PROCEED:
- If most students show A, move on.
PIVOT:
- Most likely misconception: students pick B because lanterns are often colourful.
- Reteach: gigantic is about size, not colour. Very large is the match.
- Fresh re-check: a tiny mouse - gigantic or not?

WATCH FOR:
- Students confusing gigantic with bright or colourful - re-anchor on size.""",

    20: """SAY:
- Look at the objects on the slide.
- Which ones could be described as gigantic?
- Point to the gigantic ones.

DO:
- Give 5 seconds of looking time.
- Say: 1, 2, 3, point.
- Cold call 1-2 students to explain why their object is gigantic.

TEACHER NOTES:
Students sort by size. Accept any object that is extremely large in real life. Push back if they pick something small.

WATCH FOR:
- Students pointing at every picture - remind them gigantic is extremely large, not just big.""",

    21: """SAY:
- Let's stretch this word into different sentences.
- The Eiffel Tower is gigantic because... what could come next?
- The Eiffel Tower is gigantic, but... what could come next?
- The Eiffel Tower is gigantic, so... what could come next?
- Turn and tell your partner one ending for each.

DO:
- Read each starter clearly.
- Give 60 seconds of partner talk in total.
- Take one response for each ending.
- Reveal the example answers and compare.

TEACHER NOTES:
This is sentence-stretching with conjunctions. The aim is for students to use gigantic in three different sentence shapes, not to memorise model answers.

WATCH FOR:
- Students who finish only one stem - prompt them to try a second one.
- Students who repeat the same idea - push for a different angle (cause, contrast, result).""",

    22: """SAY:
- Show me which word fits best with gigantic.
- A for huge, B for big.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which word fits best with gigantic? 3, 2, 1, show me.
- Scan for: most boards show A (huge).
PROCEED:
- If the class is on A, move on to the next section.
PIVOT:
- Most likely misconception: students pick big because both fit, but big is too small for gigantic.
- Reteach: huge is the closest match because gigantic means extremely large, not just larger than usual.
- Fresh re-check: an elephant is gigantic. Is it big, or is it huge?

WATCH FOR:
- Students who pick the easier word - keep stretching them to the strongest match.""",

    24: """SAY:
- The scenes in Under the Southern Cross are drawings of real places and real things people do at night, all around Australia.
- When we read, we can connect the story world to our own lives and say how it makes us feel.
- Today we will share what we think and feel about different places in the book.

DO:
- Hold the book so students can see one or two scenes.
- Point to the word connect on the slide.
- Echo the word connect together.

TEACHER NOTES:
This is the anchor for personal responses. Keep the focus on connecting our own experiences to scenes in the book.

WATCH FOR:
- Students who think connect means draw a line - clarify it means link our own life to the story.""",

    25: """SAY:
- Open the book to pages 3 and 4 - the Deckchair Cinema scene.
- Look carefully at the picture.
- What do you think about this scene? Have you been somewhere like this? Would you like to go here?
- Turn and tell your partner.

DO:
- Make sure every student can see pages 3-4.
- Give 30 seconds of partner talk.
- Take 2-3 responses from the room.

TEACHER NOTES:
This is the first We Do for personal response. Accept any honest answer with a brief reason.

WATCH FOR:
- Students who only describe the picture - prompt back to: would YOU like to go here? Why?""",

    26: """SAY:
- Listen to my answer for pages 3 and 4.
- I'd like to go here because it would be a great experience to see a movie under the stars.
- This reminds me of going to the drive-in movies because we sat outside to watch the movie.
- Notice how I gave a reason after because.

DO:
- Read both sentence stems clearly.
- Read the model answers with expression.
- Point to each because and ask: what is the reason?

TEACHER NOTES:
This is the I Do for sentence stems. Make the because part visible so students copy the structure.

WATCH FOR:
- Students who copy the model word for word - remind them to share their own reason.""",

    27: """SAY:
- Open the book to pages 7 and 8 - the AFL game scene.
- Would you like to be in this scene? Why or why not?
- Turn and tell your partner.

DO:
- Make sure every student can see pages 7-8.
- Give 30 seconds of partner talk.
- Take 2-3 responses.

TEACHER NOTES:
Second We Do for personal response. Football fans and non-fans both have valid responses - accept both with reasons.

WATCH FOR:
- Students who only say yes or no without a reason - prompt: tell me why.""",

    28: """SAY:
- Listen to my answer for pages 7 and 8.
- I'd like to go here because it is exciting to cheer for my team during a game.
- This reminds me of going to the ANZAC day match because I sat in the stands and cheered for my favourite team.

DO:
- Read both sentence stems clearly.
- Read the model answers with expression.
- Echo the word because.

TEACHER NOTES:
Use this to reinforce the structure. The reason is what makes a real connection.

WATCH FOR:
- Students who repeat the model exactly - prompt them for their own moment that reminds them.""",

    29: """SAY:
- Time to make your own connections to the book.
- Walk around the room and look at the scenes from Under the Southern Cross.
- Pick two scenes you want to write about.
- On one colour sticky note, write: I'd like to go here because...
- On the other colour sticky note, write: This reminds me of...
- Stick them on your chosen scenes.

DO:
- Hand out two sticky notes per student.
- Show the two sentence starters.
- Set a 5 minute timer for the gallery walk.
- Circulate and prompt students who freeze.

TEACHER NOTES:
This is the You Do. Movement plus writing keeps engagement high. Keep a list of students who finish quickly so you can stretch them.

WATCH FOR:
- Students who write only the sentence starter - prompt them to add a reason after because.
- Students copying a friend - send them to a different scene.""",

    30: """SAY:
- Look at the scenes around the room.
- Use your sticky notes to share your connection.
- I'd like to go here because... or This reminds me of...

DO:
- Show the two sentence starters again.
- Circulate and read 4-5 sticky notes to scan responses.
- Pick 2 strong examples to share with the class shortly.

TEACHER NOTES:
This is the same activity continued from the previous slide. Use the time to gather examples for the share-back.

WATCH FOR:
- Students who only share one sticky note - encourage them to try the second stem too.""",

    31: """SAY:
- Show me what making connections to a text really means.
- A - say what we like about the text. B - think about something from our own lives that is similar.
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- What does making connections mean? 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move to the share back.
PIVOT:
- Most likely misconception: students pick A because they have been giving opinions about the text.
- Reteach: making connections is about linking our OWN life to the story. Saying we like it is just an opinion.
- Fresh re-check: When I say this reminds me of... what am I doing? A - giving an opinion. B - making a connection.

WATCH FOR:
- Students who say both - push for the BEST answer based on the words on the slide.""",

    32: """SAY:
- Time to share your sticky notes with a partner.
- Read what you wrote.
- Then explain why you chose that scene.

DO:
- Pair students up.
- Give 2 minutes for sharing.
- Listen in on 2-3 pairs.
- Bring the class back and share 1-2 strong examples.

TEACHER NOTES:
Closing share-back. Use a strong example to model how a personal connection sounds when said aloud.

WATCH FOR:
- Students who freeze when sharing - allow them to read directly from the sticky note.
- Pairs going off-topic - re-cue with the sentence stems.""",
}


# =====================================================================
# LESSON 6
# =====================================================================
LESSON_6_NOTES = {
    10: """SOURCES:
Teacher Handbook, Under the Southern Cross (2En01V) - Lesson 6 read-aloud notes.

SAY:
- Today we are reading pages 29 and 30 of Under the Southern Cross.
- These pages look different to the rest of the book - they are all factual information about the Southern Cross.
- Ask: What important information does the author give us about the Southern Cross?

DO:
- Show pages 29-30 so all students can see.
- Page 29, pause after "naked eye." Ask: How many main stars make up the Southern Cross, and what shape do they form?
- Page 29, pause after "surface curves." Ask: Where can you see the Southern Cross all year round?
- Page 29, pause after "or Acrux." Ask: What alphabet are the five main stars of the Southern Cross named after?
- Page 30, pause after "true south." Ask: How could you use the Southern Cross if you ever became lost?
- Page 30, pause after "democratic protest." Ask: Can you think of a reason why the Southern Cross was used in different flags?
- Page 30, pause after "years old!" Ask: Which is older - the Earth or the Southern Cross constellation?
- Page 30, pause after "upside down." Ask: What do we now understand about why the position of the Southern Cross changes on each page of the book?

TEACHER NOTES:
Pause points come from the Teacher Handbook. The factual layer is the focus today - read slowly so students can hold the facts in mind.

WATCH FOR:
- Students who treat this like a story - remind them this is the informative layer.
- Students overwhelmed by facts - bring them back to one fact at a time after each pause.""",

    12: """SAY:
- A noun is a naming word.
- It names a person, place, animal or thing.
- We also call this a common noun.
- A common noun does not need a capital letter.
- Say these examples with me... penguin... cinema... turtles.

DO:
- Point to each example on the slide.
- Echo read the word list together.

TEACHER NOTES:
This is a quick revision from Lesson 4. Keep it brisk - students will move into proper nouns next.

WATCH FOR:
- Students who try to add a capital to penguin - remind them common nouns do not need capitals.""",

    13: """SAY:
- A proper noun names specific people, places, objects, animals or the title of something.
- A proper noun always begins with a capital letter.
- Say these examples with me... Southern Cross... Phillip Island... Milky Way.

DO:
- Point to each proper noun on the slide.
- Show the capital letter at the start of each one.

TEACHER NOTES:
Quick revision from Lesson 4. Students will use this rule when expanding sentences today.

WATCH FOR:
- Students who skip the capital letter - point to it every time on a proper noun.""",

    14: """SAY:
- Let's look at examples of common and proper nouns from our book.
- Common nouns: crocodile, bridge, beach, telescope, outback, moon. No capital letter needed.
- Proper nouns: Broome, South Pole, Swan River, The Last Post, Southern Hemisphere, Melbourne Cricket Ground. Capital letter at the start of each main word.

DO:
- Read one column, then the other.
- Point to a capital letter on each proper noun.

TEACHER NOTES:
These are the proper nouns students will hear in the sentences they expand today. Read them clearly so they recognise them later.

WATCH FOR:
- Students who think Melbourne Cricket Ground only needs one capital - point to each main word.""",

    15: """SAY:
- In Lesson 4, we expanded sentences by adding a who or what.
- Watch: They are watching a movie at night.
- We swapped the pronoun they for the families. Now we know who.
- Full sentence: The families are watching a movie at night.

DO:
- Read the original sentence.
- Point to They, then to the families.
- Read the expanded sentence with expression.

TEACHER NOTES:
This recaps Lesson 4. Use it as a launch pad for adding when and where details next.

WATCH FOR:
- Students who do not remember the move - re-anchor with: who or what is doing the action?""",

    16: """SAY:
- Today we will add even more detail by saying when and where something happens.
- Watch: At night, the families are watching a movie at the Deckchair Cinema.
- At night - that is the when. The Deckchair Cinema - that is the where.
- We often write the when detail first in an expanded sentence.

DO:
- Point to At night and name it as the when.
- Point to The Deckchair Cinema and name it as the where.
- Read the new full sentence with expression.

TEACHER NOTES:
This is the I Do for adding when and where. The when first rule is the structure students will copy.

WATCH FOR:
- Students who put when at the end - gently remind them: when usually comes first in writing.""",

    17: """SAY:
- Prepositions are little words that usually start a phrase.
- These phrases tell us where or when something happens.
- Examples include in, on, at, before, after, during.

DO:
- Read the definition aloud.
- Echo three preposition examples together.

TEACHER NOTES:
Keep this short. Students do not need a full preposition list - they need to recognise preposition phrases for when and where.

WATCH FOR:
- Students who get stuck on the term - tell them prepositions are signal words for when and where.""",

    18: """SAY:
- Here are some when words and phrases.
- Before lunch... after dinner... during the night... at sunset... in the morning... on Monday... one cold day.
- Echo each phrase after me.

DO:
- Read each phrase clearly.
- Have students echo each one.
- Point to the preposition word at the start.

TEACHER NOTES:
Build a bank of when phrases students can use later. Repetition helps the phrases stick.

WATCH FOR:
- Students who only say the preposition (e.g. before) without the noun - cue them to say the whole phrase.""",

    19: """SAY:
- Here are some where words and phrases.
- Through the grass... beside the shed... around the block... between the trees... near the farm... above the horizon... below the water.
- Echo each phrase after me.

DO:
- Read each phrase clearly.
- Have students echo each one.
- Point to the preposition word at the start.

TEACHER NOTES:
Builds a bank of where phrases. These mirror the kind of phrases used in the text.

WATCH FOR:
- Students who confuse where and when phrases - re-cue: ask, does this tell me a place or a time?""",

    20: """SAY:
- Read pages 17 and 18 with me - the Staircase to the Moon scene.
- Here is a kernel sentence: It rises.
- Who or what rises? The full moon.
- When does it rise? After sunset. Where does it rise? Over the horizon.
- Full sentence: After sunset, the full moon rises over the horizon.

DO:
- Read the kernel sentence.
- Point to It, then swap it for the full moon.
- Point to after sunset (when) and over the horizon (where).
- Read the expanded sentence with expression.

TEACHER NOTES:
This is the I Do for when and where expansion. Show the full thinking - who, when, where, then full sentence with when first.

WATCH FOR:
- Students who only add the where - prompt them to add the when as well.""",

    21: """SAY:
- Read pages 15 and 16 with me - the fireworks scene.
- Kernel sentence: They exploded in bright colours.
- Who or what exploded? The fireworks.
- When? On New Year's Eve. Where? Over the sky.
- Full sentence: On New Year's Eve, the fireworks exploded in bright colours over the sky.

DO:
- Read the kernel sentence.
- Partner talk - who, when, where?
- Take a response and reveal.

TEACHER NOTES:
First We Do. Pace will pick up over the next slides. Expect minor mistakes - keep redirecting to who, when, where.

WATCH FOR:
- Students who skip the when - point to the when phrase on the slide and re-cue.""",

    22: """SAY:
- Read pages 11 and 12 with me - the Brisbane Ferris wheel scene.
- Kernel sentence: It spins.
- Who or what spins? The Ferris wheel.
- When? In the evening. Where? Above the trees.
- Full sentence: In the evening, the Ferris wheel spins above the trees.

DO:
- Partner talk - who, when, where?
- Take a response.
- Reveal the expanded sentence and read aloud.

TEACHER NOTES:
Second We Do. Students should now be moving faster. Celebrate when they say the when first.

WATCH FOR:
- Students who put when at the end - prompt: try moving when to the start.""",

    23: """SAY:
- Read pages 19 and 20 with me - the aurora australis scene near Hobart.
- Kernel sentence: They watched the aurora australis.
- Who or what watched? The sightseers.
- When? On a dark night. Where? Near Hobart.
- Full sentence: On a dark night, the sightseers watched the aurora australis near Hobart.

DO:
- Partner talk - who, when, where?
- Take a response.
- Reveal and read aloud.

TEACHER NOTES:
Third We Do. Hobart and aurora australis link to the proper noun rule too - point out the capitals while reading.

WATCH FOR:
- Students who write hobart without a capital - re-anchor on the proper noun rule.""",

    24: """SAY:
- Read pages 25 and 26 with me - the campfire scene near Alice Springs.
- Kernel sentence: It crackles.
- Who or what crackles? The camp fire.
- When? At night. Where? In Alice Springs.
- Full sentence: At night, the camp fire crackles in Alice Springs.

DO:
- Partner talk - who, when, where?
- Take a response.
- Reveal and read aloud.

TEACHER NOTES:
Fourth We Do. By now, students should be fluent with the move. Use this one for confidence.

WATCH FOR:
- Students who add a long where phrase - keep it short and clear.""",

    25: """SAY:
- Read this sentence with me: During the night, the loggerhead turtles hatch at Mon Repos.
- Show me the when detail.
- Hold up A for the loggerhead turtles, B for during the night, C for at Mon Repos.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Letter (A, B or C)
Script:
- Which is the when detail? 3, 2, 1, show me.
- Scan for: most boards show B (during the night).
PROCEED:
- If most students show B, move on to sentence order.
PIVOT:
- Most likely misconception: students pick C (at Mon Repos) because it sounds like a time.
- Reteach: when tells us about time. Mon Repos is a place, not a time.
- Fresh re-check: at sunset - is that a when or a where?

WATCH FOR:
- Students mixing up where and when prepositions - re-cue with the wall banks of phrases.""",

    26: """SAY:
- We have practised putting the when detail at the beginning of our sentences.
- At night... On New Year's Eve... After sunset.
- But in Under the Southern Cross, Lessac wants to draw our attention to the place names. So she begins with a where detail.
- On Phillip Island... In Darwin... Near Boulia.

DO:
- Read the when starters.
- Read the where starters.
- Hold up the book and point to a page where the where comes first.

TEACHER NOTES:
This is the author's craft move - she breaks the when first rule on purpose to highlight place. Make this difference visible.

WATCH FOR:
- Students confused that we changed the rule - clarify: in our writing we usually start with when, but Lessac chose to start with where for a reason.""",

    27: """SAY:
- Read these sentences with me. Notice how the order changes the rhythm and emphasis.
- In Sydney, families watch the fireworks at midnight.
- At midnight, families watch the fireworks in Sydney.
- Both sentences mean the same thing, but the first one emphasises the place. The second emphasises the time.

DO:
- Read each pair clearly.
- After each pair, ask: which one feels like Lessac's style?
- Take a quick partner whisper for one pair.

TEACHER NOTES:
Students notice how moving the when or where changes what we focus on. Keep the discussion practical, not abstract.

WATCH FOR:
- Students who think the meaning has changed - clarify: same facts, different emphasis.""",

    28: """SAY:
- Read this sentence: In Canberra, the crowd silently stands when The Last Post is played.
- Now turn and tell your partner how you would change the order to put the when detail first.

DO:
- Give 30 seconds of partner talk.
- Cold call 1-2 students.
- Reveal: When The Last Post is played, the crowd silently stands in Canberra.

TEACHER NOTES:
Students experiment with reordering. Both versions are correct, but the second matches the when first rule for their own writing.

WATCH FOR:
- Students who only swap the words at the start - check that the rest of the sentence still reads correctly.""",

    29: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 6 - Sentence level writing.
- Add a who or what, a when and a where to expand each sentence.
- Remember to start with when first in your own writing, unless you are copying Lessac's style.

DO:
- Hand out booklets or direct students to the page.
- Circulate and check the first 2-3 answers for each student.
- Set up a share back at the end.

TEACHER NOTES:
This is the You Do for sentence expansion. Pull a small group if students still confuse when and where after the CfU.

WATCH FOR:
- Students who add only one detail - prompt them for both when and where.
- Students finishing fast - extend with: try moving the when to the end.""",

    30: """SAY:
- Time to share your expanded sentences.
- Pick one to read aloud to the class.

DO:
- Call on 3-4 students to read their best sentence.
- Affirm correct who/what, when and where as students read.

TEACHER NOTES:
Closing share-back. Use it to celebrate strong sentences and gently correct any leftover errors before moving on.

WATCH FOR:
- Students who freeze sharing - let them read straight from the booklet.""",
}


# =====================================================================
# LESSON 7
# =====================================================================
LESSON_7_NOTES = {
    10: """SOURCES:
Teacher Handbook, Under the Southern Cross (2En01V) - Lesson 7 read-aloud notes.

SAY:
- Today we are reading a new text called The Southern Cross - A Bright Star Story.
- This text will give us more information about the Southern Cross.
- Ask: Why is the Southern Cross an important part of Australia's sky?
- Some of you may remember things we learnt yesterday from pages 29 and 30. We will build on that today.

DO:
- Hold or display the new text.
- Point to the title.
- Take 2-3 quick predictions about what the text will be about.

TEACHER NOTES:
This is the launch slide for a new shared read aloud. Keep predictions short - the reading does the heavy lifting.

WATCH FOR:
- Students who do not have the text in front of them - check everyone can see before reading.""",

    11: """SAY:
- Listen as I read.
- Look up at the night sky and you will see many twinkling stars. One special group is the Southern Cross. It looks like a kite or cross and is very famous in Australia because we can see it so clearly.
- What did this part of the text tell us?

DO:
- Read the page slowly with expression.
- Pause and ask: What shape does the Southern Cross look like?
- Take 1-2 responses.

TEACHER NOTES:
First chunk of the new text. Keep the comprehension simple - shape and visibility are the key facts here.

WATCH FOR:
- Students drifting - re-cue with the kite or cross gesture using their hands.""",

    12: """SAY:
- Listen carefully.
- For thousands of years, Aboriginal and Torres Strait Islander peoples have looked at the Southern Cross and found stories and signs.
- They saw animals and patterns in the stars that helped them understand the world.
- Pause after... the land, and family. Ask: What did Aboriginal and Torres Strait Islander people use the stars for?

DO:
- Read with care and respect.
- Pause at the handbook pause point.
- Take 2-3 responses about how the stars were used.

TEACHER NOTES:
Treat First Nations content with respect. The stars taught lessons about life, the land and family. Avoid summarising Dreaming stories beyond what the text gives.

SENSITIVITY ADVISORY:
- What it is: First Nations cultural knowledge about the Southern Cross.
- Framing language: For thousands of years, Aboriginal and Torres Strait Islander peoples have looked at the night sky and shared stories from this place. Use respectful language.
- Watch for: Students treating Dreaming stories as fairy tales. Gently re-frame as stories that teach important lessons.
- Protocol: If a student shares family or community knowledge, listen, thank them, and do not press for detail.

WATCH FOR:
- Students using past tense as if these stories no longer exist - clarify these stories continue to be shared today.""",

    13: """SAY:
- Listen.
- The Southern Cross is part of the Milky Way, our galaxy.
- Its stars are very far from us, and even further from each other.
- Learning about the Southern Cross shows us how vast and wonderful space is.
- Pause after... and wonderful space is. Ask: How does the Southern Cross fit into the Milky Way?

DO:
- Read the page slowly.
- Stretch arms wide on the word vast.
- Take 1-2 responses.

TEACHER NOTES:
The word vast is taught later in the lesson. Use the gesture now so students hear it in context.

WATCH FOR:
- Students confused that the Southern Cross is in the Milky Way - clarify the Milky Way is our whole galaxy and the Southern Cross is one part of it.""",

    14: """SAY:
- Listen.
- To see the Southern Cross clearly, you need a dark sky far away from city lights.
- Australia has special areas called Dark Sky Parks or Dark Sky Sanctuaries. These places are protected so that there is no light pollution, making the skies very clear and dark.
- Pause after... amazing places for stargazing. Ask: Why is it helpful to have no light pollution when you are looking at the stars?

DO:
- Read the page slowly.
- Take 1-2 responses.
- Connect to last night's experience - has anyone seen lots of stars when away from town?

TEACHER NOTES:
Last chunk of the new text. After the read, return to the guiding question - why is the Southern Cross important to Australia?

WATCH FOR:
- Students who do not know the term light pollution - tie it to too much light from streetlights and houses making it hard to see stars.""",

    16: """SOURCES:
Teacher Handbook, Under the Southern Cross (2En01V) - Lesson 7 vocabulary.

SAY:
- In the text, the stars guided people as they travelled.
- If you guide someone in a particular direction, you lead them in that direction.
- Say the word with me... guiding!
- A teacher might guide you through new work. A map might guide you on a trip.

DO:
- Point to the image on the slide.
- Mime leading someone forward as you say guide.
- Have students repeat the word twice.

TEACHER NOTES:
Connect guide to leading. Students will use this word in writing later, so anchor it clearly.

WATCH FOR:
- Students mixing up guide with follow - clarify: a guide leads, not follows.""",

    17: """SAY:
- The stars guided people as they travelled. This means...
- A - they were too hard to see. B - there were so many stars in the sky. C - the stars gave instructions. D - the stars helped people travel in the right direction.
- Show me your answer.

DO:
- Read each option clearly.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Letter (A, B, C or D)
Script:
- What does it mean that the stars guided people? 3, 2, 1, show me.
- Scan for: most boards show D.
PROCEED:
- If most students show D, move on to the next practice slide.
PIVOT:
- Most likely misconception: students pick C because instructions sounds like leading.
- Reteach: guiding is leading someone in a direction. Stars do not give spoken instructions, but they help people find the right way.
- Fresh re-check: a sign on a path - does it guide you?

WATCH FOR:
- Students who pick A or B - they have not connected the meaning to direction.""",

    18: """SAY:
- Look at the two pictures.
- Which is an example of guiding?
- A - a father showing his child how to walk to the beach safely. B - a young girl feeding a goat.
- Show me A or B.

DO:
- Read both options.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

TEACHER NOTES:
Expected answer: A. Feeding a goat is helpful but not guiding. Guiding is leading someone where to go.

WATCH FOR:
- Students who pick B because it looks kind - clarify kind is not the same as guiding.""",

    19: """SAY:
- I will say two situations. Show me thumbs up if someone would be guiding you. Thumbs down if not.
- You eat your lunch.
- You learn to ride a bike.

DO:
- Pause between each one for thumbs.
- Scan the room.

TEACHER NOTES:
Expected answers: thumbs down for eating lunch (no one needs to lead you), thumbs up for learning to ride a bike (someone is showing you what to do).

WATCH FOR:
- Students who put thumbs up for both - re-cue: is someone leading you in a direction or showing you how to do something?""",

    20: """SAY:
- Show me which words are similar in meaning to guiding.
- Hands up for helping... following... showing... confusing.

DO:
- Pause between each one for hands.
- Scan after each.

CFU CHECKPOINT:
Technique: Hands Up for each option
Script:
- Hands up if this word is similar to guiding.
- Scan for: hands up on helping and showing. Hands down on following and confusing.
PROCEED:
- If most students get all four right, move to the next CfU.
PIVOT:
- Most likely misconception: students put hands up on following because they hear the word lead.
- Reteach: guiding is leading. Following is the opposite - going behind.
- Fresh re-check: When a teacher leads the line, are they guiding or following?

WATCH FOR:
- Students mirroring a friend - keep wait time firm before each show.""",

    21: """SAY:
- When someone is guiding you, they...
- A - show you the way. B - give you a job to do. C - tell you something kind.
- Show me your answer.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Letter (A, B or C)
Script:
- Which best completes the sentence? 3, 2, 1, show me.
- Scan for: most boards show A.
PROCEED:
- If most students show A, move on to the next vocabulary word.
PIVOT:
- Most likely misconception: students pick C because guiding can feel kind.
- Reteach: kindness is not guiding. Guiding is showing the way.
- Fresh re-check: a tour guide at the zoo - what do they do? Show the way.

WATCH FOR:
- Students mirroring fast hands - keep wait time before they show.""",

    22: """SOURCES:
Teacher Handbook, Under the Southern Cross (2En01V) - Lesson 7 vocabulary.

SAY:
- In the text, the Southern Cross shows us how vast space is.
- Something that is vast is extremely large.
- Say the word with me... vast!
- The ocean is vast. The desert is vast.

DO:
- Stretch arms wide as you say vast.
- Have students repeat the word twice with arms wide.

TEACHER NOTES:
Vast and gigantic are similar but vast is more about open spaces. Use the wide gesture to show breadth, not just size.

WATCH FOR:
- Students using vast for everyday objects - remind them it is for very large spaces.""",

    23: """SAY:
- The Southern Cross shows us how vast space is. This means that space is...
- A - really exciting. B - very big.
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- What does vast mean? 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move on.
PIVOT:
- Most likely misconception: students pick A because space feels exciting.
- Reteach: vast is about size, not feeling.
- Fresh re-check: a tiny puddle - vast or not vast?

WATCH FOR:
- Students who pick A because they read excitement into the sentence - re-cue on size.""",

    24: """SAY:
- Show me which word is the opposite of vast.
- A for small. B for dark.
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

TEACHER NOTES:
Expected answer: A (small). Dark is unrelated. This builds on the size meaning.

WATCH FOR:
- Students who pick B because dark feels like the night sky - re-cue: opposite of vast means opposite of very large.""",

    25: """SAY:
- Look at the two images.
- Which one shows a vast area?
- Point to your choice.

DO:
- Give 5 seconds of looking time.
- Say: 1, 2, 3, point.
- Cold call one student to explain why.

TEACHER NOTES:
Students pick the image that shows extremely large open space. Accept either if they can explain vast in their reason.

WATCH FOR:
- Students pointing without a reason - prompt: how can you tell it is vast?""",

    26: """SAY:
- Show me which examples could be described as vast.
- Hands up for an ocean... a library with a large range of books... a puddle... the Prime Minister's knowledge of Australia.

DO:
- Pause between each option for hands.
- Scan after each one.

TEACHER NOTES:
Expected hands up: an ocean, a library with a large range of books, the Prime Minister's knowledge. Vast can also describe knowledge or amount, not only physical size. A puddle is small.

WATCH FOR:
- Students who put hands up for puddle - clarify size again.
- Students confused by knowledge being vast - explain we can have a vast amount of knowledge, like a vast library.""",

    27: """SAY:
- Show me which words are similar to vast.
- Hands up for huge... wide... small... tiny.

DO:
- Pause between each option for hands.
- Scan after each.

CFU CHECKPOINT:
Technique: Hands Up for each option
Script:
- Hands up if this word is similar to vast.
- Scan for: hands up on huge and wide. Hands down on small and tiny.
PROCEED:
- If most students get all four right, move to the writing section.
PIVOT:
- Most likely misconception: students hesitate on wide because vast feels different to wide.
- Reteach: vast can mean very wide AND very large. Wide is part of vast.
- Fresh re-check: a vast desert - is it wide?

WATCH FOR:
- Students who mirror their neighbour - keep wait time firm.""",

    29: """SAY:
- A noun is a naming word.
- It names a person, place, animal or thing.
- We also call this a common noun.
- A common noun does not need a capital letter.

DO:
- Echo the word noun together.
- Point to each example on the slide.

TEACHER NOTES:
Quick revision before the new content on adjectives, articles and noun groups.

WATCH FOR:
- Students who try to add capitals to penguin - remind them common nouns do not.""",

    30: """SAY:
- A proper noun names specific people, places, objects, animals or the title of something.
- A proper noun begins with a capital letter.
- Say the examples with me... Southern Cross... Phillip Island... Milky Way.

DO:
- Point to each capital letter.
- Echo the words.

TEACHER NOTES:
Quick revision. Proper nouns will appear in the noun group practice today.

WATCH FOR:
- Students who skip the capital - re-cue every time.""",

    31: """SAY:
- An adjective is a word that describes a noun.
- It tells what a noun is like.
- Adjectives and verbs work together in sentences to give a clear picture.
- For example - the bright stars. Bright is the adjective. Stars is the noun.

DO:
- Read the definition aloud.
- Point to bright as the adjective and stars as the noun.

TEACHER NOTES:
This is the I Do for adjectives. Keep the example simple - bright + stars - because students will sort adjectives next.

WATCH FOR:
- Students who name the verb as the adjective - clarify: adjectives describe, verbs do.""",

    32: """SAY:
- Read these sentences with me.
- The tasty prawns are flipping in the full nets in Perth.
- The bright screen is showing the funny movie in Darwin.
- Identify the nouns and the adjectives in each sentence.
- Turn and tell your partner.

DO:
- Read each sentence twice.
- Give 30 seconds of partner talk.
- Take responses and reveal: prawns, nets, screen, movie are nouns; tasty, full, bright, funny are adjectives.

TEACHER NOTES:
Two sentences with two adjective-noun pairs each. Use a colour code if you can - one colour for nouns, one for adjectives.

WATCH FOR:
- Students who underline Perth or Darwin - those are proper nouns, also acceptable as nouns.
- Students who name verbs as adjectives - re-anchor on describes.""",

    33: """SAY:
- An article is a little word that goes in front of a noun.
- It tells us if we are talking about one or more of the noun, or if it is something special.

DO:
- Read the definition aloud.
- Echo the word article together.

TEACHER NOTES:
Definition only. The three article types are coming up next.

WATCH FOR:
- Students who think articles are sentences - clarify: articles are tiny words like a, an, the.""",

    34: """SAY:
- There are three article types - the, a and an.
- We use the for a specific or special noun already known. The moon - we know which one.
- We use a before a noun that begins with a consonant sound. A star.
- We use an before a noun that begins with a vowel sound. An island.

DO:
- Echo each example.
- Have students mouth a star and an island so they hear the sound difference.

TEACHER NOTES:
The vowel sound rule is by sound, not letter. An hour starts with a vowel sound even though the word begins with h. Keep it simple today.

WATCH FOR:
- Students who use a where they need an - cue: does the next word start with a vowel sound?""",

    35: """SAY:
- A noun group is a group of words built around a noun.
- It can have an article, an adjective, and a noun together.
- Watch how the noun group grows.
- Just a noun: stars.
- Article and noun: the stars.
- Article, adjective and noun: the ancient stars.

DO:
- Point to each level of the build.
- Echo each version together.

TEACHER NOTES:
This is the I Do for noun groups. The build grows from one word to three. Students will copy this move next.

WATCH FOR:
- Students who skip the article - point back to the article column.""",

    36: """SAY:
- Read this sentence with me: An eerie light scares travellers near Boulia.
- Find the noun group. Then highlight the article, the adjective and the noun.
- Turn and tell your partner.

DO:
- Give 30 seconds of partner talk.
- Take responses.
- Reveal: noun group is an eerie light. Article: an. Adjective: eerie. Noun: light.

TEACHER NOTES:
First We Do for finding noun groups. Vowel sound rule applies - an before eerie because eerie starts with a vowel sound.

WATCH FOR:
- Students who include scares in the noun group - clarify: scares is a verb, not part of the noun group.""",

    37: """SAY:
- Read with me: The energetic players competed at the MCG.
- Find the noun group. Then highlight the article, the adjective and the noun.

DO:
- Partner talk for 30 seconds.
- Take responses.
- Reveal: noun group is the energetic players. Article: the. Adjective: energetic. Noun: players.

TEACHER NOTES:
Second We Do. Pace should pick up here. Use this slide to celebrate quick noun group identification.

WATCH FOR:
- Students who add MCG to the noun group - clarify: MCG is a separate proper noun, not part of this noun group.""",

    38: """SAY:
- Read with me: The explosive fireworks boomed over the Sydney sky.
- Find the noun group. Highlight the article, adjective and noun.

DO:
- Partner talk for 30 seconds.
- Take responses.
- Reveal: noun group is the explosive fireworks. Article: the. Adjective: explosive. Noun: fireworks.

TEACHER NOTES:
Third We Do. Sydney sky is a separate noun phrase - not the focus today.

WATCH FOR:
- Students who pick Sydney sky - re-anchor on the noun group built around fireworks.""",

    39: """SAY:
- Read with me: At Mon Repos, you can see the tiny turtles swimming.
- Find the noun group. Highlight the article, adjective and noun.

DO:
- Partner talk for 30 seconds.
- Take responses.
- Reveal: noun group is the tiny turtles. Article: the. Adjective: tiny. Noun: turtles.

TEACHER NOTES:
Fourth We Do. Mon Repos is a proper noun setting, not part of the noun group.

WATCH FOR:
- Students who include swimming - clarify: swimming is a verb, not part of the noun group.""",

    40: """SAY:
- Read with me: A dangerous crocodile could be seen in the Daintree River.
- Identify the article, the adjective and the noun in the noun group.
- Write your answers on your whiteboard.

DO:
- Give 30 seconds of writing time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Whiteboard write, then show
Script:
- Article, adjective and noun. 3, 2, 1, show me.
- Scan for: article a, adjective dangerous, noun crocodile.
PROCEED:
- If most students have all three correct, move to creating noun groups.
PIVOT:
- Most likely misconception: students pick the as the article because it is more common.
- Reteach: a is used before a consonant sound. Crocodile starts with a c sound, so we say a crocodile.
- Fresh re-check: an apple or a apple - which one is right and why?

WATCH FOR:
- Students who write Daintree River as the noun - re-anchor on the noun group built around crocodile.""",

    41: """SAY:
- Now we are going to build our own noun groups.
- The shape is - article, plus adjective, plus noun.
- For example, crowd becomes the respectful crowd. Crocodile becomes a large crocodile. Night becomes an amazing night.

DO:
- Echo each noun group together.
- Point to each part of the build (article, adjective, noun).

TEACHER NOTES:
This is the I Do for building noun groups. Students will copy the move on the next slides.

WATCH FOR:
- Students who skip the article - prompt them with: which little word goes before?""",

    42: """SAY:
- Read pages 11 and 12 - the Brisbane Ferris wheel scene.
- Turn the noun Ferris wheel into a noun group.
- For example - the large Ferris wheel.
- Article: the. Adjective: large. Noun: Ferris wheel.

DO:
- Read the example aloud.
- Partner talk - what other adjectives could we use?
- Take 1-2 responses.

TEACHER NOTES:
First We Do for noun group creation. Accept any adjective that fits a Ferris wheel.

WATCH FOR:
- Students who use a instead of the - both can work, but the matches the example.""",

    43: """SAY:
- Read pages 26 and 27 - the War Memorial scene.
- Turn the noun War Memorial into a noun group.
- For example - the respectful War Memorial.
- Article: the. Adjective: respectful. Noun: War Memorial.

DO:
- Read the example aloud.
- Partner talk - what other adjectives could we use?
- Take 1-2 responses.

TEACHER NOTES:
War Memorial is a proper noun, so the capital W and capital M stay even when used in a noun group. Respectful is a strong adjective for this scene.

SENSITIVITY ADVISORY:
- What it is: ANZAC and military commemoration.
- Framing language: Use respectful adjectives - solemn, respectful, quiet, brave.
- Watch for: Students using silly or playful adjectives.
- Protocol: Redirect gently - this is a place where people remember soldiers, so we use respectful words.

WATCH FOR:
- Students who use playful adjectives - prompt them to choose words that match the feeling of the place.""",

    44: """SAY:
- Read pages 17 and 18 - the Staircase to the Moon scene.
- Turn the noun moon into a noun group.
- For example - the glowing moon.
- Article: the. Adjective: glowing. Noun: moon.

DO:
- Read the example aloud.
- Partner talk - what other adjectives could we use?
- Take 1-2 responses.

TEACHER NOTES:
Third We Do. By now students should be confident with the move - celebrate creative adjectives.

WATCH FOR:
- Students who pick weak adjectives like nice - push for stronger sensory words like glowing, bright, silver.""",

    45: """SAY:
- Now you write a noun group for the noun light.
- Article, plus adjective, plus noun.
- Examples - the bright light, an illuminated light, a shining light.
- Write your noun group on your whiteboard.

DO:
- Give 30 seconds of writing time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Whiteboard write, then show
Script:
- Write a noun group for light. 3, 2, 1, show me.
- Scan for: article + adjective + noun. Examples: the bright light, a shining light, an illuminated light.
PROCEED:
- If most students have all three parts, move to the booklet task.
PIVOT:
- Most likely misconception: students write only adjective and noun (e.g. bright light) without an article.
- Reteach: every noun group needs the little article word at the front.
- Fresh re-check: write a noun group for star.

WATCH FOR:
- Students using a when the adjective starts with a vowel sound (e.g. a illuminated light) - prompt them to swap to an.""",

    46: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 7 - Sentence-level writing.
- Build noun groups for each noun.
- Remember - article, adjective and noun.

DO:
- Hand out booklets or direct students to the page.
- Circulate and check the first 2-3 answers for each student.
- Set up a share-back at the end.

TEACHER NOTES:
This is the You Do. Pull a small group if students still skip the article.

WATCH FOR:
- Students who use the same adjective for every noun - prompt for variety.
- Students finishing fast - extend with two adjectives in one noun group (the bright twinkling stars).""",

    47: """SAY:
- Time to share your noun groups with the class.
- Read your favourite one.

DO:
- Call on 3-4 students.
- Affirm the article, adjective and noun in each example as students share.

TEACHER NOTES:
Closing share-back. Use it to celebrate strong adjectives and confirm the article + adjective + noun pattern.

WATCH FOR:
- Students who pick a weak adjective - praise the structure but push for a stronger word next time.""",
}


# =====================================================================
# LESSON 8
# =====================================================================
LESSON_8_NOTES = {
    10: """SAY:
- Let's revisit the personal responses you wrote on sticky notes in Lesson 5.
- Some of you wrote: I'd like to go here because... and This reminds me of...
- Today we will use those connections to build our own class version of the book.

DO:
- Display or hand out the sticky notes from Lesson 5.
- Take 2-3 quick share-backs of favourite responses.

TEACHER NOTES:
This warms up the writing task. Use student voices to set the tone for the class book.

WATCH FOR:
- Students who do not have their notes - pair them with someone who does.""",

    11: """SAY:
- Let's brainstorm the night-time activities we saw in Under the Southern Cross.
- Watching the penguins on Phillip Island... visiting an outdoor cinema in Darwin... seeing the Min Min lights near Boulia... watching an AFL game at the MCG... watching the turtles at Mon Repos... riding the Ferris wheel in Brisbane... stargazing near Adelaide... watching the fireworks in Sydney.
- Which ones do you remember most?

DO:
- Read each activity aloud.
- Have students put a hand up for any they remember from the book.
- Take 1-2 quick comments.

TEACHER NOTES:
This is recall, not new content. Speed it up - the goal is to refresh the bank of activities for writing.

WATCH FOR:
- Students who want to share a story for every activity - keep them brief and move on.""",

    12: """SAY:
- More night-time activities from the book.
- Seeing The Staircase to the Moon near Broome... watching the aurora australis near Hobart... going prawn fishing in Perth... spotting crocodiles on the Daintree River... sharing stories around the campfire near Alice Springs... attending a dawn service in Canberra.
- Which ones surprise you?

DO:
- Read each activity aloud.
- Take 1-2 surprised reactions.

TEACHER NOTES:
Continuation of the brainstorm. Keep it brisk so students stay focused on the writing to come.

WATCH FOR:
- Students confused about a place name (Broome, Mon Repos) - quickly point to it on a map or in the book.""",

    13: """SAY:
- What activities do YOU do at night?
- Cook marshmallows... go for a night walk... attend a night market... watch a movie in the backyard... go to a restaurant... camp... star gaze... look at Christmas lights.
- Have you done any of these? Turn and tell your partner one of your night-time activities.

DO:
- Give 30 seconds of partner talk.
- Take 3-4 responses from the room.
- Add any new activity ideas to the board.

TEACHER NOTES:
Personal connection sets up the writing task. Keep the list visible so students can pick from it later.

WATCH FOR:
- Students who say they do not do anything at night - prompt with simple options like reading, eating dinner, watching a show.""",

    14: """SAY:
- Today we will create our class version of Under the Southern Cross.
- Watch my model: In Newcastle, my family pitches a tent in the backyard and roasts marshmallows on the fire - under the Southern Cross.
- Notice the shape: In place, my family does activity - under the Southern Cross.

DO:
- Read the model with expression.
- Point to each part: place name, my family, activity, under the Southern Cross.
- Show the picture or sketch.

TEACHER NOTES:
This is the I Do for writing in the book's style. The structure is: where + who + what + tag line.

WATCH FOR:
- Students who skip the place name - re-anchor on Lessac's where first style.
- Students who skip under the Southern Cross - that line is the tag that ties every page together.""",

    15: """SAY:
- Now we will complete a sentence together.
- In _____, my family _____ - under the Southern Cross.
- Let's try: In Melbourne, my family watch a funny movie in our backyard - under the Southern Cross.
- Turn and tell your partner where YOU live and what your family does at night.

DO:
- Read the model.
- Give 30 seconds of partner talk.
- Take 2-3 responses and shape them into the sentence frame.
- Reveal the example.

TEACHER NOTES:
This is the We Do. Use the response routine from the slide - In Melbourne, my family... - to keep responses on structure.

WATCH FOR:
- Students who add too much detail and lose the structure - cue back to: place, family, activity.""",

    16: """SAY:
- What detail do we begin with when writing in the same style as the text?
- A for why. B for where.
- Show me your answer.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which detail comes first in Lessac's style? 3, 2, 1, show me.
- Scan for: most boards show B (where).
PROCEED:
- If most students show B, move to the booklet task.
PIVOT:
- Most likely misconception: students pick A (why) because reasons feel important.
- Reteach: Lessac wants us to notice the place names first, so where comes first.
- Fresh re-check: read one of Lessac's sentences aloud and ask: what comes first?

WATCH FOR:
- Students who confuse where with when - clarify: where is a place, when is a time.""",

    17: """SAY:
- Time to write your own page.
- Turn to the page titled: Lesson 8 - Understanding, responding to and creating text.
- Draw your night-time activity.
- Write a sentence in the style of the book - In place, my family does activity, under the Southern Cross.
- Add a small factual or descriptive detail in smaller text - just like Lessac does.

DO:
- Hand out booklets.
- Circulate and check first 2-3 students for the where-first structure.
- Time the task: 15 minutes for drawing and sentence, 5 for the smaller detail.

TEACHER NOTES:
This is the You Do. The class book becomes a real artefact - keep encouragement high.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use the sentence frame on the slide. In _____, my family _____ - under the Southern Cross. Pick one place from a word bank and one activity from the brainstorm.
- Extra Notes: Allow students to draw first and write second if they need.
EXTENDING PROMPT:
- Task: Write a second smaller-text fact about the place or the activity, in the style of Lessac's informative layer.
- Extra Notes: Encourage them to look at pages 1-2 for the model.

WATCH FOR:
- Students who freeze at the page name - prompt with their family activity from earlier.
- Students copying the model exactly - cue them to choose their own place and activity.""",

    18: """SAY:
- Use this place bank if you need help.
- Melbourne, Adelaide, Sydney, Canberra, Brisbane, Darwin, Hobart, Perth.
- Or write your own town if you live somewhere else.
- Remember the frame: In _____, my family _____ - under the Southern Cross.

DO:
- Point to each place name.
- Re-read the sentence frame.
- Circulate and confirm the frame is being used.

TEACHER NOTES:
This is a scaffold for the booklet task. Use it to support students who need a place name to start.

WATCH FOR:
- Students who pick a place but cannot link it to their family - prompt: have you ever visited there? Or pretend you live there.""",

    20: """SAY:
- A noun is a naming word.
- It names a person, place, animal or thing.
- We call this a common noun. It does not need a capital letter.

DO:
- Echo the word noun together.
- Point to each example.

TEACHER NOTES:
Quick revision before today's noun group expansion task.

WATCH FOR:
- Students who add capitals to penguin or cinema - re-anchor on common nouns.""",

    21: """SAY:
- A proper noun names specific people, places, objects, animals or the title of something.
- A proper noun begins with a capital letter.

DO:
- Point to each capital letter.
- Echo the words.

TEACHER NOTES:
Quick revision. Proper nouns will appear in today's noun groups (Phillip Island, Swan River).

WATCH FOR:
- Students who skip capital letters - re-cue.""",

    22: """SAY:
- An adjective is a word that describes a noun. It tells what a noun is like.
- Adjectives and verbs work together to give a clear picture.

DO:
- Read the definition aloud.
- Echo the word adjective.

TEACHER NOTES:
Quick revision from Lesson 7.

WATCH FOR:
- Students who confuse adjectives with verbs - clarify: adjectives describe, verbs do.""",

    23: """SAY:
- An article is a little word that goes in front of a noun.
- It tells us if we are talking about one or more of the noun, or if it is something special.

DO:
- Echo the word article.
- Read the definition.

TEACHER NOTES:
Quick revision from Lesson 7.

WATCH FOR:
- Students who skip the article - remind them noun groups need an article at the front.""",

    24: """SAY:
- There are three article types - the, a and an.
- The for a specific or special noun. The moon.
- A before a noun that begins with a consonant sound. A star.
- An before a noun that begins with a vowel sound. An island.

DO:
- Echo each example.
- Have students mouth a star and an island so they hear the difference.

TEACHER NOTES:
Quick revision. The vowel sound rule is by sound not letter.

WATCH FOR:
- Students who pick a where they need an - cue: vowel sound.""",

    25: """SAY:
- A noun group is a group of words built around a noun.
- It can have an article, an adjective, and a noun together.
- Watch how it grows. Stars... the stars... the ancient stars.

DO:
- Echo each level.
- Point to article, adjective, noun.

TEACHER NOTES:
This is the same definition from Lesson 7. Use it to launch today's expansion practice.

WATCH FOR:
- Students who skip the article - prompt: which little word goes first?""",

    26: """SAY:
- Turn the noun fireworks into a noun group.
- For example - the amazing fireworks.
- Article: the. Adjective: amazing. Noun: fireworks.
- Turn and tell your partner another adjective we could use.

DO:
- Read the example aloud.
- Give 30 seconds of partner talk.
- Take 1-2 responses.

TEACHER NOTES:
First We Do for noun group creation today. Accept any adjective that fits fireworks.

WATCH FOR:
- Students who use a instead of the - both work; the is the example here.""",

    27: """SAY:
- Turn the noun hatchling into a noun group.
- For example - a tiny hatchling.
- Article: a. Adjective: tiny. Noun: hatchling.
- Turn and tell your partner another adjective.

DO:
- Read the example aloud.
- Partner talk for 30 seconds.
- Take 1-2 responses.

TEACHER NOTES:
Second We Do for noun group creation. Hatchling starts with a consonant sound, so a not an.

WATCH FOR:
- Students who use the - that works too, but a fits the small or new feeling of a hatchling.""",

    28: """SAY:
- Now we use noun groups to expand sentences.
- Watch: They scramble.
- Who or what scramble? The desperate hatchlings.
- When do they scramble? At night. Where? Towards the beach.
- Full sentence: At night, the desperate hatchlings scramble towards the beach.
- Notice we put when first.

DO:
- Show pages 9 and 10.
- Point to They, then to the desperate hatchlings.
- Point to At night (when), towards the beach (where).
- Read the full sentence with expression.

TEACHER NOTES:
This is the I Do for sentence expansion with a noun group plus when and where. Show the full thinking.

WATCH FOR:
- Students who use only a noun, not a noun group - prompt for the article and adjective too.""",

    29: """SAY:
- Read pages 1 and 2 - the Phillip Island scene.
- Kernel sentence: They waddle.
- Who waddle? Use a noun group.
- Each night, the little penguins waddle up the beach on Phillip Island.

DO:
- Show pages 1 and 2.
- Partner talk - what is the noun group? When? Where?
- Take a response and reveal.

TEACHER NOTES:
First We Do today. The little penguins is the noun group. Each night is when. Up the beach on Phillip Island is where.

WATCH FOR:
- Students who use only penguins (no noun group) - prompt for an article and adjective.""",

    30: """SAY:
- Read pages 21 and 22 - the prawn fishing scene in Perth.
- Kernel sentence: He caught a prawn.
- Who caught the prawn? Use a noun group.
- In summer, a patient fisherman caught a prawn in the Swan River.

DO:
- Show pages 21 and 22.
- Partner talk - what is the noun group? When? Where?
- Take a response and reveal.

TEACHER NOTES:
Second We Do. A patient fisherman is the noun group. In summer is when. In the Swan River is where.

WATCH FOR:
- Students who write Swan river without a capital - re-anchor on the proper noun rule.""",

    31: """SAY:
- Read pages 11 and 12 - the Brisbane Ferris wheel.
- Kernel sentence: It spins.
- What spins? Use a noun group.
- After the sun goes down, the giant Ferris wheel spins up to the sky.

DO:
- Show pages 11 and 12.
- Partner talk - noun group? When? Where?
- Take a response and reveal.

TEACHER NOTES:
Third We Do. The giant Ferris wheel is the noun group. After the sun goes down is when. Up to the sky is where.

WATCH FOR:
- Students who write ferris wheel without capitals - clarify: Ferris is a proper noun in this title.""",

    32: """SAY:
- Read this sentence: In the evening, the excited families watched a movie at the Deckchair Cinema.
- Show me the noun group.
- A for watched a movie. B for In the evening. C for the excited families.
- Show me your answer.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Letter (A, B or C)
Script:
- Which is the noun group? 3, 2, 1, show me.
- Scan for: most boards show C.
PROCEED:
- If most students show C, move to the booklet task.
PIVOT:
- Most likely misconception: students pick B because In the evening sounds detailed.
- Reteach: a noun group is built around a noun (the thing). In the evening is a when phrase, not a noun group.
- Fresh re-check: a tiny hatchling - is that a noun group?

WATCH FOR:
- Students who pick A (watched a movie) - clarify: that is a verb phrase, not a noun group.""",

    33: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 8 - Sentence-level writing.
- Expand each sentence by adding a noun group, when and where.
- Remember to start with when first in your own writing.

DO:
- Hand out booklets or direct to the page.
- Circulate and check first 2-3 answers per student.
- Set up a share-back at the end.

TEACHER NOTES:
This is the You Do for sentence expansion with noun groups. Pull a small group if students still skip the article in their noun groups.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a word bank of articles, adjectives, nouns and prepositional phrases. Match one from each column to build the expanded sentence.
- Extra Notes: Encourage oral rehearsal before writing.
EXTENDING PROMPT:
- Task: Add an extra adjective to make an expanded noun group (e.g. the giant glowing Ferris wheel).
- Extra Notes: Look at the model on slide 28 again.

WATCH FOR:
- Students who only add when - prompt for both when and where.
- Students who skip the noun group's article - re-anchor on article + adjective + noun.""",

    34: """SAY:
- Time to share your expanded sentences.
- Read your favourite one to the class.

DO:
- Call on 3-4 students to share.
- Affirm the noun group, the when and the where as students read.

TEACHER NOTES:
Closing share-back. Celebrate strong noun groups and confirm the structure before ending the lesson.

WATCH FOR:
- Students who freeze sharing - let them read straight from the booklet.""",
}


ALL_NOTES = {
    3: LESSON_3_NOTES,
    4: LESSON_4_NOTES,
    5: LESSON_5_NOTES,
    6: LESSON_6_NOTES,
    7: LESSON_7_NOTES,
    8: LESSON_8_NOTES,
}


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    lines = text.split("\n")
    p = notes_tf.paragraphs[0]
    p.text = lines[0]
    for line in lines[1:]:
        p = notes_tf.add_paragraph()
        p.text = line


def apply(src, out, notes):
    pres = Presentation(src)
    n = 0
    skipped = []
    for i, slide in enumerate(pres.slides, 1):
        if i in notes:
            set_notes(slide, notes[i])
            n += 1
        else:
            skipped.append(i)
    pres.save(out)
    print(f"  Wrote notes to {n} teaching slides; skipped {len(skipped)} (admin/divider/credits)")
    print(f"  -> {out.name}")


def main():
    print(f"Working in: {FOLDER}")
    for lesson_num, fname in FILES.items():
        src = FOLDER / fname
        out_name = fname[:-5] + " - with notes.pptx"
        out = FOLDER / out_name
        print(f"\nLesson {lesson_num}: {fname}")
        if not src.exists():
            print(f"  SKIP - source not found")
            continue
        apply(src, out, ALL_NOTES[lesson_num])


if __name__ == "__main__":
    main()
