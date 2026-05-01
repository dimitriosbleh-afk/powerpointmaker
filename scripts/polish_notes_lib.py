"""Programmatic polish for existing teacher notes per teachernotes.md v2.0.

Functions:
    polish_deck(src_path, out_path, overrides=None)
        Reads `src_path`, applies baseline polish to every slide's notes, and
        saves to `out_path`. Per-slide overrides (1-based slide index → notes
        text) win over the programmatic polish.

Baseline polish:
    - ASCII safety: em/en dashes -> '-', curly quotes -> straight, ellipsis -> '...'
    - Banned-opener replacements in SAY lines (conservative).
    - Add brief TEACHER NOTES to slides that have no existing notes if the
      slide title matches a known admin/divider/materials pattern.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Dict, Optional

from pptx import Presentation
from pptx.util import Pt
from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------- Baseline ASCII safety ----------
ASCII_REPLACEMENTS = [
    ("—", "-"),   # em dash
    ("–", "-"),   # en dash
    ("…", "..."), # ellipsis
    ("“", '"'),    # left double quote
    ("”", '"'),    # right double quote
    ("‘", "'"),    # left single quote
    ("’", "'"),    # right single quote
    ("→", "->"),  # right arrow
    ("≥", ">="),  # greater-equal
    ("≤", "<="),  # less-equal
    ("×", "x"),   # multiplication sign
    ("•", "-"),   # bullet point -> hyphen (for note bullets)
]


def ascii_safe(text: str) -> str:
    for a, b in ASCII_REPLACEMENTS:
        text = text.replace(a, b)
    return text


# ---------- Banned opener replacements ----------
# Section 6 of v2.0 spec bans: "Today we are going to...", "Now we are going to...",
# "You will be...", "In this lesson...". Apply conservative replacements that
# preserve meaning while removing the presenter flavour.
BANNED_OPENERS = [
    (re.compile(r"^Today we are going to ", re.IGNORECASE), "We are "),
    (re.compile(r"^Today,? we will ", re.IGNORECASE), "We will "),
    (re.compile(r"^Today,? we are ", re.IGNORECASE), "We are "),
    (re.compile(r"^Today,? we read ", re.IGNORECASE), "We are reading "),
    (re.compile(r"^Today,? we'?re going to ", re.IGNORECASE), "We are "),
    (re.compile(r"^Now we are going to ", re.IGNORECASE), "We are now "),
    (re.compile(r"^Now,? we will ", re.IGNORECASE), "We will now "),
    (re.compile(r"^In this lesson,? we will ", re.IGNORECASE), "We will "),
    (re.compile(r"^In this lesson,? you will ", re.IGNORECASE), "You will "),
    (re.compile(r"^You will be learning ", re.IGNORECASE), "We are learning "),
]


def fix_banned_openers(text: str) -> str:
    out_lines = []
    for line in text.split("\n"):
        stripped = line.lstrip("- ").lstrip()
        prefix = line[: len(line) - len(line.lstrip("- "))] if line.startswith("- ") else ""
        # Strip leading straight quote (single or double) so the regex anchors
        # hit the actual word inside SAY-style quoted lines.
        leading_quote = ""
        body = stripped
        if body[:1] in ('"', "'"):
            leading_quote = body[:1]
            body = body[1:]
        replaced = body
        for pat, sub in BANNED_OPENERS:
            new = pat.sub(sub, replaced)
            if new != replaced:
                replaced = new
                break
        rebuilt = prefix + leading_quote + replaced
        if line == stripped:
            out_lines.append(leading_quote + replaced)
        else:
            out_lines.append(rebuilt)
    return "\n".join(out_lines)


# ---------- Admin slide detection ----------
ADMIN_PATTERNS = [
    # (regex over slide text, default TEACHER NOTES content)
    (
        re.compile(r"^lesson \d+", re.IGNORECASE),
        "TEACHER NOTES:\nTitle slide. Begin once materials are ready and the novel is open at the lesson's reading pages.",
    ),
    (
        re.compile(r"cultural sensitivity", re.IGNORECASE),
        "TEACHER NOTES:\nRead aloud or summarise before opening the novel for the first time.\n\nSENSITIVITY ADVISORY:\n- What it is: Storm Boy is set on Ngarrindjeri country and includes Aboriginal characters and language from the period.\n- Framing language: \"The story shows respect for the Coorong and the people who lived there. Some words from older stories sound different to how we speak today.\"\n- Watch for: students affected by names or images of deceased persons, especially Aboriginal and Torres Strait Islander students.\n- Protocol: pause if a student is upset, offer a quiet break with a peer or aide, follow up at recess and with your wellbeing lead if needed.",
    ),
    (
        re.compile(r"use of this resource", re.IGNORECASE),
        "TEACHER NOTES:\nTeacher orientation only, not for students. Read once before delivering the unit. The Literature Study Guide names the pause points and queries used through the lesson.",
    ),
    (
        re.compile(r"key principles", re.IGNORECASE),
        "TEACHER NOTES:\nTeacher reference for the I Do, We Do, You Do badges and the support and extension icons used through the deck. Not student-facing.",
    ),
    (
        re.compile(r"engagement icons", re.IGNORECASE),
        "TEACHER NOTES:\nTeacher reference for the response routines used through the deck: whiteboards, choral, thumbs, show fingers, pair share, cold call. Not student-facing.",
    ),
    (
        re.compile(r"colour coding chart|color coding chart", re.IGNORECASE),
        "TEACHER NOTES:\nTeacher reference for the sentence-element colour coding used in the modelling slides: who, what doing, when, where, why, how. Not student-facing.",
    ),
    (
        re.compile(r"^vocabulary$", re.IGNORECASE),
        "TEACHER NOTES:\nSection divider. Today's vocabulary words come from the lesson's reading pages.",
    ),
    (
        re.compile(r"text-?level reading", re.IGNORECASE),
        "TEACHER NOTES:\nSection divider. Today's reading mode is teacher choice. Have your novel pre-marked with the chosen pause points.",
    ),
    (
        re.compile(r"sentence-?level writing", re.IGNORECASE),
        "TEACHER NOTES:\nSection divider. The next slides explicitly teach the sentence-level focus for this lesson.",
    ),
    (
        re.compile(r"^you will need|in this lesson,? you will need", re.IGNORECASE),
        "SAY:\n- \"Boards out, novel out, booklet ready.\"\n- \"Texta in your hand. Lid checked.\"\n\nDO:\n- Scan the room for missing items before reading begins.\n- Pair up any students missing a board or a working texta.\n\nTEACHER NOTES:\nMaterial check. Settle this fast so the lesson can start.\n\nWATCH FOR:\n- Dry textas. Swap before the first show-me, not during it.",
    ),
    (
        re.compile(r"learning intention\s*&\s*success criteria|learning intention.*success criteria|^learning objectives$", re.IGNORECASE | re.MULTILINE),
        "SAY:\n- \"Read the learning intention with me.\"\n- \"These are the things we are practising today.\"\n- \"Ask: which one will be on your worksheet? Expected: the criterion that names today's writing or comprehension task.\"\n- \"If any of the words feel new, that is okay. We will build them together.\"\n\nDO:\n- Choral read the LI, then track each success criterion with your finger.\n- Ask one student to say SC1 in their own words.\n- Park any new vocabulary or term on the board so you can return to it during the I Do.\n\nTEACHER NOTES:\nSC1 is normally the floor. Almost every student should reach it from today's reading. The final SC is usually what the worksheet or You Do task assesses.\n\nWATCH FOR:\n- Students who cannot say SC1 in their own words. Give a concrete example before moving on.",
    ),
    (
        re.compile(r"for more resources|ochre\.org\.au|@ochreeducation", re.IGNORECASE),
        "TEACHER NOTES:\nCredits and attribution slide. Not student-facing. End the lesson on the closing reflection slide rather than this one.",
    ),
]


def slide_text(slide) -> str:
    """Concatenated text from all shapes on the slide."""
    chunks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    chunks.append(t)
        elif getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        chunks.append(t)
    return "\n".join(chunks)


def admin_default_notes(slide) -> Optional[str]:
    """Return default notes for an admin/divider slide based on its visible text, else None."""
    text = slide_text(slide)
    for pat, default in ADMIN_PATTERNS:
        if pat.search(text):
            return default
    return None


# ---------- Notes writer (mirrors write_notes.py) ----------
def _suppress_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("buChar", "buAutoNum", "buNone"):
        for el in pPr.findall("{%s}%s" % (A_NS, tag)):
            pPr.remove(el)
    etree.SubElement(pPr, "{%s}buNone" % A_NS)


def set_notes_text(slide, text: str):
    if not slide.has_notes_slide:
        _ = slide.notes_slide
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    paragraphs = text.split("\n")
    first = True
    for para_text in paragraphs:
        if first:
            p = notes_tf.paragraphs[0]
            first = False
        else:
            p = notes_tf.add_paragraph()
        is_bullet = para_text.startswith("- ")
        if is_bullet:
            para_text = para_text[2:]
        if para_text == "":
            _suppress_bullet(p)
            continue
        if not is_bullet:
            _suppress_bullet(p)
        run = p.add_run()
        run.text = para_text
        try:
            run.font.size = Pt(12)
        except Exception:
            pass


def get_notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text


# ---------- Top-level polish ----------
def polish_deck(src_path: Path, out_path: Path, overrides: Optional[Dict[int, str]] = None):
    overrides = overrides or {}
    prs = Presentation(str(src_path))
    fixed = 0
    added = 0
    overridden = 0
    for i, slide in enumerate(prs.slides, start=1):
        if i in overrides:
            new_text = overrides[i]
            new_text = ascii_safe(new_text)
            set_notes_text(slide, new_text)
            overridden += 1
            continue
        existing = get_notes_text(slide).strip()
        if not existing:
            default = admin_default_notes(slide)
            if default:
                set_notes_text(slide, ascii_safe(default))
                added += 1
            continue
        # Apply baseline polish to existing notes.
        polished = ascii_safe(existing)
        polished = fix_banned_openers(polished)
        if polished != existing:
            set_notes_text(slide, polished)
            fixed += 1
    prs.save(str(out_path))
    return {"slides": len(prs.slides), "fixed": fixed, "added": added, "overridden": overridden}


if __name__ == "__main__":
    import sys, json
    src = Path(sys.argv[1])
    out = Path(sys.argv[2])
    print(json.dumps(polish_deck(src, out), indent=2))
