#!/usr/bin/env python3
"""Merge Decimals & Fractions Lessons 1-6 into a single PPTX and gather every
companion PDF into one Resources folder.

Produces:
    output/Decimals_and_Fractions_Unit/Decimals and Fractions Unit.pptx
    output/Decimals_and_Fractions_Unit/Resources/<every PDF>.pdf
"""

import copy
import shutil
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
OUTPUT_ROOT = ROOT / "output"
UNIT_DIR = OUTPUT_ROOT / "Decimals_and_Fractions_Unit"
RESOURCES_DIR = UNIT_DIR / "Resources"

LESSON_DIRS = [
    ("DecFrac_Lesson1_Place_Value_With_Decimals",
     "DecFrac_Lesson1_Place_Value_With_Decimals.pptx", 1),
    ("DecFrac_Lesson2_Estimation_And_Rounding",
     "DecFrac_Lesson2_Estimation_And_Rounding.pptx", 2),
    ("DecFrac_Lesson3_Adding_Decimals",
     "DecFrac_Lesson3_Adding_Decimals.pptx", 3),
    ("DecFrac_Lesson4_Subtracting_Decimals",
     "DecFrac_Lesson4_Subtracting_Decimals.pptx", 4),
    ("DecFrac_Lesson5_Adding_Fractions",
     "DecFrac_Lesson5_Adding_Fractions.pptx", 5),
    ("DecFrac_Lesson6_Subtracting_Fractions_LCD",
     "DecFrac_Lesson6_Subtracting_Fractions_LCD.pptx", 6),
]

OUTPUT_PPTX = UNIT_DIR / "Decimals and Fractions Unit.pptx"


def copy_slide(src_slide, dst_prs):
    """Copy a slide from a source deck into the destination deck, preserving
    shapes, background, image relationships, and speaker notes."""
    dst_layout = dst_prs.slide_layouts[0]
    dst_slide = dst_prs.slides.add_slide(dst_layout)

    # Replace shape tree
    src_spTree = src_slide.shapes._spTree
    dst_spTree = dst_slide.shapes._spTree
    dst_spTree.getparent().replace(dst_spTree, copy.deepcopy(src_spTree))

    # Copy background
    src_bg = src_slide.background._element
    dst_bg = dst_slide.background._element
    dst_bg.getparent().replace(dst_bg, copy.deepcopy(src_bg))

    # Re-relate any image references
    rid_map = {}
    for rId, rel in list(src_slide.part.rels.items()):
        if "image" in rel.reltype:
            new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rId] = new_rId

    if rid_map:
        nsmap = {
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }
        attr_names = (
            "embed",
            "{%s}embed" % nsmap["r"],
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
        )
        for old_rId, new_rId in rid_map.items():
            if old_rId == new_rId:
                continue
            for elem in dst_slide.shapes._spTree.iter():
                for attr_name in attr_names:
                    if elem.get(attr_name) == old_rId:
                        elem.set(attr_name, new_rId)

    # Copy notes
    if src_slide.has_notes_slide:
        notes_text = src_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            dst_slide.notes_slide.notes_text_frame.text = notes_text

    return dst_slide


def merge_decks():
    UNIT_DIR.mkdir(parents=True, exist_ok=True)
    RESOURCES_DIR.mkdir(parents=True, exist_ok=True)

    first_dir, first_pptx, _ = LESSON_DIRS[0]
    base_path = OUTPUT_ROOT / first_dir / first_pptx
    base_prs = Presentation(str(base_path))
    print(f"Base deck: {base_path}  ({len(base_prs.slides)} slides)")

    for dir_name, pptx_name, session in LESSON_DIRS[1:]:
        src_path = OUTPUT_ROOT / dir_name / pptx_name
        src_prs = Presentation(str(src_path))
        print(f"Appending Lesson {session}: {src_path}  ({len(src_prs.slides)} slides)")
        for slide in src_prs.slides:
            copy_slide(slide, base_prs)

    base_prs.save(str(OUTPUT_PPTX))
    print(f"\nWrote combined deck: {OUTPUT_PPTX}  ({len(base_prs.slides)} slides total)")


def gather_resources():
    """Copy every PDF from each lesson's resources-sessionN/ folder into one
    flat Resources/ folder for the unit."""
    copied = 0
    for dir_name, _pptx_name, session in LESSON_DIRS:
        src_resources = OUTPUT_ROOT / dir_name / f"resources-session{session}"
        if not src_resources.is_dir():
            print(f"  (no resources folder for Lesson {session} at {src_resources})")
            continue
        for pdf in sorted(src_resources.glob("*.pdf")):
            dst = RESOURCES_DIR / pdf.name
            shutil.copy2(pdf, dst)
            copied += 1
            print(f"  Resources/  <-  {pdf.name}")
    print(f"\nCopied {copied} PDF(s) into {RESOURCES_DIR}")


def main():
    merge_decks()
    print()
    gather_resources()
    print(f"\nUnit folder: {UNIT_DIR}")


if __name__ == "__main__":
    main()
