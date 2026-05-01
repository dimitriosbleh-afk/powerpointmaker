"""Verify every output deck:
- has exactly one 'Review and Reflect' slide as the FINAL slide
- has at least one specific LI/SC slide (either pre-existing or newly inserted)
- new LI/SC slide (when present) sits immediately after the OCHRE objectives slide
"""
from pathlib import Path
from pptx import Presentation
import re

OUT = Path(r"C:\Users\09560329\Downloads\Challenge Slides Updated")

def lesson_no(name):
    m = re.match(r"\s*(\d+)\.", name)
    return int(m.group(1)) if m else None

def slide_text(slide):
    return "\n".join(
        "".join(r.text for r in p.runs)
        for shape in slide.shapes if shape.has_text_frame
        for p in shape.text_frame.paragraphs
    )

problems = []
for f in sorted(OUT.rglob("*.pptx")):
    p = Presentation(str(f))
    slides = list(p.slides)
    n = len(slides)
    rel = f.relative_to(OUT)
    lno = lesson_no(f.name)

    # Find Review + LI/SC + OCHRE objectives positions
    review_idxs = []
    li_sc_idxs = []  # specific LI/SC (school-style or our injected one)
    ochre_idxs = []
    new_li_idxs = []  # our injected one (title exact match)

    for i, s in enumerate(slides, 1):
        t = slide_text(s)
        low = t.lower()
        i_can = sum(1 for ln in t.splitlines() if ln.strip().lower().startswith("i can "))
        first_line = (t.strip().splitlines() or [""])[0].strip()
        if "Review and Reflect" in t:
            review_idxs.append(i)
        # Injected LI/SC has the exact title 'Learning Intention & Success Criteria'
        if first_line == "Learning Intention & Success Criteria" or t.startswith("Learning Intention & Success Criteria"):
            new_li_idxs.append(i)
        if (("learning intention" in low or "we are learning to" in low) and i_can >= 2):
            li_sc_idxs.append(i)
        if "learning objective" in low and "we are learning to" in low and i_can == 0:
            ochre_idxs.append(i)

    # Check 1: exactly one review at end
    if not review_idxs:
        problems.append(f"{rel}: NO Review slide")
    elif review_idxs != [n]:
        problems.append(f"{rel}: Review slide(s) at {review_idxs}, last slide is {n}")

    # Check 2: at least one LI/SC slide
    if not li_sc_idxs:
        problems.append(f"{rel}: NO LI/SC slide found")

    # Check 3: if injected LI/SC, it sits right after OCHRE slide
    if new_li_idxs and ochre_idxs:
        for nl in new_li_idxs:
            # find the closest preceding OCHRE
            preceding = [o for o in ochre_idxs if o < nl]
            if not preceding or preceding[-1] != nl - 1:
                problems.append(f"{rel}: new LI/SC at slide {nl} not immediately after OCHRE slide {ochre_idxs}")

print(f"Decks scanned: {sum(1 for _ in OUT.rglob('*.pptx'))}")
print(f"Problems: {len(problems)}")
for p in problems:
    print(f"  - {p}")
