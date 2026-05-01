"""Add Grade 1 teacher notes to T2 Wk 3 D1 LITERACY Planner.pptx.

Saves new file alongside source as "[original] - with teacher notes.pptx".
Preserves all slide faces, formatting, animations.
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

SRC = Path(r"c:/Users/09560329/Downloads/T2 Wk 3 D1 LITERACY Planner.pptx")
OUT = SRC.parent / (SRC.stem + " - with teacher notes.pptx")


# ============================================================
# NOTES BY SLIDE NUMBER
# ============================================================
NOTES = {

    # --- Slide 1: Title ---
    1: """TEACHER NOTES:
Week 3 literacy planner. Use this slide to settle the class while you check materials and the daily slide stack.""",

    # --- Slide 2: Week 3 Overview table ---
    2: """TEACHER NOTES:
Teacher planning overview. Skim before the week so you know which sessions need handbook pages, booklet slides, and printed craft items.

WATCH FOR:
- Mother's Day craft must go home Thursday. Print and prep early in the week.""",

    # --- Slide 3: MONDAY divider ---
    3: """TEACHER NOTES:
Monday section divider. Use to refocus the class before the daily routine.""",

    # --- Slide 4: blue bag & book box reading (Monday) ---
    4: """SAY:
- Get your blue bag.
- Find a good book nook.
- Read your readers first for at least five minutes.
- Then you can read from your book box.
- Quiet reading the whole time.

DO:
- Start the 10-minute timer on the slide.
- Move slowly around the room and check students are reading, not chatting.
- Use this time to do one or two individual reading assessments.

TEACHER NOTES:
This is independent reading. No Fluency Friends this week. Use this slot for one-on-one reading checks.

WATCH FOR:
- Students flicking pages without reading. Sit beside them and listen for 30 seconds.
- Students choosing only the book box and skipping readers. Redirect to readers first.""",

    # --- Slide 5: PAck up (Monday) ---
    5: """SAY:
- Readers back in your blue bag.
- Books back in your book box.
- Then come and sit on the floor.

DO:
- Cue the pack-up signal you normally use.
- Wait for everyone on the floor before starting the next slide.

TEACHER NOTES:
Quick transition slide. Keep it brisk so reading energy carries into the lesson.

WATCH FOR:
- Students still wandering. Use a count down to settle.""",

    # --- Slide 6: LITERACY TASK LI/SC (Monday) ---
    6: """SAY:
- Today we will look at how stories and information books are different.
- Read the learning intentions with me.
- We will also do some sentence work and check for full stops.
- Ask: who can say one thing we will learn today? Expected: about fiction and non-fiction, or sentences and fragments.

DO:
- Read each learning intention aloud, pointing to the words.
- Read each I can statement and check students are tracking.
- Tell students which success criterion the exit task will check.

TEACHER NOTES:
Two big focuses today. Keep the first chunk on text type and audience, the second chunk on sentences and fragments. Do not try to teach both at once.

WATCH FOR:
- Students who cannot say what we are learning in their own words. Restate in simple language.""",

    # --- Slide 7: English - Daily lesson 4 divider ---
    7: """TEACHER NOTES:
Lesson divider. This is the start of the Ochre TBU Lesson 4 sequence. Have the book A Year on Our Farm ready and the student booklet open at slide 13.""",

    # --- Slide 8: Section header - Understanding text (existing notes preserved) ---
    8: """TEACHER NOTES:
Section divider for slides 9 to 12. The visible notes here are the teaching guide for the upcoming text-feature slides. Show the first and last spread of A Year on Our Farm so students see how the year begins and ends in a similar way. Then move into voice, audience and purpose on the slides that follow.

WATCH FOR:
- Students grabbing for the book. Hold the book up but keep it in your hands until you ask them to look closely.""",

    # --- Slide 9: Text structure I ---
    9: """SAY:
- A Year on Our Farm is a special kind of story.
- It is called an episodic narrative because it happens over a whole year.
- The story is told in little parts, and the parts join together into one big story.
- It also has true facts about life on a farm, like an information book.
- Ask: what is the same about every part of the book? Expected: it all happens on the same farm in one year.

DO:
- Hold up the book.
- Show one or two double spreads as you talk about the little parts.
- Point to the slide picture as you say information.

TEACHER NOTES:
The threshold idea is that one book can be both a story and have true information. Keep language concrete: little parts, one big story, true facts.

WATCH FOR:
- Students confused by the word episodic. Just say little parts join together. Do not labour the term.""",

    # --- Slide 10: A Year on Our Farm - narrator W ---
    10: """SAY:
- The person telling a story is called the narrator.
- Listen as I read pages 5 and 6 again.
- Ask: who is telling this story? How do you know? Expected: a child, because he says Jess and I and we.
- Watch the slide. The narrator is one of the children. He talks about his mum and his sister Jess.

DO:
- Reread pages 5 and 6 of A Year on Our Farm.
- Point to the words Jess and I and we when they appear.
- Reveal the answer box only after students have answered.

TEACHER NOTES:
This is a We Do. Get students saying the word narrator. Use the page evidence, not just guessing.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Tell your partner one word the narrator uses that shows he is in the story.
EXTENDING PROMPT:
- Task: How do you think the narrator feels about Jess? What clue helps you?

WATCH FOR:
- Students saying the author is the narrator. Gently correct: the author wrote it, the narrator is the person inside the story telling it.""",

    # --- Slide 11: Audience W ---
    11: """SAY:
- The audience is the group of people the story is written for.
- Listen as I read pages 9 and 10 again.
- Ask: who do you think this story is written for? How do you know? Expected: children, because the words are easy and the pictures are interesting.
- The audience is children. The language is easy, the pictures are bright and the story is about a family.

DO:
- Reread pages 9 and 10.
- Point to the pictures as you talk about audience.
- Reveal the answer box only after students share.

TEACHER NOTES:
Keep the talk concrete. Audience means who the book is for. Use the pictures and the easy words as the proof.

WATCH FOR:
- Students saying the audience is one person. Correct: audience means the whole group of people, not just one reader.""",

    # --- Slide 12: Purpose W ---
    12: """SAY:
- The purpose is the reason the author wrote the book.
- Listen as I read pages 26 and 27 again.
- Ask: what does the author want us to learn or feel? Expected: what life is like on a farm across a whole year.
- The purpose is to show us a whole year on the farm.

DO:
- Reread pages 26 and 27.
- Point to the pictures as you reread.
- Reveal the answer box only after students share.

TEACHER NOTES:
Anchor purpose to learn, feel or enjoy. Help students name something specific they learn or feel from the book.

WATCH FOR:
- Students answering the purpose is to read it. Push for what we learn or how we feel.""",

    # --- Slide 13: Comparing text types I ---
    13: """SAY:
- Different kinds of books can talk about the same topic but sound different.
- A non-fiction book sounds more formal. It teaches us real facts.
- A fiction book sounds like a story. It shares ideas and feelings.
- Ask: which kind of book is A Year on Our Farm most like? Expected: it is both - mostly story, but with real facts about the farm.

DO:
- Hold up A Year on Our Farm.
- Hold up or point to the non-fiction extract from Lesson 2 if you have it.
- Compare the two side by side.

TEACHER NOTES:
This sets up the next slide. Keep examples short and visible. Do not over-explain genre theory.

WATCH FOR:
- Students saying fiction means not real. Confirm: fiction is made up, but A Year on Our Farm has true facts inside the story.""",

    # --- Slide 14: Comparing text types W (extract comparison) ---
    14: """SAY:
- Listen to two short extracts.
- One is from A Year on Our Farm: Lillypilly has...
- The other is from the non-fiction text we read in Lesson 2: On farms, people raise both plants and animals.
- Ask: how is the language different? Expected: one is friendly and personal, the other is formal and informative.
- Ask: why do you think the author chose the storytelling style for A Year on Our Farm? Expected: so children would enjoy it.

DO:
- Read both extracts slowly.
- Point to the soft, friendly words in the story extract.
- Point to the fact words in the non-fiction extract.

TEACHER NOTES:
Compare two short extracts only. Do not let the discussion run long. The point is the contrast in style.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Pick one word from each extract. Say which one sounds friendlier.
EXTENDING PROMPT:
- Task: Why might an information book about farms not work as well for young children?

WATCH FOR:
- Students who only spot length. Refocus on word choice and feeling, not length.""",

    # --- Slide 15: CFU - audience ---
    15: """SAY:
- Quick check. The group of people a story is written for is called what?
- A is purpose. B is narrator. C is audience.
- Show me on your fingers - one finger for A, two for B, three for C.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan all hands before taking the answer.

CFU CHECKPOINT:
Technique: Show Fingers (1-3)
Script:
- Ask: what do we call the group of people a story is written for? Expected: C, audience.
- Scan for: most fingers showing 3.
PROCEED:
- 80 percent or more show C. Move to the next teaching block.
PIVOT:
- Most likely: students confuse audience with narrator.
- Reteach: narrator is the one telling, audience is the ones listening or reading.
- Re-check: ask them to point to one of their classmates and say audience or narrator.

TEACHER NOTES:
Hinge check on the three terms. Do not move to sentence work until most students have this.

WATCH FOR:
- Students copying their neighbour's fingers. Eyes on you, not on them.""",

    # --- Slide 16: Sentence-level writing divider ---
    16: """TEACHER NOTES:
Section divider for slides 17 to 32. There are a lot of sentence slides here. Use teacher judgement on how many to work through before students go independent. Whiteboards at desks will keep students active.

WATCH FOR:
- Students fading after too many fragment slides. Move into the You Do as soon as you see they have it.""",

    # --- Slide 17: What is a sentence? I ---
    17: """SAY:
- A sentence is a group of words that makes a complete thought.
- It needs a who or what, and a what doing.
- If one of those parts is missing, it is a fragment, not a sentence.
- Read the checklist with me.

DO:
- Point to each part of the checklist as you say it.
- Read the bottom line about fragments slowly.

TEACHER NOTES:
This is the I Do for the whole sentence chunk. Anchor the three parts: complete thought, who or what, what doing. Use these words again on every following slide.

WATCH FOR:
- Students who think any group of words is a sentence. Push back to the checklist.""",

    # --- Slide 18: Is it a sentence? I ---
    18: """SAY:
- Watch this first.
- The first line is: the chook lays eggs.
- Is this a sentence? Let me check the list.
- It has a who - the chook. It has a what doing - lays eggs. It is a complete thought.
- But it does not have a capital letter or a full stop.
- Look at the second line. Now it has a capital T and a full stop. That is a correct sentence.

DO:
- Point to each part of the checklist as you say it.
- Point to the missing capital and the missing full stop on the first line.
- Point to the capital T and the full stop on the second line.

TEACHER NOTES:
This is modelling. Show your thinking step by step. Use the checklist out loud each time.

WATCH FOR:
- Students who only listen. Cue them to track with their finger or eyes.""",

    # --- Slide 19: Punctuate these sentences W ---
    19: """SAY:
- Sentences need a capital letter at the start and a full stop at the end.
- Look at the first line: the children picked apples.
- What is missing? Tell your partner.
- Ask: what should the first line look like? Expected: capital T at the start, full stop at the end.
- Now do the same with the cow line and the vegetables line.

DO:
- Have students fix one line at a time on their whiteboards.
- Take one student answer for each line.
- Reveal the corrected version on the slide.

TEACHER NOTES:
Three quick fixes. Move briskly. Capital and full stop only - do not introduce other punctuation here.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Just add the capital letter for now. Skip the full stop.
EXTENDING PROMPT:
- Task: Write your own sentence about an animal on a farm. Use a capital and a full stop.

WATCH FOR:
- Students writing capitals in the wrong place. Point to the very first letter of the sentence.""",

    # --- Slide 20: Fragment or sentence? W (drinks milk) ---
    20: """SAY:
- Look at the words on the slide: drinks milk.
- Is this a sentence? Let me check.
- It has a what doing - drinks. But who is drinking? We do not know.
- This is a fragment. Something is missing.
- Ask: what is missing? Expected: the who or what.

DO:
- Point to the checklist as you check each box.
- Point to where the who is missing.
- Show that we cannot tell who or what drinks.

TEACHER NOTES:
This is the model for the You Do pattern: name the fragment, identify what is missing, then we add it on the next slide.

WATCH FOR:
- Students saying the milk is the who. Correct: the milk is what is being drunk, not who is doing the drinking.""",

    # --- Slide 21: Fragment to sentence I ---
    21: """SAY:
- A fragment needs more detail to become a sentence.
- It needs both the who or what AND the what doing.
- Watch as I add a who: drinks milk becomes Daisy drinks milk.
- Now I add a capital D at the start and a full stop at the end.
- Ask: what did I add to make this a sentence? Expected: a who - Daisy.

DO:
- Point to drinks milk on the top.
- Point to Daisy drinks milk on the bottom.
- Circle the capital D and the full stop with your finger.

TEACHER NOTES:
Make the fix visible step by step. Daisy is the cow from the story - link back if students remember.

WATCH FOR:
- Students who try to add lots of words. Keep it simple - one who, then punctuate.""",

    # --- Slide 22: Convert fragments W (spreads hay / hunts mice) ---
    22: """SAY:
- Your turn with help.
- The fragment is: spreads hay.
- Whose job is it to spread hay on a farm? Show me on your boards.
- Now write your sentence: who, then the words spreads hay, then a full stop.
- Now try the next one: hunts mice.

DO:
- Give 30 seconds for each fragment.
- Scan all boards before revealing the answer.
- Reveal Dad spreads hay and Maria hunts mice.

TEACHER NOTES:
Two quick guided turns. Accept any reasonable who - Dad, the farmer, Maria, the cat. The point is naming a who and using a capital.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a name from your own family. Just write the who.
EXTENDING PROMPT:
- Task: Add a where, like Dad spreads hay in the paddock.

WATCH FOR:
- Students who forget the capital letter. Tap the start of their board and say capital.
- Students who forget the full stop. Tap the end and say full stop.""",

    # --- Slide 23: Convert fragments W (the animals / the tractor) ---
    23: """SAY:
- This time the who is given. We need to add the what doing.
- The fragment is: the animals.
- What might the animals be doing? Show me on your boards.
- Now try: the tractor.

DO:
- Give 30 seconds for each fragment.
- Scan all boards before revealing the answer.
- Reveal The animals are exploring and The tractor rumbles along.

TEACHER NOTES:
Now we are practising adding the what doing. Accept any sensible action, not just the slide answer.

WATCH FOR:
- Students who copy the slide example. Cue them to think first, then check the slide.""",

    # --- Slide 24: Convert fragments W (mixed - identify what is missing) ---
    24: """SAY:
- This time you need to work out what is missing first.
- Look at: spreads hay. What is missing? Expected: the who.
- Look at: the tractor. What is missing? Expected: the what doing.
- Now write the full sentence with a capital and a full stop.

DO:
- Walk students through each fragment.
- Let students name what is missing before they fix it.
- Reveal the example sentence after their attempts.

TEACHER NOTES:
This is a step harder. Students name the missing piece first, then fix. Slow down if many students get stuck on naming what is missing.

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
- Ask: what is missing in the tractor? Expected: the what doing.
- Scan for: students writing in the doing word, not the who.
PROCEED:
- 80 percent or more name the right missing piece. Move on.
PIVOT:
- Most likely: students add another who instead of the doing word.
- Reteach: re-anchor the checklist - who AND doing AND complete thought.
- Re-check: give them this fragment - the dog. What is missing?

WATCH FOR:
- Students who add a who when the doing was missing. Re-tap the checklist.""",

    # --- Slide 25: Sentences with extra detail (when/where) I ---
    25: """SAY:
- A sentence can have extra detail to make it more interesting.
- Watch this: the animals becomes The animals are exploring in the barn.
- I added a who and a what doing, AND I added a where - in the barn.
- Ask: what does in the barn tell us? Expected: where the animals are.

DO:
- Point to the where part - in the barn.
- Read the full sentence with feeling.

TEACHER NOTES:
This introduces when and where details. Keep the focus on these two questions only - when and where.

WATCH FOR:
- Students mixing up when and where. Use simple gestures - point to the floor for where, point to a clock for when.""",

    # --- Slide 26: Convert + when/where W ---
    26: """SAY:
- Your turn. Look at: spreads hay.
- Add a who and a when or where.
- Show me your sentence on your board.
- Now try: Our farm cat Maria. Add what doing and a where.

DO:
- Give one minute for each.
- Scan boards.
- Reveal Every morning, Dad spreads hay for the sheep, then Our farm cat Maria hunts mice in the barn.

TEACHER NOTES:
Two guided turns. Accept any sensible when or where, not just the slide example.

WATCH FOR:
- Students writing only the who and what doing. Cue them: now add the when or where.""",

    # --- Slide 27: CFU - which is a sentence? ---
    27: """SAY:
- Quick check. Look at the four choices.
- Which one is a sentence with the right capital and full stop?
- Show me 1, 2, 3 or 4 on your fingers.

DO:
- Give 10 seconds.
- Say: 3, 2, 1, show me.
- Scan all fingers before revealing.

CFU CHECKPOINT:
Technique: Show Fingers
Script:
- Ask: which option is the correctly punctuated sentence? Expected: The gander hisses.
- Scan for: most fingers on the option with capital G at the start and full stop at the end.
PROCEED:
- 80 percent or more correct. Move on to the next slide.
PIVOT:
- Most likely: students pick The Gander Hisses with capitals on every word.
- Reteach: only the first word and special names get a capital, not every word.
- Re-check: write the dog runs. Which one is correct - the dog runs, The Dog Runs, or The dog runs.

TEACHER NOTES:
Hinge check before moving to longer sentences with details. Watch for capital-on-every-word confusion.

WATCH FOR:
- Students copying neighbours. Boards down, then show.""",

    # --- Slide 28: Details in a sentence I (composite) ---
    28: """SAY:
- Watch how a sentence grows.
- Start with: the chook lays eggs.
- Add a when: every day the chook lays eggs.
- Add a where: every day the chook lays eggs in the coop.
- Now punctuate: Every day, the chook lays eggs in the coop.
- Notice the capital E and the comma after every day, and the full stop at the end.

DO:
- Point to each version on the slide as you read it.
- Point to the comma after Every day.
- Point to the capital and the full stop on the final version.

TEACHER NOTES:
This shows the growth from fragment to a longer sentence with when and where. Do not labour the comma - just point and name it.

WATCH FOR:
- Students looking overwhelmed. Cover the bottom version and reveal one line at a time.""",

    # --- Slide 29: CFU - which is a sentence (composite) ---
    29: """SAY:
- Another quick check. Look at the four choices.
- Which one is a sentence with a capital and a full stop?
- Show me 1, 2, 3 or 4.

DO:
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan all fingers before revealing.

CFU CHECKPOINT:
Technique: Show Fingers
Script:
- Ask: which option is the correctly punctuated sentence? Expected: The six goslings swim.
- Scan for: most fingers on the option with capital T and full stop, no capitals on every word.
PROCEED:
- 80 percent or more correct. Move on to the longer sentence work.
PIVOT:
- Most likely: students pick The Six Goslings Swim with extra capitals.
- Reteach: only the first letter and proper names get a capital.
- Re-check: ask if six and goslings are special names. Expected: no.

TEACHER NOTES:
Same trap as slide 27 - watch for the over-capitalised option.

WATCH FOR:
- Students sharing answers. Boards or fingers down, then show on cue.""",

    # --- Slide 30: Convert + when/where with table support W ---
    30: """SAY:
- The table on the slide gives us when words and where words.
- Pick one when from the table, and one where from the table.
- Add them to the fragment Dad spreads hay.
- Show me your sentence.

DO:
- Read the table aloud - in the morning, during spring; in the paddock, around the chicken coop.
- Give one minute for students to write.
- Reveal In the morning, Dad spreads hay around the chicken coop.

TEACHER NOTES:
The table is a scaffold. Students choose one from each column. Accept any combination as long as it makes a complete sentence.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Just choose one when from the table. Skip the where.
EXTENDING PROMPT:
- Task: Add your own when and where, not from the table.

WATCH FOR:
- Students copying the example. Cue them to mix and match from the table.""",

    # --- Slide 31: Convert + when/where + punctuation W ---
    31: """SAY:
- Two more to try together.
- The fragment is brings Daisy home. Add a who and a when. Punctuate it.
- The fragment is Jess. Add a what doing and a where. Punctuate it.

DO:
- Give one minute per fragment.
- Scan boards.
- Reveal In autumn, Dad brings Daisy home and Jess catches yabbies in the dam.

TEACHER NOTES:
Mixed practice. Some need a who, some need a doing word. Keep the checklist visible.

WATCH FOR:
- Students forgetting the punctuation. Tap the end of their board.
- Students who confuse Jess as a what doing instead of a who. Re-anchor: Jess is a person's name.""",

    # --- Slide 32: Convert + when AND where + punctuation W ---
    32: """SAY:
- Last one together. This time we add both a when and a where.
- The fragment is the children. Add a what doing, a when AND a where. Punctuate it.
- Then try ride on the truck. Add a who and a when. Punctuate it.

DO:
- Give 90 seconds per fragment.
- Scan boards.
- Reveal After school, the excited children shear sheep in the shed and In Winter, Kelly and Keeper ride on the truck on the farm.

TEACHER NOTES:
This is the hardest sentence slide. If students are tiring, take one example only and move to the You Do.

WATCH FOR:
- Students writing very long sentences that lose the meaning. Read it back and ask if it makes sense.""",

    # --- Slide 33: Your task Y (workbook) ---
    33: """SAY:
- Time to try it on your own.
- First, open your booklet to slide 13.
- Next, find the fragments and turn them into sentences.
- Then, check each sentence has a capital and a full stop.

DO:
- Direct students to booklet slide 13.
- Circulate and check the first one or two answers for each student.
- Pull a small group to the table if many are stuck.

TEACHER NOTES:
You Do for sentence-level writing. Look for capitals, full stops, and a complete who plus doing.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Just write the who and the doing. Use the slide table for ideas.
EXTENDING PROMPT:
- Task: Add a when or a where to every sentence.

WATCH FOR:
- Students starting without a capital. Stop and reset before they keep going.
- Fast finishers. Send them to add a where to each sentence.""",

    # ============================================================
    # TUESDAY (slides 34-58)
    # ============================================================

    # --- Slide 34: TUESDAY divider ---
    34: """TEACHER NOTES:
Tuesday section divider. Use to refocus class for the new daily routine.""",

    # --- Slide 35: blue bag & book box (Tuesday) ---
    35: """SAY:
- Get your blue bag.
- Find a good book nook.
- Read your readers first for at least five minutes.
- Then you can read from your book box.
- Quiet reading the whole time.

DO:
- Start the 10-minute timer.
- Move slowly around the room.
- Use this time for one or two reading assessments.

TEACHER NOTES:
Same independent reading routine as Monday.

WATCH FOR:
- Students who never finish a book. Note for follow-up reading group.""",

    # --- Slide 36: PAck up (Tuesday) ---
    36: """SAY:
- Readers back in your blue bag.
- Books back in your book box.
- Then sit on the floor.

DO:
- Cue the pack-up signal.
- Wait until everyone is settled.

TEACHER NOTES:
Quick transition. Keep it brisk.

WATCH FOR:
- Students still wandering. Use a count down.""",

    # --- Slide 37: LITERACY TASK LI/SC (Tuesday) ---
    37: """SAY:
- Today we will look at how things change over time.
- Read the learning intention with me.
- Read the success criteria with me.
- We are going to choose an image and show how it changes, then make our own drawing.

DO:
- Read the learning intention and each I can statement.
- Point to the words.
- Tell students which success criterion the You Do task will check.

TEACHER NOTES:
LI and SC for Lesson 5. Today is about change over time and the idea of symbols.

WATCH FOR:
- Students unsure what change over time means. Give one quick example - a baby growing up.""",

    # --- Slide 38: Daily lesson 5 divider ---
    38: """TEACHER NOTES:
Lesson divider for Lesson 5. Have A Year on Our Farm and the student booklet open at slide 15.""",

    # --- Slide 39: In this lesson you will need ---
    39: """SAY:
- Check you have your gear.
- Mini whiteboard and texta.
- Booklet.
- Pencil.

DO:
- Hold up each item as you say it.
- Wait until everyone has all three before moving on.

TEACHER NOTES:
Resource check slide. Sort gear gaps now so they do not slow the read-aloud.

WATCH FOR:
- Students missing a texta. Have spare textas ready.""",

    # --- Slide 40: Shared read aloud divider ---
    40: """TEACHER NOTES:
Section divider for the read-aloud. Use to refocus students before reading begins.""",

    # --- Slide 41: Re-read the Text (existing notes preserved as guide) ---
    41: """SAY:
- We are going to read A Year on Our Farm again.
- This is our third time, so today we will think about why things happen.
- Listen carefully and watch the pictures.
- Ask: what big ideas do you think the author wants us to understand about farm life? Hold this question in your head as we read.

DO:
- Pre-mark pages 14, 15, 27 and 28 as pause points.
- At page 14 (after baby bottle), pause and ask: what happens when a lamb is born? How does the family help?
- At page 15 (after help Mum), pause and ask: how does the family feel about working on the farm? What pictures tell you that?
- At page 27 (after baling hay), pause and ask: why do they prepare the land in certain months?
- At page 28 (after them alive), pause and ask: what might happen if they did not do the jobs at the right time?

TEACHER NOTES:
Third reading. Focus is cause and effect, plus the big ideas about farm life. Suggested big-idea responses include: farm life is hard work but everyone helps; you have to care for animals every day; farm life changes with the seasons; you have to be responsible.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Listen for one feeling word at each pause point and whisper it to your partner.
EXTENDING PROMPT:
- Task: Pick one big idea and find one picture in the book that shows it.

WATCH FOR:
- Students drifting during long read-aloud. Stand near them and point to the page.
- Students who answer the surface only. Push for what does this picture make you think.""",

    # --- Slide 42: Vocabulary divider ---
    42: """TEACHER NOTES:
Section divider. Two words today: freezing and allowed. Have whiteboards ready.""",

    # --- Slide 43: freezing - vocab I ---
    43: """SAY:
- Our first word is freezing.
- Say the word with me. Freezing.
- Freezing means very cold.
- In the story, the weather is freezing in July.
- Listen to my example: The man is freezing because the icy wind is blowing through his thin jacket.
- Ice cream is not a good idea on a freezing cold day.

DO:
- Point to the picture on the slide.
- Have students say freezing two times.
- Mime shivering as you say the word.

TEACHER NOTES:
Vocab I Do. Use the story link first, then your own example. Keep it brisk.

WATCH FOR:
- Students who say cold but not very cold. Push for the very part - freezing is a strong word.""",

    # --- Slide 44: Think and move - freezing W ---
    44: """SAY:
- In the story, Kelly and Keeper sleep outside even when it is freezing.
- Ask: why do you think they do not get cold? Expected: they have thick fur, they curl up together, they are working dogs.
- Now show me what you do when you are freezing.

DO:
- Take 2 quick answers to the why question.
- Cue everyone to stand and mime being freezing - shiver, hug arms, blow on hands.
- Sit students back down.

TEACHER NOTES:
Quick movement break tied to the word. Do not let it run more than one minute.

WATCH FOR:
- Students who do not move at all. Model first, then they copy.""",

    # --- Slide 45: Sort objects - freezing / not freezing W ---
    45: """SAY:
- Look at the pictures.
- Help me sort them - which ones are freezing, and which are not freezing?
- Point to one and show me thumbs up for freezing or thumbs down for not freezing.

DO:
- Point to each picture in turn.
- Take a class show of thumbs.
- Call on one student to explain a tricky one.

TEACHER NOTES:
Quick sort. Tie answers back to the meaning - very cold.

WATCH FOR:
- Students who confuse cool with freezing. Re-anchor: freezing is the very cold end.""",

    # --- Slide 46: Which drink would be freezing? W ---
    46: """SAY:
- Look at drink A and drink B.
- Which drink would be freezing?
- Show me A or B on your fingers - 1 for A, 2 for B.

DO:
- Give 5 seconds.
- Say: 3, 2, 1, show me.
- Take one student to explain the clue - ice, frost, cold colour.

TEACHER NOTES:
Quick choice. The answer depends on the slide images. Confirm with students why the freezing drink looks freezing.

WATCH FOR:
- Students choosing on which drink they like. Re-anchor: which one is freezing cold.""",

    # --- Slide 47: Which words fit best with freezing? CFU ---
    47: """SAY:
- Show me which words fit with the word freezing.
- Hold up one finger for each word that matches.
- Choices are: icy, hot, sticky, snowy.

DO:
- Read each word aloud.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers
Script:
- Ask: which words fit with freezing? Expected: icy and snowy.
- Scan for: most students showing two fingers, with the right two words named.
PROCEED:
- 80 percent or more get icy and snowy. Move to the next word.
PIVOT:
- Most likely: students pick sticky because it sounds like a weather word.
- Reteach: sticky is about touch, not temperature. Hot is the opposite of freezing.
- Re-check: which fits freezing - chilly or warm? Expected: chilly.

TEACHER NOTES:
Hinge check on freezing. Looking for students who can sort by meaning, not just by word feel.

WATCH FOR:
- Students who pick all four. Re-cue: only the ones that fit freezing.""",

    # --- Slide 48: allowed - vocab I ---
    48: """SAY:
- Our second word is allowed.
- Say the word with me. Allowed.
- If you are allowed to do something, it is okay to do it and you will not get in trouble.
- In the story, Kelly and Keeper are not allowed in the house because they are working dogs.
- Listen to my example: You are allowed to take photos here. You are allowed to exit this way.

DO:
- Point to the picture on the slide.
- Have students say allowed two times.
- Use a thumbs-up gesture for allowed.

TEACHER NOTES:
Vocab I Do. Allowed is an everyday word students will recognise. Anchor it to a school rule for fast meaning.

WATCH FOR:
- Students who confuse allowed with aloud. Say the meaning, not the spelling.""",

    # --- Slide 49: Think and move - allowed W ---
    49: """SAY:
- Kelly and Keeper are not allowed in the house.
- Pretend to be the farmer.
- What would you say or do if Kelly and Keeper tried to come inside?
- Tell your partner.

DO:
- Give partners 30 seconds.
- Take 2 student examples - words they would say, actions they would do.

TEACHER NOTES:
Short partner-talk. Keep it under one minute. Look for students using the word allowed in their answer.

WATCH FOR:
- Students who copy the slide picture. Push for their own words.""",

    # --- Slide 50: Are you allowed to... W ---
    50: """SAY:
- I will read three things. You show me thumbs up if you are allowed, thumbs down if you are not allowed.
- Are you allowed to ride in a car without a seat belt? Show me.
- Are you allowed to play outside during play time? Show me.
- Are you allowed to bring your pet into school every day? Show me.

DO:
- Read each one slowly.
- Give 3 seconds for thumbs each time.
- Confirm the answer with the slide.

TEACHER NOTES:
Three quick safety and school rule checks. Reinforces the meaning of allowed and links to classroom rules.

WATCH FOR:
- Students who get the seat belt one wrong. Stop and reteach - we are not allowed because it is not safe.""",

    # --- Slide 51: Section divider - understanding text ---
    51: """TEACHER NOTES:
Section divider before the symbols teaching block. Use to refocus students.""",

    # --- Slide 52: A Year on Our Farm I - flip through (existing notes preserved) ---
    52: """SAY:
- Look at the cover of A Year on Our Farm.
- Notice the tree.
- Watch as I flip through the pages and stop at each season.
- Ask: how does the tree change in each picture? Expected: leaves change colour, leaves fall off, leaves come back, the tree looks bare in winter.

DO:
- Hold up the cover.
- Flip through the book and stop once in each season.
- Pause at each season for students to look and describe the tree.

TEACHER NOTES:
This sets up the next slide on symbols. The tree is the same tree, but it changes across the year. That is the threshold idea.

WATCH FOR:
- Students who name colours but not change. Push: how is this tree different from the one before.""",

    # --- Slide 53: Symbols I (existing notes preserved) ---
    53: """SAY:
- Sometimes a picture in a story means more than just what we see.
- A symbol is a picture or thing that stands for something else.
- In this book, the tree does not just show us a tree. It shows us how time is passing and how the seasons are changing.
- Look at the picture on the slide. Do you know what this symbol stands for? Expected: power.

DO:
- Point to the symbol on the slide.
- Take one or two student answers.
- Confirm: this symbol means power.

TEACHER NOTES:
Introduce the word symbol. Anchor the meaning with the slide image first, then back to the tree in A Year on Our Farm.

WATCH FOR:
- Students who treat symbol as just any picture. Stress: a symbol stands for something else, not just itself.""",

    # --- Slide 54: Symbols CFU I ---
    54: """SAY:
- Read the slide with me.
- The tree is a symbol because it shows how time is passing and how the seasons change.
- Each time we see the tree again, it tells us another part of the year has gone by.
- Ask: what does the tree help us understand? Expected: time passing, the seasons changing.

DO:
- Read each bullet on the slide.
- Pause after each one and have students nod or thumb up if they agree.
- Take one student to explain in their own words.

TEACHER NOTES:
Anchor the symbol idea before the next slide. The tree stands for time passing.

WATCH FOR:
- Students who repeat the slide word for word. Push for their own words.""",

    # --- Slide 55: Other symbols for time CFU W ---
    55: """SAY:
- We can use other symbols to show time passing.
- Look at the pictures on the slide. What other symbol could show time passing?
- Talk to your partner.
- Take some student ideas.

DO:
- Point to the slide pictures.
- Give partners 30 seconds.
- Take 3 student ideas and write one or two on the board.

TEACHER NOTES:
Brainstorm slide. Accept ideas like a clock, a calendar, leaves falling, a baby growing.

WATCH FOR:
- Students who name objects but not symbols of time. Re-cue: what does this picture show changing.""",

    # --- Slide 56: Other symbols for time (existing notes - Think Pair Share) ---
    56: """SAY:
- Now use Think, Pair, Share.
- Think on your own first - what could be a symbol of time on the farm?
- Pair with the person next to you.
- Share your idea with your partner.

DO:
- Give 30 seconds think time, eyes closed if helpful.
- Cue Pair - 30 seconds.
- Take 3 to 4 ideas and record on the board.

TEACHER NOTES:
Examples to listen for: a lamb growing into a sheep, hay bales through the year, the dam filling and drying, fruit in the orchard. Record student ideas on the board for the You Do task.

WATCH FOR:
- Students stuck for an idea. Use the orchard or the lamb examples to seed thinking.""",

    # --- Slide 57: Which set of images shows passing of time? CFU ---
    57: """SAY:
- Quick check. Which set of images shows passing of time most clearly?
- Show me on your fingers - 1 for the first set, 2 for the second.

DO:
- Give 10 seconds.
- Say: 3, 2, 1, show me.
- Scan all fingers before revealing.

CFU CHECKPOINT:
Technique: Show Fingers
Script:
- Ask: which set of images shows the passing of time more clearly? Expected: the set that shows clear change between each image.
- Scan for: most students choosing the set with obvious change.
PROCEED:
- 80 percent or more correct. Move to the You Do.
PIVOT:
- Most likely: students pick the set they like the look of.
- Reteach: time passing means we can see something change between each picture.
- Re-check: ask them to point to the change between picture 1 and picture 4.

TEACHER NOTES:
Hinge before the You Do drawing task. Need students to spot which set shows real change before they create their own.

WATCH FOR:
- Students who change their answer to match a friend. Boards down, then show.""",

    # --- Slide 58: Your task Y - four panel drawing ---
    58: """SAY:
- Time to make your own drawing.
- Watch first as I draw an example - eggs to chicks to hens, in four panels.
- I am thinking out loud about how the picture changes from one panel to the next.
- Now your turn. Choose one idea from our list and draw it in four panels in your booklet.

DO:
- Model the four-panel drawing on the board first - eggs, cracking, baby chicks, full hens.
- Think aloud as you draw - first the eggs are still, then a crack, then the chicks come out, then they grow.
- Direct students to booklet slide 15.
- Circulate and check the first panel for each student.

TEACHER NOTES:
You Do drawing task. Push for clear change between panels, not four of the same picture.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use the lamb to sheep example. Just draw two big changes - small and big.
EXTENDING PROMPT:
- Task: Add a label to each panel - season, age, or time.

WATCH FOR:
- Students drawing four pictures with no change. Stop and ask what is changing.
- Students stuck on what to draw. Send them back to the class brainstorm list.""",

    # ============================================================
    # WEDNESDAY (slides 59-67)
    # ============================================================

    # --- Slide 59: WEDNESDAY divider ---
    59: """TEACHER NOTES:
Wednesday section divider. Use to refocus class.""",

    # --- Slide 60: blue bag (Wed) ---
    60: """SAY:
- Get your blue bag.
- Find a good book nook.
- Read your readers first for at least five minutes.
- Then you can read from your book box.
- Quiet reading the whole time.

DO:
- Start the timer.
- Move around the room.
- Use this time for one or two reading assessments.

TEACHER NOTES:
Same routine. Tag students you have not assessed yet this week.

WATCH FOR:
- Students who skip readers and only choose book box. Redirect to readers first.""",

    # --- Slide 61: PAck up (Wed) ---
    61: """SAY:
- Readers back in your blue bag.
- Books back in your book box.
- Then sit on the floor.

DO:
- Cue the pack-up signal.
- Wait until everyone is settled.

TEACHER NOTES:
Quick transition.

WATCH FOR:
- Students still wandering. Count down to settle.""",

    # --- Slide 62: LITERACY TASK LI/SC (Wed) - proper nouns ---
    62: """SAY:
- Today we are learning about a special kind of noun called a proper noun.
- Read the learning intentions with me.
- Read the I can statements with me.
- Proper nouns always start with a capital letter.

DO:
- Read each learning intention aloud.
- Read each I can statement, pointing to the words.
- Tell students the writing task will check that they can find a proper noun and use a capital.

TEACHER NOTES:
LI and SC for the noun lesson. Two ARC slides follow, then the writing activity.

WATCH FOR:
- Students who think a proper noun means a long word. Re-anchor: it is a special name for a person, place or thing.""",

    # --- Slide 63: ARC Lesson - Identify common and proper nouns ---
    63: """SAY:
- We are going to use a short ARC lesson to learn about common and proper nouns.
- Watch the slides carefully.
- I will pause to ask you questions.

DO:
- Make sure the file is open in PowerPoint, not Google Slides, so the transitions work.
- Run the ARC lesson Level 1, Lesson 4 in slideshow mode.
- Pause at the end to recap: a common noun is a general name. A proper noun is a special name and starts with a capital letter.

TEACHER NOTES:
This is a third-party ARC lesson. The transitions only work when the file is opened in PowerPoint. Make sure students hear the term proper noun several times.

WATCH FOR:
- Students who confuse common and proper. Stop and give two examples - dog vs Rex, school vs Greenlands Primary.""",

    # --- Slide 64: ARC Lesson - Capitalise proper nouns ---
    64: """SAY:
- This second ARC lesson shows how to capitalise proper nouns in sentences.
- Watch and listen.
- We will practise after the lesson.

DO:
- Run the ARC lesson Level 1, Lesson 5 in slideshow mode.
- Pause at the end to recap: every proper noun starts with a capital letter, even if it is in the middle of a sentence.

TEACHER NOTES:
Second ARC lesson. Recap with one or two examples on the board before students go to the writing slide.

WATCH FOR:
- Students who only capitalise the first word in a sentence. Stop and reteach.""",

    # --- Slide 65: Activity - write a sentence with a common and proper noun ---
    65: """SAY:
- Look at the picture together.
- Watch how I write a sentence about it.
- I need a common noun, a proper noun, and capital letters in the right places.
- Listen: My dog Pepper is very small.
- The common noun is dog. The proper noun is Pepper - capital P.
- The very first word My also has a capital letter.

DO:
- Display the slide picture.
- Write the example sentence on the board.
- Circle the common noun.
- Underline the proper noun.
- Point to each capital letter as you say it.

TEACHER NOTES:
This is the I Do for the noun-sentence task. Be explicit about both kinds of capitals - sentence start AND proper noun.

WATCH FOR:
- Students who skip the proper noun. Cue: every sentence today needs a name.""",

    # --- Slide 66: Your Turn - 3 pictures ---
    66: """SAY:
- Time to try it on your own.
- Look at picture 1, 2 and 3.
- Write a sentence in your writing book about each picture.
- Each sentence needs a common noun and a proper noun.
- Circle the common noun. Underline the proper noun.
- Use capitals at the start of the sentence and for any proper nouns.

DO:
- Direct students to their writing books.
- Circulate and check the first sentence for each student.
- Pull a small group to the table if many students are stuck.

TEACHER NOTES:
You Do. Three sentences. Looking for capital at sentence start, capital on proper noun, common and proper noun marked.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Write one sentence only. Use a sentence frame: My ____ is called ____.
EXTENDING PROMPT:
- Task: Use more than one proper noun in a sentence, like My friends Sam and Jack play soccer.

WATCH FOR:
- Students using lower-case names. Stop and tap the first letter of the name.
- Students who forget to circle and underline. Re-cue the marking task.""",

    # --- Slide 67: Optional Activity - noun sort ---
    67: """SAY:
- If we have time, you can do the noun sort activity.
- You can do this on your own or with a partner.
- Sort the words into common nouns and proper nouns.

DO:
- Hand out the Twinkl noun sort activity.
- Set partners or solo.
- Circulate and check sorting.

TEACHER NOTES:
Optional activity. Use only if pace allows. Skip if students are still on the writing task.

WATCH FOR:
- Students who sort by topic instead of common vs proper. Re-anchor: we are sorting on capitals, not on topic.""",

    # ============================================================
    # THURSDAY (slides 68-75)
    # ============================================================

    # --- Slide 68: THURSDAY divider ---
    68: """TEACHER NOTES:
Thursday section divider. Mother's Day craft must go home today - have all printed pieces ready.""",

    # --- Slide 69: blue bag (Thu) ---
    69: """SAY:
- Get your blue bag.
- Find a good book nook.
- Read your readers first for at least five minutes.
- Then you can read from your book box.
- Quiet reading the whole time.

DO:
- Start the timer.
- Move around the room.
- Use this time for one or two reading assessments.

TEACHER NOTES:
Same independent reading routine.

WATCH FOR:
- Students who finish a reader. Have a system for swapping to a new one.""",

    # --- Slide 70: PAck up (Thu) ---
    70: """SAY:
- Readers back in your blue bag.
- Books back in your book box.
- Then sit on the floor.

DO:
- Cue the pack-up signal.
- Wait until everyone is settled.

TEACHER NOTES:
Quick transition.

WATCH FOR:
- Students still wandering. Count down to settle.""",

    # --- Slide 71: LITERACY TASK LI/SC (Thursday) - Mother's Day ---
    71: """SAY:
- Today we are learning to follow steps to make a Mother's Day present.
- Read the learning intention with me.
- Read each I can statement.
- Following steps in order is called a procedure.

DO:
- Read the learning intention and success criteria.
- Hold up the Mother's Day craft so students see what they will make.
- Tell students the writing must be done first, then the craft.

TEACHER NOTES:
LI and SC for the Mother's Day session. Note: this task may run into Friday's slot if needed. The craft MUST go home today.

WATCH FOR:
- Students worried they cannot finish in time. Reassure - they will have time across two sessions if needed.""",

    # --- Slide 72: Mother's Day Writing - Why is Mother's Day special? ---
    72: """SAY:
- Today is a special day. We are making something for Mother's Day.
- Ask: why is Mother's Day a special day? Take a few ideas.
- Now look at the craft activity. We are going to follow the steps to make a present for our mum or someone special.

DO:
- Take 3 student answers about why Mother's Day is special.
- Display the Mother's Day craft pages.
- Show the instruction document briefly.

TEACHER NOTES:
Optional: read a Mother's Day story now or in library. The lesson is also a procedure lesson - keep that focus alongside the craft.

SENSITIVITY ADVISORY:
- What it is: Some students may not have a mum, may have lost their mum, or may have a complicated home situation.
- Framing language: Use mum or someone special who looks after you. Allow students to make the gift for any female caregiver - aunty, grandma, foster carer, family friend.
- Watch for: Students going quiet, looking sad, or saying they have no one to give it to.
- Protocol: Quietly check in with that student. Help them choose a person to make it for. Speak to your wellbeing lead if needed.

WATCH FOR:
- Students excited but distracted. Settle before moving to the next slide.""",

    # --- Slide 73: Mother's Day Writing - Why is your Mum special? ---
    73: """SAY:
- Now think about your own mum, or the special person you are making this for.
- What makes them special?
- We will write down our ideas.
- Then we will use these in our craft.

DO:
- Give 30 seconds of thinking time.
- Take 3 to 4 student ideas aloud.
- Direct students to write their ideas where the craft instructions show.

TEACHER NOTES:
Brainstorm and write step. Keep examples simple - she gives me hugs, she cooks dinner, she helps me read.

SENSITIVITY ADVISORY:
- What it is: Same as previous slide.
- Framing language: Use the special person who looks after you.
- Watch for: Students who cannot start writing.
- Protocol: Sit beside them and help them name one person and one reason.

WATCH FOR:
- Students writing very long answers. Cue them to keep ideas short for the craft.""",

    # --- Slide 74: Mother's Day Writing - procedure ---
    74: """SAY:
- When we read instructions and follow them, that is called a procedure.
- Look at the steps for our craft.
- Ask: why do we need to follow the steps in order? Expected: so the craft turns out right.
- Ask: what might happen if we do not? Expected: it might not work, we might miss a piece.

DO:
- Walk through each step of the chosen craft once.
- Point to each step in order.
- Show the students the writing step.

TEACHER NOTES:
Reinforce the genre - procedure. The teacher must select which bouquet template to use and the writing step that goes with it. Decide before the lesson and have the steps printed or displayed.

WATCH FOR:
- Students rushing into the craft before they understand the order. Hold them at the steps slide a little longer.""",

    # --- Slide 75: You Do - Mother's Day Writing ---
    75: """SAY:
- Now it is your turn to follow the steps.
- Start at step 1.
- Do each step in order.
- Put your hand up if you get stuck.

DO:
- Hand out craft materials.
- Circulate.
- Check the writing step is done before students move to the next step.
- Stop the class if many are getting stuck at the same step.

TEACHER NOTES:
You Do for the craft. The writing must be done as part of the procedure, not skipped. Remind students this MUST go home today.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Sit with a buddy who is one step ahead and check together.
EXTENDING PROMPT:
- Task: Once finished, write a card or extra message to go with the gift.

WATCH FOR:
- Students skipping the writing step and going straight to craft. Stop and reset.
- Students who finish quickly. Direct them to write a card or help a partner with the next step.""",

    # ============================================================
    # FRIDAY (slides 76-86)
    # ============================================================

    # --- Slide 76: FRIDAY divider ---
    76: """TEACHER NOTES:
Friday section divider. Plan to skip the text-connections piece if Mother's Day craft needs more time.""",

    # --- Slide 77: blue bag (Fri) ---
    77: """SAY:
- Get your blue bag.
- Find a good book nook.
- Read your readers first for at least five minutes.
- Then you can read from your book box.
- Quiet reading the whole time.

DO:
- Start the timer.
- Move around the room.
- Use this time for one or two reading assessments.

TEACHER NOTES:
Final independent reading slot for the week. Aim to have all students individually assessed by today.

WATCH FOR:
- Students you have not yet heard read this week. Prioritise them now.""",

    # --- Slide 78: PAck up (Fri) ---
    78: """SAY:
- Readers back in your blue bag.
- Books back in your book box.
- Then sit on the floor.

DO:
- Cue the pack-up signal.
- Wait until everyone is settled.

TEACHER NOTES:
Quick transition.

WATCH FOR:
- Students still wandering. Count down to settle.""",

    # --- Slide 79: WHOLE CLASS TEXT - Read Chapter 3 ---
    79: """SAY:
- It is class novel time.
- Today we are reading Chapter 3 of George's Marvellous Medicine.
- Listen carefully. You can read along with me if you want.

DO:
- Hold up the book so all students can see.
- Read Chapter 3 aloud at a steady pace.
- Pause briefly at the funny moments.

TEACHER NOTES:
Class novel. Read aloud, students follow with you. Keep the focus on enjoyment and understanding before the comprehension check.

WATCH FOR:
- Students drifting off. Glance at them and lift your reading energy.""",

    # --- Slide 80: Comprehension check in - George's rule ---
    80: """SAY:
- Quick check.
- What was George's one rule about his medicine?
- A: only brown things went in.
- B: it had to be healthy.
- C: runny, powdery or gooey things went in.
- Show me thumbs up for the answer you think is right when I read it.

DO:
- Read each option aloud, one at a time.
- Watch for thumbs at the right option.
- Confirm the answer from the text.

CFU CHECKPOINT:
Technique: Thumbs Up
Script:
- Ask: which was George's rule? Expected: C - runny, powdery or gooey things went in.
- Scan for: most thumbs up at option C.
PROCEED:
- 80 percent or more get C. Move to the inferential check.
PIVOT:
- Most likely: students pick B because it sounds sensible.
- Reteach: George's rule was about how things go in - not about being healthy.
- Re-check: ask one student to find the line in the chapter.

TEACHER NOTES:
Comprehension hinge after Chapter 3. The answer should be findable in the chapter just read.

WATCH FOR:
- Students copying neighbours. Eyes on the slide, then thumbs.""",

    # --- Slide 81: Inferential check in - bathroom cupboard ---
    81: """SAY:
- This question makes us think hard.
- Why was the medicine cupboard in the bathroom the only thing in the entire house that George was forbidden to touch?
- Turn and tell your partner.
- Take some ideas.

DO:
- Give 30 seconds for partner talk.
- Take 3 student answers.
- Connect to text-to-self - rules in your house, things you are not allowed to touch.

TEACHER NOTES:
Inferential question - the answer is not stated. Look for ideas like: medicines can be dangerous, they can make you very sick, the parents wanted to keep him safe.

WATCH FOR:
- Students who say I do not know. Use a sentence frame: maybe because... .""",

    # --- Slide 82: Share time - Text to Self Connection ---
    82: """SAY:
- Today we are practising a Text to Self Connection.
- That is when something we read reminds us of our own life.
- We are going to think about rules in our house and connect them to George's story.

DO:
- Read the slide aloud slowly.
- Have students repeat: Text to Self Connection.
- Tell students the next slide will give them the question to share.

TEACHER NOTES:
Quick teaching slide. Anchor the term Text to Self Connection. Keep this short and move into the share.

WATCH FOR:
- Students confused by the term. Use the example: George could not touch the cupboard - that reminds me of my house rule about the heater.""",

    # --- Slide 83: Share time - rules in your house ---
    83: """SAY:
- George was not allowed to touch anything in the bathroom cupboard. That was a rule his parents had.
- Have you got rules like that in your house?
- Are there places or things you are not allowed to touch?
- Share your Text to Self Connection with your POD.

DO:
- Give 30 seconds thinking time.
- Cue students to turn to their POD and share.
- Walk around and listen.
- Take 2 students to share with the whole class.

TEACHER NOTES:
Share time. Look for connections that link a story rule to a home rule.

SENSITIVITY ADVISORY:
- What it is: Home rules can sometimes touch on hard topics - locked rooms, alcohol cupboards, safety issues.
- Framing language: Things our family says we should not touch to keep us safe.
- Watch for: Students sharing something that feels worrying or unsafe.
- Protocol: Listen, do not react in front of the class. Follow up quietly later and refer to the wellbeing lead if needed.

WATCH FOR:
- Students who go off topic. Re-cue: a rule from your house, like George's rule.""",

    # --- Slide 84: WHOLE CLASS TEXT - Chapter 4 ---
    84: """SAY:
- Now we read Chapter 4 of George's Marvellous Medicine.
- Listen carefully.
- I will pause for a question after we finish.

DO:
- Hold up the book.
- Read Chapter 4 aloud at a steady pace.

TEACHER NOTES:
Plan note: only read Chapter 4 today, not Chapter 5. Move Chapter 5 to next session.

WATCH FOR:
- Students who lose focus mid-chapter. Slow your pace and add some expression.""",

    # --- Slide 85: Inferential check in - George in the shed ---
    85: """SAY:
- Inferential question.
- How do you think George felt when he went into the shed to find more ingredients for his medicine?
- Turn and tell your partner.
- Use a feeling word and a reason.

DO:
- Give 30 seconds for partner talk.
- Take 3 student answers.
- Push for a feeling and a reason: I think George felt ___ because ___.

TEACHER NOTES:
Inferential question. Possible answers: excited because he had a plan, nervous because he might get caught, brave because he was determined. Accept any feeling that the student can back with a reason.

WATCH FOR:
- Students who give a feeling with no reason. Re-prompt: why do you think that.""",

    # --- Slide 86: You Do - Mother's Day Writing finish ---
    86: """SAY:
- Now we finish our Mother's Day craft.
- This MUST go home today.
- Pick up where you stopped yesterday and finish all the steps.

DO:
- Hand back unfinished crafts.
- Circulate and help students at the steps where they got stuck.
- Make sure each finished gift is named and ready to go in bags.

TEACHER NOTES:
Finishing slot for the Mother's Day craft. Priority is that every student takes their gift home today.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Sit with a buddy. Work through the next step together.
EXTENDING PROMPT:
- Task: Once finished, decorate or add an extra detail to your card.

WATCH FOR:
- Students who have not started. Sit with them and start step 1 together.
- Students who finish early. Help a friend with the step they are stuck on.""",
}


# ============================================================
# APPLIER
# ============================================================
def suppress_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for tag_name in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag_name)):
            pPr.remove(el)
    pPr.append(etree.SubElement(pPr, qn("a:buNone")))


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    parts = []
    for line in text.split("\n"):
        if line.startswith("- "):
            parts.append((line[2:], True))
        else:
            parts.append((line, False))

    p = notes_tf.paragraphs[0]
    p.text = parts[0][0]
    if not parts[0][1]:
        suppress_bullet(p)
    for content, is_bullet in parts[1:]:
        p = notes_tf.add_paragraph()
        p.text = content
        if not is_bullet:
            suppress_bullet(p)


def main():
    pres = Presentation(str(SRC))
    slide_count = len(pres.slides)
    n = 0
    missing = []
    for i, slide in enumerate(pres.slides, 1):
        if i in NOTES:
            set_notes(slide, NOTES[i])
            n += 1
        else:
            missing.append(i)
    pres.save(str(OUT))
    print(f"Source has {slide_count} slides")
    print(f"Applied notes to {n} slides -> {OUT.name}")
    if missing:
        print(f"Slides without notes: {missing}")
    else:
        print("All slides received notes.")


if __name__ == "__main__":
    main()
