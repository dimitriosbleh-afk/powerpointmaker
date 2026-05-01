"""Inspect the OCHRE 'Learning objectives' slide in lesson 16 to understand structure."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

p = Presentation(r"C:\Users\09560329\Downloads\Challenge Slides\4\16. literature_presentation Sentence Combining - Compound Sentences 4En06V Storm Boy - Information report.pptx")

# Slide 7 is the objectives slide
slide = list(p.slides)[6]

print(f"Slide layout name: {slide.slide_layout.name}")
print(f"Number of shapes: {len(slide.shapes)}")
print()
for i, shape in enumerate(slide.shapes):
    info = []
    info.append(f"shape[{i}] type={shape.shape_type}")
    if shape.has_text_frame:
        tf = shape.text_frame
        info.append(f"text:\n  {tf.text!r}"[:200])
        info.append(f"#paragraphs={len(tf.paragraphs)}")
        for j, p in enumerate(tf.paragraphs):
            info.append(f"  para[{j}] level={p.level} text={p.text!r}"[:200])
    info.append(f"left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
    info.append(f"name={shape.name}")
    print("\n".join(info))
    print("-" * 60)

