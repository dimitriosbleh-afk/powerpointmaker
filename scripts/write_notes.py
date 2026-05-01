"""Write teacher notes into a PPTX from a JSON {slide_num (1-based): notes_text} mapping.

Usage:
  python write_notes.py <input.pptx> <notes.json> <output.pptx>

Notes are written as plain text into the notes slide. Existing notes are replaced
entirely. Slides not present in the JSON are left untouched.
"""
import sys
import io
import json
import copy
from pptx import Presentation
from pptx.util import Pt
from lxml import etree

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def set_notes_text(slide, text):
    """Replace the notes slide body text frame with `text` (preserving paragraphs)."""
    if not slide.has_notes_slide:
        # Force the notes slide to exist
        _ = slide.notes_slide
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()

    paragraphs = text.split("\n")
    # First paragraph uses the existing first paragraph
    first = True
    for para_text in paragraphs:
        if first:
            p = notes_tf.paragraphs[0]
            first = False
        else:
            p = notes_tf.add_paragraph()
        # Track whether the original line was a hyphen bullet so PowerPoint
        # only renders bullets on real bullet lines, not on headers or prose.
        is_bullet = para_text.startswith("- ")
        if is_bullet:
            para_text = para_text[2:]
        if para_text == "":
            # Blank line: suppress bullet so empty separator paragraphs
            # do not show a stray dot.
            _suppress_bullet(p)
            continue
        if not is_bullet:
            _suppress_bullet(p)
        run = p.add_run()
        run.text = para_text
        try:
            run.font.size = Pt(12)
        except Exception:
            pass


def _suppress_bullet(p):
    """Add <a:buNone/> to the paragraph's pPr so PowerPoint renders no bullet."""
    pPr = p._p.get_or_add_pPr()
    # Remove any existing bullet-related child to avoid duplicates.
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall("{%s}%s" % (A_NS, tag.split(":", 1)[1])):
            pPr.remove(el)
    bu_none = etree.SubElement(pPr, "{%s}buNone" % A_NS)


def main():
    in_path = sys.argv[1]
    json_path = sys.argv[2]
    out_path = sys.argv[3]

    with open(json_path, "r", encoding="utf-8") as f:
        notes_map = json.load(f)
    # Keys come back as strings — normalise to int
    notes_map = {int(k): v for k, v in notes_map.items()}

    prs = Presentation(in_path)
    written = 0
    for i, slide in enumerate(prs.slides, start=1):
        if i in notes_map:
            set_notes_text(slide, notes_map[i])
            written += 1

    prs.save(out_path)
    print(f"Wrote notes to {written} slides -> {out_path}")


if __name__ == "__main__":
    main()
