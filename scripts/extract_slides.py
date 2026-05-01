"""Extract slide text + existing notes from a PPTX for teacher-notes authoring."""
import sys
import io
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")


def extract(path):
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            elif shape.shape_type == 19 or getattr(shape, "has_table", False):
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = cell.text.strip()
                            if t:
                                texts.append(t)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        out.append({"index": i, "text": texts, "notes": notes})
    return out


if __name__ == "__main__":
    path = sys.argv[1]
    data = extract(path)
    for s in data:
        print(f"===== SLIDE {s['index']} =====")
        for line in s["text"]:
            print(f"  {line}")
        if s["notes"]:
            print("  --- EXISTING NOTES ---")
            for line in s["notes"].splitlines():
                print(f"  {line}")
        print()
