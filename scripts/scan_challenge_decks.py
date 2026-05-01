"""Scan all Challenge Slides decks. Per-lesson, find:
- whether a specific 'Learning Intention & Success Criteria' slide exists (i.e. has 'I can' bullets)
- whether an OCHRE generic 'Learning objectives' slide exists
- index of each, and the resulting verdict.
"""
from pathlib import Path
from pptx import Presentation
import re
import sys

ROOT = Path(r"C:\Users\09560329\Downloads\Challenge Slides")

def lesson_no(name):
    m = re.match(r"\s*(\d+)\.", name)
    return int(m.group(1)) if m else None

def slide_text(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                if t.strip():
                    out.append(t)
    return "\n".join(out)

def classify(p):
    """Return tuple (li_sc_slide_idx_or_None, ochre_slide_idx_or_None)."""
    li_sc = None
    ochre = None
    for i, s in enumerate(p.slides, 1):
        t = slide_text(s).lower()
        # Specific LI/SC: at least 2 specific 'I can ...' bullets (not the boilerplate)
        i_can_lines = [line for line in t.splitlines() if line.strip().startswith("i can ")]
        has_li_phrase = "learning intention" in t or "we are learning to" in t
        if has_li_phrase and len(i_can_lines) >= 2 and li_sc is None:
            li_sc = i
        # Generic OCHRE objectives slide: has 'learning objectives' header and the
        # boilerplate phrase but no 'I can' bullets
        if ("learning objective" in t and "we are learning to" in t and
            len(i_can_lines) == 0 and ochre is None):
            ochre = i
    return li_sc, ochre

decks = sorted(ROOT.rglob("*.pptx"))
records = []
for f in decks:
    lno = lesson_no(f.name)
    if lno is None:
        continue
    p = Presentation(str(f))
    li_sc, ochre = classify(p)
    needs_new = li_sc is None
    records.append((lno, "WN" if "with notes" in f.name.lower() else "OG", li_sc, ochre, needs_new, f))

# Pretty print
print(f"{'Lesson':<7}{'Ver':<4}{'LI/SC':<7}{'OCHRE':<7}{'Need new':<9} File")
print("-"*120)
for lno, ver, li_sc, ochre, needs, f in records:
    rel = f.relative_to(ROOT)
    print(f"{lno:<7}{ver:<4}{str(li_sc or '-'):<7}{str(ochre or '-'):<7}{'YES' if needs else 'no':<9} {rel}")

# Lessons that need new LI/SC across both versions
needs_lessons = sorted({lno for lno, _, _, _, n, _ in records if n})
print()
print(f"Lessons needing new LI/SC: {needs_lessons}")
print(f"Total decks: {len(records)}")
