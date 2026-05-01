"""Add teacher notes (megapromptlean_3.0 format) to Storm Boy lessons 21-25.

Saves a new file alongside each source: "[original name] - with notes.pptx".
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

SRC_DIR = Path(r"c:/Users/09560329/Downloads/Challenge Slides/5")

L21_FILE = "21. literature_presentation Punctuate direct speech 4En06V Storm Boy - Information report.pptx"
L22_FILE = "22. literature_presentation Single paragraph outline (SPO) to summarise a text 4En06V Storm Boy - Information report.pptx"
L23_FILE = "23. literature_presentation Note taking 4En06V Storm Boy - Information report.pptx"
L24_FILE = "24. literature_presentation Plan a concluding paragraph for an information report 4En06V Storm Boy - Information report.pptx"
L25_FILE = "25. literature_presentation Write a concluding paragraph for an information report 4En06V Storm Boy - Information report.pptx"


# ============================================================
# LESSON 21 - Punctuate direct speech
# ============================================================
LESSON_21_NOTES = {

    9: """SAY:
- Quick reset before we read.
- Today we will read pages 69 to 72 of Storm Boy.
- Listen carefully and watch the pictures - we will pause to talk.

DO:
- Pre-mark your novel with the chosen pause points.
- Settle students for read-aloud with novels open or in view.

TEACHER NOTES:
This is the read-aloud setup slide. Mode of reading is teacher choice. Pre-marking the novel before students arrive prevents stop-and-search during the read.

WATCH FOR:
- Students who do not have a copy ready - sort this before the first page.
- Students still in transition energy - reset with a short routine before opening.

Stage: Text-Level Reading | VTLM: Engagement""",


    10: """SAY:
- We are going to read pages 69 to 72.
- Your job is to listen carefully and look at the pictures.
- I will pause at certain points to ask questions.
- Ask: what do we do at a pause point? Expected: think first, then share.

DO:
- Read pages 69 to 72 of Storm Boy aloud, or run the chosen reading mode.
- Pause at the points you have pre-marked using the Literature Study Guide.
- Take 1-2 responses at each pause - then keep moving.
- Gloss any incidental vocabulary in one short sentence.

TEACHER NOTES:
Pause points and queries come from the Literature Study Guide - choose the ones that match your students' needs, not all of them. Pages 69 to 72 cover the moment Mr Percival is shot, so the chapter ends emotionally.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Listening partner - the supported student listens, then their partner whispers a one-word summary at each pause.
EXTENDING PROMPT:
- Task: Note one literary device the author uses on these pages and how it shapes the mood.

WATCH FOR:
- Students drifting off - point to the page and ask a quick concrete question.
- Confused looks at unknown words - one-sentence gloss, then keep reading.

Stage: Text-Level Reading | VTLM: Explicit Teaching""",


    11: """SAY:
- We are going to teach three new words from today's pages.
- Watch and listen first, then we will practise together.

DO:
- Have whiteboards and textas ready.
- Display the first vocabulary word.

TEACHER NOTES:
Brief divider slide. The Vocabulary guide on the Ochre website has more detail if needed.

WATCH FOR:
- Students who need wait time - protect it across the next set of slides.

Stage: Vocabulary Divider | VTLM: Engagement""",


    12: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our first word is blazing.
- Watch this first - I am going to think aloud about what this word means.
- Blazing means burning brightly.
- I notice the story uses it for the sun - Storm Boy was watching the sun blazing like a round, shiny coin. So blazing is not only fire - it can be very bright sun too.
- Say the word with me... blazing!
- Ask: which word family does it belong to? Expected: blaze, blazes, blazed, blazingly.

DO:
- Point to the image on the slide.
- Use a strong, bright voice for blazing so students hear the feeling.
- Have students repeat the word twice.

TEACHER NOTES:
First meeting with the word. Linking it back to the sun in the story gives students a familiar anchor before the practice slides.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Repeat-after-me on the word, then point to the image and say one short example sentence.
EXTENDING PROMPT:
- Task: Add one more example sentence using blazing for speed (e.g. a blazing sprinter).

WATCH FOR:
- Students who think blazing only means on fire - confirm a very bright sun also counts.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    13: """SAY:
- Look at the picture.
- If this image can be described as blazing, say blazing.
- If not, stay silent.

DO:
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.
- Listen for a strong choral response.

TEACHER NOTES:
Expected response: blazing. The image shows a bright, fiery sun.

WATCH FOR:
- A few students saying it before others - reset wait time on the next image.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    14: """SAY:
- Look at the picture.
- If this image can be described as blazing, say blazing.
- If not, stay silent.

DO:
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.

TEACHER NOTES:
Expected response: silence (not blazing). The image is dim or cool, not bright or fiery.

WATCH FOR:
- Students saying blazing because the picture is interesting - link back to the meaning: blazing is brightness, heat or speed.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    15: """SAY:
- Look at the picture.
- If this image can be described as blazing, say blazing.
- If not, stay silent.

DO:
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.

TEACHER NOTES:
Expected response: blazing. Confirm the link to brightness or fire on the reveal.

WATCH FOR:
- Students who hesitate - re-anchor the meaning before the next slide.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    16: """SAY:
- Two pictures - A and B.
- Which image can be described as blazing? Show me A or B on your whiteboard.
- Ask: why did you choose that one? Expected: it is very bright.

DO:
- Read both options aloud.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.
- Cold call one student to explain their choice.

TEACHER NOTES:
Expected response: car headlights are blazing - they are very bright. Use the language of brightness on the reveal.

WATCH FOR:
- Students choosing the brighter colour photo for the wrong reason - prompt: is it bright enough to dazzle you?

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    17: """SAY:
- Blazing has three meanings we use a lot - brightness, heat and speed.
- Watch how I think about each one. A blazing sun is bright. A blazing fire is hot. An olympic runner is fast - they have a blazing pace.
- Ask: can you think of something else that is blazing in each of these ways? Turn and tell your partner.

DO:
- Give 30 seconds of partner talk.
- Take 1-2 ideas for each meaning from the room.
- Build a quick word bank on the board if helpful.

TEACHER NOTES:
This is the move from definition to flexible use. Three meanings are the core knowledge for the rest of the cycle.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Choose one meaning only and give one example for it.
EXTENDING PROMPT:
- Task: Write one sentence that uses blazing in two of the three meanings (e.g. The sun was blazing on the blazing sprinter).

WATCH FOR:
- Students who can only give the bright meaning - prompt with: what about something hot? Or something fast?

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    18: """SAY:
- Show me which word fits best with blazing.
- Hold up 1 for snow, 2 for fire, 3 for wind, 4 for water.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which word fits best with blazing? On three, show me 1, 2, 3 or 4.
- Scan for: fingers showing 2 (fire).
PROCEED:
- If most students are on 2, move on to sentence rephrase practice.
PIVOT:
- Most likely misconception: students pick wind because of blazing speed.
- Reteach: blazing fits best with fire because fire shows brightness AND heat together. Wind has speed but not brightness.
- Fresh re-check: which fits better with blazing - sun or rain?

TEACHER NOTES:
Multiple-choice CFU. Fire is the strongest match because it carries two of the three meanings.

WATCH FOR:
- Students mirroring the fastest hand - keep wait time firm.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    19: """SAY:
- Watch how I use blazing to make a sentence stronger.
- Original: The car sped around the track at an incredibly fast pace.
- I am going to swap incredibly fast for blazing.
- New version: The car sped around the track at a blazing pace.
- Ask: which version sounds stronger? Expected: the blazing version.

DO:
- Read the original sentence.
- Read the rephrased version with emphasis on blazing.
- Have students echo the new sentence with you.

TEACHER NOTES:
Students see how blazing replaces a longer descriptive phrase. Highlights the speed meaning.

WATCH FOR:
- Students who lose the speed connection - remind them blazing has three meanings, and speed is one.

Stage: Vocabulary I Do | VTLM: Modelling""",


    20: """SAY:
- Listen: The chef cooked the steak on a burning hot stove, making it sizzle.
- Watch how I rephrase this with blazing.
- I look for a word that means very hot - that is blazing.
- New version: The chef cooked the steak on a blazing stove, making it sizzle.
- Ask: which meaning of blazing did I use here? Expected: heat.

DO:
- Read both versions aloud.
- Point to blazing each time.

TEACHER NOTES:
This time blazing replaces burning hot - the heat meaning. Two of three meanings have now been used in rephrase practice.

WATCH FOR:
- Students who think blazing always means on fire - confirm a stove can be blazing because it is very hot.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    21: """SAY:
- Build a sentence using the word blazing.
- Sentence frame: The ____ was blazing, making it ____.
- Watch my example: The sun was blazing, making it hard to see.
- Now write your own. Pick a different blazing thing.

DO:
- Give 60 seconds for students to write.
- Circulate and check 3-4 sentences.
- Take 1-2 examples to share.

TEACHER NOTES:
Active sentence building. Accept brightness, heat or speed examples - reject sentences that use blazing wrongly.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use the sentence frame exactly with one word from the board word bank.
EXTENDING PROMPT:
- Task: Write a sentence that does NOT use the frame. Use blazing for speed.

WATCH FOR:
- Students stuck for an idea - prompt with sun, fire, headlights, oven, sprinter.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    22: """SAY:
- Which scenario describes something blazing? A or B?
- A - The campfire crackled, sending sparks into the night with intense heat.
- B - The leaves swayed gently in the breeze, casting shadows below.
- Show me A or B.

DO:
- Read both scenarios.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Write A or B. 3, 2, 1, show me.
- Scan for: most boards show A.
PROCEED:
- If most students show A, move on to splayed.
PIVOT:
- Most likely misconception: students pick B because the language sounds nice.
- Reteach: blazing needs brightness, heat or speed. A has intense heat and sparks - that is blazing. B is gentle, not blazing.
- Fresh re-check: which is blazing - a calm pond or a bonfire?

TEACHER NOTES:
End-of-word CFU. Strong response to A signals readiness for the next vocabulary word.

WATCH FOR:
- Students who copy a neighbour - scan quickly and protect wait time.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    23: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our next word is splayed.
- Watch this first - I am going to think aloud.
- Splayed means spread out.
- In the story, Mr Percival's wing was splayed out after he was shot. So I picture wings stretched wide, not folded.
- Say the word with me... splayed!
- Ask: what does this gesture show? (spread fingers wide). Expected: splayed.

DO:
- Point to the image on the slide.
- Spread your fingers wide and hold them up to model splayed.
- Have students repeat the word twice and copy the gesture.

TEACHER NOTES:
The hand gesture is sticky for students - keep using it across the practice slides.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Repeat the word twice and copy the spread-finger gesture.
EXTENDING PROMPT:
- Task: Name two things in the classroom that could be splayed (e.g. fan, deck of cards).

WATCH FOR:
- Students who muddle splayed with closed - re-show the wide-spread gesture.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    24: """SAY:
- Read the sentence with me.
- The bird's wings were splayed out, ready to take flight.
- Ask: what is splayed in this sentence? Expected: the bird's wings.
- Use your hands to show splayed wings.

DO:
- Choral read.
- Pause to gesture splayed wings together.

TEACHER NOTES:
Connecting the word to a clear picture. Use the gesture to anchor the meaning.

WATCH FOR:
- Students who do not gesture - encourage everyone to try the spread shape.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    25: """SAY:
- Look at the picture.
- If the image depicts something splayed, say splayed.
- If not, stay silent.

DO:
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.

TEACHER NOTES:
Expected response: splayed. Confirm with the spread-shape gesture on reveal.

WATCH FOR:
- Students hesitating - say the meaning again before the next image.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    26: """SAY:
- Look at the picture.
- If the image depicts something splayed, say splayed.
- If not, stay silent.

DO:
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.

TEACHER NOTES:
Expected response: silence (not splayed). The image shows something closed, neat or together.

WATCH FOR:
- Students who say splayed because they want to participate - protect wait time and the silent option.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    27: """SAY:
- Look at the picture.
- If the image depicts something splayed, say splayed.
- If not, stay silent.

DO:
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.

TEACHER NOTES:
Expected response: splayed.

WATCH FOR:
- Students who name a different feature of the picture - bring focus back to the spread-out shape.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    28: """SAY:
- Which words can be paired with splayed?
- Read with me: spread, stretched, clumped, fan out, closed, extended.
- Show thumbs up for each word that pairs with splayed and thumbs down if it does not.

DO:
- Read each word clearly.
- Give 3 seconds for thumbs after each one.
- Scan the room.
- Reveal: spread, stretched, fan out, extended (yes); clumped, closed (no).

CFU CHECKPOINT:
Technique: Thumbs Up / Thumbs Down
Script:
- Thumbs up if the word pairs with splayed, thumbs down if it does not.
- Scan for: matching thumbs for each word in turn.
PROCEED:
- If most students are correct on each, move to scene comparisons.
PIVOT:
- Most likely misconception: students give thumbs up to closed because it sounds related.
- Reteach: splayed means spread apart. Closed is the opposite.
- Fresh re-check: thumbs up for spread, thumbs down for closed.

TEACHER NOTES:
Sorting CFU - reveals whether students hold the spread meaning firmly.

WATCH FOR:
- Students copying a neighbour - change order and pace.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    29: """SAY:
- Two scenes - which one describes something splayed?
- A pile of books, or clothes spread all over the floor.
- Ask: which one is spread out? Expected: clothes spread.

DO:
- Partner talk for 20 seconds.
- Cold call one pair.
- Reveal: clothes spread all over the floor.

TEACHER NOTES:
Spread is the key clue. A pile is clumped, not splayed.

WATCH FOR:
- Students who pick the pile because it has more items - return to the meaning: spread out.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    30: """SAY:
- Two scenes - which one describes something splayed?
- A finger pointing, or fingers spread apart.
- Show me with your hand which one matches splayed.

DO:
- Have students raise the gesture they think matches.
- Take a quick scan and reveal: fingers spread apart.

TEACHER NOTES:
The hand gesture is now embedded. Use it to confirm the answer.

WATCH FOR:
- Students still pointing one finger after the reveal - re-cue: splayed needs spread.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    31: """SAY:
- Look at the picture.
- Write a sentence using the word splayed to describe it.
- Watch my example: The Pokemon cards were splayed out in the person's hand.
- Now write your own.

DO:
- Give 60 seconds.
- Circulate and check 3-4 sentences.
- Share 1-2 examples.

TEACHER NOTES:
Independent sentence using the word. Accept any sentence where splayed correctly describes spread out.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use the model sentence frame: The ____ were splayed out.
EXTENDING PROMPT:
- Task: Use splayed in a sentence with another splayed pair (e.g. wings and feathers).

WATCH FOR:
- Students who write a sentence where splayed is in the wrong context - prompt them to use it for something spread.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    32: """SAY:
- I will read four sentences. For each one, decide - has the word splayed been used in the correct context?
- Thumbs up if yes, thumbs down if no.
- One: The shoes were side by side, splayed out.
- Two: The blanket was folded neatly and splayed on the bed.
- Three: The book was splayed open on the desk, showing all its pages.
- Four: The peacock's tail feathers were splayed in a vibrant fan.

DO:
- Read each sentence twice.
- Give 3 seconds for thumbs after each.
- Scan and reveal.

CFU CHECKPOINT:
Technique: Thumbs Up / Thumbs Down per sentence
Script:
- Thumbs up if splayed is used correctly, thumbs down if it is not.
- Scan for: thumbs down on 1 and 2, thumbs up on 3 and 4.
PROCEED:
- If most students are correct across all four, move on to ashamed.
PIVOT:
- Most likely misconception: students miss that side by side and folded neatly are the opposite of splayed.
- Reteach: splayed means spread apart, not neat or together.
- Fresh re-check: a fan opened wide - splayed or not? Thumbs up or down.

TEACHER NOTES:
End-of-word CFU - four-item version pushes for accuracy on context.

WATCH FOR:
- Students who keep thumbs sideways - require a clear up or down.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    33: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our next word is ashamed.
- Watch this first - I am thinking about the body language.
- If you are ashamed, you feel guilt or sorrow.
- In the story, the shooters ran away ashamed after hurting Mr Percival - so I picture them not making eye contact.
- Say the word with me... ashamed!

DO:
- Point to the image on the slide.
- Lower your eyes and slump your shoulders to model ashamed.
- Have students repeat the word twice.

TEACHER NOTES:
Feeling word. Body language clues (head down, shoulders slumped) help students recognise it.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Repeat the word and copy the head-down body language.
EXTENDING PROMPT:
- Task: Name a sentence using ashamed about a character in another story they have read.

SENSITIVITY ADVISORY:
- What it is: Some students may have recent experiences of feeling ashamed.
- Framing language: Keep examples external and gentle - stay with the shooters or the dancer, not personal student stories.
- Watch for: Students going quiet or looking upset.
- Protocol: If a student looks distressed, move on and check in privately afterwards.

WATCH FOR:
- Students who confuse ashamed with embarrassed - confirm they are close, but ashamed is stronger and includes guilt.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    34: """SAY:
- Read the sentence with me.
- I was ashamed because I got caught cheating.
- Ask: what feeling does ashamed show? Expected: guilt or sorrow.

DO:
- Choral read.
- Lower head as you say ashamed to model body language.

TEACHER NOTES:
Connect the feeling to a familiar context. Cheating is a clear, age-appropriate trigger for shame without being personal.

WATCH FOR:
- Students who laugh at the example - redirect: this is a serious feeling word.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    35: """SAY:
- I will read an example.
- If the example can be described as feeling ashamed, say ashamed.
- If not, stay silent.
- Example: feeling embarrassed and wanting to hide.

DO:
- Read the example.
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.

TEACHER NOTES:
Expected response: ashamed. Embarrassed and wanting to hide is the body language of shame.

WATCH FOR:
- Students who stay silent because they are unsure - reread and gesture the head-down shape.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    36: """SAY:
- I will read an example.
- If the example can be described as feeling ashamed, say ashamed.
- If not, stay silent.
- Example: feeling like you let your friends down.

DO:
- Read the example.
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.

TEACHER NOTES:
Expected response: ashamed. Letting friends down is a guilt feeling.

WATCH FOR:
- Students who say the wrong word - re-anchor: ashamed is guilt or sorrow.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    37: """SAY:
- I will read an example.
- If the example can be described as feeling ashamed, say ashamed.
- If not, stay silent.
- Example: feeling confident and proud of your achievements.

DO:
- Read the example.
- Give 5 seconds wait time.
- Cue: 3, 2, 1, say it.

TEACHER NOTES:
Expected response: silence (not ashamed). Confident and proud is the opposite feeling.

WATCH FOR:
- Students who say ashamed automatically - re-cue silence as a valid answer.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    38: """SAY:
- Choose: would or wouldn't.
- I would feel ashamed if I dropped my friend's phone, or I wouldn't.
- Show me with W or N on your whiteboard.
- Ask: why? Expected: a reason that fits guilt or sorrow.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Take 1-2 reasons.

TEACHER NOTES:
Answers may vary. Accept either if the reason matches the meaning of ashamed.

WATCH FOR:
- Students who give a reason that is not about guilt or sorrow - prompt: would you feel guilty?

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    39: """SAY:
- Choose: would or wouldn't.
- I would feel ashamed if I tried a new sport and I wasn't very good at it.
- Or I wouldn't feel ashamed.
- Show me W or N.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Take 1-2 reasons.

TEACHER NOTES:
Answers may vary. The growth-mindset answer is wouldn't feel ashamed - trying something new is brave. Accept either with a reason.

WATCH FOR:
- Students who feel ashamed of trying - gently model that trying is something to be proud of.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    40: """SAY:
- Choose: would or wouldn't.
- I would feel ashamed if I spilt my drink all over the kitchen table.
- Or I wouldn't.
- Show me W or N.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Take 1-2 reasons.

TEACHER NOTES:
Answers may vary. An accident is not usually shame - confirm shame is for things you did on purpose or knew were wrong.

WATCH FOR:
- Students who feel ashamed for accidents - help them name embarrassment instead.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    41: """SAY:
- Show me which word is unlike ashamed.
- Hold up 1 for embarrassed, 2 for confident, 3 for disgraced, 4 for humiliated.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which word is unlike ashamed? Show me on three.
- Scan for: fingers showing 2 (confident).
PROCEED:
- If most students are on 2, move to picture practice.
PIVOT:
- Most likely misconception: students pick humiliated because the word is unfamiliar.
- Reteach: embarrassed, disgraced and humiliated all sit near ashamed. Confident is the opposite.
- Fresh re-check: is proud like ashamed? Thumbs up or down.

TEACHER NOTES:
Synonym vs antonym CFU. Catches confusion with new words.

MISCONCEPTIONS:
- Misconception: humiliated and ashamed are opposites.
  Why: students do not know what humiliated means.
  Impact: students will reject correct synonyms in future tasks.
  Quick correction: explain humiliated means made to feel ashamed in front of others - same family.

WATCH FOR:
- Students mirroring a neighbour - reset wait time before scanning.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    42: """SAY:
- Look at the two people - A and B.
- Which person looks ashamed?
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.
- Take one student to explain the body-language clues.

TEACHER NOTES:
Confirm body-language clues - head down, shoulders slumped, eyes lowered.

WATCH FOR:
- Students choosing on smile alone - prompt them to look at posture, not just face.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    43: """SAY:
- Look at the two dogs - A and B.
- Which dog looks ashamed?
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Yes, dogs can show ashamed body language - eyes down, body low. Same clues as for people.

WATCH FOR:
- Students who get distracted by which dog is cuter - bring focus back to the body-language clues.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    44: """SAY:
- Complete the sentences.
- I was ashamed because I forgot ____.
- I was ashamed because I missed ____.
- Choose any reason that makes sense.

DO:
- Give 60 seconds for students to write both sentences.
- Circulate and check.
- Reveal example completions: my friend's birthday, the goal.

CFU CHECKPOINT:
Technique: Independent write, then cold call for answers
Script:
- Complete both sentences. Write your reason.
- Scan for: reasons that fit guilt or sorrow.
PROCEED:
- If most students have a sensible reason for both, move to the booklet task.
PIVOT:
- Most likely misconception: students write reasons that show embarrassment, not shame (e.g. tripped over).
- Reteach: ashamed is a stronger feeling - guilt or sorrow because of something you did.
- Fresh re-check: which fits better - I was ashamed because I lied, or I was ashamed because I sneezed?

TEACHER NOTES:
End-of-word CFU. Independent write surfaces if students confuse shame with embarrassment.

WATCH FOR:
- Students writing extreme or distressing reasons - redirect to lighter examples like forgetting or missing something.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    45: """SAY:
- Time to try this on your own.
- First, turn to the page titled Lesson 21: Vocabulary.
- Next, complete the booklet tasks for blazing, splayed and ashamed.
- Then, read your sentences back to a partner.

DO:
- Direct students to the correct booklet page.
- Circulate and check the first 2 to 3 answers for each student.
- Pull a small group if many are still unsure.

TEACHER NOTES:
You Do practice for the three vocabulary words. Aim for accuracy on meaning, not full sentences.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Complete only the meaning-match questions for now. Use the word wall on the board.
EXTENDING PROMPT:
- Task: Write one fresh sentence for each word that uses the meaning differently from the slide examples.

WATCH FOR:
- Students who skip definitions and only do the easy questions - direct them back to the word meanings.
- Students finishing fast - ask them to write a fresh sentence with each word.

Stage: Vocabulary You Do | VTLM: Independent Practice""",


    46: """SAY:
- Quick reset.
- Now we move to sentence-level writing.
- Today our focus is punctuating direct speech.

DO:
- Have whiteboards and textas ready.

TEACHER NOTES:
Brief divider.

WATCH FOR:
- Energy dipping - a quick stand-stretch can reset before the next chunk.

Stage: Sentence-Level Divider | VTLM: Engagement""",


    47: """SAY:
- Watch this first.
- Direct speech is when we write the exact words a character says.
- Punctuation helps the reader see who is speaking.
- I am going to read this example carefully: The men shouted with rage, "The ducks have gone!"
- I notice three parts. The speaker tag - the men shouted with rage - tells me who said it.
- The exact words - The ducks have gone - sit inside the inverted commas.
- The speaker tag is NOT inside the speech marks because those words are not being said out loud.
- Ask: where do the inverted commas go? Expected: around the words spoken.

DO:
- Point to the speaker tag on the slide.
- Point to the inverted commas around the spoken words.
- Read the example with extra punch on the words inside the speech marks.

TEACHER NOTES:
Anchor concept for the next twelve slides. The three labels (speaker tag, inverted commas, exact words spoken) repeat throughout - keep them clear.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Echo each label as the teacher points to it.
EXTENDING PROMPT:
- Task: Find one example of direct speech in pages already read and name the three parts.

MISCONCEPTIONS:
- Misconception: the speaker tag belongs inside the inverted commas.
  Why: students think everything to do with speaking is part of the speech.
  Impact: students put their tag inside the speech marks in their own writing.
  Quick correction: re-point to the example - the speaker tag sits OUTSIDE because the speaker did not say their own name.

WATCH FOR:
- Students who try to put the speaker tag inside the speech marks - re-show the example.

Stage: I Do | VTLM: Explicit Teaching""",


    48: """SAY:
- Watch how I punctuate this sentence step by step.
- Original: Storm Boy bellowed don't shoot Mr Percival.
- Step 1: I need to find the speaker tag. That is Storm Boy bellowed.
- Step 2: I am going to put a comma after the speaker tag.
- Step 3: I need a capital letter to start the words spoken aloud, and inverted commas around them.
- Step 4: I need to put the final punctuation mark BEFORE the closing inverted comma.
- Final: Storm Boy bellowed, "Don't shoot Mr Percival!"

DO:
- Point to each step on the slide as you say it.
- Read the original, then each version, then the final.
- Use a strong voice to model an exclamation mark.

TEACHER NOTES:
Core I Do. Students will follow this four-step sequence in We Do and You Do. Make every step visible on the board.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Echo each step number as the teacher names it. Hold up fingers to count 1, 2, 3, 4.
EXTENDING PROMPT:
- Task: Identify the speaking verb (bellowed) and explain why an exclamation mark fits.

MISCONCEPTIONS:
- Misconception: the final punctuation goes after the closing inverted comma.
  Why: students think punctuation belongs at the very end of the sentence on paper.
  Impact: incorrect punctuation in independent writing.
  Quick correction: re-show step 4 - inside the speech marks, not outside.

WATCH FOR:
- Students missing the comma between speaker tag and speech - point explicitly.
- Students forgetting the capital letter at the start of the spoken words - re-cue.

Stage: I Do | VTLM: Modelling""",


    49: """SAY:
- Look at the sentence: Hide-Away whispered, "He is badly hurt."
- I want you to find each part.
- Ask: where is the speaker tag? Expected: Hide-Away whispered.
- Ask: where are the inverted commas? Expected: around He is badly hurt.
- Ask: where are the exact words spoken? Expected: He is badly hurt.

DO:
- Read the sentence twice.
- Have students point to each feature on the slide as you name it.
- Reveal each feature one at a time.

TEACHER NOTES:
Identification practice. Students must be able to name speaker tag, inverted commas, and exact words spoken before they punctuate independently.

WATCH FOR:
- Students mixing up speaker tag and exact words - re-anchor: the words inside the speech marks are the spoken ones.

Stage: We Do | VTLM: Guided Practice""",


    50: """SAY:
- Same job - identify the parts.
- Storm Boy sobbed, "You're the best, best friend I ever had!"
- Ask: where is the speaker tag? The inverted commas? The exact words?
- Ask: what kind of sentence is this? Expected: an exclamation.

DO:
- Read the sentence twice.
- Reveal each feature.
- Confirm the exclamation mark on the reveal.

TEACHER NOTES:
Repeat identification routine, with the sentence-type clue added as a bridge to the next slide.

WATCH FOR:
- Students who treat the comma as part of the spoken words - point to the comma sitting outside the speech marks.

Stage: We Do | VTLM: Guided Practice""",


    51: """SAY:
- Three sentences with the punctuation mark missing.
- The men snapped, "It's that pelican again ___"
- Fingerbone inquired, "What can I do to help ___"
- Storm Boy stated, "I can see some men ___"
- I am going to look at the speaking verb each time - the verb tells me the sentence type.
- Ask: snapped means? Expected: an exclamation.
- Ask: inquired means? Expected: a question.
- Ask: stated means? Expected: a statement.

DO:
- Read each sentence aloud.
- Have students write the correct mark on their whiteboards.
- Reveal: !, ?, .
- Discuss the speaking verb each time.

TEACHER NOTES:
Speaking-verb-to-sentence-type rule. Every speaking verb has an emotional fingerprint.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Match each verb to its mark using a quick T-chart on the whiteboard.
EXTENDING PROMPT:
- Task: Add a fourth example using a different speaking verb (whispered, demanded, asked).

WATCH FOR:
- Students who guess marks without checking the verb - prompt: what does the verb tell us?

Stage: We Do | VTLM: Guided Practice""",


    52: """SAY:
- Read this sentence: Storm Boy screamed, "They've shot Mr Percival"!
- Has this been punctuated correctly? Show me thumbs up or thumbs down.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Cold call a non-volunteer to explain the error.
- Reveal: Storm Boy screamed, "They've shot Mr Percival!"

CFU CHECKPOINT:
Technique: Thumbs Up / Thumbs Down with cold-call explanation
Script:
- Is this sentence punctuated correctly? Thumbs up or down.
- Scan for: thumbs down (it is not correct).
PROCEED:
- If most students show thumbs down, move to step-by-step correction practice.
PIVOT:
- Most likely misconception: students think the exclamation mark belongs after the closing speech marks.
- Reteach: the final punctuation mark goes BEFORE the closing inverted comma, not after. Step 4 of our routine.
- Fresh re-check: He shouted, "Run away" - where does the exclamation mark go?

TEACHER NOTES:
Hinge CFU - exposes the most common error in direct-speech writing.

MISCONCEPTIONS:
- Misconception: the final punctuation goes after the closing inverted comma.
  Why: students think the mark closes the whole sentence.
  Impact: independent writing has the mark in the wrong place.
  Quick correction: re-show step 4 with finger - INSIDE the speech marks.

WATCH FOR:
- Students who only put thumbs sideways - require a clear answer.

Stage: CFU | VTLM: Formative Assessment""",


    53: """SAY:
- Let's correctly punctuate this sentence together.
- Original: Storm Boy sobbed poor Mr Percival.
- Step 1: Place a comma after the speaker tag.
- Step 2: Begin the speech with a capital letter and place inverted commas around the spoken words.
- Step 3: Use a final punctuation mark BEFORE the second inverted comma.
- I am going to listen to the speaking verb sobbed - that sounds sad. So an exclamation mark fits.
- Final: Storm Boy sobbed, "Poor Mr Percival!"

DO:
- Read the original aloud.
- Build each version on the board, step by step.
- Use a sad voice for sobbed so students choose the right end mark.

TEACHER NOTES:
We Do correction. The verb-to-mark rule is the key teaching move here.

WATCH FOR:
- Students who write a full stop - prompt them with the speaking verb: how does sobbed sound?

Stage: We Do | VTLM: Guided Practice""",


    54: """SAY:
- Punctuate this sentence: Storm Boy enquired will Mr percival be alright.
- Step 1: Comma after the speaker tag.
- Step 2: Capital letter and inverted commas.
- Step 3: Final punctuation mark before the second inverted comma.
- Ask: what mark fits enquired? Expected: a question mark.
- Final: Storm Boy enquired, "Will Mr Percival be alright?"

DO:
- Build the corrected version on the board.
- Double-check the capital P on Percival.
- Read the final sentence with rising intonation for the question.

TEACHER NOTES:
Question mark practice. Note: the original had a lowercase p on Percival - confirm it is a proper noun and needs a capital.

WATCH FOR:
- Students who miss the capital on Percival - point to proper-noun capitalisation.

Stage: We Do | VTLM: Guided Practice""",


    55: """SAY:
- Punctuate this sentence: Hide-Away asked would you like some breakfast.
- Same three steps.
- Ask: what mark fits asked? Expected: a question mark.
- Final: Hide-Away asked, "Would you like some breakfast?"

DO:
- Build the corrected sentence on the board.
- Read aloud with question intonation.

TEACHER NOTES:
Question mark again. Asked is a clear question signal.

WATCH FOR:
- Students adding extra punctuation inside the speech - remind them only one final mark per sentence.

Stage: We Do | VTLM: Guided Practice""",


    56: """SAY:
- Punctuate this sentence: Storm Boy cried Mr Percival Oh Mr Percival.
- This one needs commas inside the spoken words to show the pause.
- Watch how I add them: "Mr Percival, Oh, Mr Percival!"
- Final: Storm Boy cried, "Mr Percival, Oh, Mr Percival!"

DO:
- Build the corrected version.
- Read it aloud with feeling so the commas are heard.

TEACHER NOTES:
Slightly harder - students must add commas INSIDE the speech, not just at the end. The exclamation mark fits the speaking verb cried.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Read aloud first - if you pause, that is where a comma goes.
EXTENDING PROMPT:
- Task: Add another short interjection inside the speech (e.g. Oh no) and punctuate it correctly.

WATCH FOR:
- Students who try to use full stops inside the speech - confirm commas hold the pauses.

Stage: We Do | VTLM: Guided Practice""",


    57: """SAY:
- The speaker tag is Storm Boy sobbed.
- Ask: which option fits best with sobbed? Show me 1, 2, 3 or 4.
- Option 1: "Poor Mr Percival!"
- Option 2: "There you are, Mr Percival!"
- Option 3: "Where are the hunters?"
- Option 4: "Look, there's the pelican!"

DO:
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which sentence fits best with sobbed? On three, show me.
- Scan for: fingers showing 1 (Poor Mr Percival!).
PROCEED:
- If most students are on 1, move to picture-stimulus writing.
PIVOT:
- Most likely misconception: students pick option 2 because it has an exclamation mark.
- Reteach: sobbed shows sadness. Option 1 fits sadness. Option 2 sounds happy.
- Fresh re-check: would you sob "Where are the hunters?" - thumbs up or down.

TEACHER NOTES:
Verb-meaning CFU. Students must match the speaking verb to the emotional fit, not just the punctuation.

WATCH FOR:
- Students who guess on punctuation alone - bring focus back to the meaning of the speaking verb.

Stage: CFU | VTLM: Formative Assessment""",


    58: """SAY:
- Use the picture to punctuate the direct speech.
- Watch how I think about this.
- The speaker is Storm Boy. The exact words spoken are: You're the best friend I've ever had.
- I am going to choose a speaking verb that fits the picture - remarked seems calm and warm.
- Step 1: Comma after the speaker tag.
- Step 2: Capital letter and inverted commas.
- Step 3: Final punctuation mark before the closing inverted comma.
- Final: Storm Boy remarked, "You're the best friend I've ever had!"

DO:
- Point to the picture and ask students what Storm Boy might be feeling.
- Build the punctuated sentence step by step.
- Confirm the exclamation mark matches the strong feeling.

TEACHER NOTES:
Picture-stimulus modelling. Students see a teacher pick the speaking verb, then apply the four-step routine.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use the suggested speaker tag and exact words from the slide. Just punctuate.
EXTENDING PROMPT:
- Task: Write the same direct speech with a different speaker (e.g. Hide-Away).

WATCH FOR:
- Students who go for a full stop - ask: how does this picture make Storm Boy feel?

Stage: I Do | VTLM: Modelling""",


    59: """SAY:
- Look at the picture.
- Think about what Storm Boy might say in this moment.
- Write a sentence with direct speech.
- First: speaker tag followed by a comma.
- Next: exact words inside inverted commas, starting with a capital.
- Then: final punctuation mark before the closing inverted comma.
- For example: Storm Boy gasped, "Watch out Mr Percival!"

DO:
- Show the picture for 10 seconds.
- Give 90 seconds for students to write a direct speech sentence.
- Take 2-3 examples from the room.

TEACHER NOTES:
We Do moving toward independent. Accept any speaker tag and spoken words that fit the picture, as long as the four-step routine is followed.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a sentence frame: ____ said, "____."
EXTENDING PROMPT:
- Task: Add a second sentence showing what Mr Percival might say back.

WATCH FOR:
- Students who forget the comma after the speaker tag - check this first.
- Students who place the final mark outside the speech marks - re-anchor step 4.

Stage: We Do | VTLM: Guided Practice""",


    60: """SAY:
- New picture, same job.
- Write a sentence with direct speech.
- First, speaker tag, comma. Next, capital, inverted commas, exact words. Then, final mark before the closing inverted comma.
- For example: Storm Boy remarked, "You'll be okay, Mr Percival."

DO:
- Show the picture.
- Give 90 seconds for students to write.
- Circulate and check 4-5 sentences.
- Share 1-2 examples.

TEACHER NOTES:
Second picture-stimulus. By now most students should be applying the four steps confidently.

WATCH FOR:
- Students who match a speaking verb that does not fit the picture - prompt them to choose a verb that matches the mood.

Stage: We Do | VTLM: Guided Practice""",


    61: """SAY:
- Which sentence is punctuated correctly? Show me 1, 2, 3 or 4.
- Option 1: "The men shouted with rage, The ducks have gone!"
- Option 2: "Storm Boy bellowed," Don't shoot Mr Percival.
- Option 3: Storm Boy screamed, "They've shot Mr Percival!"
- Option 4: Storm Boy sobbed "They've shot Mr Percival!"

DO:
- Read each option aloud.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which sentence is punctuated correctly? On three, show me.
- Scan for: fingers showing 3.
PROCEED:
- If most students are on 3, move to the booklet task.
PIVOT:
- Most likely misconception: students pick option 4 - they miss the missing comma after sobbed.
- Reteach: every direct-speech sentence needs a comma between the speaker tag and the inverted commas.
- Fresh re-check: She whispered "Be quiet" - what is missing?

TEACHER NOTES:
Pre-You-Do hinge CFU. Strong response to 3 signals readiness.

WATCH FOR:
- Students who pick option 1 - the speech words are missing inverted commas.
- Students who pick option 2 - the comma is in the wrong place.

Stage: CFU | VTLM: Formative Assessment""",


    62: """SAY:
- Time to apply this on your own.
- First, turn to the page titled Lesson 21: Sentence level writing.
- Next, complete the tasks on punctuating direct speech.
- Then, read each sentence back aloud after punctuating.

DO:
- Direct students to the correct booklet page.
- Circulate and check the first 2 to 3 sentences for each student.
- Pull a small group if many students still place the final mark outside the inverted commas.

TEACHER NOTES:
You Do practice. The most common errors are missing comma after speaker tag and final mark outside the speech marks.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a four-step checklist card alongside the booklet for each sentence.
EXTENDING PROMPT:
- Task: After completing the booklet, write three of your own direct-speech sentences using different speaking verbs.

WATCH FOR:
- Students who skip the comma after the speaker tag - flag this first.
- Students who use lowercase to start the speech - re-cue capital letter.
- Students finishing fast - ask them to swap and check a partner's work.

Stage: You Do | VTLM: Independent Practice""",
}


# ============================================================
# LESSON 22 - Single paragraph outline (SPO)
# ============================================================
LESSON_22_NOTES = {

    9: """SAY:
- Quick reset before we read.
- Today we read pages 72 to 73 of Storm Boy.
- These are emotional pages - listen carefully.

DO:
- Have your novel pre-marked with the chosen pause points.
- Settle students for read aloud.

TEACHER NOTES:
Read-aloud setup. Mode of reading is teacher choice.

SENSITIVITY ADVISORY:
- What it is: These pages cover Mr Percival's death and Storm Boy's grief.
- Framing language: Acknowledge the pages are sad. Use a calm, steady voice.
- Watch for: Students with recent loss may find this hard.
- Protocol: Allow a quiet exit if needed. Have a pre-agreed signal with students who may struggle.

WATCH FOR:
- Students who go very quiet - check in privately at a natural pause.

Stage: Text-Level Reading | VTLM: Engagement""",


    10: """SAY:
- We are going to read pages 72 to 73.
- Listen carefully and watch the pictures.
- I will pause to ask questions.

DO:
- Read pages 72 to 73 of Storm Boy aloud, or run the chosen reading mode.
- Pause at the points you have pre-marked.
- Take 1-2 responses at each pause.
- Gloss any incidental vocabulary in one short sentence.

TEACHER NOTES:
Pause points come from the Literature Study Guide. These pages are emotionally heavy - hold the silences and do not rush.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Listening partner - the supported student listens, their partner whispers a one-word feeling at each pause.
EXTENDING PROMPT:
- Task: Note one literary device (simile, imagery, repetition) the author uses to convey grief.

SENSITIVITY ADVISORY:
- What it is: Mr Percival dies in this section.
- Framing language: Read calmly and steadily. Do not act out grief.
- Watch for: Students who go quiet or look upset.
- Protocol: Pause briefly, then continue. Check in privately afterwards.

WATCH FOR:
- Students who look upset - keep reading calmly and check in privately afterwards.

Stage: Text-Level Reading | VTLM: Explicit Teaching""",


    11: """SAY:
- Three new words from today's pages.
- Watch and listen first, then we practise together.

DO:
- Have whiteboards and textas ready.
- Display the first vocabulary word.

TEACHER NOTES:
Brief divider.

WATCH FOR:
- Students who need wait time - protect it across the next set of slides.

Stage: Vocabulary Divider | VTLM: Engagement""",


    12: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our first word is flung.
- Watch this first - I am going to think aloud.
- To fling something is to throw it, usually forcefully or angrily.
- In the story, Storm Boy flung himself down on his bunk and sobbed softly. So flung is a strong, overwhelmed motion - he was so upset he did not just lie down.
- Say the word with me... flung!

DO:
- Point to the image on the slide.
- Use a quick, sharp throwing gesture as you say flung.
- Have students repeat the word twice.

TEACHER NOTES:
Connect the strong-throw motion to the story moment - Storm Boy is overwhelmed when he flings himself down.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Repeat the word and copy the throwing gesture.
EXTENDING PROMPT:
- Task: Identify one other moment in Storm Boy where someone could fling something.

WATCH FOR:
- Students who think flung means a gentle throw - confirm it is forceful or careless.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    13: """SAY:
- Read the sentence with me.
- Alex flung the ball across the field to his teammate, hoping for a quick goal.
- Ask: what does flung tell us about how Alex threw the ball? Expected: forcefully and quickly.

DO:
- Choral read.
- Use the strong-throw gesture together.

TEACHER NOTES:
Use the gesture to anchor the meaning before independent practice.

WATCH FOR:
- Students who do not gesture - encourage everyone to copy the motion.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    14: """SAY:
- Two sentences. Which one is correct?
- A: In a fit of frustration, Jane flung her backpack onto the couch, scattering her books everywhere.
- B: He gently flung the delicate vase onto the table.
- Show me A or B.

DO:
- Read both clearly.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.
- Reveal: A.

TEACHER NOTES:
Expected response: A. Flung implies force, so it does not fit with gently or delicate.

WATCH FOR:
- Students who pick B because gently sounds nicer - re-anchor the meaning.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    15: """SAY:
- Think of a time you have flung something.
- Stay silent and act it out.
- I might choose two volunteers to show their flinging.

DO:
- Give 30 seconds of silent thinking.
- Invite 2 volunteers to act.
- Ask the class: what was the volunteer flinging?

TEACHER NOTES:
Movement task. Force, frustration or carelessness should be visible in the motion.

WATCH FOR:
- Students who throw very lightly - prompt: flung is forceful.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    16: """SAY:
- Show me which option fits best.
- When you fling something, you...
- Hold up 1 for throw it delicately, 2 for throw it forcefully and angrily, 3 for throw it with careful aim, 4 for throw it carelessly.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which option fits best with fling? Show me on three.
- Scan for: fingers showing 2 or 4 (forcefully and angrily, OR carelessly).
PROCEED:
- If most students are on 2 or 4, move on to active sentence writing.
PIVOT:
- Most likely misconception: students pick 1 (delicately).
- Reteach: flung means strong, careless or angry - not delicate or careful.
- Fresh re-check: would you fling a feather? Thumbs up or down.

TEACHER NOTES:
End-of-word CFU. Both 2 and 4 are correct because flung covers both meanings.

WATCH FOR:
- Students who pick 3 (careful aim) - confirm flung is not careful.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    17: """SAY:
- Write a sentence using the word flung.
- Pick something you might fling.
- Watch my example: During the snowball fight, Sophie flung a snowball at her friend, laughing as it exploded into powder.

DO:
- Give 90 seconds.
- Circulate and check 4-5 sentences.
- Take 1-2 examples to share.
- Reveal the example sentence on click.

TEACHER NOTES:
Active practice. Accept any sentence where flung shows a forceful or careless throw.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a sentence frame: ____ flung the ____.
EXTENDING PROMPT:
- Task: Use flung in a sentence that shows frustration without using the word frustrated.

WATCH FOR:
- Students who use flung for a delicate action - prompt them to choose a different verb.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    18: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our next word is sobbed.
- Watch this first - I am going to think aloud.
- If someone sobs, they cry with so much emotion that they cannot control the sound of their breathing.
- I notice this word is right next to flung in the story - Storm Boy flung himself down and sobbed. The two words tell us how big his feelings are.
- Say the word with me... sobbed!

DO:
- Point to the image on the slide.
- Use a soft, breathy voice for sobbed so students hear the feeling.
- Have students repeat the word twice.

TEACHER NOTES:
First meeting with sobbed. The word ties directly to the chapter we just read.

SENSITIVITY ADVISORY:
- What it is: Some students may have recent grief experiences.
- Framing language: Use the story example. Avoid asking for personal stories of loss.
- Watch for: Students going quiet or tearful.
- Protocol: Allow a student to step out for a moment. Check in privately.

WATCH FOR:
- Students who confuse sobbed with cried - confirm sobbed is harder, with breath caught.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    19: """SAY:
- Read the sentence with me.
- After falling and scraping her knee, Mia sobbed loudly until her mum came over to comfort her.
- Ask: why might Mia sob? Expected: the pain and shock made it hard to control her breathing.

DO:
- Choral read.
- Use a slightly sad tone.

TEACHER NOTES:
Connecting sobbed to a familiar context that students can imagine.

WATCH FOR:
- Students who laugh - redirect: this is a serious feeling word.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    20: """SAY:
- When you sob, you cry with so much emotion that you cannot control your breathing.
- Think of a time you might have sobbed - keep it gentle.
- For example: when I thought I lost my parents at the park.
- Tell your partner an example.

DO:
- Give 30 seconds of partner talk.
- Take 1-2 voluntary responses only.

TEACHER NOTES:
Personal connection. Do not push reluctant students to share. Light examples are fine - dropped ice-cream, missed friends, lost-in-the-park feeling.

SENSITIVITY ADVISORY:
- What it is: Personal sharing about crying.
- Framing language: Keep examples light.
- Watch for: Students who share something heavy.
- Protocol: Listen, acknowledge, redirect, follow up after the lesson.

WATCH FOR:
- Students who share something heavy - thank them and follow up privately.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    21: """SAY:
- Build a sentence using the word sobbed.
- Sentence frame: Siara sobbed uncontrollably when ____, watching it ____.
- Watch my example: Siara sobbed uncontrollably when she dropped her ice cream on the hot cement, watching it melt slowly.

DO:
- Read the example.
- Give 60 seconds for students to build their own.
- Take 1-2 examples to review.

TEACHER NOTES:
Sentence frame supports students who need a starter.

WATCH FOR:
- Students who use sobbed for happy moments - re-anchor: sobbed is sad, not happy.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    22: """SAY:
- Which words have a similar meaning to sobbed?
- Read with me: cried, frowned, wept, bawled, sniffled, wailed.
- Show thumbs up if it is similar, thumbs down if not.

DO:
- Read each word.
- Give 3 seconds for thumbs.
- Reveal: cried, wept, bawled, wailed (yes); frowned, sniffled (sniffled is borderline - lighter than sobbed).

TEACHER NOTES:
Cried, wept, bawled and wailed are all strong-crying synonyms. Frowned is a face only - no crying. Sniffled is too light to count.

WATCH FOR:
- Students unsure on bawled or wept - quick gloss: same family of strong crying.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    23: """SAY:
- Two children - A and B.
- Which child is sobbing?
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.
- Take a non-volunteer to explain the body-language clues.

TEACHER NOTES:
Body language - shoulders shaking, mouth open, hands to face. These are the visual clues.

WATCH FOR:
- Students who choose on facial expression alone - prompt them to look at the whole body.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    24: """SAY:
- Sentence building. Read each one with me.
- She was sobbing, but her friend quickly brought her a new balloon.
- She was sobbing, so her dad gave her a big hug.
- She was sobbing because she accidentally let go of her balloon.
- Ask: what does each connector tell us? Expected: but flips the idea, so shows a result, because gives the reason.

DO:
- Read each sentence aloud.
- Point to each connector word.

CFU CHECKPOINT:
Technique: Class discussion plus targeted question
Script:
- Which connector shows a reason? Cold call.
- Scan for: students saying because.
PROCEED:
- If students name the connectors correctly, move to clutching.
PIVOT:
- Most likely misconception: students think but and so are interchangeable.
- Reteach: but flips the idea, so is the result, because is the reason.
- Fresh re-check: She was sobbing ____ she missed her friend - which connector?

TEACHER NOTES:
End-of-word CFU using the connector grammar - sets up sentence-construction work later.

WATCH FOR:
- Students who only read the sentences without thinking about the connectors - prompt them to focus on the joining word.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    25: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our next word is clutching.
- Watch this first - I am going to think aloud.
- If you clutch something, you hold it tightly.
- In the story, all night Storm Boy lay clutching his cold wet pillow. So clutching shows me he could not let go - the feeling was too strong.
- Say the word with me... clutching!

DO:
- Point to the image on the slide.
- Make a tight fist to model clutching.
- Have students repeat the word twice and copy the gesture.

TEACHER NOTES:
First meeting with clutching. The strong-grip gesture is the visual anchor.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Repeat the word twice and copy the tight-fist gesture.
EXTENDING PROMPT:
- Task: Compare clutching with grasping - same or different?

WATCH FOR:
- Students who confuse clutching with holding - confirm clutching is tighter, often when scared or upset.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    26: """SAY:
- Read the sentence with me.
- Clutching his new book tightly, Jake eagerly found a cozy corner to start reading.
- Ask: why is Jake clutching the book? Expected: he loves it and does not want to let it go.

DO:
- Choral read.
- Use the tight-fist gesture.

TEACHER NOTES:
Connect the meaning to a happy context this time. Clutching can also be from excitement.

WATCH FOR:
- Students who think clutching is always sad - clarify it can be from any strong feeling.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    27: """SAY:
- Two sentences. Which one is correct?
- A: She was gently clutching the newborn baby, careful not to drop it.
- B: The little girl was clutching her favourite teddy bear tightly as she went to bed.
- Show me A or B.

DO:
- Read both clearly.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Expected response: B. Clutching is a tight grip, so it does not fit with gently.

WATCH FOR:
- Students who pick A because gently sounds caring - clarify clutching is firm, not gentle.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    28: """SAY:
- When you clutch something, you hold it tightly.
- Ask: when have you clutched something tightly? Expected: a teddy at night, parents' hand crossing the road, balloon string.
- Tell your partner one example.

DO:
- Give 30 seconds of partner talk.
- Take 1-2 examples.

TEACHER NOTES:
Personal connection. All age-appropriate examples are fine.

WATCH FOR:
- Students stuck for an idea - prompt with crossing the road, holding a teddy or carrying a hot drink.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    29: """SAY:
- Two images - A and B.
- Which one shows clutching?
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Look for the tight grip - that is the visual clue.

WATCH FOR:
- Students who pick the picture with the more complete object - return them to the grip itself.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    30: """SAY:
- Show me which word fits best with clutching.
- Hold up 1 for touching, 2 for gripping, 3 for patting, 4 for crushing.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which word fits best with clutching? Show me on three.
- Scan for: fingers showing 2 (gripping).
PROCEED:
- If most students are on 2, move to the booklet task.
PIVOT:
- Most likely misconception: students pick 4 (crushing) because both are tight.
- Reteach: clutching is a tight hold but it does not damage. Crushing breaks the thing.
- Fresh re-check: would you crush a teddy bear or clutch it?

TEACHER NOTES:
End-of-word CFU. Distinguishes tight grip (clutching) from harmful pressure (crushing).

WATCH FOR:
- Students mirroring a fast neighbour - keep wait time firm.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    31: """SAY:
- Time to apply this on your own.
- First, turn to the page titled Lesson 22: Vocabulary.
- Next, complete the booklet tasks for flung, sobbed and clutching.
- Then, read your sentences back to a partner.

DO:
- Direct students to the booklet page.
- Circulate and check the first 2 to 3 answers per student.
- Pull a small group if needed.

TEACHER NOTES:
You Do practice for the three vocabulary words.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Complete only the meaning-match questions for now. Use the word wall.
EXTENDING PROMPT:
- Task: Write one fresh sentence for each word in the context of Storm Boy.

WATCH FOR:
- Students writing definitions instead of the booklet answers - direct them back to the prompts.
- Students finishing fast - ask them to write a sentence with each word.

Stage: Vocabulary You Do | VTLM: Independent Practice""",


    32: """SAY:
- Quick reset.
- Now we move to sentence-level writing.
- Today we are going to summarise the story so far using a Single Paragraph Outline.

DO:
- Have whiteboards and textas ready.

TEACHER NOTES:
Brief divider.

WATCH FOR:
- Energy dipping - quick stand-stretch can help reset.

Stage: Sentence-Level Divider | VTLM: Engagement""",


    33: """SAY:
- Watch this first - I am thinking about why we summarise.
- A summary is a brief statement. It does not include extra details - just the main ideas.
- Summaries help us in four ways - boost comprehension, focus on the main idea and supporting details, retain information, and analyse information.
- Ask: which one matters most for your learning? Turn and tell your partner.

DO:
- Read each bullet on the slide.
- Give 20 seconds of partner talk.
- Take 1-2 responses.

TEACHER NOTES:
I Do for why we summarise. Keep it short - the point is to motivate the SPO work that follows.

WATCH FOR:
- Students who do not see the value - link it to retelling the chapter without retelling every detail.

Stage: I Do | VTLM: Explicit Teaching""",


    34: """SAY:
- Let's summarise what happened in this section of the story.
- Use the question prompts to guide you - who, what doing, when, where, why, how.
- Ask: where are they? What does Storm Boy do to comfort Mr Percival? Why couldn't Hide-Away help? How does Storm Boy feel about losing Mr Percival? What characters are in this section?
- Turn and tell your partner one main idea.

DO:
- Read the question prompts on the slide.
- Give 30 seconds of partner talk.
- Take 2-3 responses.
- Keep the focus on main ideas, not every detail.

TEACHER NOTES:
We Do for the summary discussion. Cap responses at one or two sentences each.

WATCH FOR:
- Students who retell every detail - prompt them to give just the main idea.

Stage: We Do | VTLM: Guided Practice""",


    35: """SAY:
- Big ideas from the story so far.
- Read with me - Storm Boy saves three baby pelicans... he nurses Mr Percival back to health... after letting them go, Mr P returns... Mr P and Storm Boy become inseparable... Mr P learns to play fetch... Mr P's skills save the crew's lives... the shooters return and Mr P warns the ducks... Mr P gets shot... Storm Boy comforts Mr P until he passes away.
- Ask: which big idea will we focus on first?

DO:
- Choral read each big idea.
- Pause at each one and select with the class.

TEACHER NOTES:
Knowledge-anchor slide. Each big idea is a candidate main idea for the SPO.

WATCH FOR:
- Students who get lost in the details of each event - bring focus back to the headline.

Stage: We Do | VTLM: Guided Practice""",


    36: """SAY:
- A Single Paragraph Outline - or SPO - puts your writing in an order that makes it easy for readers to understand.
- Today we will build an SPO together to summarise key ideas in Storm Boy.

DO:
- Read the slide aloud as choral reading.
- Make eye contact across the room as you say SPO so the term sticks.

TEACHER NOTES:
Introduce SPO as the planning tool for paragraph writing.

WATCH FOR:
- Students who confuse SPO with a finished paragraph - confirm it is the plan, not the writing.

Stage: I Do | VTLM: Explicit Teaching""",


    37: """SAY:
- A paragraph has three parts.
- Topic sentence - expresses the main idea of the paragraph.
- Supporting details - they relate to and support the topic sentence.
- Concluding sentence - summarises the paragraph, echoing the topic sentence without repeating it.
- Say with me: topic, supporting details, concluding.

DO:
- Point to each part as you name it.
- Choral repeat the three parts.

TEACHER NOTES:
Second I Do anchor. Students will use these labels for the rest of the lesson.

WATCH FOR:
- Students who muddle topic sentence and concluding sentence - clarify the topic introduces, the concluding wraps up.

Stage: I Do | VTLM: Explicit Teaching""",


    38: """SAY:
- What would be a suitable topic sentence for the main idea: Mr Percival warns the ducks?
- A topic sentence introduces the main idea and tells the reader who, what doing, when, where, why or how.
- Turn and tell your partner one possible topic sentence.

DO:
- Give 30 seconds of partner talk.
- Take 2 responses.
- Co-build the best version on the board.

TEACHER NOTES:
Modelled topic sentence writing. Aim for: Mr Percival heroically warned the ducks of the shooters' approach.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a sentence frame: Mr Percival ____ the ducks ____.
EXTENDING PROMPT:
- Task: Write two different topic sentences using different verbs and compare which is stronger.

WATCH FOR:
- Students who write a full paragraph - re-cue: just one topic sentence, not the whole paragraph.

Stage: We Do | VTLM: Guided Practice""",


    39: """SAY:
- Watch this first.
- Now I am going to find supporting details for the topic sentence.
- I will read an excerpt and pull out details that show how Mr Percival warned the ducks.
- I am looking for actions in the text that match the topic sentence - that is the criterion.

DO:
- Display the next slide with the excerpt.
- Talk through your thinking before reading - what am I looking for?

TEACHER NOTES:
I Do for finding supporting details. Make the search criterion explicit before reading.

WATCH FOR:
- Students who do not know what to look for - re-state the criterion: details about how Mr P warns the ducks.

Stage: I Do | VTLM: Modelling""",


    40: """SAY:
- Listen as I read the extract from page 67.
- I am looking for details that describe Mr Percival warning the ducks.
- I notice three details - he stared rudely until he was chased away, he swam annoyingly near their hidden boats, he flew round and round their hiding places.
- Each detail describes how he warned the ducks.

DO:
- Read the extract aloud.
- Underline or circle the three details on the slide as you find them.
- Use a clear think-aloud voice.

TEACHER NOTES:
Reveal each underlined detail in turn so students can see the thinking process.

WATCH FOR:
- Students who try to copy the whole extract - prompt them to pick only the details that match the topic sentence.

Stage: I Do | VTLM: Modelling""",


    41: """SAY:
- Watch how I convert each detail into note form for the SPO.
- Stared rudely - chased away.
- Swam near boats - made known.
- Flew around hiding place to warn ducks.
- I am using key words and short forms, not full sentences.

DO:
- Read the notes on the slide.
- Point to each abbreviation - the arrow means leads to. The number 2 is short for to.

TEACHER NOTES:
I Do for note form. Sets up the convention students will use - notes, not full sentences.

WATCH FOR:
- Students who write full sentences - remind them notes are quicker to use later.

Stage: I Do | VTLM: Modelling""",


    42: """SAY:
- Now we add a concluding sentence.
- A concluding sentence summarises the paragraph, echoing the topic sentence but not repeating it.
- For example: Through his clever and persistent actions, Mr Percival made sure the ducks had a chance to escape.
- Tell your partner another way to conclude this paragraph.

DO:
- Read the topic sentence and supporting details aloud first.
- Give 30 seconds of partner talk.
- Take 1-2 ideas.
- Co-build the concluding sentence on the board.

TEACHER NOTES:
Completes the first SPO model.

WATCH FOR:
- Students who repeat the topic sentence word for word - prompt them to say it differently.

Stage: We Do | VTLM: Guided Practice""",


    43: """SAY:
- When planning a summary paragraph, you...
- Show me 1, 2, 3 or 4.
- Option 1: write the points in chronological order.
- Option 2: focus on one main idea.
- Option 3: add your own opinion.
- Option 4: write about everything you can remember.

DO:
- Read each option.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- When planning a summary, what do you do? Show me on three.
- Scan for: fingers showing 2 (focus on one main idea).
PROCEED:
- If most students are on 2, move to the second SPO build.
PIVOT:
- Most likely misconception: students pick 1 (chronological order) - that fits a recount, not a summary.
- Reteach: a summary picks one main idea - it does not retell every step in order.
- Fresh re-check: a summary is - one main idea, OR everything I can remember? Choose one.

TEACHER NOTES:
Hinge CFU before the second SPO build.

MISCONCEPTIONS:
- Misconception: a summary should include everything you can remember.
  Why: students confuse summary with full retell.
  Impact: their booklet summaries become long, unfocused recounts.
  Quick correction: a summary is short and focused on the main idea only.

WATCH FOR:
- Students who pick 4 - re-anchor: summary means short.

Stage: CFU | VTLM: Formative Assessment""",


    44: """SAY:
- Now we will build an SPO for a different main idea.
- Main idea: Hide-Away comforts Storm Boy.
- A topic sentence introduces the main idea and tells the reader who, what doing, when, where, why or how.
- We will read an excerpt to find supporting details.

DO:
- Read the main idea aloud.
- Frame the next two slides as the source text.
- Set students up to listen for the supporting details.

TEACHER NOTES:
We Do build. The main idea is preselected so students focus on detail extraction.

WATCH FOR:
- Students drifting off - bring them back with the search criterion: how does Hide-Away comfort Storm Boy?

Stage: We Do | VTLM: Guided Practice""",


    45: """SAY:
- Listen as I read the extract from pages 72 to 73.
- I am looking for details that show how Hide-Away comforts Storm Boy.
- Hide-Away didn't light the lantern... the three of them stayed on in front of the fireplace... and at nine o'clock Mr Percival died... only then did Hide-Away move... gently took Mr Percival from Storm Boy... Storm Boy gave him up... he flung himself down on his bunk and sobbed softly... Hide-Away came over and put a hand on his shoulder.

DO:
- Read the extract slowly.
- Pause at the points where Hide-Away acts.
- Ask: what is Hide-Away doing here?

TEACHER NOTES:
Pause at action points so students can name them. Supporting details: stayed silently, took Mr P gently, put a hand on Storm Boy's shoulder.

SENSITIVITY ADVISORY:
- What it is: Mr Percival dies in this extract.
- Framing language: Read with calm, steady voice. Do not act out grief.
- Watch for: Students who go quiet or look upset.
- Protocol: Hold the silence, then move on. Check in privately afterwards.

WATCH FOR:
- Students who struggle emotionally - allow a quiet moment before moving on.

Stage: We Do | VTLM: Guided Practice""",


    46: """SAY:
- Listen to the second part of the extract.
- It's right that you should cry for Mr Percival for a while, he said... but don't keep on brooding... but why did they shoot Mr Percival? He wasn't hurting anyone... in the world there will always be men who are cruel, just as there will always be men who are lazy or stupid or wise or kind... today you've seen what cruel and stupid men can do... he pulled a blanket over Storm Boy and said quietly, now try to get some sleep.

DO:
- Read the extract slowly.
- Point to each new comfort action - words of permission, words of explanation, putting the blanket over Storm Boy.

TEACHER NOTES:
Hide-Away comforts with words and actions. The supporting details extend - permission to grieve, explanation of the world, tucking Storm Boy in.

WATCH FOR:
- Students who try to summarise too quickly - direct them back to the actions in the text.

Stage: We Do | VTLM: Guided Practice""",


    47: """SAY:
- Now let's list the supporting details in note form together.
- Stayed with Storm Boy silently.
- Storm Boy sobbed - Hide-Away put a hand on his shoulder.
- Explained the cruel world to Storm Boy.
- Notice - we are using abbreviations and short forms.

DO:
- Read each note on the slide.
- Point to the abbreviations - w/ for with, S Boy for Storm Boy, H-A for Hide-Away.

TEACHER NOTES:
We Do for note form. Abbreviations should match the ones used earlier.

WATCH FOR:
- Students who do not understand the abbreviations - clarify each one out loud.

Stage: We Do | VTLM: Guided Practice""",


    48: """SAY:
- Now we add the topic sentence to the SPO.
- A topic sentence introduces the main idea.
- For example: As the sun set, Hide-Away lovingly comforted Storm Boy as he grieved.
- Cold call - what is another way we could write the topic sentence?

DO:
- Read the example topic sentence.
- Give 30 seconds of think-time.
- Cold call 2 non-volunteers.
- Co-build a strong topic sentence on the board.

TEACHER NOTES:
Use cold call to involve students who do not usually volunteer. Aim for varied openings.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a frame: As ____, Hide-Away comforted Storm Boy.
EXTENDING PROMPT:
- Task: Write a topic sentence that hints at the emotional tone without using the word sad.

WATCH FOR:
- Students who give a fragment - check it has a who and a what doing.

Stage: We Do | VTLM: Guided Practice""",


    49: """SAY:
- Add a concluding sentence to the SPO.
- A concluding sentence summarises the paragraph, echoing the topic sentence but not repeating it.
- For example: Hide-Away, a loving father, did everything he could to help Storm Boy deal with Mr Percival's passing.
- Read with me.

DO:
- Read the concluding sentence aloud.
- Point to the topic sentence at the top of the SPO.
- Show that the conclusion echoes but does not repeat.

TEACHER NOTES:
Full SPO is now visible - topic, supporting details, concluding. Use this as the model for the You Do.

WATCH FOR:
- Students copying the topic sentence as the conclusion - re-anchor: echo, do not repeat.

Stage: We Do | VTLM: Guided Practice""",


    50: """SAY:
- Read the SPO on the slide carefully.
- The main idea is: Storm Boy runs to find Hide-Away.
- One supporting detail does not match the main idea.
- Find the one that does not fit and explain why.

DO:
- Give 30 seconds reading time.
- Cold call one student to identify the wrong detail.
- Reveal: supporting detail 2 focuses on what the shooters did - it does not match Storm Boy running to find Hide-Away.

CFU CHECKPOINT:
Technique: Independent read, then cold call
Script:
- Which supporting detail does not fit the main idea? Cold call.
- Scan for: students naming detail 2.
PROCEED:
- If students name detail 2 with the right reason, move to the You Do.
PIVOT:
- Most likely misconception: students pick a detail because it sounds less interesting.
- Reteach: each detail must match the main idea. Detail 2 is about the shooters, not Storm Boy running.
- Fresh re-check: which is the main idea here? And does each detail match it?

TEACHER NOTES:
Hinge CFU before You Do. Tests whether students can spot off-topic supporting details.

WATCH FOR:
- Students who pick a detail without giving a reason - require a reason.

Stage: CFU | VTLM: Formative Assessment""",


    51: """SAY:
- This is your scaffold for the next SPO.
- Main idea: Storm Boy comforting Mr Percival.
- Use the SPO frame to build a topic sentence, supporting details, and a concluding sentence.

DO:
- Direct students to the booklet task next.
- Use this slide for students who need extra scaffolding.

TEACHER NOTES:
Optional support slide. Unhide if needed.

WATCH FOR:
- Students who copy the example without thinking - prompt them to use their own words.

Stage: Scaffold (optional) | VTLM: Differentiation""",


    52: """SAY:
- Time to apply this on your own.
- First, turn to the page titled Lesson 22: Sentence level writing.
- Next, complete the SPO task with a topic sentence, supporting details in note form, and a concluding sentence.
- Then, read your SPO aloud to a partner to check it makes sense.

DO:
- Direct students to the booklet page.
- Circulate and check first 2 to 3 supporting details for each student.
- Remind students that supporting details are notes, not full sentences.

TEACHER NOTES:
You Do practice. Pull a small group if many students are still writing full sentences for supporting details.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use the scaffolded slide as a worked example. Replace the words but keep the structure.
EXTENDING PROMPT:
- Task: After completing the booklet, expand your SPO into a full paragraph.

WATCH FOR:
- Students who skip the topic sentence - require it to be written first.
- Students who write full sentences for supporting details - re-cue note form.
- Students finishing fast - ask them to extend the SPO into a full paragraph.

Stage: You Do | VTLM: Independent Practice""",
}


# ============================================================
# LESSON 23 - Note taking (KPAS)
# ============================================================
LESSON_23_NOTES = {

    8: """SAY:
- Quick reset before we read.
- Today we read pages 73 to 76 of Storm Boy.
- Listen carefully and watch for the changes Storm Boy is going through.

DO:
- Pre-mark your novel with chosen pause points.

TEACHER NOTES:
Read-aloud setup. Mode of reading is teacher choice.

WATCH FOR:
- Students who do not have a copy ready - sort this out before reading begins.

Stage: Text-Level Reading | VTLM: Engagement""",


    9: """SAY:
- We are going to read pages 73 to 76.
- Your job is to listen carefully and look at the pictures.
- I will pause at certain points to ask questions.

DO:
- Read pages 73 to 76 of Storm Boy aloud, or run the chosen reading mode.
- Pause at the points you have pre-marked.
- Take 1-2 responses at each pause.

TEACHER NOTES:
Pause points come from the Literature Study Guide. Choose what matches your students' needs.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Listening partner - whisper a one-word feeling at each pause.
EXTENDING PROMPT:
- Task: Note one change Storm Boy is going through and the evidence in the text.

WATCH FOR:
- Students drifting off - point to the page and ask a quick concrete question.

Stage: Text-Level Reading | VTLM: Explicit Teaching""",


    10: """SAY:
- Three new words from today's pages.
- Watch and listen first, then we practise together.

DO:
- Have whiteboards and textas ready.

TEACHER NOTES:
Brief divider.

WATCH FOR:
- Students who need wait time - protect it across the next set of slides.

Stage: Vocabulary Divider | VTLM: Engagement""",


    11: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our first word is arrange.
- Watch this first - I am going to think aloud.
- If you arrange something, you make plans for it to happen.
- In the story, the sailors will arrange to have Mr Percival put into the museum. So arrange is not just moving - it is planning ahead.
- Say the word with me... arrange!

DO:
- Point to the image on the slide.
- Mime placing things in order with your hands.
- Have students repeat the word twice.

TEACHER NOTES:
First meeting with arrange. The word has two related meanings - to plan, and to put in order.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Repeat the word twice and mime arranging chairs.
EXTENDING PROMPT:
- Task: Use arrange in two different sentences - one for planning, one for ordering.

WATCH FOR:
- Students who only think arrange means putting flowers in order - confirm it also means making plans.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    12: """SAY:
- Read the sentence with me.
- Can you help me arrange the chairs in a circle for the group discussion? the teacher asked.
- Ask: what does arrange mean here? Expected: put them in order.

DO:
- Choral read.
- Mime placing chairs in a circle.

TEACHER NOTES:
Connect to a familiar classroom example.

WATCH FOR:
- Students who treat arrange as just move - confirm it has the planning or ordering element.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    13: """SAY:
- Two sentences. Which one is correct?
- A: They decided to arrange a surprise party for their friend's birthday.
- B: She was surprised to see the pillows carefully arranged in a messy pile.
- Show me A or B.

DO:
- Read both clearly.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Expected response: A. Arranged means organised - so a messy pile cannot be arranged.

WATCH FOR:
- Students who pick B because pillows sounds like decorating - clarify that messy and arranged contradict each other.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    14: """SAY:
- When you arrange something, you make plans for it to happen.
- Ask: when have you arranged something? Expected: a sleepover, a party, a sport game.
- Tell your partner one example.

DO:
- Give 30 seconds of partner talk.
- Take 1-2 examples.

TEACHER NOTES:
Personal connection. Common examples: sleepovers, parties, family outings.

WATCH FOR:
- Students stuck for an idea - prompt with sleepover, party, sport game, family outing.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    15: """SAY:
- Write a sentence using the word arrange.
- Pick something you might arrange.
- Watch my example: Can you arrange the books on the shelf by colour? I kindly asked my sister.

DO:
- Give 90 seconds.
- Circulate and check.
- Take 1-2 examples to share.
- Reveal example sentence on click.

TEACHER NOTES:
Active practice. Accept any sentence where arrange shows planning or organising.

WATCH FOR:
- Students who write arrange the room - prompt them to specify how the room is being arranged.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    16: """SAY:
- Show me which option fits best.
- When you arrange something you...
- Hold up 1 for forget it is happening, 2 for make plans for it to happen, 3 for ask your parents to organise it, 4 for hope it will happen.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which option fits best with arrange? Show me on three.
- Scan for: fingers showing 2.
PROCEED:
- If most students are on 2, move on to tossing.
PIVOT:
- Most likely misconception: students pick 4 (hope) because hope feels related.
- Reteach: arrange is making the plans, not just hoping.
- Fresh re-check: did you arrange the trip if you only hoped it would happen? Thumbs up or down.

TEACHER NOTES:
End-of-word CFU.

WATCH FOR:
- Students mirroring the fastest hand - keep wait time firm.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    17: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our next word is tossing.
- Watch this first - I am going to think aloud.
- If you toss something, you throw it lightly and carelessly.
- In the story, Storm Boy looks out at the tossing trees - so tossing is not just throwing, it can be moving lightly back and forth.
- Say the word with me... tossing!

DO:
- Point to the image on the slide.
- Mime a light underarm toss.
- Have students repeat the word twice.

TEACHER NOTES:
Tossing has two meanings - throwing lightly, and moving lightly back and forth (like trees in the wind).

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Repeat the word and copy the light-toss gesture.
EXTENDING PROMPT:
- Task: Compare tossing with flinging - which is gentler?

WATCH FOR:
- Students who confuse tossing with flinging from yesterday - confirm flinging is forceful, tossing is light.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    18: """SAY:
- Read the sentence with me.
- The kids were happily tossing a beach ball back and forth in the pool.
- Ask: what does tossing tell us? Expected: throwing it lightly and playfully.

DO:
- Choral read.
- Mime the light back-and-forth toss.

TEACHER NOTES:
Connect to a happy, playful context.

WATCH FOR:
- Students who think tossing is hard - clarify it is light, not strong.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    19: """SAY:
- Think of a time you have tossed something by throwing it lightly and carefully.
- Stay silent and act it out.
- I might choose two volunteers to show their tossing.

DO:
- Give 30 seconds of silent thinking.
- Invite 2 volunteers to act.
- Ask the class: what was the volunteer tossing?

TEACHER NOTES:
Movement task. Light, easy motion is the visual signature.

WATCH FOR:
- Students who throw hard - prompt: tossing is light, not forceful.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    20: """SAY:
- Two sentences. Which one is correct?
- A: He decided to go tossing after lunch to catch a fish.
- B: The chef was busy tossing the salad with a light dressing.
- Show me A or B.

DO:
- Read both.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Expected response: B. Tossing is not a fishing word.

WATCH FOR:
- Students who pick A because it sounds like a sport - clarify that tossing is not the right word for fishing.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    21: """SAY:
- Two pictures of food - A and B.
- Which one is being tossed?
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Look for the food in motion - usually salad or vegetables being tossed in a pan.

WATCH FOR:
- Students who pick the food that looks tastier - bring focus back to the action.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    22: """SAY:
- Show me which word fits best with tossing.
- Hold up 1 for fling, 2 for boot, 3 for punt, 4 for hit.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which word fits best with tossing? Show me on three.
- Scan for: fingers showing 1 (fling).
PROCEED:
- If most students are on 1, move to silently.
PIVOT:
- Most likely misconception: students pick boot, punt or hit - all stronger actions.
- Reteach: tossing is light. Fling is the closest match because both involve throwing.
- Fresh re-check: would you toss a feather or boot a feather?

TEACHER NOTES:
End-of-word CFU.

WATCH FOR:
- Students confused by the technical difference between fling and toss - acknowledge they overlap, but tossing is the lighter action.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    23: """SOURCES:
Macquarie Dictionary, 2024.

SAY:
- Our next word is silently.
- Watch this first - I am going to think aloud.
- If something happens silently, it happens without any noise.
- In the story, Storm Boy stood for a long time looking silently all around him - so silently is not just quiet, it is making no sound at all.
- Say the word with me... silently!

DO:
- Point to the image on the slide.
- Hold a finger to your lips and tip-toe to model silently.
- Have students repeat the word twice.

TEACHER NOTES:
Silently is an adverb - tells us how something happens. Use the finger-to-lips gesture as an anchor.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Repeat the word and copy the finger-to-lips gesture.
EXTENDING PROMPT:
- Task: Use silently in a sentence about a character from a different book.

WATCH FOR:
- Students who confuse silently with quietly - clarify silently means absolutely no sound, quietly means soft sound.

Stage: Vocabulary I Do | VTLM: Explicit Teaching""",


    24: """SAY:
- Read the sentence with me.
- The children watched the fireworks silently, mesmerised by the bright colours in the night sky.
- Ask: what does silently tell us? Expected: they watched without making any sound.

DO:
- Choral read.
- Use a hushed voice.

TEACHER NOTES:
Connect to a vivid familiar context.

WATCH FOR:
- Students who lose the meaning of mesmerised - gloss it as completely focused.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    25: """SAY:
- When you do something silently, you do it without a noise.
- Ask: when have you needed to do something silently? Expected: library, sneaking past a sleeping baby, watching a movie at night.
- Tell your partner one example.

DO:
- Give 30 seconds of partner talk.
- Take 1-2 examples.

TEACHER NOTES:
Personal connection. Library, sleeping baby, late-night movie are common examples.

WATCH FOR:
- Students stuck - prompt with library, sleeping baby, movie at night.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    26: """SAY:
- Build a sentence using the word silently.
- Sentence frame: The ____ moved silently through the garden, looking for ____.
- Watch my example: The cat moved silently through the garden, looking for a mouse to catch.

DO:
- Read the example.
- Give 60 seconds for students to build their own.
- Take 1-2 examples to review.

TEACHER NOTES:
Sentence frame supports students who need a starter.

WATCH FOR:
- Students who write silently for a noisy thing - re-anchor: silently means no sound.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    27: """SAY:
- Which words have a similar meaning to silently?
- Read with me: noiselessly, carefully, quietly, seriously, stealthily, hushed.
- Show thumbs up if it is similar, thumbs down if not.

DO:
- Read each word.
- Give 3 seconds for thumbs.
- Reveal: noiselessly, quietly, stealthily, hushed (yes); carefully, seriously (no).

TEACHER NOTES:
Carefully is about taking care, not being silent. Seriously is about tone, not sound.

WATCH FOR:
- Students unsure on stealthily - explain it means moving quietly to avoid being noticed.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    28: """SAY:
- Two groups of people - A and B.
- Which group is moving silently?
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Look for the body language clues - careful steps, fingers to lips, no open mouths.

WATCH FOR:
- Students who pick the more interesting picture - bring focus back to the body language.

Stage: Vocabulary We Do | VTLM: Guided Practice""",


    29: """SAY:
- Sentence building.
- She worked silently, but found it hard to concentrate without any background noise.
- She worked silently, so that she wouldn't disturb her sleeping classmates.
- She worked silently because she was focused and didn't want to make any mistakes.
- Ask: what does each connector tell us? Expected: but for contrast, so for result, because for reason.

DO:
- Read each sentence aloud.
- Point to each connector.

CFU CHECKPOINT:
Technique: Class discussion plus targeted question
Script:
- Which connector shows a reason? Cold call.
- Scan for: students saying because.
PROCEED:
- If students name the connectors correctly, move to the booklet task.
PIVOT:
- Most likely misconception: students mix up so and because.
- Reteach: so introduces the result that follows. Because introduces the reason that came first.
- Fresh re-check: She walked silently ____ the baby was asleep - so or because?

TEACHER NOTES:
End-of-word CFU using the connector grammar.

WATCH FOR:
- Students who skip the connector and only read the sentence - prompt them to focus on the joining word.

Stage: Vocabulary CFU | VTLM: Formative Assessment""",


    30: """SAY:
- Time to apply this on your own.
- First, turn to the page titled Lesson 23: Vocabulary.
- Next, complete the booklet tasks for arrange, tossing and silently.
- Then, read your sentences back to a partner.

DO:
- Direct students to the booklet page.
- Circulate and check first 2 to 3 answers per student.
- Pull a small group if needed.

TEACHER NOTES:
You Do practice for the three vocabulary words.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Complete only the meaning-match questions for now.
EXTENDING PROMPT:
- Task: Write one fresh sentence for each word in a Storm Boy context.

WATCH FOR:
- Students writing definitions instead of booklet answers - direct them back to the prompts.
- Students finishing fast - ask them to write a fresh sentence with each word.

Stage: Vocabulary You Do | VTLM: Independent Practice""",


    31: """SAY:
- Quick reset.
- Now we move to sentence-level writing.
- Today we are going to take notes using the KPAS strategy.

DO:
- Have whiteboards and textas ready.

TEACHER NOTES:
Brief divider.

WATCH FOR:
- Energy dipping - quick stand-stretch can reset before the next chunk.

Stage: Sentence-Level Divider | VTLM: Engagement""",


    32: """SAY:
- Watch this first - I am thinking about why hand-written notes work.
- In previous lessons, we worked on note taking.
- Taking notes by hand helps us in three ways - it helps us work out what is essential, retain information better, and understand what we read.
- Studies show this is true for handwritten notes more than typed notes.

DO:
- Read the bullets on the slide.
- Hold up a notebook to anchor handwritten note taking.

TEACHER NOTES:
I Do anchor for why we take notes by hand.

WATCH FOR:
- Students who think typing is faster and better - acknowledge but reframe with the research.

Stage: I Do | VTLM: Explicit Teaching""",


    33: """SAY:
- Why do we take notes? Four reasons.
- Planning and preparing for writing.
- Organising thoughts and ideas.
- Staying focused on the task.
- Clarifying and remembering information.
- Tell your partner one reason that matters most for you.

DO:
- Read the four reasons.
- Give 30 seconds of partner talk.
- Take 2-3 responses.

TEACHER NOTES:
We Do for the why.

WATCH FOR:
- Students stuck - prompt with: how do you remember after a holiday what you learned before?

Stage: We Do | VTLM: Guided Practice""",


    34: """SAY:
- Taking notes is done by using the KPAS strategy.
- K stands for Key words.
- P stands for Phrases.
- A stands for Abbreviations.
- S stands for Symbols.
- Say with me: K, P, A, S.

DO:
- Point to each letter on the slide.
- Choral repeat KPAS.
- Use your fingers to count off K, P, A, S.

TEACHER NOTES:
Core strategy anchor. Students will use these four labels for the rest of the lesson.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Choral repeat KPAS three times until it sticks.
EXTENDING PROMPT:
- Task: Identify one abbreviation and one symbol they already use in maths or science.

WATCH FOR:
- Students who muddle the order - choral repeat several times.

Stage: I Do | VTLM: Explicit Teaching""",


    35: """SAY:
- Watch this first.
- Keywords and phrases are the parts of a sentence that tell us the main idea.
- I am going to find them by asking - who, what doing, when, where, why and how.
- They are the most important words.

DO:
- Read the slide.
- Point to who, what doing, when, where, why, how.

TEACHER NOTES:
I Do for keywords and phrases.

WATCH FOR:
- Students who try to highlight every word - prompt them to choose only the most important.

Stage: I Do | VTLM: Modelling""",


    36: """SAY:
- Two examples - which one shows keywords and phrases highlighted?
- Example 1 highlights everything.
- Example 2 highlights only the keywords - Mr Percival lived on in their mind's eye because a bird like Mr Percival never really dies.
- Show me 1 or 2.

DO:
- Read both examples.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Expected response: example 2. Keywords and phrases mean only the most important words.

MISCONCEPTIONS:
- Misconception: highlighting more words means better notes.
  Why: students think more equals more thinking.
  Impact: notes become as long as the original sentence.
  Quick correction: keywords are the few most important - the ones you would say if you only had three seconds.

WATCH FOR:
- Students who pick example 1 because more highlighted seems better - reteach: keywords are the few most important.

Stage: We Do | VTLM: Guided Practice""",


    37: """SAY:
- Watch this first.
- Abbreviations are shortened versions of words.
- We use them to save time when taking notes.
- I notice we already use them every day - AM and PM, Jan and Feb and Dec, NSW and Vic and WA, Rd and St and Pl, cm and m and km, tsp and ml and kg.

DO:
- Read the abbreviations.
- Point to each one.
- Ask: what does each one stand for?

TEACHER NOTES:
I Do for abbreviations. Use the familiar ones to anchor the concept.

WATCH FOR:
- Students unsure about NSW or Vic - clarify they are state names.

Stage: I Do | VTLM: Explicit Teaching""",


    38: """SAY:
- Now let's apply abbreviations to our novel.
- Ask: what might these abbreviations mean in Storm Boy?
- Mr P - what does that stand for? Crng? Adl? H-A?

DO:
- Cold call for each abbreviation.
- Reveal: Mr Percival, the Coorong, Adelaide, Hide-Away.

TEACHER NOTES:
Connect abbreviations to story characters and places. Students will use these in the rest of the lesson.

WATCH FOR:
- Students who do not recognise Crng - prompt: it is a place from the story.

Stage: We Do | VTLM: Guided Practice""",


    39: """SAY:
- Symbols are shorthand techniques to take notes efficiently.
- We have explored these before.
- Today we will use them when converting notes to sentences.

DO:
- Read the slide.

TEACHER NOTES:
I Do anchor for symbols. Quick activation - the next slide unpacks the most common ones.

WATCH FOR:
- Students who do not know what shorthand is - clarify: short symbols that stand for full words.

Stage: I Do | VTLM: Explicit Teaching""",


    40: """SAY:
- Ask: what do these symbols mean?
- Forward slash - new idea.
- Up arrow - up, over, increase.
- Equals sign - equals or leads to.
- Ampersand (the and symbol) - and.
- Tell your partner one symbol you might use today.

DO:
- Read each symbol.
- Give 20 seconds of partner talk.
- Take 2 responses.

TEACHER NOTES:
We Do for symbols.

WATCH FOR:
- Students confused by the up arrow - explain: it means more, higher, or increase.

Stage: We Do | VTLM: Guided Practice""",


    41: """SAY:
- Which set of notes best represents the sentence?
- Sentence: Mr Percival is a very intelligent bird.
- Option 1: Mr P equals very intelligent bird.
- Option 2: Mr P slash intelligent bird.
- Option 3: Mr P arrow intelligent bird.
- Show me 1, 2 or 3.

DO:
- Read each option.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Expected response: option 1. Equals shows that Mr P is the bird.

WATCH FOR:
- Students who pick option 3 because arrows look advanced - reteach: arrow is for cause or sequence, not equals.

Stage: We Do | VTLM: Guided Practice""",


    42: """SAY:
- Which set of notes best represents the sentence?
- Sentence: Mr Percival was a brave bird who heroically saved sailors' lives.
- Option 1: Mr P was brave equals who heroically saved arrow lives.
- Option 2: Mr P equals brave bird slash saved sailors' lives.
- Show me 1 or 2.

DO:
- Read both.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

TEACHER NOTES:
Expected response: option 2. The equals shows Mr P is the brave bird, and the slash separates the new idea about saving sailors' lives.

WATCH FOR:
- Students who pick option 1 because longer looks like more thinking - reteach: notes use fewer words, not more.

Stage: We Do | VTLM: Guided Practice""",


    43: """SAY:
- Watch this first.
- Now we convert KPAS notes back into full sentences.
- Notes: w/out Mr P arrow S Boy no stay.
- I read this as: without Mr Percival, Storm Boy could not stand to stay.
- I am going to decode each piece - w slash out is short for without. Mr P is Mr Percival. Arrow means leads to. S Boy is Storm Boy. No stay means could not stay.

DO:
- Point to each abbreviation and symbol.
- Read aloud the full sentence.

TEACHER NOTES:
I Do for converting notes to sentences. Make every step visible.

WATCH FOR:
- Students who lose track of the symbols - re-anchor each one slowly.

Stage: I Do | VTLM: Modelling""",


    44: """SAY:
- Notes: Mr P buried down arrow Lookout Pst.
- Decode it. Mr P - Mr Percival. Buried. Down arrow - below or under. Lookout Pst - Lookout Post.
- Full sentence: Mr Percival was buried below the Lookout Post.

DO:
- Point to each piece.
- Read the full sentence aloud.
- Have students echo it back.

TEACHER NOTES:
We Do practice. Build the sentence step by step on the board if helpful.

WATCH FOR:
- Students who skip the down arrow - confirm it means below or under.

Stage: We Do | VTLM: Guided Practice""",


    45: """SAY:
- Notes: H-A equals alone again.
- Decode it. H-A - Hide-Away. Equals - is. Alone again - by himself once more.
- Full sentence: Hide-Away was all alone once again.

DO:
- Point to each piece.
- Read the full sentence aloud.

TEACHER NOTES:
Short notes - quick conversion practice.

WATCH FOR:
- Students who try to add details that are not in the notes - re-anchor: convert what is there.

Stage: We Do | VTLM: Guided Practice""",


    46: """SAY:
- Notes: Mr P gone arrow S Boy rdy go 2 school in ADL.
- Decode it. Mr P gone - Mr Percival is gone. Arrow - means or leads to. S Boy rdy - Storm Boy is ready. Go 2 school - go to school. In ADL - in Adelaide.
- Full sentence: Since Mr Percival was gone, Storm Boy was ready to go to school in Adelaide.

DO:
- Point to each piece.
- Read the full sentence aloud.

TEACHER NOTES:
Multiple abbreviations and a symbol. Build it up step by step.

WATCH FOR:
- Students stuck on rdy - confirm it means ready.

Stage: We Do | VTLM: Guided Practice""",


    47: """SAY:
- Notes: S Boy lost Mr P b slash c cruel men.
- Decode it. S Boy - Storm Boy. Lost Mr P - lost Mr Percival. B slash c - because. Cruel men - the cruel men.
- Full sentence: Storm Boy lost Mr Percival because of cruel men.

DO:
- Point to each piece.
- Read the full sentence aloud.

TEACHER NOTES:
B slash c is the abbreviation for because. Confirm this for students who may not recognise it.

WATCH FOR:
- Students who do not know b slash c - explain it stands for because.

Stage: We Do | VTLM: Guided Practice""",


    48: """SAY:
- Notes: S Boy left arrow H-A sad.
- Decode it. S Boy - Storm Boy. Left. Arrow - leads to or causes. H-A - Hide-Away. Sad.
- Full sentence: When Storm Boy left, Hide-Away was sad.

DO:
- Point to each piece.
- Read the full sentence aloud.

TEACHER NOTES:
The arrow shows that one event caused another - Storm Boy leaving caused Hide-Away to feel sad.

WATCH FOR:
- Students who write was Storm Boy left - prompt them to start with When or After.

Stage: We Do | VTLM: Guided Practice""",


    49: """SAY:
- Look at the notes and the sentence.
- Notes: S Boy arrow Adelaide w slash out H-A.
- Sentence: Storm Boy saw Adelaide when out of Hide-Away's humpy.
- Does the sentence reflect the notes? Show me thumbs up or thumbs down.

DO:
- Read both.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Reveal correct sentence: Storm Boy went to Adelaide without Hide-Away.

CFU CHECKPOINT:
Technique: Thumbs Up / Thumbs Down with reveal
Script:
- Does the sentence match the notes? Thumbs up or down.
- Scan for: thumbs down.
PROCEED:
- If most students show thumbs down, move to the next CfU.
PIVOT:
- Most likely misconception: students confuse arrow with saw.
- Reteach: arrow means went to, not saw. W slash out means without, not when out.
- Fresh re-check: read the notes again - what does arrow mean here?

TEACHER NOTES:
CFU testing whether students decode symbols accurately.

WATCH FOR:
- Students putting thumbs sideways - require a clear up or down.

Stage: CFU | VTLM: Formative Assessment""",


    50: """SAY:
- Notes: H-A look up arrow at sky slash S Boy look arrow window.
- Sentence: While Hide-Away looked up at the night sky, Storm Boy looked out through his window.
- Does the sentence reflect the notes? Show me thumbs up or thumbs down.

DO:
- Read both.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Thumbs Up / Thumbs Down
Script:
- Does the sentence match the notes? Thumbs up or down.
- Scan for: thumbs up.
PROCEED:
- If most students show thumbs up, move to the booklet task.
PIVOT:
- Most likely misconception: students miss the slash symbol meaning new idea, and treat it as one continuous action.
- Reteach: the slash separates the two ideas - Hide-Away looking at the sky, AND Storm Boy looking through his window.
- Fresh re-check: read the notes again - what does the slash mean?

TEACHER NOTES:
Hinge CFU before You Do. Tests slash-symbol decoding.

WATCH FOR:
- Students who hesitate - reread the sentence aloud.

Stage: CFU | VTLM: Formative Assessment""",


    51: """SAY:
- Time to apply this on your own.
- First, turn to the page titled Lesson 23: Sentence level writing.
- Next, complete the tasks on note taking using KPAS.
- Then, swap with a partner and check each other's notes make sense.

DO:
- Direct students to the booklet page.
- Circulate and check the first 2 to 3 answers for each student.
- Remind students of the four KPAS letters.

TEACHER NOTES:
You Do practice. Pull a small group if many students still write full sentences instead of using notes.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a KPAS reference card alongside the booklet. Convert one sentence at a time.
EXTENDING PROMPT:
- Task: Take notes from a paragraph in a different chapter using the full KPAS strategy.

WATCH FOR:
- Students who skip abbreviations - prompt them to use them where they make sense.
- Students who overuse symbols and lose meaning - re-anchor: notes must still make sense to you later.
- Students finishing fast - ask them to convert their notes back into full sentences as a self-check.

Stage: You Do | VTLM: Independent Practice""",
}


# ============================================================
# LESSON 24 - Plan a concluding paragraph
# ============================================================
LESSON_24_NOTES = {

    9: """SAY:
- Today we are using a non-fiction text to help us plan our concluding paragraph.
- We need to know about the Coorong before we write about it.

DO:
- Have the supplementary text ready - either the slides or the linked articles.

TEACHER NOTES:
Brief divider. The Literature Study Guide lists optional reading sources.

WATCH FOR:
- Students who do not know what the Coorong is - the next four slides will introduce it.

Stage: Background Knowledge Divider | VTLM: Engagement""",


    10: """SAY:
- Listen as I read about the Coorong National Park.
- The Coorong is a special place in South Australia's southeast.
- It is a protected area that preserves unique habitats along the coast.
- It is located 200 km southeast of Adelaide.
- It is long and narrow, stretching from Goolwa Barrage to nearly Kingston.
- It includes the Goolwa and Coorong Channels, the northern and southern lagoons, Ocean Beach, sand dunes and the surrounding land.

DO:
- Read the slide aloud.
- Point to each fact as you say it.
- Pause to ask: where is the Coorong?

TEACHER NOTES:
Background-knowledge slide 1 of 4. Keep facts visible - students will use them when planning their paragraph.

WATCH FOR:
- Students who lose the location detail - rewind: 200 km southeast of Adelaide.

Stage: Background Knowledge | VTLM: Explicit Teaching""",


    11: """SAY:
- More about the Coorong.
- It is home to water birds and migratory species - they have safe shelter there.
- It shows how rocks and landforms change over time.
- It holds significant historical and archaeological sites.
- It supports activities like fishing and tourism.
- Many groups care for the region - the Department for Environment and Water, the Ngarrindjeri people, recreational boaters and fishers, conservation groups, and local councils.

DO:
- Read the slide aloud.
- Point to each group as you name it.

TEACHER NOTES:
Background-knowledge slide 2 of 4.

WATCH FOR:
- Students who lose track of the groups - re-cue: many groups care for it.

Stage: Background Knowledge | VTLM: Explicit Teaching""",


    12: """SAY:
- An Internationally Important Wetland.
- The Coorong is home to 115 species that need protection.
- That includes 79 birds, 2 amphibians, 15 mammals, 4 reptiles and 15 plants.
- It is the most important waterbird wetland in the Murray-Darling Basin system.
- It is internationally recognised - it has been listed as a Ramsar Wetland because it provides habitats for water birds.

DO:
- Read the slide aloud.
- Point to the numbers.

TEACHER NOTES:
Background-knowledge slide 3 of 4.

WATCH FOR:
- Students who get bogged down in the species count - keep the focus on importance.

Stage: Background Knowledge | VTLM: Explicit Teaching""",


    13: """SAY:
- The Coorong is crucial for both local species and migratory shorebirds.
- Some birds travel all the way from Alaska and Siberia.
- The Ramsar Convention is an international agreement that aims to protect wetlands worldwide.
- South Australia has six other Ramsar sites besides the Coorong.

DO:
- Read the slide aloud.
- Use a finger to trace the migration journey from Alaska or Siberia to Australia.

TEACHER NOTES:
Background-knowledge slide 4 of 4.

WATCH FOR:
- Students who do not know where Alaska or Siberia is - quick gloss: very far north.

Stage: Background Knowledge | VTLM: Explicit Teaching""",


    14: """SAY:
- What does the Coorong help preserve?
- Show me 1, 2, 3 or 4.
- Option 1: the Ngarrindjeri people.
- Option 2: unique habitats.
- Option 3: historical and archaeological sites.
- Option 4: Department for Environment and Water.

DO:
- Read each option.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- What does the Coorong help preserve? Show me on three.
- Scan for: fingers showing 2 or 3 (both are valid - unique habitats AND historical sites).
PROCEED:
- If most students are on 2 or 3, move to text-level writing.
PIVOT:
- Most likely misconception: students pick 1 (Ngarrindjeri people) or 4 (Department) - they are caretakers, not what is preserved.
- Reteach: the Coorong preserves places and habitats. People and organisations look after it.
- Fresh re-check: which is preserved - the people, or the wetlands?

TEACHER NOTES:
End-of-knowledge CFU.

WATCH FOR:
- Students who pick 1 - clarify: the people are not preserved, the lands and culture they care about are.

Stage: CFU | VTLM: Formative Assessment""",


    15: """SAY:
- Now we move to writing.
- Today we plan a concluding paragraph for our information report.

DO:
- Have whiteboards and textas ready.
- Display the next slide.

TEACHER NOTES:
Brief divider. Mentor text is provided - keep it visible during planning.

WATCH FOR:
- Energy dipping - quick stand-stretch can help reset before the writing chunk.

Stage: Text-Level Writing Divider | VTLM: Engagement""",


    16: """SAY:
- Watch this first.
- The purpose of an information report is to inform the reader.
- An information report classifies or describes factual information about a specific topic or event.
- Today our topic is the Coorong.

DO:
- Read the slide aloud.
- Point to the words classify and describe.

TEACHER NOTES:
I Do anchor for purpose. Sets up the rest of the lesson.

WATCH FOR:
- Students who confuse information report with persuasive writing - clarify it is about facts, not opinions.

Stage: I Do | VTLM: Explicit Teaching""",


    17: """SAY:
- We have already written our introduction and our three body paragraphs.
- Today we are planning the conclusion.
- A conclusion concludes the report and summarises the main points discussed.

DO:
- Read the slide aloud.
- Point to the link to the mentor text.

TEACHER NOTES:
I Do for purpose of conclusion. Keep it tight.

WATCH FOR:
- Students who think the conclusion is for new information - we will correct this on the next slide.

Stage: I Do | VTLM: Explicit Teaching""",


    18: """SAY:
- A factual recount has three body paragraphs.
- Ours have been about the Coorong - what it is and where, cultural significance, and flora and fauna.
- In the concluding paragraph we will - summarise key points, link to the introduction, avoid introducing any new information.
- Notice the third one - we do not add new facts in a conclusion.

DO:
- Read the slide aloud.
- Point to each rule.
- Highlight the third rule with extra emphasis.

TEACHER NOTES:
We Do for what a conclusion does. The no-new-information rule is the most-missed.

WATCH FOR:
- Students who plan to add new facts - re-anchor the rule.

Stage: We Do | VTLM: Guided Practice""",


    19: """SAY:
- True or false?
- Statement 1: The concluding paragraph summarises the factual information from the body paragraphs.
- Statement 2: The concluding paragraph introduces new facts, events and information.
- Statement 3: The concluding paragraph is an opportunity to recap what has been discussed in the introduction and body paragraphs.
- Show me thumbs up for true, thumbs down for false, on each one.

DO:
- Read each statement.
- Give 3 seconds for thumbs.
- Reveal: 1 true, 2 false, 3 true.

CFU CHECKPOINT:
Technique: Thumbs Up / Thumbs Down per statement
Script:
- True or false? Thumbs up or down for each.
- Scan for: thumbs up on 1 and 3, thumbs down on 2.
PROCEED:
- If most students are correct on all three, move to TSG planning.
PIVOT:
- Most likely misconception: students get statement 2 wrong - they think a conclusion can include new info.
- Reteach: a conclusion only summarises what is already in the report. New information goes in the body paragraphs.
- Fresh re-check: should we add a new fact about pelicans in the conclusion? Thumbs up or down.

TEACHER NOTES:
Three-statement CFU. Statement 2 is the trap that catches the most-common error.

MISCONCEPTIONS:
- Misconception: a conclusion can introduce new facts.
  Why: students think a strong ending needs something new.
  Impact: their conclusions read like a fourth body paragraph.
  Quick correction: a conclusion only summarises what is already there.

WATCH FOR:
- Students putting thumbs sideways - require a clear answer.

Stage: CFU | VTLM: Formative Assessment""",


    20: """SAY:
- Watch this first.
- A concluding paragraph has three parts - the TSG formula.
- Thesis statement - a statement of what the information report discussed.
- Specific statement - reiterates the significance of the wetland.
- General statement - concludes the topic of the factual recount.

DO:
- Read the slide aloud.
- Point to each part.
- Use your fingers to count off T, S, G.

TEACHER NOTES:
I Do for TSG structure. Students will use these labels for the rest of the lesson. The notes section gives the example mentor text.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Choral repeat T, S, G three times until it sticks.
EXTENDING PROMPT:
- Task: Compare TSG with another paragraph formula they have used (e.g. PEEL).

WATCH FOR:
- Students who confuse TSG with TPSAC or other paragraph formulas - clarify TSG is for the conclusion.

Stage: I Do | VTLM: Explicit Teaching""",


    21: """SAY:
- The introduction uses GST - General, Specific, Thesis.
- The conclusion uses TSG - Thesis, Specific, General.
- I notice these mirror each other.
- Ask: why might that be? Turn and tell your partner.

DO:
- Point to GST in the introduction column and TSG in the conclusion column.
- Give 30 seconds of partner talk.
- Take 2 responses.
- Confirm: a clear bookend - the report starts wide and zooms in, then ends by zooming back out.

TEACHER NOTES:
We Do for the structural connection.

WATCH FOR:
- Students who do not see the mirror - point to each pair: General matches General, Specific matches Specific, Thesis matches Thesis.

Stage: We Do | VTLM: Guided Practice""",


    22: """SAY:
- Watch this first.
- The thesis statement says what the information report discussed.
- For our Coorong report, the body paragraphs covered location, cultural significance, and flora and fauna.
- I am going to list those into a thesis: This report detailed the location of the Coorong and its cultural and ecological significance.

DO:
- Read the example.
- Point to the brainstormed words on the slide - location, conservation, cultural significance, animals and plants.

TEACHER NOTES:
We Do for the thesis statement. Students should use the brainstormed words to build their own.

WATCH FOR:
- Students who add new ideas to the thesis - re-anchor: thesis covers what we already wrote.

Stage: We Do | VTLM: Modelling""",


    23: """SAY:
- Use the brainstormed ideas and the introduction's thesis statement to build your own thesis statement.
- For example: This report will discuss exactly what the Coorong is, where it is located, its cultural significance, and the flora and fauna that rely on it for their survival.
- Now write your own thesis statement.
- Share with your partner when you are done.

DO:
- Give 90 seconds for students to write.
- Circulate and check 4-5 sentences.
- Take 2 examples to share.

TEACHER NOTES:
We Do practice. Accept any thesis statement that names the three or four topic areas.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use a frame: This report discussed ____, ____ and ____.
EXTENDING PROMPT:
- Task: Write two different thesis sentences - one short, one with more detail.

WATCH FOR:
- Students who write a generic statement (e.g. This report is about the Coorong) - prompt them to name the topics.

Stage: We Do | VTLM: Guided Practice""",


    24: """SAY:
- The specific statement addresses the topic directly and uses key terms.
- Three guiding questions - who or what is the Coorong important to? Why is it important? Does it remain important today?
- Key words on the slide - Ngarrindjeri people, fauna, supports life, biodiversity, wetland, tourists, flora.

DO:
- Read the guiding questions.
- Cold call non-volunteers to share answers.
- Reveal the suggested key terms.

TEACHER NOTES:
I Do for the specific statement. Cold-call non-volunteers to ensure all students engage.

WATCH FOR:
- Students who only give one answer - prompt them to address all three questions.

Stage: I Do | VTLM: Modelling""",


    25: """SAY:
- Use the key words and the specific statement from the introduction to build your own specific statement.
- For example: It is an expansive area of wetlands and waterways that is both culturally and ecologically important.
- Now write your own.

DO:
- Give 90 seconds for students to write.
- Circulate and check.
- Take 2 examples to share.

TEACHER NOTES:
We Do practice. Aim for a statement that captures importance and key terms.

WATCH FOR:
- Students who repeat the introduction's specific statement word for word - prompt them to rephrase.

Stage: We Do | VTLM: Guided Practice""",


    26: """SAY:
- The general statement summarises the topic and concludes the information report.
- For example: The Coorong, a long shallow lagoon filled with life, is truly an iconic Australian landscape!
- I notice two features - it uses an appositive, and it is an exclamation.
- An appositive is a phrase that renames or describes the noun.

DO:
- Read the example sentence.
- Point to the appositive - a long shallow lagoon filled with life.
- Point to the exclamation mark.

TEACHER NOTES:
I Do for the general statement. The appositive plus exclamation mark is the modelled craft move.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Identify the appositive in the example sentence by underlining it.
EXTENDING PROMPT:
- Task: Find another example of an appositive in their reading and share it.

WATCH FOR:
- Students who do not know what an appositive is - re-explain: a phrase that adds detail to the noun, set off with commas.

Stage: I Do | VTLM: Modelling""",


    27: """SAY:
- Appositives rename, explain or define the noun - the Coorong.
- Some examples - a natural ancient escape, a long shallow lagoon, a spectacular saline lagoon, a popular tourist destination.
- Pick one and read it aloud.
- Tell your partner which one fits the Coorong best and why.

DO:
- Read each appositive.
- Give 30 seconds of partner talk.
- Take 2 responses.

TEACHER NOTES:
We Do for choosing an appositive.

WATCH FOR:
- Students who pick the longest one without thinking - prompt: which fits the message of your conclusion best?

Stage: We Do | VTLM: Guided Practice""",


    28: """SAY:
- Use one of the appositives to build your own general statement.
- For example: The Coorong, a spectacular saline lagoon, is one of South Australia's most visited tourist attractions!
- Now build your own with a partner.

DO:
- Give 90 seconds of partner work.
- Circulate and check 3-4 pairs.
- Take 2 examples to share.

TEACHER NOTES:
We Do practice. Encourage exclamation marks for emphasis but do not force them.

WATCH FOR:
- Students who skip the appositive - prompt them to choose one from the list.

Stage: We Do | VTLM: Guided Practice""",


    29: """SAY:
- Here is the model conclusion building up.
- First the Thesis - This report detailed the location of the Coorong and its cultural and ecological significance.
- Then the Specific - The importance of this expansive wetland area is enormous for both past and present-day people, flora, and fauna.
- Then the General - The Coorong, a long shallow lagoon filled with life, is truly an iconic Australian landscape!

DO:
- Read each part as it builds up.
- Point to T, S, G as you read each one.
- Choral read the full conclusion.

TEACHER NOTES:
Full mentor text. Use it as a reference for the You Do.

WATCH FOR:
- Students who lose track of the structure - point to each label as you build.

Stage: We Do | VTLM: Guided Practice""",


    30: """SAY:
- Ask: what words might you see in a Thesis Statement?
- Show me 1, 2 or 3.
- Option 1: Once upon a time.
- Option 2: This text explained...
- Option 3: I hope you liked this information report...

DO:
- Read each option.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-3)
Script:
- Which words might you see in a thesis statement? Show me on three.
- Scan for: fingers showing 2 (This text explained).
PROCEED:
- If most students are on 2, move to the booklet task.
PIVOT:
- Most likely misconception: students pick 1 (Once upon a time) - that is a narrative opener, not an information report thesis.
- Reteach: thesis statements name what the report discussed - common openers are This text explained, This report discussed, This report detailed.
- Fresh re-check: which fits an information report - This report explained, or Once upon a time?

TEACHER NOTES:
Pre-You-Do hinge CFU.

WATCH FOR:
- Students who pick 3 - clarify: hope and like are personal, not factual.

Stage: CFU | VTLM: Formative Assessment""",


    31: """SAY:
- Time to apply this on your own.
- First, turn to the page titled Lesson 24: Plan a concluding paragraph for an information report.
- Next, plan your Thesis, Specific and General statement.
- Then, share your plan with a partner.
- You will write the full concluding paragraph in Lesson 25.

DO:
- Direct students to the booklet page.
- Circulate and check the first thesis statement for each student.
- Remind students - planning is in note form, full sentences come next lesson.

TEACHER NOTES:
You Do practice. Pull a small group if many students struggle to start the thesis.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use the model conclusion as a worked example. Replace the words but keep the structure.
EXTENDING PROMPT:
- Task: Draft a second appositive option for the general statement and choose between them.

WATCH FOR:
- Students who write the full paragraph instead of planning notes - re-cue: notes only today.
- Students who add new information - re-anchor the no-new-info rule.
- Students finishing fast - ask them to draft a second appositive.

Stage: You Do | VTLM: Independent Practice""",


    33: """SAY:
- Reference - the introduction model text.
- The Coorong is a narrow, shallow lagoon that stretches across the Limestone coast in South Australia. It is an expansive area of wetlands and waterways that is both culturally and ecologically important. This report will discuss exactly what the Coorong is, where it is located, its cultural significance, and the flora and fauna that rely on it for their survival.
- Notice the GST formula - General, Specific, Thesis.

DO:
- Display the text.
- Point to each part - General, Specific, Thesis.

TEACHER NOTES:
Reference slide for the modelled introduction. Use it to compare with the conclusion.

WATCH FOR:
- Students who confuse GST with TSG - clarify: introduction starts general and zooms in.

Stage: Reference (Mentor Text) | VTLM: Modelling""",


    34: """SAY:
- Reference - body paragraph 1.
- Topic: what the Coorong is and where it is located.
- This paragraph names the Murray-Darling basin connection, the length, the Younghusband Peninsula, and the mix of fresh, brackish and saltwater.

DO:
- Display the text.
- Point to the topic-introducing sentence.

TEACHER NOTES:
Reference slide. Knowledge anchor when planning the thesis statement.

WATCH FOR:
- Students who try to copy this whole paragraph - re-cue: this is for reference, not copying.

Stage: Reference (Mentor Text) | VTLM: Modelling""",


    35: """SAY:
- Reference - body paragraph 2.
- Topic: cultural significance.
- The paragraph names the Ngarrindjeri people, the Kungun Ngarrindjeri Yunnan agreement, and the joint management of the parks.

DO:
- Display the text.
- Point to the topic-introducing sentence.

TEACHER NOTES:
Reference slide. The Ngarrindjeri name and KNY agreement are key facts.

SENSITIVITY ADVISORY:
- What it is: First Nations cultural content.
- Framing language: Use the Ngarrindjeri name with respect. Pronounce it carefully.
- Watch for: Students who do not know how to say Ngarrindjeri - model the pronunciation.
- Protocol: Refer back to the Cultural Sensitivity slide if needed.

WATCH FOR:
- Students who skip the Ngarrindjeri reference - re-anchor that cultural significance is a key topic.

Stage: Reference (Mentor Text) | VTLM: Modelling""",


    36: """SAY:
- Reference - body paragraph 3.
- Topic: flora and fauna.
- The paragraph names vulnerable species, local birds and migratory species from Alaska and Siberia, and the Bonney Upwelling.

DO:
- Display the text.
- Point to the topic-introducing sentence.

TEACHER NOTES:
Reference slide. The Bonney Upwelling and Alaska or Siberia migration are vivid facts.

WATCH FOR:
- Students who do not know what fauna or flora means - quick gloss: animals and plants.

Stage: Reference (Mentor Text) | VTLM: Modelling""",


    37: """SAY:
- Reference - the modelled conclusion.
- This report detailed the location of the Coorong and its cultural and ecological significance. The importance of this expansive wetland area is enormous for both past and present-day people, flora, and fauna. The Coorong, a long shallow lagoon filled with life, is truly an iconic Australian landscape!

DO:
- Display the text.
- Point to T, S, G as you read.

TEACHER NOTES:
Reference slide. Use it as the model for the You Do task.

WATCH FOR:
- Students who try to copy this exactly - prompt them to use it as a structure model, not a script.

Stage: Reference (Mentor Text) | VTLM: Modelling""",
}


# ============================================================
# LESSON 25 - Write a concluding paragraph
# ============================================================
LESSON_25_NOTES = {

    9: """SAY:
- Today we keep building our information report.
- We will use a non-fiction text to refresh our knowledge of the Coorong.

DO:
- Have the supplementary text ready - either the slides or the linked articles.

TEACHER NOTES:
Brief divider. Background-knowledge slides repeat from Lesson 24 to remind students of the facts.

WATCH FOR:
- Students who tuned out yesterday - keep the read-aloud short and active.

Stage: Background Knowledge Divider | VTLM: Engagement""",


    10: """SAY:
- Listen as I read about the Coorong National Park.
- The Coorong is a special place in South Australia's southeast.
- It is a protected area that preserves unique habitats along the coast.
- It is located 200 km southeast of Adelaide.
- It is long and narrow, stretching from Goolwa Barrage to nearly Kingston.
- It includes the Goolwa and Coorong Channels, the northern and southern lagoons, Ocean Beach, sand dunes and the surrounding land.

DO:
- Read the slide aloud.
- Point to each fact as you say it.

TEACHER NOTES:
Background-knowledge slide 1 of 4. Same content as Lesson 24 - quick refresh.

WATCH FOR:
- Students who already remember from yesterday - acknowledge and read briskly.

Stage: Background Knowledge | VTLM: Explicit Teaching""",


    11: """SAY:
- More about the Coorong.
- It is home to water birds and migratory species.
- It shows how rocks and landforms change over time.
- It holds significant historical and archaeological sites.
- It supports activities like fishing and tourism.
- Many groups care for the region - the Department for Environment and Water, the Ngarrindjeri people, recreational boaters and fishers, conservation groups, and local councils.

DO:
- Read the slide aloud.
- Point to each group.

TEACHER NOTES:
Background-knowledge slide 2 of 4.

WATCH FOR:
- Students who lose track of the groups - quick re-cue: many groups care for it.

Stage: Background Knowledge | VTLM: Explicit Teaching""",


    12: """SAY:
- An Internationally Important Wetland.
- The Coorong is home to 115 species that need protection.
- That includes 79 birds, 2 amphibians, 15 mammals, 4 reptiles and 15 plants.
- It is the most important waterbird wetland in the Murray-Darling Basin system.
- It is internationally recognised as a Ramsar Wetland.

DO:
- Read the slide aloud.

TEACHER NOTES:
Background-knowledge slide 3 of 4.

WATCH FOR:
- Students who get bogged down in numbers - keep the focus on importance.

Stage: Background Knowledge | VTLM: Explicit Teaching""",


    13: """SAY:
- The Coorong is crucial for both local species and migratory shorebirds.
- Some birds travel from Alaska and Siberia.
- The Ramsar Convention is an international agreement to protect wetlands.
- South Australia has six other Ramsar sites besides the Coorong.

DO:
- Read the slide aloud.

TEACHER NOTES:
Background-knowledge slide 4 of 4.

WATCH FOR:
- Students unsure of where Alaska or Siberia is - quick gloss.

Stage: Background Knowledge | VTLM: Explicit Teaching""",


    14: """SAY:
- What activities does the park support for tourists? Select all that apply.
- Show me 1, 2, 3 or 4 - or hold up multiple fingers if you think more than one is correct.
- Option 1: fishing.
- Option 2: bird watching.
- Option 3: whale watching.
- Option 4: kayaking and boating.

DO:
- Read each option.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.
- Reveal: 1, 2, 4 (whale watching is not mentioned in the texts).

CFU CHECKPOINT:
Technique: Show Fingers - select all that apply
Script:
- Which activities does the park support? Hold up the numbers - you can show more than one.
- Scan for: fingers showing 1, 2 and 4.
PROCEED:
- If most students have 1, 2 and 4 (and not 3), move to text-level writing.
PIVOT:
- Most likely misconception: students include 3 (whale watching) because the Coorong is on the coast.
- Reteach: re-read the text - whale watching is not mentioned. The text mentioned fishing, tourism, bird watching and boating.
- Fresh re-check: did the text mention whale watching? Thumbs up or down.

TEACHER NOTES:
Multi-answer CFU. Trains students to verify against the text.

WATCH FOR:
- Students who only show one answer when more apply - re-cue: select all that apply.

Stage: CFU | VTLM: Formative Assessment""",


    15: """SAY:
- Now we move to writing.
- Today we are writing the concluding paragraph for our information report.

DO:
- Have the planning notes from Lesson 24 ready.
- Have whiteboards and textas ready.

TEACHER NOTES:
Brief divider.

WATCH FOR:
- Students who do not have their planning notes - direct them to find Lesson 24 booklet pages first.

Stage: Text-Level Writing Divider | VTLM: Engagement""",


    16: """SAY:
- Watch this first.
- A conclusion concludes the report and summarises the main points discussed.
- We have already planned our conclusion using the TSG formula.
- Today we are going to review and edit our work, then write the full paragraph.

DO:
- Read the slide aloud.
- Point to the link to the mentor text.

TEACHER NOTES:
I Do anchor for today's task.

WATCH FOR:
- Students who confuse plan with write - clarify: today we write the full paragraph using yesterday's plan.

Stage: I Do | VTLM: Explicit Teaching""",


    17: """SAY:
- In the concluding paragraph we will - summarise key points, link to the introduction, avoid introducing any new information.
- The third rule is the most important - no new facts in a conclusion.

DO:
- Read the slide aloud.
- Highlight the third rule.

TEACHER NOTES:
We Do for what a conclusion does. The no-new-information rule needs another visit before students start writing.

WATCH FOR:
- Students who plan to add new facts - re-anchor the rule.

Stage: We Do | VTLM: Guided Practice""",


    18: """SAY:
- The introduction uses GST - General, Specific, Thesis.
- The conclusion uses TSG - Thesis, Specific, General.
- Read both with me - introduction first, then conclusion.
- Notice the mirror - the introduction starts wide and zooms in. The conclusion starts on what we discussed and zooms back out.
- Ask: what is the purpose of the thesis statement? What does the specific statement tell us? What do you notice about the general statement?

DO:
- Read both paragraphs as choral reading.
- Point to G, S, T in the introduction.
- Point to T, S, G in the conclusion.
- Cold call non-volunteers for the questions.

TEACHER NOTES:
We Do for linking introduction and conclusion. The general statement uses an appositive and an exclamation - point that out.

WATCH FOR:
- Students who do not see the mirror - point to each pair: General with General, Specific with Specific, Thesis with Thesis.

Stage: We Do | VTLM: Guided Practice""",


    19: """SAY:
- The purpose of the concluding paragraph is to...
- Show me 1, 2, 3 or 4.
- Option 1: introduce new information on the Coorong.
- Option 2: summarise the key points.
- Option 3: give an opinion on the Coorong as a travel destination.
- Option 4: form links with the introduction.

DO:
- Read each option.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- What is the purpose of the concluding paragraph? Show me on three.
- Scan for: fingers showing 2 or 4 (both are valid).
PROCEED:
- If most students are on 2 or 4, move to converting ideas into a conclusion.
PIVOT:
- Most likely misconception: students pick 1 (introduce new information).
- Reteach: a conclusion summarises and links - it never introduces new information.
- Fresh re-check: should we add a new fact? Thumbs up or down.

TEACHER NOTES:
Pre-modelling CFU. Catches the no-new-info confusion before writing.

WATCH FOR:
- Students who pick 3 (give an opinion) - clarify information reports are factual, not opinion-based.

Stage: CFU | VTLM: Formative Assessment""",


    20: """SAY:
- Imagine we have written an information report about pelicans in the Coorong.
- Watch this first - over the next slides, I am going to show you how to convert ideas into a conclusion.

DO:
- Read the slide aloud.
- Set up the next three slides as the modelled walkthrough.

TEACHER NOTES:
I Do anchor for the modelled walkthrough.

WATCH FOR:
- Students who switch off - the next slides involve them more actively.

Stage: I Do | VTLM: Modelling""",


    21: """SAY:
- The thesis statement makes clear what the body paragraphs of the information report discussed.
- Body paragraphs - habitat and ecological importance, diet and behaviours, conservation of pelicans.
- Sentence starter: This report discussed...
- Watch how I build the thesis: This report discussed the habitat, behaviour, and conservation of pelicans in the Coorong.

DO:
- Read the body topics on the slide.
- Build the thesis sentence aloud.
- Point to each body topic as it goes into the thesis.

TEACHER NOTES:
I Do for converting body topics into a thesis statement.

WATCH FOR:
- Students who add new information - re-cue: thesis only covers what we already wrote.

Stage: I Do | VTLM: Modelling""",


    22: """SAY:
- The specific statement addresses the topic directly and uses key terms.
- Main idea: pelicans are important to the ecosystem of the Coorong.
- Key details: control fish populations, their waste enriches the shoreline soil, they indicate environmental health.
- Watch how I build the specific statement: Pelicans are really important in the Coorong because they help keep the ecosystem healthy by eating fish, adding nutrients to the water and land, and showing us how healthy the environment is.

DO:
- Read each detail on the slide.
- Build the specific statement aloud.
- Point to each detail as it goes into the sentence.

TEACHER NOTES:
We Do for converting main ideas into a specific statement. The student version may be shorter - that is fine.

WATCH FOR:
- Students who only mention one detail - prompt them to combine the three.

Stage: We Do | VTLM: Modelling""",


    23: """SAY:
- The general statement summarises the topic and concludes the report.
- Most important thing to remember: pelicans are guardians of the ecosystem.
- Watch how I build the general statement with an appositive and an exclamation.
- Pelicans, the vital guardians of the Coorong's ecosystem, what incredible helpers they are!

DO:
- Read the most-important-thing on the slide.
- Build the general statement.
- Point to the appositive - the vital guardians of the Coorong's ecosystem.
- Point to the exclamation mark.

TEACHER NOTES:
We Do for the general statement. Highlight the appositive and exclamation mark.

WATCH FOR:
- Students who skip the appositive - re-cue: the appositive renames or describes the noun.

Stage: We Do | VTLM: Modelling""",


    24: """SAY:
- Here is the full pelicans conclusion built up.
- This report discussed the habitat, behaviour, and conservation of pelicans in the Coorong.
- Pelicans are really important in the Coorong because they help keep the ecosystem healthy by eating fish, adding nutrients to the water and land, and showing us how healthy the environment is.
- Pelicans, the vital guardians of the Coorong's ecosystem, what incredible helpers they are!

DO:
- Read each part aloud.
- Point to T, S, G as you read.
- Choral read the full conclusion.

TEACHER NOTES:
Full modelled pelicans conclusion. Use it as a reference for the You Do.

WATCH FOR:
- Students who lose track of the structure - re-anchor with T, S, G labels.

Stage: We Do | VTLM: Guided Practice""",


    25: """SAY:
- Ask: what would you find in a concluding paragraph of an information report?
- Show me 1, 2, 3 or 4 - hold up multiple fingers if more than one applies.
- Option 1: thesis statement.
- Option 2: summary.
- Option 3: This report discussed...
- Option 4: new details.

DO:
- Read each option.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers - select all that apply
Script:
- What would you find in a concluding paragraph? Show me on three - more than one is allowed.
- Scan for: fingers showing 1, 2 and 3 (and NOT 4).
PROCEED:
- If most students show 1, 2 and 3, move to the booklet task.
PIVOT:
- Most likely misconception: students pick 4 (new details).
- Reteach: a conclusion uses a thesis statement, summarises, and uses sentence starters like This report discussed - it does NOT add new details.
- Fresh re-check: are new details in a conclusion? Thumbs up or down.

TEACHER NOTES:
Pre-You-Do hinge CFU.

MISCONCEPTIONS:
- Misconception: a conclusion can include new details.
  Why: students think a strong ending needs something new.
  Impact: their conclusions read like a fourth body paragraph.
  Quick correction: a conclusion only summarises what is already there.

WATCH FOR:
- Students who only pick one answer - re-cue: select all that apply.

Stage: CFU | VTLM: Formative Assessment""",


    26: """SAY:
- Quick recap before you write.
- In Lesson 24, we used the TSG formula when planning.
- Thesis statement - what the information report discussed in our three body paragraphs.
- Specific statement - reiterates the significance of the wetland - who or what it matters to, why, and today.
- General statement - concludes the topic, using an appositive.

DO:
- Read each part on the slide.
- Use your fingers to count off T, S, G.

TEACHER NOTES:
Recap before writing.

WATCH FOR:
- Students who do not have their plan ready - direct them to find their Lesson 24 booklet page.

Stage: We Do | VTLM: Guided Practice""",


    27: """SAY:
- Share your planning notes with a partner.
- First, locate your key words and ideas for your thesis statement.
- Next, practise saying it aloud in a full sentence.
- Then, get positive feedback or helpful advice from your partner.

DO:
- Direct students to their Lesson 24 plan.
- Give 5 minutes for partner sharing.
- Circulate and listen for clear or unclear ideas.
- Optional - repeat with the specific and general statements before writing.

TEACHER NOTES:
We Do for partner rehearsal before writing. Catches gaps before they hit the paper.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Read your plan aloud word for word - the partner gives a thumbs up if it sounds like a complete sentence.
EXTENDING PROMPT:
- Task: Suggest one stronger word choice for your partner's specific statement.

WATCH FOR:
- Students who only read their notes silently - prompt them to say it aloud as a full sentence.
- Students who give only positive feedback - prompt them to also offer one helpful improvement.

Stage: We Do | VTLM: Guided Practice""",


    28: """SAY:
- Time to apply this on your own.
- First, turn to the page titled Lesson 25: Write a concluding paragraph for an information report.
- Next, write your thesis, specific and general statement.
- Then, use the checklist on the next slide and in your workbook to check your work.

DO:
- Direct students to the booklet page.
- Circulate and check the first thesis statement for each student.
- Remind students - full sentences this time, not notes.

TEACHER NOTES:
You Do practice. Pull a small group if many students are still struggling with the thesis statement.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: Use the model conclusion as a worked example. Replace the words but keep the structure.
EXTENDING PROMPT:
- Task: Add a second appositive option for the general statement and choose between them when revising.

WATCH FOR:
- Students who copy the pelicans example - prompt them to use Coorong content from their plan.
- Students who add new information - re-anchor the no-new-info rule.
- Students finishing fast - ask them to read aloud and check for the appositive in the general statement.

Stage: You Do | VTLM: Independent Practice""",


    29: """SAY:
- Use the checklist while you write.
- First, make sure you have a thesis, specific and general statement.
- Next, make sure you have an appositive in the general statement.
- Then, check your spelling and punctuation.

DO:
- Direct students to the checklist in the workbook.
- Circulate and reference the checklist as you check.

TEACHER NOTES:
Continuation of the booklet task. Use the checklist on the slide and in the workbook.

WATCH FOR:
- Students who skip the checklist - direct them back to it before they say done.

Stage: You Do | VTLM: Independent Practice""",


    30: """SAY:
- Optional next step - publishing.
- First, put all the pieces together.
- Next, use lined paper or a Word or Google document.
- Then, combine the five paragraphs into one finished piece of writing.

DO:
- Direct students to the publishing format.
- Hand out lined paper or open the document.

TEACHER NOTES:
Optional task. Use this if there is time.

WATCH FOR:
- Students who copy slowly - acknowledge the time it takes and pace accordingly.
- Students who skip paragraph breaks - re-cue: each paragraph starts on a new line.

Stage: You Do (Publishing) | VTLM: Independent Practice""",


    32: """SAY:
- Reference - the introduction model text.
- The Coorong is a narrow, shallow lagoon that stretches across the Limestone coast in South Australia. It is an expansive area of wetlands and waterways that is both culturally and ecologically important. This report will discuss exactly what the Coorong is, where it is located, its cultural significance, and the flora and fauna that rely on it for their survival.
- Notice the GST formula.

DO:
- Display the text.
- Point to each part - General, Specific, Thesis.

TEACHER NOTES:
Reference slide for the modelled introduction. Use it to compare with the conclusion.

WATCH FOR:
- Students who confuse GST with TSG - clarify: introduction starts general and zooms in.

Stage: Reference (Mentor Text) | VTLM: Modelling""",


    33: """SAY:
- Reference - body paragraph 1.
- Topic: what the Coorong is and where it is located.
- This paragraph names the Murray-Darling basin connection, the length, the Younghusband Peninsula, and the mix of fresh, brackish and saltwater.

DO:
- Display the text.
- Point to the topic-introducing sentence.

TEACHER NOTES:
Reference slide. Use it as a knowledge anchor.

WATCH FOR:
- Students who try to copy this whole paragraph - re-cue: this is for reference only.

Stage: Reference (Mentor Text) | VTLM: Modelling""",


    34: """SAY:
- Reference - body paragraph 2.
- Topic: cultural significance.
- The paragraph names the Ngarrindjeri people, the Kungun Ngarrindjeri Yunnan agreement, and joint management.

DO:
- Display the text.
- Point to the topic-introducing sentence.

TEACHER NOTES:
Reference slide. Pronounce Ngarrindjeri and KNY clearly.

SENSITIVITY ADVISORY:
- What it is: First Nations cultural content.
- Framing language: Use the Ngarrindjeri name with respect.
- Watch for: Students unsure how to pronounce the names.
- Protocol: Refer back to the Cultural Sensitivity slide if needed.

WATCH FOR:
- Students who skip the Ngarrindjeri reference - re-anchor.

Stage: Reference (Mentor Text) | VTLM: Modelling""",


    35: """SAY:
- Reference - body paragraph 3.
- Topic: flora and fauna.
- The paragraph names vulnerable species, local birds and migratory species from Alaska and Siberia, and the Bonney Upwelling.

DO:
- Display the text.
- Point to the topic-introducing sentence.

TEACHER NOTES:
Reference slide. The Bonney Upwelling and migration distance are vivid facts.

WATCH FOR:
- Students who do not know what fauna or flora means - quick gloss: animals and plants.

Stage: Reference (Mentor Text) | VTLM: Modelling""",


    36: """SAY:
- Reference - the modelled conclusion.
- This report detailed the location of the Coorong and its cultural and ecological significance. The importance of this expansive wetland area is enormous for both past and present-day people, flora, and fauna. The Coorong, a long shallow lagoon filled with life, is truly an iconic Australian landscape!

DO:
- Display the text.
- Point to T, S, G as you read.

TEACHER NOTES:
Reference slide. Use it as the model for the You Do task.

WATCH FOR:
- Students who try to copy this exactly - prompt them to use it as a structure model, not a script.

Stage: Reference (Mentor Text) | VTLM: Modelling""",
}


# ============================================================
# APPLIER
# ============================================================
def suppress_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for tag_name in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag_name)):
            pPr.remove(el)
    pPr.append(etree.SubElement(pPr, qn("a:buNone")))


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    parts = []
    for line in text.split("\n"):
        if line.startswith("- "):
            parts.append((line[2:], True))
        else:
            parts.append((line, False))

    p = notes_tf.paragraphs[0]
    p.text = parts[0][0]
    if not parts[0][1]:
        suppress_bullet(p)
    for content, is_bullet in parts[1:]:
        p = notes_tf.add_paragraph()
        p.text = content
        if not is_bullet:
            suppress_bullet(p)


def apply(src_name, notes):
    src = SRC_DIR / src_name
    out_name = src_name.replace(".pptx", " - with notes.pptx")
    out = SRC_DIR / out_name
    pres = Presentation(str(src))
    n = 0
    for i, slide in enumerate(pres.slides, 1):
        if i in notes:
            set_notes(slide, notes[i])
            n += 1
    pres.save(str(out))
    print(f"Applied notes to {n} slides -> {out.name}")


def main():
    apply(L21_FILE, LESSON_21_NOTES)
    apply(L22_FILE, LESSON_22_NOTES)
    apply(L23_FILE, LESSON_23_NOTES)
    apply(L24_FILE, LESSON_24_NOTES)
    apply(L25_FILE, LESSON_25_NOTES)


if __name__ == "__main__":
    main()
