"""
Add teacher notes to: Discovery 1 Maths Planner Term 2, Week 3 - Additive Thinking and 2D shape.pptx

Year 1 numeracy. Preserves slide faces and existing meaningful notes; fills empties.
Saves a new copy with ' - with teacher notes' suffix.
"""

from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
from lxml import etree
import os

SRC = r"C:/Users/09560329/Downloads/use this one Discovery 1 Maths Planner Term 2, Week 3 - Additive Thinking and 2D shape.pptx"
DST = r"C:/Users/09560329/Downloads/use this one Discovery 1 Maths Planner Term 2, Week 3 - Additive Thinking and 2D shape - with teacher notes.pptx"


# ---------------------------------------------------------------------------
# Notes per slide (1-indexed). Plain text, ASCII-safe, no em dashes.
# ---------------------------------------------------------------------------

NOTES = {

# ============================================================
# Front matter (admin slides)
# ============================================================

1: """TEACHER NOTES:
Title slide for Term 2 Week 3, Discovery 1. Use as you walk into the room while you set up materials. No student action needed.
""",

2: """TEACHER NOTES:
Teacher reference. Use this slide before the week begins to align planning with the Victorian Curriculum 2.0 Learning Continuum links shown. Not a student-facing slide.

SOURCES:
Scope and Sequence: Foundation, Discovery 1, Discovery 2, Challenge and Enrichment Learning Continuums.
""",

3: """TEACHER NOTES:
Weekly overview slide. Read through with co-teachers before lessons. The week covers two main areas: additive thinking with money (Lessons 1 to 2) and 2D shape (Lessons 3 to 5).
""",

4: """TEACHER NOTES:
Diamond Creek East Numeracy Learning Model. Reference only. Each lesson follows this shape: Daily Review, Number Fluency, Launch, Explore, Summarise, with Hinge checks placed across the lesson. Use the timing as a guide, not a strict cap.
""",

5: """TEACHER NOTES:
Icon legend slide. Use early in the week if students are new to the icons. The four icons shown here are: speak whole-class, write on a template or in a book, sign in Auslan, and Turn and Talk with a table partner.
""",

6: """TEACHER NOTES:
Icon legend continued. Difficulty bars let students choose their entry level. Other icons here: teacher teaching (explicit instruction) and mini-whiteboards. Skip if students already know the icons.
""",

7: """TEACHER NOTES:
Final icon legend slide. Thumbs up or down is a quick check signal. Concrete materials icon means manipulatives can be used. The skip icon means the slide is optional if students already understand.
""",

8: """TEACHER NOTES:
Lesson 1 divider. Topic is Additive Thinking. Use as a clear visual reset before starting the lesson. Have play money or images of $5, $10, $1 ready.
""",

# ============================================================
# Lesson 1 - Counting and Comparing Money (Session 11)
# ============================================================

9: """TEACHER NOTES:
Lesson 1 title slide. Session 11 of 11 in the Additive Thinking sequence. Today's focus is counting collections of money up to $20 and deciding 'can I afford it'. Have play money and the worksheet ready.
""",

10: """SAY:
- Let's warm up with some ten frame problems
- Look at the ten frames. Work out the total on your whiteboard
- Remember: a full ten frame is 10

DO:
- Read each problem with the class
- Allow 30 to 40 seconds per problem
- Circulate to check students are recognising the full ten frame as 10

TEACHER NOTES:
Daily Review activates ten frame and partitioning fluency from prior weeks. The 'ten plus more' pattern feeds directly into starting with the largest denomination later in the lesson.

WATCH FOR:
- Students who say 19 or 14 quickly: ready for today's counting
- Students counting all dots one by one from 1: prompt them to see 10 first, then count on

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

11: """SAY:
- Now we are partitioning. Split the number into tens and ones
- Show me on your board. Ten and how many more?
- The first part is 10. The second part is the leftover ones

DO:
- Read each partition prompt
- Give 20 seconds for boards
- Reveal the partition in a tens-and-ones format

TEACHER NOTES:
Partitioning to tens and ones underpins the way students will count money: the $10 note is the 'ten', then add the smaller coins.

WATCH FOR:
- Students writing 18 = 10 + 8 confidently: ready
- Students splitting unevenly (e.g. 18 = 9 + 9): re-anchor on the ten frame and try again

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

12: """SAY:
- Quick skip count. Eyes on the screen
- Count by 2s with me to 20
- Count by 5s with me to 50
- Count by 10s with me to 100

DO:
- Lead a brisk choral count. Track the numerals with your finger
- Run each count once, then once more if shaky
- Keep the pace fast. This is fluency, not new teaching

TEACHER NOTES:
Quick skip counting fluency. The 5s and 10s links directly to counting $5 and $10 notes later. Do not stop to teach a strategy here.

WATCH FOR:
- Students mouthing rather than voicing: prompt 'voices on'
- Students stuck on the bridge (29 to 30, 49 to 50): point to the numeral and chorus once more

[Stage 1: Number Fluency | VTLM 2.0: Retention and Recall]
""",

13: """SAY:
- Here is what we are learning today. Read it with me
- Read from slide: 'We are learning to count collections of money and decide if we can afford something'
- Now our success criteria

DO:
- Choral read the LI then each SC
- Hold up a real or play $10, $5 and $1 as you read SC1

TEACHER NOTES:
SC1 is the floor (count money), SC2 is the core target (compare to cost), SC3 extends to the 'can I afford it' decision. The exit ticket checks SC1 and SC3 together.

WATCH FOR:
- Students unsure of 'denominations': say it means the different coins and notes (gold and silver coins, paper notes)

[General: LI/SC | VTLM 2.0: Clear Learning Intention]
""",

14: """SAY:
- Watch first. I have a $10 note, a $5 note and a $1 coin
- I always start with the BIGGEST. So I start with $10
- $10 plus $5 makes $15. Now I add the $1. $15 plus $1 makes $16
- I have $16 altogether

DO:
- Hold up the matching note or coin as you say each amount
- Point to the strategy box on the slide as you count
- Slide your finger from $10 to $15 to $16 to show the running total

TEACHER NOTES:
This is the modelling slide for the lesson's key strategy: start with the largest. Keep your think-aloud short and tied to the visual on the slide.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: model again with two amounts only ($10 + $5) before adding the $1
EXTENDING PROMPT:
- Task: ask 'what if my coin was $2 instead of $1? What would I have?'

WATCH FOR:
- Students who try to start with the smallest: bring them back and re-count starting at $10

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

15: """SAY:
- I have $16. Can I afford each item on the menu?
- Pizza is $16. I have $16. They are the same. So I CAN afford it
- Fried rice is $20. I have $16. $16 is less than $20. I CANNOT afford it
- Juice is $4. I have $16. I have more than enough. I CAN afford it

DO:
- Point to each menu item as you read it
- Hold your $16 up against each price
- Use a clear voice for 'can' and 'cannot' so students hear the difference

TEACHER NOTES:
This builds SC2 and SC3. Make the comparison visible and physical: hold the money beside the price. Stay with same/less/more language before any maths symbols.

MISCONCEPTIONS:
- Misconception: students think 'I can buy it as long as I have some money'
  Why: they do not yet compare two numbers
  Impact: they say YES to everything
  Quick correction: line up the two amounts and ask 'is mine bigger, smaller or the same?'

WATCH FOR:
- Students who say 'yes' to fried rice: re-count to $16, then point to $20 and say 'is $16 bigger or smaller?'

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

16: """SAY:
- Quick check on your boards
- I have a $5 note and a $2 coin. How much altogether?
- Start with the BIGGEST first. 15 seconds. Show me

DO:
- Hold up the $5 note then the $2 coin as you read
- Count down 15 seconds. Boards up on 'show me'
- Scan all boards from back to front before revealing

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
- Ask: $5 plus $2. How much altogether? Expected: $7
- Scan for: $7 written clearly with the dollar sign
PROCEED:
- >=80% show $7. Move to the We Do
PIVOT:
- Most likely: students count from $2 first and get muddled
- Reteach: hold the $5 first, say 'five', then count on 'six, seven' as you move the $2 across
- Re-check: $5 + $1 on boards. Show me

TEACHER NOTES:
Hinge check before guided practice. Do not move on if many students start from the smaller coin or count all from 1.

WATCH FOR:
- Boards showing $7: ready
- Boards showing $52 or 25: students have written the digits side by side. Re-anchor on counting on

[Stage 2: I Do | VTLM 2.0: Check for Understanding]
""",

17: """SAY:
- Let's do this one together. Collection 1: $10, $2, $2
- Start with the biggest. Whisper to your partner: which one first?
- Now count it on your board

DO:
- Read Collection 1 aloud
- Cue Turn and Whisper for the first step
- Take a board scan for the total before revealing
- Then move to compare with each menu price together

TEACHER NOTES:
Stay in supported practice. Build the answer with the class. Use the menu to check 'can I afford?' for each total. Steps in existing notes apply: count from biggest, write the total, compare to each price, decide YES or NO.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: only do Collection 1 plus one menu item to start
EXTENDING PROMPT:
- Task: ask 'which menu items can I afford with both collections?'

WATCH FOR:
- Students confidently starting with $10: ready for You Do
- Students hesitating on which is biggest: ask them to point to the largest note

[Stage 3: We Do | VTLM 2.0: Guided Practice]
""",

18: """SAY:
- Your turn. Independent practice
- Step 1: count the money. Start with the biggest
- Step 2: write the total
- Step 3: compare to the prices
- Step 4: write YES or NO

DO:
- Direct students to the worksheet (Session 11 Worksheet)
- Set a 10 minute window to start
- Circulate to the support group first. Use the Enabling Scaffold sheet if needed

TEACHER NOTES:
This is the SC2/SC3 evidence task. Students must show counting and a YES or NO decision. Early finishers move to the Extension sheet (same amount different ways).

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: use the Enabling Scaffold (denominations pre-sorted with counting guide)
EXTENDING PROMPT:
- Task: use the Extension sheet (make the same amount different ways)

WATCH FOR:
- Students who count and write a total but skip the compare step: prompt 'now check the prices'
- Students who write YES with no working: ask them to show the count

[Stage 4: You Do | VTLM 2.0: Independent Practice]
""",

19: """SAY:
- Show me what you can do on your own
- Question 1: count the money. $10, $5 and $2. Write the total
- Question 2: a skateboard costs $18. Can you afford it? Circle YES or NO
- Show your counting on the page

DO:
- Hand out exit tickets or direct students to the page
- Set 5 minutes
- Collect or sight every ticket. No support during the ticket
- Mark students who need re-teaching for the next lesson

TEACHER NOTES:
Exit ticket assesses SC1 (count to $17) and SC3 (compare $17 to $18 and decide NO). Students who write YES and $17 have not yet compared correctly.

WATCH FOR:
- Total $17 with NO and visible counting: secure
- Total wrong (e.g. $16 or $20): SC1 not yet secure
- Total $17 with YES: comparison error. Mark for follow up

[Stage 5: Exit Ticket | VTLM 2.0: Check for Understanding]
""",

20: """SAY:
- Let's check our success criteria
- SC1: I can count money in different denominations. Thumbs?
- SC2: I can compare my money to the cost of something. Thumbs?
- SC3: I can decide whether I can afford something. Thumbs?
- Turn and Talk: when counting money, which one do you start with? Why?

DO:
- Display SC on screen
- Run thumbs check for each
- Allow 30 seconds for Turn and Talk
- Cold call 2 to 3 students. Listen for 'start with the biggest because it is easier to add on'

TEACHER NOTES:
The Turn and Talk targets the counting strategy (start with the largest denomination). This is the key takeaway for efficient money counting.

WATCH FOR:
- Students who say 'start with the biggest': they have the strategy
- Students who are unsure: revisit in tomorrow's launch

[General: Closing | VTLM 2.0: Review and Reflect]
""",

21: """SAY:
- These are the resources for today's lesson

DO:
- Confirm all four PDFs are printed before the lesson
- Have the Enabling Scaffold ready for the support group

TEACHER NOTES:
Resource slide with clickable PDF links. Worksheet and Answer Key are core. Enabling Scaffold has denominations pre-sorted. Extension asks students to make the same amount different ways.

WATCH FOR:
- Teacher reference only

[General: Resources | VTLM 2.0: Planning]
""",

# ============================================================
# Lesson 2 - Adding and Subtracting Money (Session 10)
# ============================================================

22: """TEACHER NOTES:
Lesson 2 divider. Topic continues with Additive Thinking. Have play money and the toy catalogue worksheet ready before starting.
""",

23: """TEACHER NOTES:
Lesson 2 title slide. Session 10 in the Additive Thinking sequence. Today's focus is adding two prices and subtracting from $20 to find change.
""",

24: """SAY:
- Let's warm up with some ten frame problems
- Look at the ten frames. Work out the total on your whiteboard
- Remember: a full ten frame is 10

DO:
- Read each problem with the class
- Allow 30 to 40 seconds per problem
- Circulate to check students are recognising the full ten frame as 10

TEACHER NOTES:
Daily Review activates ten plus more. The 'ten plus a small number' pattern supports today's two-step money problems where the total often crosses a friendly ten.

WATCH FOR:
- Students writing 17 and 13 quickly: ready
- Students counting all dots from 1: prompt 'see the 10 first, then count on'

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

25: """SAY:
- Now we are partitioning. Split the number into tens and ones
- Show me on your board
- The first part is 10. The second part is the leftover ones

DO:
- Read each partition prompt
- Give 20 seconds for boards
- Confirm the partition before moving on

TEACHER NOTES:
Partitioning prepares students for crossing the ten when adding two prices today (e.g. $9 + $6 = $15 by partitioning the 6 into 1 and 5).

WATCH FOR:
- 14 = 10 + 4 and 17 = 10 + 7: ready
- Students writing two equal parts (e.g. 14 = 7 + 7): re-anchor on the ten frame

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

26: """SAY:
- Show me the partition on your board
- 15 is 10 plus how many? 18 is 10 plus how many? 12? 20?
- Partition means to split into parts. 15 splits into 10 and 5

DO:
- Point to the partition definition on the slide
- Run all four prompts. Boards up after each
- Reveal the partition for any prompt students missed

TEACHER NOTES:
Fluency on tens-and-ones partitioning. This is the same move students will use in the I Do when bridging through ten ($9 + $6 by taking 1 from the 6).

WATCH FOR:
- 20 = 10 + 10 (or 20 + 0): both correct, accept
- Students writing 12 = 6 + 6: not the tens partition. Re-anchor

[Stage 1: Number Fluency | VTLM 2.0: Retention and Recall]
""",

27: """SAY:
- Here is what we are learning today. Read it with me
- Read from slide: 'We are learning to add and subtract money'
- Now our success criteria

DO:
- Choral read the LI then each SC
- Hold up play money as you read each SC

TEACHER NOTES:
SC1 is the floor (add dollar amounts), SC2 is the core target (subtract to find change), SC3 extends to choosing a strategy. The exit ticket checks SC1 and SC2 together.

WATCH FOR:
- Students unsure about 'change': say it means the money you get back when you pay too much

[General: LI/SC | VTLM 2.0: Clear Learning Intention]
""",

28: """SAY:
- Let's buy two toys. A toy car is $9 and a soccer ball is $6
- I want the total cost. The parts are 9 and 6. The whole is what I am finding
- I will use bridging to ten. $9 plus $1 is $10. Take 1 from the 6, that leaves 5
- $10 plus $5 is $15. The total cost is $15

DO:
- Point to the number bond on the slide
- Trace the bridging-to-ten move with your finger
- Slow down on the partition: 6 = 1 + 5

TEACHER NOTES:
Modelling the bridging-to-ten strategy with money. The number bond shows parts and whole. Tie the language to the partitioning practised in Number Fluency.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: model with smaller numbers first ($5 + $3) before bridging
EXTENDING PROMPT:
- Task: ask 'what other strategy could we use? Could we count on?'

MISCONCEPTIONS:
- Misconception: students count all from 1 ($1, $2, ..., $9, $10... $15)
  Why: they have not internalised counting on
  Impact: slow and error-prone with larger amounts
  Quick correction: cover the 9 with your hand and prompt 'start with 9 in your head, count on'

WATCH FOR:
- Students who copy the answer with no strategy: ask one to retell the bridge step

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

29: """SAY:
- The total cost was $15. I paid with $20. How much change do I get?
- I subtract: $20 take away $15
- I will count on. 15... 16, 17, 18, 19, 20. That is 5 more
- My change is $5

DO:
- Point to $20, then to $15 on the slide
- Tap each number as you count on: 16, 17, 18, 19, 20
- Use fingers held up to track the 5 more

TEACHER NOTES:
Subtraction by counting on. This is a Year 1-friendly strategy. Avoid formal column subtraction at this stage.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: count on with a number line first ($20 - $18 = ?)
EXTENDING PROMPT:
- Task: ask 'what change would I get if I paid with $50?'

MISCONCEPTIONS:
- Misconception: students count $15 itself when counting on
  Why: they include the start number
  Impact: change is one too many
  Quick correction: model 'we start FROM 15, the next number is 16'

WATCH FOR:
- Students who say 'change is 6': they counted 15 itself. Reset and count on from 16

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

30: """SAY:
- Quick check on your boards
- A teddy bear is $7 and a puzzle is $8. What is the total cost?
- 20 seconds. Show me

DO:
- Read both prices clearly
- Count down 20 seconds. Boards up on 'show me'
- Scan all boards before revealing the answer

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
- Ask: $7 plus $8. Total cost? Expected: $15
- Scan for: $15 written clearly with the dollar sign and a visible strategy if shown
PROCEED:
- >=80% show $15. Move to the We Do
PIVOT:
- Most likely: students answer $14 or $16 (off by one)
- Reteach: model bridging with doubles: 7 + 7 = 14, plus 1 more = 15
- Re-check: $7 + $6 on boards. Show me

TEACHER NOTES:
Hinge check before guided practice. The numbers cross ten so this checks the bridging strategy.

WATCH FOR:
- Boards showing $15: ready
- Boards showing 78 or 87: students have stuck the digits together. Re-model addition

[Stage 2: I Do | VTLM 2.0: Check for Understanding]
""",

31: """SAY:
- Let's do this together. Kite $8 plus crayons $4. What's the total?
- Now you spent $12. From $20, what's the change?
- Next: doll $11 plus book $5. Total? Change from $20?

DO:
- Read the first row aloud
- Cue 'show me' for the total
- Count on together for the change. 12... 13, 14... 20
- Repeat for row 2 (doll plus book)

TEACHER NOTES:
Co-build the answers with the class. Existing notes apply: add the prices first, then subtract from $20. Use the count-on strategy modelled in I Do.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: do only the total step. Do not also do the change step
EXTENDING PROMPT:
- Task: ask 'which item left more change? Why?'

WATCH FOR:
- Students who add the prices but stop there: prompt 'now find the change'
- Students confidently chaining add then subtract: ready for You Do

[Stage 3: We Do | VTLM 2.0: Guided Practice]
""",

32: """SAY:
- Shopping role play. One of you is the customer. One of you is the shopkeeper
- Customer: you start with $20. Choose 2 items
- Shopkeeper: add the prices. Then take that off $20 to find change
- Then swap roles

DO:
- Pair students up. Hand out the role play sheet (linked above)
- If you have play money, give each pair a $20 collection
- Circulate. Listen for the count-on strategy
- Stop the class halfway and prompt the swap

TEACHER NOTES:
Use real or play money if available. The role play sheet works without money. Either way, students must record the total and change in writing.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: work with the teacher group. Buy ONE item only and find change first
EXTENDING PROMPT:
- Task: buy 3 items within $20. Find the change

WATCH FOR:
- Customers spending more than $20: they need a check 'do you have enough?'
- Shopkeepers not subtracting: redirect to the steps

[Stage 4: You Do | VTLM 2.0: Independent Practice]
""",

33: """SAY:
- Show me what you can do on your own
- Question 1: a robot is $7 and a yo-yo is $5. What is the total cost?
- Question 2: you pay with $20. How much change?
- Show your working

DO:
- Hand out exit tickets
- Set 5 minutes
- Collect or sight every ticket. No support during the ticket

TEACHER NOTES:
Exit ticket assesses SC1 (total $12) and SC2 (change $8). Working should show either bridging for the total or count-on for the change. Mark for follow up if either step is missing.

WATCH FOR:
- Total $12 and change $8 with working: secure
- Total correct, change wrong: SC2 not yet secure. Re-teach count-on
- Both steps wrong: SC1 not yet secure

[Stage 5: Exit Ticket | VTLM 2.0: Check for Understanding]
""",

34: """SAY:
- Let's check our success criteria
- SC1: I can add dollar amounts to find the total. Thumbs?
- SC2: I can subtract from dollar amounts to find change. Thumbs?
- SC3: I can use a strategy to solve a money problem. Thumbs?
- Turn and Talk: which step do you do first when you go shopping? Add or subtract?

DO:
- Display the SC on screen
- Run thumbs check for each
- Allow 30 seconds for Turn and Talk
- Cold call 2 to 3 students. Listen for 'add the prices first, then take it off the $20'

TEACHER NOTES:
The Turn and Talk targets the order of operations in a shopping problem (add first, then subtract). This is today's key takeaway.

WATCH FOR:
- Students who say 'add first': they have the order
- Students who only say 'subtract': re-anchor with a quick example tomorrow

[General: Closing | VTLM 2.0: Review and Reflect]
""",

35: """SAY:
- These are the resources for today's lesson

DO:
- Confirm all four PDFs are printed before the lesson
- Have the Enabling Scaffold ready for the support group

TEACHER NOTES:
Resource slide with clickable PDF links. Worksheet and Answer Key are core. Enabling Scaffold reduces to a single item from $10 or $20 with visual supports. Extension is buying 3 items within $20.

WATCH FOR:
- Teacher reference only

[General: Resources | VTLM 2.0: Planning]
""",

# ============================================================
# Lesson 3 - 2D Shape Identification
# ============================================================

36: """TEACHER NOTES:
Lesson 3 divider. Topic switches to 2D Shape. Have shape cards or wooden 2D shapes ready (circle, square, rectangle, triangle).
""",

37: """SAY:
- Skip count by 2s on the screen
- Count with me. 2, 4, 6...
- Now count by 2s to find the total

DO:
- Point to each pair as you count
- Lead a brisk choral count
- Use the picture to anchor the skip count (each pair adds 2)

TEACHER NOTES:
Daily Review on skip counting by 2s. Tracks the totals 4, 6, 8, 10, 12, 14, 16, 18, 20, 22, 24. Keep brisk. This is fluency, not new teaching.

WATCH FOR:
- Students mouthing rather than voicing: prompt 'voices on'
- Students stuck at 18 to 20: chorus once more

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

38: """SAY:
- 15 is the whole. Partition it into two parts on your board
- Answers may vary. Show me one way

DO:
- Read 15 aloud
- Give 30 seconds for boards
- Take 3 student examples (e.g. 10 + 5, 8 + 7, 14 + 1)
- Reveal a few partitions side by side

TEACHER NOTES:
Open partition task. Multiple correct answers. Use this to revise that 15 can be split many ways.

WATCH FOR:
- Students writing 15 = 10 + 5 first: secure with tens partitioning
- Students who write 15 = 15 + 0: accept once, then prompt 'try a different split'
- Students whose two parts do not sum to 15: re-count together

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

39: """SAY:
- 11 is the whole. Now partition it into three parts
- Answers may vary. Show me one way

DO:
- Read 11 aloud
- Give 30 to 40 seconds for boards (three parts is harder)
- Take 3 student examples (e.g. 5 + 5 + 1, 4 + 4 + 3, 10 + 1 + 0 if accepted)
- Reveal a few partitions side by side

TEACHER NOTES:
Three-part partitioning extends the previous slide. Multiple correct answers. Some Year 1 students will need teacher prompting to find a third part.

WATCH FOR:
- Students who give two parts only: prompt 'now find a third part'
- Students whose three parts do not sum to 11: count the parts together using fingers

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

40: """SAY:
- Number Fluency. Skip count by 10s
- Today we start at numbers other than 10
- I'll say a start number. You count on by 10. Ready?

DO:
- Choose a start number (e.g. 4): count 4, 14, 24, 34...
- Repeat with another start (e.g. 7, 12, 25)
- Run each chain to about 5 jumps
- Boards optional

TEACHER NOTES:
Existing note applies: skip count by 10s starting at numbers other than 10. This builds the place value pattern that the ones digit stays the same when adding 10.

WATCH FOR:
- Students who keep the pattern (4, 14, 24): secure
- Students who restart at 10 (e.g. 4, 14, 20): re-anchor 'we keep the ones, change the tens'

[Stage 1: Number Fluency | VTLM 2.0: Retention and Recall]
""",

41: """SAY:
- What we want to LEARN today. Read with me
- 'We are learning to name, identify and describe two-dimensional shapes'
- Now our success criteria

DO:
- Choral read the LI then each SC
- Show a shape card (e.g. a triangle) as you read 'name'
- Run a finger along the sides as you read 'features'

TEACHER NOTES:
The blank in the title is for 'learn'. SC1 is the floor (name), SC2 is the core (label features such as sides and vertices), SC3 is the extension (identify polygons by criteria).

WATCH FOR:
- Students unsure about 'features': say it means the parts of the shape, like sides and corners
- Students unsure about 'polygon': you'll teach this word in the Launch (closed shape with straight sides)

[General: LI/SC | VTLM 2.0: Clear Learning Intention]
""",

42: """SAY:
- We are going to learn this vocabulary today

DO:
- Pause briefly on this vocabulary header

TEACHER NOTES:
Vocabulary section header. Sets up the next few slides. No student action.
""",

43: """SAY:
- Let's activate what you already know

DO:
- Pause briefly on this header

TEACHER NOTES:
Activate prior learning header. Sets up the next three identifying slides (circle, rectangle, triangle).
""",

44: """SAY:
- Look at the shape on the screen. What is its name?
- Some of you may remember. Show me on your whiteboard
- The name is circle. A circle has 1 curved line

DO:
- Hold up a circle shape card or trace a circle in the air
- Boards up after 10 seconds. Reveal name
- Trace the curved line on the shape with your finger

TEACHER NOTES:
First identification slide. The W icon indicates whiteboards. Activates prior knowledge of the circle and introduces the feature language '1 curved line'.

WATCH FOR:
- Students who write 'oval' or 'ball': ask 'is it flat or 3D? Is it round all the way?'
- Students who write 'circle': secure

[Stage 2: I Do | VTLM 2.0: Activate Prior Knowledge]
""",

45: """SAY:
- Look at this shape. What is its name?
- Show me on your whiteboard
- The name is rectangle. A rectangle has 4 sides and 4 straight lines

DO:
- Hold up a rectangle shape card
- Boards up after 10 seconds. Reveal name
- Touch each of the 4 sides as you count: 1, 2, 3, 4

TEACHER NOTES:
Activates the rectangle and the side-counting routine. Two short ends and two long sides. Some students may say 'square' for a long thin one.

WATCH FOR:
- Students who write 'square': prompt 'are all 4 sides the same? If two are long and two are short, it's a rectangle'
- Students who count more than 4 sides: re-trace with them

[Stage 2: I Do | VTLM 2.0: Activate Prior Knowledge]
""",

46: """SAY:
- Look at this shape. What is its name?
- Show me on your whiteboard
- The name is triangle. A triangle has 3 sides and 3 straight lines

DO:
- Hold up a triangle shape card
- Boards up after 10 seconds. Reveal name
- Touch each of the 3 sides as you count: 1, 2, 3

TEACHER NOTES:
Activates triangle. Tri = three. Make the link between the prefix and the side count if helpful.

WATCH FOR:
- Students who write 'pyramid': pyramid is 3D. Triangle is flat
- Students who write 'triangle': secure

[Stage 2: I Do | VTLM 2.0: Activate Prior Knowledge]
""",

47: """SAY:
- Time to learn something new. For this Launch you will need...

DO:
- Read the materials list shown on the slide
- Confirm whiteboards and markers ready
- Have a circle, rectangle, triangle and square card ready to hold up

TEACHER NOTES:
Launch setup slide. Use to check materials before teaching. No student response needed.
""",

48: """SAY:
- Look at these shapes. They are 2D shapes
- 2D shapes are FLAT and CLOSED
- I can lay them on the table because they are flat
- They are closed because the lines join up

DO:
- Hold up a real 2D shape card and a 3D object (e.g. a ball or a block)
- Press the 2D card flat on the table
- Trace each shape on the slide with your finger to show 'closed'

TEACHER NOTES:
First definition slide. Two key features: flat and closed. Hold up a 3D object as a non-example so students can see the difference.

WATCH FOR:
- Students confusing flat shapes with 3D objects: hold the ball next to the circle to compare

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

49: """SAY:
- 2D shapes have length and width
- Length is how long. Width is how wide
- Watch as I trace the length and the width

DO:
- Hold up a rectangle card. Trace the long way (length) and the short way (width)
- Repeat for the circle (across one way and across the other)
- Have students trace the air with two fingers: one finger for length, one for width

TEACHER NOTES:
Builds on the previous slide. Length and width are introduced as the two flat measurements. Avoid the word 'depth' or 'thickness' here. That belongs to 3D.

WATCH FOR:
- Students who can show length and width on the air: secure
- Students who only trace one way: re-model both directions

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

50: """SAY:
- Hinge question. Which shapes are two-dimensional?
- Look at A and look at B. Choose carefully
- Show me on your fingers. 1 finger for A. 2 fingers for B. Both fingers if both

DO:
- Give 20 seconds of silent thinking
- Cue 'show me' for the finger response
- Scan every student before taking responses

CFU CHECKPOINT:
Technique: Finger Show Me
Script:
- Ask: which shapes are 2D? Expected: only the flat ones (whichever option is the flat shape)
- Scan for: most students choosing the flat option
PROCEED:
- >=80% choose the correct option. Move on
PIVOT:
- Most likely: students choose the 3D object because it 'looks like a shape'
- Reteach: hold the option up. Can it lay flat on the table? If not, it's 3D
- Re-check: hold up a new pair (e.g. a ball vs a circle card). Show me

TEACHER NOTES:
Hinge slide. Decides whether to move to the 'shapes around us' slide or stop and re-teach 'flat and closed'.

WATCH FOR:
- Students who pick the 3D shape: they need the flat-vs-3D contrast again

[Stage 2: I Do | VTLM 2.0: Check for Understanding]
""",

51: """SAY:
- Look at the picture. What 2D shapes can you see?
- Turn and Talk to your partner. Find at least 2 shapes
- I'll cold call after the talk

DO:
- Allow 30 seconds for Turn and Talk
- Cold call 3 students for examples
- Point to each shape they name on the picture

TEACHER NOTES:
Existing note applies: squares on checkerboards and pieces of paper are 2D shapes. Ask students to name 2D shapes they can see in the classroom too.

WATCH FOR:
- Students naming 3D objects (e.g. 'a ball'): say 'a ball is 3D. What 2D shape do you see on the ball?'
- Students naming 2D shapes from the picture: secure

[Stage 2: I Do | VTLM 2.0: Connect to Real World]

SOURCES:
Image credit: Pixabay and Pexels
""",

52: """SAY:
- Look at the circle. A circle is one curved line that joins up
- The line goes all the way around and meets itself
- Show me thumbs up if you can see one curved line

DO:
- Trace the curved line on the slide with your finger
- Hold up a circle card and trace it again
- Run the thumbs check

TEACHER NOTES:
Existing note applies: all points on the curved line are the same distance from the centre. Keep the language simple for Year 1: 'one curved line that joins up'.

WATCH FOR:
- Thumbs up: secure with the definition
- Students who say 'two lines': they may be seeing top and bottom as separate. Re-trace as one line

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

53: """SAY:
- This is an oval. An oval is also one curved line that joins up
- Watch this. If I fold it in the middle, both halves look the same
- But only one way. An oval is squashed at one end

DO:
- Trace the oval on the slide
- Use a paper oval cut-out and fold it once down the long way (halves match)
- Try to fold it the other way (halves do not match)

TEACHER NOTES:
Existing note applies: an oval is narrower at one end than an ellipse. It can only be folded into equal halves one way. Year 1 students do not need to distinguish oval from ellipse formally; the fold-test is enough.

WATCH FOR:
- Students who say 'it's a circle': trace the squashed end and ask 'is it round all the way?'
- Students who try to fold the cut-out the other way and notice it doesn't match: secure

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

54: """SAY:
- Hinge. Which shape is a circle? A or B?
- Look carefully. Choose. Show me on your fingers

DO:
- Give 15 seconds of silent thinking
- Cue 'show me'
- Scan every student before revealing

CFU CHECKPOINT:
Technique: Finger Show Me
Script:
- Ask: which is the circle? Expected: A (the truly circular one)
- Scan for: students choosing the option whose curved line stays the same distance from the centre
PROCEED:
- >=80% choose A. Move on
PIVOT:
- Most likely: students choose B (an oval or wonky shape)
- Reteach: existing note applies: discuss why option B is not a circle. Some of the curved line is closer to the middle than other parts
- Re-check: show two new options. Show me

TEACHER NOTES:
Hinge slide. Confirms students can spot a true circle vs a near-circle.

WATCH FOR:
- Students who point to the oval: they need the 'same distance from middle' check again

[Stage 2: I Do | VTLM 2.0: Check for Understanding]
""",

55: """SAY:
- This is a square. Write the number of sides and vertices on your board
- Sides are the straight lines. Vertices are the corners where lines meet
- Show me

DO:
- Hold up a square card. Touch each side as you count: 1, 2, 3, 4
- Touch each vertex (corner) as you count: 1, 2, 3, 4
- Existing note applies: this can be done on the floor or at table using mini-whiteboards

TEACHER NOTES:
Introduces the word 'vertices' (plural of vertex). For Year 1, 'corner' and 'vertex' can be used together. The square has 4 sides and 4 vertices.

WATCH FOR:
- Students writing 4 and 4: secure
- Students confused about vertices: hold up the shape and tap each corner as you count

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

56: """SAY:
- This is a triangle. Write the number of sides and vertices
- Show me on your board

DO:
- Hold up a triangle card
- Touch each side and each vertex as students count along
- Reveal 3 sides and 3 vertices

TEACHER NOTES:
Triangle has 3 sides and 3 vertices. The number of sides always equals the number of vertices in a polygon. You may name this pattern if students are ready.

WATCH FOR:
- Students writing 3 and 3: secure
- Students writing different numbers: re-count together

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

57: """SAY:
- Which shape has 4 vertices? Look at A, B and C
- Show me on your fingers

DO:
- Give 15 seconds of silent thinking
- Cue 'show me' for the finger response
- Scan every student before revealing

CFU CHECKPOINT:
Technique: Finger Show Me
Script:
- Ask: which shape has 4 vertices? Expected: the 4-cornered shape (square or rectangle)
- Scan for: most students choosing the 4-cornered option
PROCEED:
- >=80% choose correctly. Move on
PIVOT:
- Most likely: existing note applies: students pick option B because they confuse 3 and 4. Option B is a triangle (3 vertices)
- Reteach: hold up the triangle. Tap and count each corner: 1, 2, 3
- Re-check: show a new pair. Which has 4 vertices?

TEACHER NOTES:
Hinge check on counting vertices. Decides whether to move to the rectangle or pause to re-teach.

WATCH FOR:
- Students who pick the triangle: they need the corner-count again

[Stage 2: I Do | VTLM 2.0: Check for Understanding]
""",

58: """SAY:
- This is a rectangle. Write the number of sides and vertices
- Show me on your board

DO:
- Hold up a rectangle card
- Touch each side as you count: 1, 2, 3, 4
- Touch each vertex as you count: 1, 2, 3, 4

TEACHER NOTES:
Rectangle has 4 sides and 4 vertices. Same numbers as a square. Name the difference: the square has all 4 sides equal, the rectangle has 2 long and 2 short.

WATCH FOR:
- Students writing 4 and 4: secure
- Students who say 'this is just a square': prompt 'are the sides all the same length?'

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

59: """SAY:
- Time for the main lesson. For Explore you will need...

DO:
- Read the materials list shown on the slide
- Confirm workbooks and pencils ready
- Have wooden 2D shapes available for the Enabling group

TEACHER NOTES:
Main session setup slide. Confirm materials before students start.
""",

60: """SAY:
- In your workbook, draw as many 2D shapes as you know
- Label each shape: name, number of sides, number of vertices
- Watch me first. I'll do a square

DO:
- Existing note applies: complete in workbooks. Model drawing and labelling on the whiteboard
- Draw a square. Write 'Square. 4 sides. 4 vertices.'
- Then release students to try
- Enablers may use wooden shapes to trace around. Have shape names written on the board for them to copy

TEACHER NOTES:
Independent practice for SC1 and SC2. The model on the slide (square example) is the format students should copy.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: trace around wooden shapes. Copy the shape name from the board
EXTENDING PROMPT:
- Task: include a shape with more than 4 sides (e.g. pentagon or hexagon) and label it

WATCH FOR:
- Students drawing and labelling at least 3 shapes correctly: secure
- Students stuck: prompt 'start with a circle' and walk them through the labelling

[Stage 4: You Do | VTLM 2.0: Independent Practice]
""",

61: """TEACHER NOTES:
NOTE: this slide does not match Year 1 content. The prompt is about multiplying by 1000 and the place value chart, which is upper-primary content. Likely a leftover slide from another deck.

If you reach this slide, skip it. Do not teach multiplying by 1000 to a Year 1 class. Move directly to slide 62.

If you want to use a True/False Hinge here instead, ask: 'A square has 4 sides. True or false?' Expected: True. Then move on.
""",

62: """SAY:
- Lesson focus check. Show me what you've learned today
- Existing note applies: complete the sheet linked above, or write on mini-whiteboards

DO:
- Display the lesson focus task
- Set 5 to 8 minutes
- Circulate and check shape names, sides and vertices
- Sight every response

TEACHER NOTES:
End of Lesson 3 hinge or focus task. Use to confirm students can name a 2D shape and state its sides and vertices.

WATCH FOR:
- Students who name AND label features: secure with SC1 and SC2
- Students who name only: SC2 not yet secure. Mark for follow up

[Stage 5: Hinge | VTLM 2.0: Check for Understanding]
""",

# ============================================================
# Lesson 4 - 2D Shapes in the Environment
# ============================================================

63: """TEACHER NOTES:
Lesson 4 divider. Today's focus: identifying 2D shapes around us, polygons vs not-polygons. Print or have ready: Shape Hunt BLM and the Shape Mystery Bag materials (shape cards in opaque bags).
""",

64: """SAY:
- Daily Review. Doubles and near doubles
- On your whiteboard, work them out: 3+3, 5+5, 10+10, 4+4
- Now near doubles: 5+6, 10+11, 2+3, 4+5

DO:
- Existing note applies: on the whiteboard, review addition strategies, doubles and near doubles
- Read each fact aloud
- Boards up after 30 seconds per row
- Reveal the answer and the strategy (e.g. 5+6 = 5+5+1)

TEACHER NOTES:
Doubles and near doubles fluency. Near doubles uses 'one more than a double'. This is bridging support for tomorrow's geoboard counting.

WATCH FOR:
- Students using doubles to find near doubles: secure
- Students counting all from 1: model 5+5=10, then add 1

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

65: """SAY:
- Number Fluency. Counting forwards and backwards by 1s
- I'll say a start number. You count on. Then we'll go backwards

DO:
- Existing note applies: count forwards and backwards by 1s starting at different numbers
- Choose start numbers: e.g. 23 (count to 30), 47 (count to 53)
- Then go backwards: e.g. 12 (count back to 5)
- Boards optional

TEACHER NOTES:
Quick fluency with counting on and back. Backwards counting is harder for Year 1; pause and re-chorus if shaky.

WATCH FOR:
- Students confident on bridges (e.g. 29 to 30, 19 to 20): secure
- Students stuck on backwards: chorus once more, slowly

[Stage 1: Number Fluency | VTLM 2.0: Retention and Recall]
""",

66: """SAY:
- Read the LI with me: 'We are learning to identify two-dimensional shapes in our environment'
- Now our success criteria
- Today we will name, label and record 2D shapes we see

DO:
- Choral read the LI then each SC
- Hold up a real environmental example (e.g. a square tile or rectangular book) as you read SC1

TEACHER NOTES:
Environment is a new word. Say it means 'the world around you, like the classroom or the playground'. SC4 (identify polygons) carries from yesterday's lesson into today's checklist.

WATCH FOR:
- Students unsure about 'environment': use the simple translation 'the world around us'
- Students unsure about 'polygon': remind 'closed shape with straight sides' (you'll re-teach this in the launch)

[General: LI/SC | VTLM 2.0: Clear Learning Intention]
""",

67: """SAY:
- For this Launch you will need...

DO:
- Read the materials list shown on the slide
- Confirm whiteboards and pencils ready

TEACHER NOTES:
Launch setup slide. No student action.
""",

68: """SAY:
- New word today: polygon
- A polygon is a closed shape with straight sides
- Closed means the lines join up. Straight means no curves
- Say it with me: po-ly-gon

DO:
- Trace the shape on the slide
- Run your finger along each side: 'closed... closed... closed... straight... straight'
- Have students chorus 'polygon' twice

TEACHER NOTES:
This is the key vocabulary slide for Lesson 4. Two features: closed and straight sides. The next slides will use a checklist with these features.

WATCH FOR:
- Students who can repeat 'polygon': good first step
- Students who can also say 'closed and straight sides': stronger understanding

[Stage 2: I Do | VTLM 2.0: Vocabulary]
""",

69: """SAY:
- Look at the parallelogram. Use the checklist
- Closed shape? Straight sides? Two-dimensional?
- I'll model. Closed: yes. Straight sides: yes. 2D: yes
- So this IS a polygon

DO:
- Point to the parallelogram on the slide
- Tick each box as you say yes
- Trace the closed shape and the straight sides

TEACHER NOTES:
First example of using the polygon checklist. Model the thinking aloud. The parallelogram is a polygon. Point out it's the new shape word, but no need to define it formally for Year 1.

WATCH FOR:
- Students following the checklist: secure
- Students confused by the new shape name: focus on the checklist, not the name

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

70: """SAY:
- Your turn on whiteboards. Look at the kite
- Use the checklist. Closed? Straight sides? 2D?
- Tick or cross each box on your board

DO:
- Read the three checklist items aloud
- Boards up after 30 seconds
- Reveal: closed yes, straight sides yes, 2D yes. The kite IS a polygon

TEACHER NOTES:
First We Do with the checklist. The kite is a polygon. The W icon means whiteboards.

WATCH FOR:
- Students who tick all three: secure
- Students who untick 'straight sides' for the kite: re-trace the sides together

[Stage 2: We Do | VTLM 2.0: Guided Practice]
""",

71: """SAY:
- Look at the cube. Use the checklist
- Closed? Straight sides? Two-dimensional?
- Boards up. Show me

DO:
- Read the three checklist items
- Boards up after 30 seconds
- Reveal: closed yes, straight sides yes, BUT 2D NO. A cube is 3D
- So a cube is NOT a polygon

TEACHER NOTES:
Trick example. The cube is 3D, not 2D. This is the most likely place students will tick everything without checking the 2D feature. Use this slide to make sure students apply ALL three criteria.

MISCONCEPTIONS:
- Misconception: students see straight sides and tick polygon
  Why: they only check the 'sides' feature
  Impact: they call 3D objects polygons
  Quick correction: hold up a real cube. 'Can you flatten it on the table? No. So it's not 2D. Not a polygon.'

WATCH FOR:
- Students ticking all three: re-anchor on 2D vs 3D
- Students who untick 2D: secure with the full checklist

[Stage 2: We Do | VTLM 2.0: Check for Understanding]
""",

72: """SAY:
- Look at the arrow. Use the checklist
- Closed? Straight sides? Two-dimensional?
- Boards up. Show me

DO:
- Read the three checklist items
- Boards up after 30 seconds
- Reveal: depends on the arrow. If the outline is a closed shape with straight sides, yes. If it has a curve or is open, no
- Use the actual arrow shown on the slide

TEACHER NOTES:
The polygon answer depends on whether the arrow is drawn with all straight sides and is closed. If unsure, trace the outline together to check. A typical arrow icon is a closed shape with straight sides, so it IS a polygon.

WATCH FOR:
- Students who untick 'closed' if there is an open end: prompt them to trace again
- Students who tick all three for a closed straight-sided arrow: secure

[Stage 2: We Do | VTLM 2.0: Guided Practice]
""",

73: """SAY:
- Find the error. Look at the ticks
- Existing note applies: identify which feature shouldn't have a tick
- Look at the shape. Look at the ticks. Which one is wrong?

DO:
- Read the checklist with the ticks shown
- Give 30 seconds of silent thinking
- Take responses. Have a student come to the board to point to the wrong tick

TEACHER NOTES:
Error analysis slide. Students must check the actual shape against each ticked feature. This builds critical thinking with the checklist.

WATCH FOR:
- Students who can point to the wrong tick: secure
- Students who say 'they're all right': re-trace the shape and the features one by one

[Stage 2: We Do | VTLM 2.0: Error Analysis]
""",

74: """SAY:
- We have a hands-on activity. Shape Mystery Bag
- Existing note applies: each POD will get a mystery bag with shape cards
- One partner picks a shape without showing. Describe it without saying its name
- Your partner has to guess. Then swap

DO:
- Existing note applies: pre-prepare shape mystery bags before the lesson
- Distribute one mystery bag per partner pod
- Model first: pick a shape, describe its sides and vertices
- Set 5 minutes. Circulate and prompt feature language

TEACHER NOTES:
Hands-on activity to consolidate naming and describing 2D shapes. Listen for feature language: number of sides, number of vertices, curved or straight.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: provide a feature card (sides, vertices) to read from when describing
EXTENDING PROMPT:
- Task: describe two shapes at once and ask the partner to find both

WATCH FOR:
- Pairs using feature language ('4 sides, 4 vertices'): secure
- Pairs giving non-feature clues ('it's red', 'it's big'): redirect to features

[Stage 3: We Do | VTLM 2.0: Hands-On Practice]
""",

75: """SAY:
- Time for the main lesson. For Explore you will need...

DO:
- Read the materials list shown on the slide
- Confirm whiteboards, pencils and worksheets ready

TEACHER NOTES:
Main session setup slide. No student action.
""",

76: """SAY:
- Look at the picture. Draw and count the 2D shapes
- I'll model the first one
- I see 5 of one shape, 2 of another, 6 of another

DO:
- Point to each shape group on the picture
- Count aloud with the class
- Have students draw the shape and write the count beside it

TEACHER NOTES:
Modelling the count-and-record routine. Use the actual shapes shown in the picture (e.g. windows, tiles, doors) and record the count.

WATCH FOR:
- Students drawing the shape and writing the count: secure
- Students writing the count only: prompt 'draw the shape too'

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]

SOURCES:
Image credit: Pixabay
""",

77: """SAY:
- True or false? There are two squares in this picture
- Look carefully. Count them
- Show me. Thumb up for true. Thumb down for false

DO:
- Give 20 seconds of silent thinking
- Cue thumbs
- Scan all thumbs before revealing

CFU CHECKPOINT:
Technique: Thumbs Up Thumbs Down
Script:
- Ask: are there two squares in the picture? Expected: depends on the picture (use the actual count)
- Scan for: most students choosing the same answer with confidence
PROCEED:
- >=80% agree on the correct answer. Move on
PIVOT:
- Most likely: students count rectangles as squares
- Reteach: a square has 4 EQUAL sides. Hold up a square card. Compare
- Re-check: how many SQUARES are in the picture?

TEACHER NOTES:
Hinge check on identifying squares. The trick is squares vs rectangles. Both have 4 sides and 4 vertices, but only squares have all sides equal.

WATCH FOR:
- Students who count rectangles as squares: re-anchor on 'all sides equal'

[Stage 2: I Do | VTLM 2.0: Check for Understanding]

SOURCES:
Image credit: Pixabay
""",

78: """SAY:
- On your whiteboard. Draw and count the 2D shapes you see
- Show me when you're done

DO:
- Read the slide
- Boards up after 1 minute
- Reveal the counts (14 and 3)

TEACHER NOTES:
We Do with whiteboards. Students count 2 shape types in the picture. Numbers shown on slide: 14 and 3 for the two shapes.

WATCH FOR:
- Students writing 14 and 3 with shape drawings: secure
- Students writing only one count: prompt 'count the second shape too'

[Stage 3: We Do | VTLM 2.0: Guided Practice]
""",

79: """SAY:
- On your whiteboard. Draw and count the 2D shapes
- Show me

DO:
- Read the slide
- Boards up after 1 minute
- Reveal the counts (10 and 5)

TEACHER NOTES:
Second We Do count. Numbers shown on slide: 10 and 5.

WATCH FOR:
- Students writing 10 and 5: secure
- Students who skip a shape: re-anchor and re-count

[Stage 3: We Do | VTLM 2.0: Guided Practice]

SOURCES:
Image credit: Pixabay
""",

80: """SAY:
- On your whiteboard. Draw and count the 2D shapes
- Show me

DO:
- Read the slide
- Boards up after 1 minute
- Reveal the counts (1 and 13)

TEACHER NOTES:
Third We Do count. Numbers shown on slide: 1 and 13. The big difference between the two counts may help students notice each shape distinctly.

WATCH FOR:
- Students writing 1 and 13: secure
- Students unsure where the 13 are: trace through them together

[Stage 3: We Do | VTLM 2.0: Guided Practice]

SOURCES:
Image credit: Pixabay
""",

81: """SAY:
- We're going on a 2D Shape Hunt
- Take your worksheet, board and pencil
- We're recording how many circles, squares, rectangles and triangles we can find
- Tally each one as you find it

DO:
- Existing note applies: use the Shape Hunt BLM
- Read the instructions and notes to teacher
- Take the class outside (or hunt in the classroom if preferred)
- Prompt support students to a specific area for obvious examples
- Set 10 to 15 minutes
- Bring students back together. Share counts

TEACHER NOTES:
Hands-on environment activity. Students apply naming and counting from the lesson. Stay close to the support group early.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: prompt students to a specific area (e.g. play structure with obvious shapes)
EXTENDING PROMPT:
- Task: existing note applies: students record 2D shapes on the school building, including rhombus, parallelograms and kites

WATCH FOR:
- Students using tallies efficiently: secure
- Students recording 3D objects (e.g. ball, bin): redirect to 2D shapes ON the object

[Stage 4: You Do | VTLM 2.0: Independent Practice]
""",

# ============================================================
# Lesson 5 - Making 2D shapes on geoboards
# ============================================================

82: """TEACHER NOTES:
Lesson 5 divider. Topic is a Challenging Task / Game using geoboards. Have geoboards and rubber bands ready (or set up digital geoboards). Have number lines available for the Daily Review.
""",

83: """SAY:
- Daily Review. Counting on and counting back
- Watch me model first on the TV
- Then take your number line back to your table

DO:
- Existing note applies: model count on and count back on the screen
- Then send students back with number lines
- Ask: 14+2, 9+6, 8-2, 19-3 (or similar)
- Boards or number line responses

TEACHER NOTES:
Daily Review activates count on and count back. This supports counting pegs on the geoboard later in the lesson.

WATCH FOR:
- Students using the number line correctly: secure
- Students counting all from 1: re-anchor 'start at the bigger number, count on'

[Stage 1: Daily Review | VTLM 2.0: Retention and Recall]
""",

84: """SAY:
- Number Fluency. Quick mental maths
- I'll call out, you answer

DO:
- Choose a fluency target (e.g. count by 5s to 50, doubles to 10, count back from 20)
- Run as a brisk choral or board response
- Keep the pace fast

TEACHER NOTES:
Generic Number Fluency slide. Pick a fluency relevant to the geoboard work today: counting pegs in 2s or 5s would be useful.

WATCH FOR:
- Students chorussing confidently: secure
- Students drifting: chorus once more

[Stage 1: Number Fluency | VTLM 2.0: Retention and Recall]
""",

85: """SAY:
- Read the LI with me: 'We are learning to make 2D shapes'
- Now our success criteria
- Today we will use a geoboard

DO:
- Choral read the LI then each SC
- Hold up a geoboard with rubber bands as you read SC1

TEACHER NOTES:
SC1 is the floor (use a geoboard), SC2 is the core (regular shapes), SC3 extends (irregular shapes). 'Regular' means all sides equal. 'Irregular' means not all sides equal.

WATCH FOR:
- Students unsure about 'regular' and 'irregular': you'll teach these on slides 88 to 89

[General: LI/SC | VTLM 2.0: Clear Learning Intention]
""",

86: """SAY:
- Let's activate what you already know

DO:
- Pause briefly on this header

TEACHER NOTES:
Activate prior learning header. Sets up the next few slides on 2D, regular polygons, irregular polygons.
""",

87: """SAY:
- Two-dimensional. We learned this word earlier in the week
- Two-dimensional means it has length and width. It is flat
- Watch me press it flat on the table

DO:
- Hold up a 2D shape card
- Press it flat against the desk
- Existing note applies: definition from VIC Mathematics glossary

TEACHER NOTES:
Re-activates 2D from earlier lessons. Existing note applies: an object with width and length is two-dimensional. A polygon is an example of a 2D geometric object.

WATCH FOR:
- Students who can press a 2D shape flat: secure
- Students who say 3D objects are 2D: re-show the flat-vs-3D contrast

[Stage 2: I Do | VTLM 2.0: Vocabulary]
""",

88: """SAY:
- Regular polygons. New word today
- Regular means all sides are the same length. All angles are the same too
- A square is a regular polygon. All 4 sides are equal

DO:
- Hold up a square
- Trace each side and say 'same'
- Existing note applies: regular polygons have all sides of equal length and equal angles. Source: VIC Mathematics glossary

TEACHER NOTES:
Introduces 'regular' as 'all the same'. For Year 1, focus on equal sides. Equal angles is a stretch concept.

WATCH FOR:
- Students who can say 'all sides the same': secure
- Students confused: re-trace and chant 'same... same... same... same'

[Stage 2: I Do | VTLM 2.0: Vocabulary]
""",

89: """SAY:
- Irregular polygons. The opposite of regular
- Irregular means NOT all sides are the same. Some sides are longer than others
- A long thin rectangle is irregular. A wonky triangle is irregular

DO:
- Hold up an irregular shape (e.g. a long thin rectangle, a scalene triangle)
- Trace the sides. Show that some are longer
- Existing note applies: a polygon with not all sides or angles equal is irregular. Source: VIC Mathematics glossary

TEACHER NOTES:
Pairs with regular polygons. The 'ir-' prefix means 'not'. Tie this to other 'ir-' words students might know if helpful.

WATCH FOR:
- Students who can name a sample as irregular: secure
- Students unsure: re-show the regular vs irregular pair side by side

[Stage 2: I Do | VTLM 2.0: Vocabulary]
""",

90: """SAY:
- Name and describe this shape
- It is a square. It has 4 sides and 4 vertices
- And the sides are equal. So it is a regular polygon

DO:
- Hold up a square card
- Tap each side and each vertex as you count
- Read each fact on the slide

TEACHER NOTES:
Modelling the full description: name, sides, vertices, regular feature. The I icon means I Do.

WATCH FOR:
- Students following the description: secure
- Students who only get the name: re-anchor on the full description

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

91: """SAY:
- Name and describe this shape
- It is a triangle. It has 3 sides and 3 vertices

DO:
- Hold up a triangle card
- Tap each side and each vertex
- Read each fact on the slide

TEACHER NOTES:
Triangle description. Note the slide does not specify regular or irregular for the triangle, so just name and feature-count.

WATCH FOR:
- Students who can repeat 3 and 3: secure

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

92: """SAY:
- Name and describe this shape
- It is a rectangle. It has 4 sides and 4 vertices

DO:
- Hold up a rectangle card
- Tap each side and each vertex
- Read each fact on the slide

TEACHER NOTES:
Rectangle description. Note that the rectangle and square both have 4 and 4. The difference is whether the sides are equal.

WATCH FOR:
- Students who say 'this is a square too': prompt 'are all 4 sides equal? If two long and two short, it's a rectangle'

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

93: """SAY:
- Describe this shape on your board
- Write: name, sides, vertices

DO:
- Read the prompts on the slide
- Boards up after 30 seconds
- Reveal: rectangle, 4, 4

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
- Ask: name, sides, vertices? Expected: rectangle, 4, 4
- Scan for: all three pieces of information on the board
PROCEED:
- >=80% have all three correct. Move on
PIVOT:
- Most likely: students write 'square' instead of 'rectangle'
- Reteach: hold up the slide. Are all 4 sides the same? No. So it's a rectangle
- Re-check: show another rectangle. Describe it

TEACHER NOTES:
First independent CfU on full description. The CfU icon appears twice on the slide; take one scan.

WATCH FOR:
- Students with all three correct: secure
- Students writing 'square': re-anchor on equal sides

[Stage 2: We Do | VTLM 2.0: Check for Understanding]
""",

94: """SAY:
- Describe this shape on your board
- Write: name, sides, vertices

DO:
- Read the prompts
- Boards up after 30 seconds
- Reveal: square, 4, 4

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
- Ask: name, sides, vertices? Expected: square, 4, 4
- Scan for: all three pieces of information on the board
PROCEED:
- >=80% have all three correct. Move on
PIVOT:
- Most likely: students write 'rectangle' instead of 'square'
- Reteach: hold up the slide. Are all 4 sides the same? Yes. So it's a square
- Re-check: show another square (or a rectangle for contrast). Describe it

TEACHER NOTES:
Second independent CfU. Pairs with slide 93 to contrast square and rectangle.

WATCH FOR:
- Students with all three correct: secure
- Students writing 'rectangle': re-anchor on equal sides

[Stage 2: We Do | VTLM 2.0: Check for Understanding]
""",

95: """SAY:
- For this Launch you will need...

DO:
- Read the materials list shown on the slide
- Confirm geoboards and rubber bands ready (or digital geoboards)

TEACHER NOTES:
Launch setup slide. Confirm materials before students start.
""",

96: """SAY:
- This is a geoboard. Look at the parts
- The board has pegs in a grid
- We use rubber bands. We stretch them around the pegs to make 2D shapes

DO:
- Hold up a real geoboard
- Point to the board, the pegs and a rubber band
- Stretch one band around 4 pegs to model a small square

TEACHER NOTES:
Tool introduction slide. Make sure every student can see the geoboard parts. Year 1 students will need a quick safety reminder about rubber bands (don't flick them).

WATCH FOR:
- Students who can name the parts: secure
- Students playing with the bands: redirect with 'eyes on me first'

[Stage 2: I Do | VTLM 2.0: Tool Introduction]

SOURCES:
Image credit: Pixabay
""",

97: """SAY:
- Watch me make a square. I stretch the band around 4 pegs
- The shape is closed. It has 4 sides
- Let me count: 1, 2, 3, 4

DO:
- Hold up your geoboard
- Stretch the band around 4 pegs to make a square
- Trace each side and count

TEACHER NOTES:
First square modelling. Focus on 'closed' and '4 sides'. Use a small square to start (2x2 pegs is enough).

WATCH FOR:
- Students watching attentively: ready
- Students fiddling: collect bands until your eye contact is on the model

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]

SOURCES:
Image credit: Pixabay
""",

98: """SAY:
- Now I'll make a bigger square. Each side is 7 pegs long
- Watch me count: 1, 2, 3, 4, 5, 6, 7
- Each side has the SAME number of pegs because it is a square

DO:
- Existing note applies: point and count each peg on the top side: 1 to 7
- Existing note applies: highlight that each side is equal, 7 pegs in length
- Stretch the band around 4 corners to enclose 7 pegs per side

TEACHER NOTES:
This slide names the equal-sides feature of a square. The 'each side 7 pegs' fact is the key takeaway.

WATCH FOR:
- Students who can repeat 'each side is the same': secure
- Students focused on the band: redirect to the count

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]
""",

99: """SAY:
- Hinge. Which geoboard shows a square? A or B?
- Look at the sides. Count if you need to
- Show me on your fingers

DO:
- Give 20 seconds of silent thinking
- Cue 'show me'
- Scan every student before revealing

CFU CHECKPOINT:
Technique: Finger Show Me
Script:
- Ask: which is the square? Expected: the one with 4 equal sides (whichever option shows a true square)
- Scan for: most students choosing the same option
PROCEED:
- >=80% choose correctly. Move on
PIVOT:
- Most likely: students choose the rectangle, thinking 4 sides means square
- Reteach: count the pegs on each side. Are they all the same?
- Re-check: show two new options. Show me

TEACHER NOTES:
Hinge slide. Decides whether to move to the rectangle or stop and re-teach 'all sides equal'.

WATCH FOR:
- Students who pick the rectangle: re-anchor on equal sides

[Stage 2: I Do | VTLM 2.0: Check for Understanding]
""",

100: """SAY:
- What error has been made when making this square?
- Look carefully. What's wrong?
- Existing note applies: the sides are not equal

DO:
- Trace each side on the slide
- Count the pegs on each side aloud
- Take 2 student responses

TEACHER NOTES:
Existing note applies: answers may vary. The main error is that the sides are not equal. A square needs all 4 sides the same length.

WATCH FOR:
- Students who say 'sides are different lengths': secure
- Students who say 'it's a rectangle': also accept as a related observation

[Stage 2: I Do | VTLM 2.0: Error Analysis]
""",

101: """SAY:
- How many pegs should the next side of the square be?
- Look at the side that's already made. Count its pegs
- Show me on your board

DO:
- Trace the existing side
- Boards up after 20 seconds
- Reveal: 6 pegs (because a square needs all sides equal)

TEACHER NOTES:
Apply 'all sides equal' to predict the next side. The answer 6 only makes sense if the existing side is also 6 pegs long.

WATCH FOR:
- Students writing 6: secure with the equal-sides rule
- Students writing other numbers: re-anchor on the rule

[Stage 2: We Do | VTLM 2.0: Guided Practice]
""",

102: """SAY:
- Now I'll make a rectangle. A rectangle has 4 sides
- But two sides are LONG and two sides are SHORT
- Watch me

DO:
- Stretch a band around the pegs to make a rectangle
- Trace the long sides, then the short sides
- Make the contrast clear: long, long, short, short

TEACHER NOTES:
Rectangle modelling slide. Contrast with square: same number of sides, but two pairs of equal sides instead of all four equal.

WATCH FOR:
- Students who can spot 'two long, two short': secure

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]

SOURCES:
Image credit: Pixabay
""",

103: """SAY:
- What error has been made when making this rectangle?
- Look carefully
- Existing note applies: the shape is not closed

DO:
- Trace the lines on the slide
- Show where the gap is
- Take 2 student responses

TEACHER NOTES:
Existing note applies: answers may vary. The main error is that the shape is not closed. A polygon must have all sides joined up.

WATCH FOR:
- Students who say 'the shape isn't closed': secure
- Students who say 'sides are different lengths': also acceptable depending on the picture

[Stage 2: I Do | VTLM 2.0: Error Analysis]
""",

104: """SAY:
- How many pegs should the last side of the rectangle be?
- Look at the side opposite. The two long sides match. The two short sides match
- Show me on your board

DO:
- Point to the side opposite the missing one
- Boards up after 20 seconds
- Reveal: 6 pegs

TEACHER NOTES:
Apply 'opposite sides are equal' to predict the last side. The answer 6 makes sense if the opposite side is also 6.

WATCH FOR:
- Students writing 6: secure with the opposite-sides rule
- Students writing other numbers: re-anchor on opposite sides matching

[Stage 2: We Do | VTLM 2.0: Guided Practice]
""",

105: """SAY:
- Time for the main lesson. For Explore you will need...

DO:
- Read the materials list
- Confirm geoboards and rubber bands distributed

TEACHER NOTES:
Main session setup slide. No student action.
""",

106: """SAY:
- Hands-on activity. Make these on your geoboard:
- A small square and a large square
- A small rectangle and a large rectangle

DO:
- Existing note applies: do this for 10 minutes
- Existing note applies: there will be more slides afterwards to work through
- Existing note applies: check that squares and rectangles have been made using the correct lengths of sides
- Circulate. Use 'count the pegs on each side' to check

TEACHER NOTES:
First Explore activity. Students apply equal sides for square and opposite sides for rectangle. Time-boxed at 10 minutes.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: provide a peg-count card (e.g. 'square: 3 pegs each side')
EXTENDING PROMPT:
- Task: make a square AND a rectangle that use the same number of pegs

WATCH FOR:
- Students who make squares with unequal sides: pause and re-count
- Students who make rectangles with all sides equal: that's a square. Stretch one pair longer

[Stage 4: You Do | VTLM 2.0: Independent Practice]
""",

107: """SAY:
- Look at the rectangle on the screen. It is not finished
- Create and complete the rectangle on your geoboard
- Existing note applies: answers may vary

DO:
- Read the prompt
- Allow 1 to 2 minutes for students to complete
- Take 2 student examples and show them on the screen

TEACHER NOTES:
Existing note applies: answers may vary. Students must complete the rectangle so opposite sides match.

WATCH FOR:
- Students who close the rectangle with matching opposite sides: secure
- Students who close it with all sides different: re-anchor on opposite sides

[Stage 3: We Do | VTLM 2.0: Guided Practice]
""",

108: """SAY:
- Now I'll make a triangle. A triangle has 3 sides
- Watch me. I stretch the band around 3 pegs

DO:
- Stretch a band around 3 pegs on your geoboard
- Trace each of the 3 sides
- Count: 1, 2, 3

TEACHER NOTES:
Triangle modelling. The simplest polygon. Make sure students see exactly 3 corners (vertices).

WATCH FOR:
- Students who can repeat 3 sides: secure

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]

SOURCES:
Image credit: Pixabay
""",

109: """SAY:
- Now I'll make an irregular triangle
- Irregular means the sides are not all the same length
- Watch. This one has a long side, a medium side and a short side

DO:
- Stretch a band around 3 pegs in different positions
- Trace each side. Compare the lengths
- Say 'all 3 sides different' as you trace

TEACHER NOTES:
Introduces the irregular triangle. The sides are not equal. This contrasts with the regular shapes from earlier slides.

WATCH FOR:
- Students who notice the sides are different: secure
- Students confused: re-anchor on 'irregular = not all the same'

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]

SOURCES:
Image credit: Pixabay
""",

110: """SAY:
- Hinge. Which geoboard shows an irregular triangle? A or B?
- Look at the sides. Are they all the same? If not, it's irregular
- Show me on your fingers

DO:
- Give 20 seconds of silent thinking
- Cue 'show me'
- Scan every student before revealing

CFU CHECKPOINT:
Technique: Finger Show Me
Script:
- Ask: which is the irregular triangle? Expected: the one with sides of different lengths
- Scan for: most students choosing the same option
PROCEED:
- >=80% choose correctly. Move on
PIVOT:
- Most likely: students choose the regular (equal-sided) triangle
- Reteach: count the pegs on each side of both options. Are they the same?
- Re-check: show two new triangles. Which is irregular?

TEACHER NOTES:
Hinge slide for SC3 (irregular shapes). Decides whether to move on or re-teach.

WATCH FOR:
- Students who pick the equal-sided one: re-anchor on 'irregular = not equal'

[Stage 2: We Do | VTLM 2.0: Check for Understanding]
""",

111: """SAY:
- Now I'll make a quadrilateral. New word today
- Quadrilateral means 4 sides. Quad means 4
- Watch me

DO:
- Stretch a band around 4 pegs to make a 4-sided shape
- Trace each of the 4 sides as you count
- Note: this could be a square, rectangle, or any 4-sided shape

TEACHER NOTES:
Introduces 'quadrilateral' as the family name for any 4-sided polygon. For Year 1, this is a stretch word. The number connection (quad = 4) may help.

WATCH FOR:
- Students who can repeat 'quadrilateral': good first step
- Students who say 'square' or 'rectangle': also correct, since both are quadrilaterals

[Stage 2: I Do | VTLM 2.0: Vocabulary]

SOURCES:
Image credit: Pixabay
""",

112: """SAY:
- Make another quadrilateral. A different one this time
- It still has 4 sides, but it doesn't have to be a square or a rectangle
- Watch me

DO:
- Stretch a band around 4 pegs in an irregular pattern
- Trace each of the 4 sides
- Compare to the previous quadrilateral

TEACHER NOTES:
Second quadrilateral example. Shows that quadrilaterals can look very different but all have 4 sides.

WATCH FOR:
- Students who can name 4 sides: secure
- Students who think only squares and rectangles count: re-emphasise 4 sides is the rule

[Stage 2: I Do | VTLM 2.0: Explicit Teaching]

SOURCES:
Image credit: Pixabay
""",

113: """SAY:
- What error has been made when making this quadrilateral?
- Look carefully
- Existing note applies: the shape is not closed

DO:
- Trace the lines on the slide
- Show the gap
- Take 2 student responses

TEACHER NOTES:
Existing note applies: answers may vary. The main error is the shape is not closed. A quadrilateral must have all 4 sides joined up.

WATCH FOR:
- Students who say 'the shape isn't closed': secure
- Students who say 'only 3 sides': also acceptable if that matches the picture

[Stage 2: We Do | VTLM 2.0: Error Analysis]
""",

114: """SAY:
- Final hands-on. Make these on your geoboard:
- 2 different triangles
- 2 different quadrilaterals

DO:
- Existing note applies: ask students to point to and name the shapes they have made
- Distribute geoboards if not already done
- Set 5 to 8 minutes
- Circulate. Ask 'how many sides? Is it regular or irregular?'

TEACHER NOTES:
Final consolidation. Students apply triangle and quadrilateral knowledge with variation. Encourages creativity within the rules of each shape family.

ENABLING & EXTENDING:
ENABLING PROMPT:
- Task: provide a side-count card (e.g. 'triangle: 3 sides')
EXTENDING PROMPT:
- Task: make a triangle that is regular and one that is irregular. Same for quadrilaterals

WATCH FOR:
- Students who make 2 different triangles: secure with SC3
- Students who make 2 the same: prompt 'change the size or the side lengths'

[Stage 4: You Do | VTLM 2.0: Independent Practice]
""",

}


# ---------------------------------------------------------------------------
# Apply notes
# ---------------------------------------------------------------------------

def set_notes(slide, text):
    """Replace the notes text frame with the given text. Plain text, no markdown."""
    if not slide.has_notes_slide:
        notes_slide = slide.notes_slide
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # Clear existing
    tf.clear()
    # First paragraph already exists from clear()
    lines = text.split('\n')
    if not lines:
        return
    p0 = tf.paragraphs[0]
    if lines:
        p0.text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line


def main():
    p = Presentation(SRC)
    total = len(p.slides)
    applied = 0
    skipped = []
    for i, slide in enumerate(p.slides, start=1):
        text = NOTES.get(i)
        if text is None:
            skipped.append(i)
            continue
        set_notes(slide, text.strip())
        applied += 1
    p.save(DST)
    print(f"Total slides: {total}")
    print(f"Notes applied: {applied}")
    if skipped:
        print(f"Skipped (no notes defined): {skipped}")
    print(f"Saved: {DST}")


if __name__ == "__main__":
    main()
