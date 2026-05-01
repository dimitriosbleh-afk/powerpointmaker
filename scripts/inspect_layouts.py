"""List slide layouts available in OCHRE deck."""
from pathlib import Path
from pptx import Presentation

decks = [
    r"C:\Users\09560329\Downloads\Challenge Slides\4\16. literature_presentation Sentence Combining - Compound Sentences 4En06V Storm Boy - Information report.pptx",
    r"C:\Users\09560329\Downloads\Challenge Slides\1\1. Storm Boy - sentence expansion.pptx",
]
for path in decks:
    print("=" * 60)
    print(Path(path).name)
    print("=" * 60)
    pres = Presentation(path)
    print(f"Slide W x H: {pres.slide_width}, {pres.slide_height}")
    for i, layout in enumerate(pres.slide_layouts):
        # count placeholders
        ph_names = [ph.placeholder_format.idx for ph in layout.placeholders]
        print(f"  [{i}] {layout.name}  (placeholder idxs: {ph_names})")
    print()
    print(f"Slide 7 layout: {list(pres.slides)[6].slide_layout.name}")
    print(f"Slide 1 layout: {list(pres.slides)[0].slide_layout.name}")
    print(f"Last slide layout: {list(pres.slides)[-1].slide_layout.name}")
