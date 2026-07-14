#!/usr/bin/env python3
"""Merge a multi-session unit into a single PPTX with a flat Resources/ folder.

Reads a JSON manifest describing the unit and its ordered lesson folders, then:
  1. Combines the per-lesson PPTX decks (in manifest order) into one PPTX.
  2. Copies every PDF from each lesson's `resources-session{N}/` folder into
     a single `Resources/` folder alongside the combined PPTX.

Manifest format (see docs/resource-system.md):

    {
      "unit_folder": "Decimals_and_Fractions_Unit",
      "unit_pptx_name": "Decimals and Fractions Unit.pptx",
      "lessons": [
        {"folder": "DecFrac_Lesson1_Place_Value_With_Decimals", "session": 1},
        {"folder": "DecFrac_Lesson2_Estimation_And_Rounding",   "session": 2}
      ]
    }

Usage:
    python scripts/merge_unit.py builds/manifests/<unit>.json

Each lesson folder is expected under `output/<folder>` and must contain
exactly one `*.pptx` file. Resources are picked up from
`output/<folder>/resources-session{session}/` if that folder exists.
"""

import argparse
import copy
import json
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
OUTPUT_ROOT = ROOT / "output"
TEACHER_BRIEF_SCRIPT = ROOT / "scripts" / "build_teacher_brief.js"

HYPERLINK_RELTYPE_SUFFIX = "/hyperlink"


def load_manifest(manifest_path: Path) -> dict:
    with manifest_path.open("r", encoding="utf-8") as fh:
        manifest = json.load(fh)

    required = ("unit_folder", "unit_pptx_name", "lessons")
    missing = [k for k in required if k not in manifest]
    if missing:
        raise SystemExit(f"Manifest missing keys: {missing}")

    if not manifest["lessons"]:
        raise SystemExit("Manifest has no lessons.")

    if "teacher_brief" in manifest and not isinstance(manifest["teacher_brief"], dict):
        raise SystemExit("Manifest 'teacher_brief' must be an object when supplied.")

    seen_sessions = {}
    for i, lesson in enumerate(manifest["lessons"]):
        if "folder" not in lesson or "session" not in lesson:
            raise SystemExit(
                f"Lesson {i} missing 'folder' or 'session': {lesson}"
            )
        session = lesson["session"]
        if session in seen_sessions:
            raise SystemExit(
                f"Duplicate session number in manifest: session {session} "
                f"appears in lessons {seen_sessions[session]} and {i}. "
                "Session numbers must be unique so resource folders cannot collide."
            )
        seen_sessions[session] = i

    return manifest


def build_teacher_brief(manifest_path: Path, manifest: dict) -> int:
    """Generate the optional one-page weekly teacher preparation brief."""
    if "teacher_brief" not in manifest:
        return 0

    try:
        manifest_arg = str(manifest_path.resolve().relative_to(ROOT))
    except ValueError:
        manifest_arg = str(manifest_path.resolve())

    result = subprocess.run(
        ["node", str(TEACHER_BRIEF_SCRIPT), manifest_arg],
        cwd=str(ROOT),
        capture_output=True,
        text=True,
    )
    if result.stdout:
        print(result.stdout.rstrip())
    if result.returncode != 0:
        detail = result.stderr.strip() or result.stdout.strip()
        raise SystemExit(f"Teacher brief generation failed: {detail}")
    return 1


def find_lesson_pptx(lesson_dir: Path) -> Path:
    pptx_files = sorted(lesson_dir.glob("*.pptx"))
    if not pptx_files:
        raise SystemExit(f"No PPTX found in {lesson_dir}")
    if len(pptx_files) > 1:
        raise SystemExit(
            f"Multiple PPTX files in {lesson_dir}: {[p.name for p in pptx_files]}"
        )
    return pptx_files[0]


def copy_slide(src_slide, dst_prs):
    """Copy a slide from src into dst, preserving shapes, background,
    image and hyperlink relationships, and speaker notes. Hyperlink targets
    are carried over unchanged here; rewrite_resource_hyperlinks() does the
    `resources-session{N}/` to `Resources/` rewrite once on the saved file
    so base-deck slides and appended-deck slides go through one code path."""
    dst_layout = dst_prs.slide_layouts[0]
    dst_slide = dst_prs.slides.add_slide(dst_layout)

    src_sp_tree = src_slide.shapes._spTree
    dst_sp_tree = dst_slide.shapes._spTree
    dst_sp_tree.getparent().replace(dst_sp_tree, copy.deepcopy(src_sp_tree))

    src_bg = src_slide.background._element
    dst_bg = dst_slide.background._element
    dst_bg.getparent().replace(dst_bg, copy.deepcopy(src_bg))

    rid_map = {}
    for r_id, rel in list(src_slide.part.rels.items()):
        if "image" in rel.reltype:
            new_r_id = dst_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[r_id] = new_r_id
        elif rel.reltype.endswith(HYPERLINK_RELTYPE_SUFFIX):
            new_r_id = dst_slide.part.relate_to(
                rel.target_ref, rel.reltype, is_external=True
            )
            rid_map[r_id] = new_r_id

    if rid_map:
        attr_names = (
            "embed",
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id",
        )
        for old_r_id, new_r_id in rid_map.items():
            if old_r_id == new_r_id:
                continue
            for elem in dst_slide.shapes._spTree.iter():
                for attr_name in attr_names:
                    if elem.get(attr_name) == old_r_id:
                        elem.set(attr_name, new_r_id)

    if src_slide.has_notes_slide:
        notes_text = src_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            dst_slide.notes_slide.notes_text_frame.text = notes_text


# Matches `Target="resources-sessionN/<file>"` inside slide rels XML.
# Quoting can be `"` or `'`; we capture the quote so we can reproduce it.
_HYPERLINK_TARGET_RE = re.compile(
    r'Target=(["\'])resources-session\d+/([^"\']+)\1'
)


def rewrite_resource_hyperlinks(pptx_path: Path) -> int:
    """Rewrite hyperlink targets in slide rels files from
    `resources-session{N}/<file>` to `Resources/<file>` so in-slide links
    resolve from the unit folder. Returns the number of targets rewritten."""
    tmp_path = pptx_path.with_suffix(pptx_path.suffix + ".tmp")
    rewrites = 0

    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(
        tmp_path, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if (
                item.filename.startswith("ppt/slides/_rels/")
                and item.filename.endswith(".rels")
            ):
                text = data.decode("utf-8")
                new_text, n = _HYPERLINK_TARGET_RE.subn(
                    lambda m: f"Target={m.group(1)}Resources/{m.group(2)}{m.group(1)}",
                    text,
                )
                if n:
                    rewrites += n
                    data = new_text.encode("utf-8")
            zout.writestr(item, data)

    shutil.move(str(tmp_path), str(pptx_path))
    return rewrites


def validate_no_duplicate_zip_members(pptx_path: Path) -> None:
    with zipfile.ZipFile(pptx_path, "r") as zf:
        names = zf.namelist()
    duplicates = sorted({name for name in names if names.count(name) > 1})
    if duplicates:
        raise SystemExit(
            f"Merged PPTX contains duplicate zip entries: {duplicates[:10]}. "
            "This usually means a relationship or media part was copied unsafely."
        )


def merge_decks(manifest: dict, unit_dir: Path) -> Path:
    lessons = manifest["lessons"]
    out_pptx = unit_dir / manifest["unit_pptx_name"]

    first_dir = OUTPUT_ROOT / lessons[0]["folder"]
    base_path = find_lesson_pptx(first_dir)
    base_prs = Presentation(str(base_path))
    print(f"Base deck: {base_path}  ({len(base_prs.slides)} slides)")

    for lesson in lessons[1:]:
        src_dir = OUTPUT_ROOT / lesson["folder"]
        src_path = find_lesson_pptx(src_dir)
        src_prs = Presentation(str(src_path))
        print(
            f"Appending session {lesson['session']}: {src_path}  "
            f"({len(src_prs.slides)} slides)"
        )
        for slide in src_prs.slides:
            copy_slide(slide, base_prs)

    base_prs.save(str(out_pptx))
    print(f"\nWrote combined deck: {out_pptx}  ({len(base_prs.slides)} slides total)")

    rewrites = rewrite_resource_hyperlinks(out_pptx)
    if rewrites:
        print(
            f"Rewrote {rewrites} hyperlink target(s) "
            f"from resources-session{{N}}/ to Resources/"
        )
    validate_no_duplicate_zip_members(out_pptx)

    return out_pptx


def gather_resources(manifest: dict, resources_dir: Path) -> int:
    copied = 0
    seen_names = {}
    for lesson in manifest["lessons"]:
        session = lesson["session"]
        src_resources = OUTPUT_ROOT / lesson["folder"] / f"resources-session{session}"
        if not src_resources.is_dir():
            print(f"  (no resources folder for session {session} at {src_resources})")
            continue
        for pdf in sorted(src_resources.glob("*.pdf")):
            dst = resources_dir / pdf.name
            if pdf.name in seen_names and seen_names[pdf.name] != session:
                raise SystemExit(
                    f"Duplicate resource filename across sessions: {pdf.name} "
                    f"(session {seen_names[pdf.name]} and session {session}). "
                    f"Resource filenames must be unique across the unit."
                )
            seen_names[pdf.name] = session
            shutil.copy2(pdf, dst)
            copied += 1
            print(f"  Resources/  <-  {pdf.name}")
    return copied


def prepare_unit_output(unit_dir: Path, resources_dir: Path) -> None:
    """Prepare a clean delivery folder without touching per-lesson outputs."""
    unit_dir.mkdir(parents=True, exist_ok=True)

    for pptx in unit_dir.glob("*.pptx"):
        pptx.unlink()

    if resources_dir.exists():
        shutil.rmtree(resources_dir)
    resources_dir.mkdir(parents=True, exist_ok=True)


def merge_unit(manifest_path: Path) -> Path:
    manifest = load_manifest(manifest_path)
    unit_dir = OUTPUT_ROOT / manifest["unit_folder"]
    resources_dir = unit_dir / "Resources"

    prepare_unit_output(unit_dir, resources_dir)

    merge_decks(manifest, unit_dir)
    print()
    copied = gather_resources(manifest, resources_dir)
    copied += build_teacher_brief(manifest_path, manifest)
    print(f"\nCopied {copied} PDF(s) into {resources_dir}")
    print(f"\nUnit folder: {unit_dir}")
    return unit_dir


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "manifest",
        type=Path,
        help="Path to a unit manifest JSON file (see docstring for format).",
    )
    args = parser.parse_args()

    if not args.manifest.is_file():
        print(f"Manifest not found: {args.manifest}", file=sys.stderr)
        return 2

    merge_unit(args.manifest)
    return 0


if __name__ == "__main__":
    sys.exit(main())
