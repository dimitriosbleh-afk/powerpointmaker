"""Add teacher notes (megaprompt format) to the Under the Southern Cross lesson 3 and 4 decks."""
from pptx import Presentation
from pathlib import Path

DL = Path(r"c:/Users/09560329/Downloads")
L3_SRC = DL / "Copy of 3. Lesson Slides Extended Daily lesson 3 2En01V Under the Southern Cross - Informative.pptx"
L3_OUT = DL / "Copy of 3. Lesson Slides Extended Daily lesson 3 2En01V Under the Southern Cross - Informative - with notes.pptx"
L4_SRC = DL / "Copy of 4. Lesson Slides Extended Daily lesson 4 2En01V Under the Southern Cross - Informative.pptx"
L4_OUT = DL / "Copy of 4. Lesson Slides Extended Daily lesson 4 2En01V Under the Southern Cross - Informative - with notes.pptx"


LESSON_3_NOTES = {

    10: """SOURCES:
Teacher Handbook, Under the Southern Cross (2En01V) - Lesson 3 read-aloud notes.

SAY:
- Today we are going to keep exploring our book, Under the Southern Cross.
- I am going to read the cover and pages 1 to 28 out loud. Your job is to listen carefully and look at the pictures.
- Ask: How do the illustrations help us to understand more about each place and the Southern Cross? Expected: the pictures show what each place looks like at night, who lives there, and where the stars are.

DO:
- Hold the book so all students can see the illustrations.
- On each page, ask students to represent the orientation of the Southern Cross with their hands in a diamond shape.
- When Banjo is found on each page, ask: What do you notice about what Banjo is doing in the picture? Does he help show us something about the place?
- Page 6, pause after "outback Queensland." Ask: What do you see in this picture and how do the faces in the car match the text?
- Page 10, pause after "ping-pong balls." Ask: How does this picture help us understand the role of the man and woman helping the hatchlings?
- Page 20, pause after "mainland too." Ask: What can you tell about the people's feelings by looking at their faces in this picture?

TEACHER NOTES:
Pause points, the Southern Cross hand gesture and the Banjo-spotting routine are taken directly from the Teacher Handbook. Keep them intact and pace the read-aloud so each pause is a real pause, not a rushed aside.

WATCH FOR:
- Students drifting off - bring them back by pointing to an illustration and asking what they can see.
- Students calling out - remind the class we listen first, then share.
- Students forgetting the diamond hand shape - model it again and re-cue on the next page.""",


    12: """SOURCES:
Collins Primary School Dictionary, https://schools.collinsdictionary.com/dictionary/primary, accessed 19.10.25

SAY:
- In the text, the Min Min lights in the Australian outback are described as eerie.
- Something that is eerie is strange or frightening.
- Say the word with me... eerie!
- An old, empty house might look eerie. An owl might sound eerie when it hoots at night.

DO:
- Point to the image on the slide.
- Use a slightly spooky voice so students hear the feeling of the word.
- Have students repeat the word twice.

TEACHER NOTES:
This is the first meeting with the word. Keep it concrete and memorable - students will practise using it across the next few slides.

WATCH FOR:
- Students mixing up eerie with movie-scary - confirm eerie can feel strange or uneasy, not only frightening.""",


    13: """SAY:
- Listen: travellers have seen eerie lights in their travels near Boulia.
- Which answer fits? A - the lights are unusual and frightening. B - the lights are beautiful to look at.
- Show me A or B on your whiteboard.

DO:
- Read both options clearly.
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Write A or B. 3, 2, 1, show me.
- Scan for: most boards show A.
PROCEED:
- If most students show A, move on to the next practice slide.
PIVOT:
- Most likely misconception: students pick B because lights can look pretty.
- Reteach: re-read the eerie definition - strange or frightening. Remind students that unusual and frightening is the match.
- Fresh re-check: Say the meaning of eerie with me.

WATCH FOR:
- Students guessing quickly - protect wait time before they show.""",


    14: """SAY:
- Look carefully at the pictures on the slide.
- Point to the picture that could be described as eerie.
- Remember - eerie means strange or frightening.

DO:
- Give 5 seconds wait time.
- Say: 1, 2, 3, point.
- Cold call one student to explain why their picture looks eerie.

TEACHER NOTES:
Active image choice helps students link the word to a real-looking scene. Ask why, not just which.

WATCH FOR:
- Students pointing to the brightest picture just because it stands out - ask them to match to the meaning.""",


    15: """SAY:
- Let's finish this sentence together.
- The forest looked eerie because... what could come next?
- Turn and tell your partner one idea.

DO:
- Give 20 seconds of partner talk.
- Take 2-3 ideas from the room.
- Reveal the example answer and compare with student ideas.

TEACHER NOTES:
Students are stretching the word into new sentences. Accept any reason that matches strange or frightening.

WATCH FOR:
- Answers that only describe the forest (e.g. the forest was big) without a strange or frightening reason - redirect with: strange or frightening?""",


    16: """SAY:
- I will say a feeling. Thumbs up if an eerie thing might make you feel that way, thumbs down if not.
- happy?... scared?... uncomfortable?... safe?

DO:
- Pause between each word so thumbs come up together.
- Scan the room on each one.

TEACHER NOTES:
This checks whether students connect eerie to a feeling, not a shape or size. scared and uncomfortable are thumbs up; happy and safe are thumbs down.

WATCH FOR:
- Students putting thumbs up for all four to avoid choosing - re-cue: only if eerie could make you feel that way.""",


    17: """SAY:
- Show me which word fits best with eerie.
- Hold up 1 for normal, 2 for strange, 3 for bright, 4 for interesting.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which word fits best with eerie? On three, show me 1, 2, 3 or 4.
- Scan for: fingers showing 2 (strange).
PROCEED:
- If the class is on 2, move on to the next vocabulary word.
PIVOT:
- Most likely misconception: students pick bright because some eerie things have unusual light (Min Min lights).
- Reteach: eerie is about the feeling, not whether there is light.
- Fresh re-check: Say the meaning of eerie with me... strange or frightening.

WATCH FOR:
- Students mirroring the fastest hand - keep wait time firm.""",


    18: """SOURCES:
Collins Primary School Dictionary, https://schools.collinsdictionary.com/dictionary/primary, accessed 29.10.25

SAY:
- In the text, aspiring astronomers are keen to use their telescopes to gaze at the Milky Way.
- If you are keen to do something, or for something to happen, you want very much to do it or for it to happen.
- Say the word with me... keen!
- You might be keen to go to the beach. The children were keen to play a game of soccer.

DO:
- Point to the image on the slide.
- Use an enthusiastic voice for keen so students hear the feeling.
- Have students repeat the word twice.

TEACHER NOTES:
Keen is a strong want-to word. Connect it to something the class genuinely wants to do so it sticks.

WATCH FOR:
- Students mixing up keen with good at - keen is about wanting to, not being skilled.""",


    19: """SAY:
- Aspiring astronomers are keen to use their telescopes to look at the Milky Way.
- A - they really want to use their telescopes to look at the Milky Way. B - they are scared to use their telescopes to look at the Milky Way.
- Show me A or B on your whiteboard.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Write A or B. 3, 2, 1, show me.
- Scan for: most boards show A.
PROCEED:
- If the class shows A, move on to the next practice slide.
PIVOT:
- Most likely misconception: students confuse keen with nervous.
- Reteach: keen means you want to do it very much.
- Fresh re-check: Are you keen for the weekend? Thumbs up or down.

WATCH FOR:
- Students guessing from the word aspiring - remind them the key word today is keen.""",


    20: """SAY:
- Which activity would you be keen to do?
- Jump on a trampoline, or clean your room?
- Point to your choice.

DO:
- Give wait time.
- Students point.
- Take 1-2 explanations - why are you keen to do that one?

TEACHER NOTES:
Accept either answer if students can explain keen in the reason. The goal is using the word, not agreeing with a favourite.

WATCH FOR:
- Students changing their choice when a friend picks differently - protect independent thinking with wait time first.""",


    21: """SAY:
- Look at the two children. Which child looks keen to go on the rollercoaster?
- Show me A or B.

DO:
- Wait time.
- Show Me Boards.
- Ask one student how they could tell - what were the clues?

TEACHER NOTES:
Students read facial and body clues to match the word to a feeling. Name the clues (big smile, eyes wide, hands up) when you reveal.

WATCH FOR:
- Students picking the scared child - rewind: keen means want to do it.""",


    22: """SAY:
- I will read four sentences. Thumbs up if the sentence shows the person is keen. Thumbs down if not.
- I don't want to get my hair cut.
- I can't wait to go to the beach!
- I am so excited to ride my bike!
- I am too tired to go to swimming lessons.

DO:
- Pause between sentences for thumbs.
- Scan each time.

TEACHER NOTES:
Students are matching the word to language clues (can't wait, so excited, don't want to, too tired). Name the clue words when you reveal each answer.

WATCH FOR:
- Students copying their neighbour - scan quickly on each one.""",


    23: """SAY:
- Show me which word fits best with keen.
- 1 for bored, 2 for scared, 3 for excited, 4 for surprised.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me. Scan.

CFU CHECKPOINT:
Technique: Show Fingers (1-4)
Script:
- Which word fits best with keen? On three, show me.
- Scan for: fingers showing 3 (excited).
PROCEED:
- If the class is on 3, move on to the sentence-level writing section.
PIVOT:
- Most likely misconception: students pick surprised because excited and surprised feel similar.
- Reteach: keen is wanting to do something before it happens.
- Fresh re-check: Are you keen for something after school? Thumbs up or down.

WATCH FOR:
- Students mirroring the fastest hand - keep wait time firm.""",


    25: """SAY:
- A sentence is a complete thought that makes sense.
- Every sentence has a who or what - the subject.
- Every sentence has a verb - a doing or did word.
- Every sentence starts with a capital letter and ends with a punctuation mark.

DO:
- Point to each icon as you say each rule.
- Have students echo each rule back.

TEACHER NOTES:
This is the anchor for the rest of the lesson. Keep it short and clear - students will practise soon.

WATCH FOR:
- Students getting stuck on the word subject - tie it to who or what every time.""",


    26: """SAY:
- Listen to this sentence: The penguins swim through the icy water to the shore.
- Who or what is doing the action? The penguins. That is the subject.
- What are they doing? They swim. That is the verb.
- Does it start with a capital letter? Yes. Does it end with a full stop? Yes.
- So this is a complete sentence.

DO:
- Point to The penguins and swim as you name each part.
- Read the explanation at the bottom aloud to confirm.

TEACHER NOTES:
Model the thinking aloud step by step so students can copy it on the next slides.

WATCH FOR:
- Students who name icy water as the subject - point back to who is doing the action.""",


    27: """SAY:
- If a group of words is missing a subject or a verb, we call it a fragment.
- A fragment is not a complete sentence.
- Listen: watch the dazzling fireworks. This has a verb, but no subject. Who is watching? We do not know. This is a fragment.
- Listen: In Sydney, the excited families. This has a subject, but no verb. What are they doing? We do not know. This is a fragment.

DO:
- Point to each example as you read it.
- Ask: Is this a complete thought? Answer: no.

TEACHER NOTES:
Keep hammering who/what + what doing. The missing piece is what makes it a fragment.

WATCH FOR:
- Students thinking any short group of words is a fragment - link it back to missing subject or missing verb.""",


    28: """SAY:
- We can turn a fragment into a sentence by adding the missing piece.
- If the fragment has a who or what, we add a what doing.
- If the fragment has a what doing, we add a who or what.
- Watch: In Sydney, the excited families plus watch the dazzling fireworks makes a full sentence.
- In Sydney, the excited families watch the dazzling fireworks.

DO:
- Point to the subject box, then the verb box, then the joined sentence.
- Read the full sentence aloud with expression.

TEACHER NOTES:
This is the core I Do for the lesson. Emphasise that we add only the missing part - we do not rewrite everything.

WATCH FOR:
- Students trying to replace the fragment instead of adding to it - redirect: what is missing?""",


    29: """SAY:
- I will show you a group of words. You decide - is it a sentence or a fragment?
- Hands on head for a sentence. Hands on shoulders for a fragment.
- Let's try each one: the excited supporters... the crocodiles' eyes glow in the torchlight... the mysterious lights shine at night... rises above the ocean... the crowd quietly listens to The Last Post... hatch from their eggs and scuttle towards the ocean.

DO:
- Read each group clearly.
- Give wait time before the signal.
- Scan for hands on head or shoulders.
- Reveal the answer and the corrected sentence each time.

CFU CHECKPOINT:
Technique: Movement (hands on head for sentence, hands on shoulders for fragment)
Script:
- Sentence - hands on head. Fragment - hands on shoulders. Choose now.
- Scan for: correct response for each one.
PROCEED:
- If 80% or more are correct on each, move to the I Do on subject vs verb fragments.
PIVOT:
- Most likely misconception: students think a capital letter or a full stop makes it a sentence.
- Reteach: say the rule again - a sentence needs a subject AND a verb.
- Fresh re-check: the happy dogs - sentence or fragment?

WATCH FOR:
- Students copying their neighbour's movement - scan quickly and change order next time.""",


    30: """SAY:
- The subject is the person or thing doing the action - the who or what.
- A subject fragment has a subject, but is missing a verb.
- Here is an example: In Brisbane, the magical Ferris wheel.
- We have a what - the Ferris wheel - but no doing word. So this is a subject fragment.

DO:
- Point to the magical Ferris wheel and name it as the subject.
- Ask: Is there a doing word? Shake head.

TEACHER NOTES:
Make the name stick - subject fragment has a subject, missing a verb.

WATCH FOR:
- Students reading extra detail words as verbs (e.g. magical) - confirm that verbs show action.""",


    31: """SAY:
- A verb is an action word. It tells what the subject is doing.
- A verb fragment has a verb, but is missing a subject.
- Here is an example: spins rapidly up to the stars.
- There is a doing word - spins - but no who or what. So this is a verb fragment.

DO:
- Point to spins as the verb.
- Ask: Who or what spins? Shake head - we don't know.

TEACHER NOTES:
Mirror the slide before. Say the matching names aloud - verb fragment has a verb, missing a subject.

WATCH FOR:
- Students naming the missing piece as an adjective or noun - steer back to who or what is doing the action.""",


    32: """SAY:
- Listen: The waddling penguins with short legs and large feet.
- Is this a sentence or a fragment? What is missing?
- It has a subject - the waddling penguins - but no verb. This is a subject fragment.
- Let's add a verb: are heading for the dunes.
- Now our sentence is complete: The waddling penguins with short legs and large feet are heading for the dunes.

DO:
- Read the fragment, then ask the question.
- Think aloud as you identify the subject and the missing verb.
- Read the whole sentence aloud when it is complete.

TEACHER NOTES:
This is the last teacher-led example before We Do. Students should see the full thinking before they try one.

WATCH FOR:
- Students who think the description (short legs, large feet) is the verb - point back to what doing.""",


    33: """SAY:
- Read with me: twinkle like stars in the night sky.
- Is this a subject fragment or a verb fragment? What is missing?
- Turn and tell your partner what is missing.
- Now let's add it together - The city lights twinkle like stars in the night sky.

DO:
- Give 20 seconds of partner talk.
- Take 1-2 responses.
- Build the full sentence on the board, or reveal the completed one.

CFU CHECKPOINT:
Technique: Turn and Tell, then cold call
Script:
- Is this a subject fragment or a verb fragment?
- Scan for: students saying verb fragment.
PROCEED:
- If most pairs say verb fragment, move to the next example.
PIVOT:
- Most likely misconception: students say subject fragment because it looks short.
- Reteach: name the verb (twinkle), then ask who or what twinkles.
- Fresh re-check: float on the sea - subject or verb fragment?

WATCH FOR:
- Students adding extra detail but not a subject - redirect: who or what twinkles?""",


    34: """SAY:
- Read with me: In the evening, the excited children with huge smiles on their faces.
- What is missing? Is this a subject fragment or a verb fragment?
- This is a subject fragment. It needs a verb.
- Let's add one together - go prawn fishing on the beach.
- Full sentence: In the evening, the excited children with huge smiles on their faces go prawn fishing on the beach.

DO:
- Read the fragment.
- Partner talk: what is missing?
- Take a response and write the verb on the board.
- Read the completed sentence with expression.

TEACHER NOTES:
Keep the focus on identifying the fragment type first, then adding the missing piece.

WATCH FOR:
- Students who try to replace the long description - remind them only the missing piece needs adding.""",


    35: """SAY:
- Read with me: Near the ocean, the tiny hatchlings.
- Is this a subject fragment or a verb fragment?
- This is a subject fragment. It needs a verb.
- Let's add one - break free from their eggs.
- Full sentence: Near the ocean, the tiny hatchlings break free from their eggs.

DO:
- Read the fragment.
- Partner talk - what doing word could we add?
- Take a response, write the verb, read the full sentence.

TEACHER NOTES:
Third We Do fragment - pace should pick up as confidence grows.

WATCH FOR:
- Students adding only are - push for a stronger action verb.""",


    36: """SAY:
- Read with me: cheer loudly for their favourite team.
- Is this a subject fragment or a verb fragment?
- This is a verb fragment. It needs a subject.
- Let's add one - At the Melbourne Cricket Ground, the excited fans.
- Full sentence: At the Melbourne Cricket Ground, the excited fans cheer loudly for their favourite team.

DO:
- Read the fragment.
- Partner talk - who or what cheers?
- Take a response and complete the sentence.

TEACHER NOTES:
Students should now notice verb fragments quickly. Celebrate correct identification before rewriting.

WATCH FOR:
- Students using a pronoun only (they cheer) - accept it, then stretch: can we add more detail about who?""",


    37: """SAY:
- Look at the starting words: In the evening, local families.
- Which example shows this fragment correctly turned into a complete sentence - A or B?
- A: In the evening, local families with big smiles and excited faces.
- B: In the evening, local families watch exciting movies at the outdoor cinema.
- Show me A or B.

DO:
- Give 5 seconds wait time.
- Say: 3, 2, 1, show me.
- Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Write A or B. 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move to the booklet task.
PIVOT:
- Most likely misconception: students pick A because it sounds longer and more detailed.
- Reteach: check for a verb. A has no verb. B has watch - that is the verb that completes the sentence.
- Fresh re-check: The tired koala sleeps in the tree - sentence or fragment?

WATCH FOR:
- Students who assume more words always means a sentence - pull back to subject plus verb.""",


    38: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 3 - Sentence level writing.
- Read each fragment carefully.
- Decide if it is a subject fragment or a verb fragment.
- Add the missing piece to make a complete sentence.

DO:
- Hand out booklets or direct students to the page.
- Circulate and check the first 2-3 answers for each student.
- Remind students of the rule - every sentence needs a subject and a verb.

TEACHER NOTES:
Use this time to scan for students who still confuse subject and verb fragments. Pull a small group if needed.

WATCH FOR:
- Students who write new sentences that still miss a verb - prompt: what is the doing word?
- Students finishing fast - ask them to read their sentences aloud and check for a capital letter and full stop.""",
}


LESSON_4_NOTES = {

    10: """SAY:
- Authors write texts for different reasons.
- Some texts are written to entertain - like a funny story or a picture book we read for fun.
- Some texts are written to inform - like a fact book about animals.
- Some texts are written to persuade - like a poster that tries to get you to buy something.

DO:
- Point to each purpose as you name it.
- Give one quick familiar example for each (e.g. a picture book we have read, a non-fiction animal book, a school fundraiser poster).

TEACHER NOTES:
This is the anchor for the next few slides. Keep the three purposes short and clear so students can sort texts by purpose.

WATCH FOR:
- Students who think all books are to entertain - show the non-fiction example again.""",


    11: """SAY:
- Why do you think Frane Lessac wrote Under the Southern Cross?
- Was it to entertain us, to inform us, or to persuade us?
- Turn and tell your partner what you think and why.

DO:
- Give 20 seconds of partner talk.
- Take 2-3 responses from the room.
- Affirm that the book is a mix - we learn facts about Australia AND enjoy the journey. We will unpack this next.

TEACHER NOTES:
Accept to entertain, to inform, or both, as long as students can give a reason. This sets up the idea of a hybrid text on the next slide.

WATCH FOR:
- Students who only say to entertain - prompt back to the factual parts of the book (the facts about each place).""",


    12: """SOURCES:
Image credit: Free SVG

SAY:
- The word hybrid means something that is a mix of two different things.
- For example, a hybrid car has two ways of getting power - electricity and petrol.
- Say the word with me... hybrid!

DO:
- Point to the hybrid car image.
- Use your fingers to show two parts coming together.

TEACHER NOTES:
Use the car as a concrete anchor before applying the word to a text. Students need the everyday meaning first.

WATCH FOR:
- Students who hear hybrid for the first time - have them repeat it twice.""",


    13: """SAY:
- A hybrid text combines features of two types of texts together into one text.
- Under the Southern Cross combines a narrative journey layer with an informative layer.
- The narrative journey layer tells a story as we travel around Australia at night.
- The informative layer gives us facts about each place.

DO:
- Point to the word narrative, then to the word informative.
- Hold up the book and flick through a page so students can see the two different print sizes.

TEACHER NOTES:
Name both layers clearly. Students will be asked to spot them on the next slides.

WATCH FOR:
- Students who think the two layers are two different books - re-show the single page with both types of text.""",


    14: """SOURCES:
Image credit: Free SVG

SAY:
- Open the book to pages 1 and 2.
- Can you see two layers on the page? One tells the story, one gives us facts.
- Why do you think the two layers are printed and positioned differently?
- What is different about the language features used in both layers?

DO:
- Make sure every student can see pages 1-2 (show the book, or direct them to their copies).
- Partner talk for 30 seconds.
- Take 2-3 responses and highlight the bigger print (narrative) versus the smaller factual text (informative).

TEACHER NOTES:
The layout itself signals the two layers. Name the signals: bigger print and story-style is narrative; smaller print with facts is informative.

WATCH FOR:
- Students who focus only on the illustrations - pull them back to the two text layers.""",


    15: """SAY:
- Let's look at pages 15 and 16.
- On these pages, we can see the language features that belong to each layer.
- Narrative layer uses adjectives, rhythmic phrases, repetition of Under the Southern Cross, and imagery.
- Informative layer uses factual information, technical language, proper nouns, geographical references, and explanatory paragraphs.

DO:
- Read a short example of narrative language from the page.
- Read a short example of informative language from the page.
- Point to each feature on the slide as you name it.

TEACHER NOTES:
This is a big list - do not teach every feature deeply. Name them and use the examples on the page to anchor the key difference: story-sounding words versus fact-sounding words.

WATCH FOR:
- Students overwhelmed by the vocabulary - focus on two contrasts: descriptive words for narrative, and proper nouns or facts for informative.""",


    16: """SAY:
- Listen to this sentence: Across Australia, the Southern Cross is a constellation that can be seen in the night sky year-round.
- Is this narrative or informative?
- I can hear a fact - year-round - and a technical word - constellation. This is the informative layer.

DO:
- Read the sentence aloud twice.
- Think aloud as you identify the clues (fact, technical word).
- Point to the word informative as you confirm.

TEACHER NOTES:
This is the I Do for sorting by layer. Model the thinking students will copy on the next slides.

WATCH FOR:
- Students who just vote without clues - name the clue word (constellation, year-round) each time.""",


    17: """SAY:
- Listen: In the enchanted rainforest, eyes gleam in the dark while frogs sing their midnight songs.
- Is this the narrative layer or the informative layer?
- Turn and tell your partner your answer and one clue word.

DO:
- Read the sentence clearly.
- Partner talk for 20 seconds.
- Cold call 1-2 pairs to share.
- Reveal and name the clues (enchanted, gleam, sing - descriptive story words).

TEACHER NOTES:
Expected answer: narrative layer. Students should point to descriptive or story-sounding words.

WATCH FOR:
- Students pointing to rainforest as a fact - remind them that facts sound like information, not a story scene.""",


    18: """SAY:
- Listen: In Broome, the Staircase to the Moon appears when the full moon rises over the ocean at low tide.
- Is this narrative or informative?
- Show me with your hand - thumbs up if informative, thumbs down if narrative.

DO:
- Read the sentence twice.
- Give 5 seconds wait time.
- Scan thumbs.
- Reveal and name the clues (Broome as proper noun, factual description of when it happens).

TEACHER NOTES:
Expected answer: informative. Even though it sounds lovely, the clues are proper noun and factual explanation.

WATCH FOR:
- Students who pick narrative because of the word moonlight-type images - re-cue on the fact that it explains when and where.""",


    19: """SAY:
- Read with me.
- A: In the outback, glowing balls called Min Min lights are sometimes seen at night.
- B: The desert glows softly as strange orbs of light float gently across the sand. A dingo watches quietly nearby.
- Which sentence is written with narrative features? Show me A or B.

DO:
- Read both options clearly.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which has narrative features? 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move on to the next CfU.
PIVOT:
- Most likely misconception: students pick A because Min Min lights sound interesting.
- Reteach: narrative uses descriptive story words - glows softly, float gently, quietly. A uses facts - are sometimes seen.
- Fresh re-check: Read a short narrative sentence and a short informative sentence. Which is which?

WATCH FOR:
- Students who choose on length - remind them the clue is in the type of words, not how long the sentence is.""",


    20: """SAY:
- Read with me.
- A: Under the Southern Cross, the waves glitter as sea turtles crawl towards the moonlit beach, while a young girl watches in silence from the dunes.
- B: The Aurora Australis is a natural light display in the sky, seen in the southern parts of the world.
- Which sentence is written with informative features? Show me A or B.

DO:
- Read both options clearly.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which has informative features? 3, 2, 1, show me.
- Scan for: most boards show B.
PROCEED:
- If most students show B, move to the booklet task.
PIVOT:
- Most likely misconception: students pick A because it mentions the Southern Cross and feels factual.
- Reteach: B explains what the Aurora Australis is and where it is seen - that is an informative explanation. A describes a scene with characters.
- Fresh re-check: Say one clue word for an informative sentence - expected: explain, fact, describes what something is.

WATCH FOR:
- Students who mix up the two layers - re-anchor: narrative = story, informative = facts.""",


    21: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 4 - Understanding, responding to and creating text.
- Read each sentence and decide which layer it belongs to.

DO:
- Hand out booklets or direct students to the page.
- Circulate and check the first 2-3 answers for each student.
- Remind students to look for the clue words each time.

TEACHER NOTES:
This is an independent practice moment. Pull a small group if many students are still unsure after the CfU checkpoints.

WATCH FOR:
- Students who pick randomly - ask them to underline one clue word in each sentence before they decide.""",


    23: """SAY:
- A noun is a naming word.
- It names a person, place, animal or thing.
- We also call this a common noun.
- A common noun does not need a capital letter.
- Say these examples with me... penguin... cinema... turtles.

DO:
- Point to each example on the slide.
- Echo read the word list together.

TEACHER NOTES:
Keep the noun definition ultra simple before introducing proper nouns on the next slide.

WATCH FOR:
- Students who try to add a capital to penguin - remind them common nouns are general names, no capital needed.""",


    24: """SOURCES:
The Grammar Book: Understanding and teaching primary grammar, Zoe Paramour and Timothy Paramour (2020).

SAY:
- A proper noun is special - it names a specific person, place, object, animal or title.
- A proper noun always begins with a capital letter.
- Say these examples with me... Southern Cross... Phillip Island... Milky Way.

DO:
- Point to each proper noun on the slide.
- Show the capital letter at the start of each one.

TEACHER NOTES:
This is the key rule of the lesson - proper nouns always start with a capital letter. Students will practise this next.

WATCH FOR:
- Students who think all nouns need capitals - remind them penguin is a common noun, Phillip Island is a proper noun.""",


    25: """SAY:
- Let's look at common and proper noun examples from our book.
- Common noun examples: penguin, cinema, lights, turtle, fireworks, bush - no capital letter.
- Proper noun examples: Southern Cross, Phillip Island, Queensland, Min Min Hotel, Australian Football League, Milky Way - capital letter at the start of each word.

DO:
- Read one column, then the other.
- Point to a capital letter on each proper noun.

TEACHER NOTES:
Connect the grammar content back to Under the Southern Cross so students see these nouns in context.

WATCH FOR:
- Students who think Min Min Hotel needs only one capital - point to each word that needs a capital.""",


    26: """SAY:
- Read the words on the screen.
- Which words are proper nouns? Which ones need a capital letter?
- Animals, eggs, November, Darwin, ocean, South Pacific.
- Turn and tell your partner which ones are proper nouns.

DO:
- Give 20 seconds of partner talk.
- Cold call 2-3 students to share.
- Reveal: November, Darwin, South Pacific.

TEACHER NOTES:
Students pick out the proper nouns from a mixed list. The clue is that proper nouns name specific people, places or things.

WATCH FOR:
- Students who pick ocean or animals - remind them these are general names (common nouns), not specific ones.""",


    27: """SAY:
- Read the words.
- Gallipoli, fire, South Pole, dawn, Perth, crocodiles.
- Which ones are proper nouns? Write them down with a capital letter at the start.

DO:
- Give 30 seconds of independent writing time.
- Circulate and scan 3-4 books.
- Reveal: Gallipoli, South Pole, Perth.

CFU CHECKPOINT:
Technique: Independent write, then cold call for answers
Script:
- Write the proper nouns on your whiteboard with a capital letter at the start.
- Scan for: Gallipoli, South Pole, Perth.
PROCEED:
- If most students have all three, move on to capital letter correction.
PIVOT:
- Most likely misconception: students capitalise crocodiles or fire because they sound important.
- Reteach: Is it a specific place, person, or title? Crocodiles names a type of animal - common noun.
- Fresh re-check: Is Sydney a proper noun? Thumbs up or down.

WATCH FOR:
- Students who miss Perth because it is a single word - remind them that cities are proper nouns.""",


    28: """SAY:
- Read this sentence: The southern cross can be seen all year round.
- What is the proper noun? Southern Cross - a specific name.
- So we need capital letters - S for Southern, C for Cross.
- The corrected sentence is: The Southern Cross can be seen all year round.

DO:
- Read the original sentence aloud.
- Point to southern cross and circle it with your finger.
- Read the corrected sentence with the capital letters emphasised.

TEACHER NOTES:
Model the thinking: find the proper noun, then check every word in the name for a capital letter.

WATCH FOR:
- Students who only capitalise Southern and leave cross lowercase - point out that the full name needs capitals on every word.""",


    29: """SAY:
- Read this sentence with me: There are 100 to 400 billion stars in the milky way.
- What is the proper noun that needs a capital? The Milky Way.
- Rewrite the sentence with a capital M for Milky and a capital W for Way.
- Turn and tell your partner the corrected version.

DO:
- Give 20 seconds of partner talk.
- Take one response.
- Reveal: There are 100 - 400 billion stars in the Milky Way.

TEACHER NOTES:
Keep the focus on the proper noun only. Other words in the sentence do not change.

WATCH FOR:
- Students who capitalise stars or billion - redirect to what is a specific name here.""",


    30: """SAY:
- Read this sentence with me: Loggerhead turtles hatch at mon repos near bundaberg.
- There are two proper nouns that need capitals - Mon Repos and Bundaberg.
- Rewrite the sentence with the capital letters in the right places.

DO:
- Partner talk for 30 seconds.
- Take a response.
- Reveal: Loggerhead turtles hatch at Mon Repos near Bundaberg.

TEACHER NOTES:
This one has two proper nouns, which is the stretch. Prompt students to find both names before writing.

WATCH FOR:
- Students who only fix one name - ask: is there another specific place here?""",


    31: """SAY:
- Look at the two sentences.
- The Deckchair Cinema is home to great movies and lots of animals.
- The wheel of brisbane looks over the Brisbane River.
- Are these sentences punctuated correctly? Show me thumbs up if yes, thumbs down if no.

DO:
- Read each sentence clearly.
- Give 10 seconds of thinking time.
- Scan thumbs.
- Reveal the correction: The Wheel of Brisbane looks over the Brisbane River.

CFU CHECKPOINT:
Technique: Thumbs up or thumbs down
Script:
- Is sentence 2 correct? Thumbs up or down.
- Scan for: thumbs down (it is not correct - wheel of brisbane needs capitals).
PROCEED:
- If most students show thumbs down, move on to sentence expansion.
PIVOT:
- Most likely misconception: students miss that wheel of brisbane is the name of a specific thing.
- Reteach: the Wheel of Brisbane is a specific landmark, just like the MCG - it is a title, so every main word needs a capital.
- Fresh re-check: the sydney opera house - correct or not?

WATCH FOR:
- Students who only capitalise Wheel - remind them of, and, the stay lowercase but Brisbane is a proper noun.""",


    32: """SAY:
- We can use what we know about common and proper nouns to expand sentences and add more detail.
- Listen: They are watching a movie at night.
- Who is doing it? We do not know. Let's add a who/what.
- The families are watching a movie at night.
- Now we can read pages 3 and 4 to get even more detail.

DO:
- Point to They, then to the families.
- Point to each part of the build as you name it.
- Read the final expanded sentence aloud with expression.

TEACHER NOTES:
This is the I Do for expanding sentences. The move students need to copy: swap or add a specific who/what to add detail.

WATCH FOR:
- Students who rewrite the whole sentence - remind them we only add or swap the who/what piece.""",


    33: """SAY:
- Read pages 7 and 8 with me.
- Here is a simple sentence: He plays football at the MCG.
- Who is he? Can we add more detail to the who/what?
- Let's try: the athlete plays football.
- Full expanded sentence: The athlete plays football at the MCG.

DO:
- Point to He, then swap it for the athlete.
- Partner talk - can we add more detail about who?
- Take a response and reveal the expanded sentence.

TEACHER NOTES:
The MCG is a great proper noun example - link this back to capital letters for proper nouns.

WATCH FOR:
- Students who add more detail after football - gently redirect: we are expanding the who/what piece today.""",


    34: """SAY:
- Read pages 15 and 16 with me.
- Here is a simple sentence: They exploded in bright colours over the sky.
- Who or what is they?
- Let's add detail: the fireworks exploded in bright colours.
- Full sentence: The fireworks exploded in bright colours over the sky.

DO:
- Point to They on the slide.
- Partner talk - who or what exploded?
- Take a response, then reveal.

TEACHER NOTES:
Students should be faster at this move now. Celebrate the swap - they becomes the fireworks.

WATCH FOR:
- Students who say the word they again - prompt: can we say exactly what kind of thing?""",


    35: """SAY:
- Read pages 21 and 22 with me.
- Simple sentence: They can be caught in the Swan River.
- Who or what is they?
- Let's add: prawns can be caught.
- Full sentence: Prawns can be caught in the Swan River.

DO:
- Point to They.
- Partner talk to identify the who/what.
- Take a response and reveal.

TEACHER NOTES:
Swan River is another proper noun in context - point to the capital S and capital R while reading.

WATCH FOR:
- Students who forget the capital letter on Prawns at the start of the sentence - remind them every sentence starts with a capital letter.""",


    36: """SAY:
- Read the sentence: It tells us about our military's history.
- We need a who/what to expand this sentence. Read pages 26 and 27 to help.
- Which detail fits best? A - The Australian War Memorial. B - The Last Post.
- Show me A or B.

DO:
- Read both options clearly.
- Give 10 seconds wait time.
- Say: 3, 2, 1, show me. Scan all boards.

CFU CHECKPOINT:
Technique: Show Me Boards (A or B)
Script:
- Which best expands the sentence? 3, 2, 1, show me.
- Scan for: most boards show A (The Australian War Memorial).
PROCEED:
- If most students show A, move to the booklet task.
PIVOT:
- Most likely misconception: students pick B because The Last Post appears in the book and is memorable.
- Reteach: The Last Post is a tune, not a place that tells us about history. The Australian War Memorial is the place that tells us about our military's history.
- Fresh re-check: point to a proper noun on the book cover - is it a thing or a tune?

WATCH FOR:
- Students who forget to look at the original sentence meaning - re-read: it tells us about our military's history. Which one does that?""",


    37: """SAY:
- Time to try this on your own.
- Turn to the page titled: Lesson 4 - Sentence level writing.
- Remember to use a capital letter for every proper noun.
- Add a who/what detail to expand each sentence.

DO:
- Hand out booklets or direct students to the page.
- Circulate and check the first 2-3 answers for each student.
- Remind students that common nouns do not need a capital, but proper nouns do.

TEACHER NOTES:
Scan for two errors: missing capital letters on proper nouns, and replacing instead of adding detail. Pull a small group if needed.

WATCH FOR:
- Students who leave proper nouns lowercase - point back to the capital letter rule.
- Students finishing fast - ask them to read their sentences aloud and check that every proper noun has a capital letter.""",
}


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    lines = text.split("\n")
    p = notes_tf.paragraphs[0]
    p.text = lines[0]
    for line in lines[1:]:
        p = notes_tf.add_paragraph()
        p.text = line


def apply(src, out, notes):
    pres = Presentation(src)
    n = 0
    for i, slide in enumerate(pres.slides, 1):
        if i in notes:
            set_notes(slide, notes[i])
            n += 1
    pres.save(out)
    print(f"Applied notes to {n} slides -> {out}")


def main():
    apply(L3_SRC, L3_OUT, LESSON_3_NOTES)
    apply(L4_SRC, L4_OUT, LESSON_4_NOTES)


if __name__ == "__main__":
    main()
