"""Add teacher notes to T2 Wk 4 Maths Discovery 1 Planner 2025.pptx.

Slide faces, theme, layout, animations, transitions and hyperlinks are unchanged.
Notes are written in plain paragraphs (a:buNone) to match the deck's existing
notes style. Section headers are uppercase + colon. Blank lines separate
sections. ASCII-safe; no em dashes; straight quotes only.
"""

from pathlib import Path
from copy import deepcopy
from lxml import etree
from pptx import Presentation

SRC = Path(r"C:\Users\09560329\Downloads\T2 Wk 4 Maths Discovery 1 Planner 2025.pptx")
DST = SRC.with_name(SRC.stem + " - with teacher notes.pptx")

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
A = "{%s}" % A_NS
NSMAP = {"a": A_NS}


def make_paragraph(text: str) -> etree._Element:
    """Build a single <a:p> with buNone (matches the deck's existing notes style).

    Empty text -> blank separator paragraph.
    """
    p = etree.Element(A + "p")
    pPr = etree.SubElement(p, A + "pPr")
    pPr.set("indent", "0")
    pPr.set("marL", "0")
    pPr.set("rtl", "0")
    pPr.set("algn", "l")
    lnSpc = etree.SubElement(pPr, A + "lnSpc")
    spcPct = etree.SubElement(lnSpc, A + "spcPct")
    spcPct.set("val", "100000")
    spcBef = etree.SubElement(pPr, A + "spcBef")
    etree.SubElement(spcBef, A + "spcPts").set("val", "0")
    spcAft = etree.SubElement(pPr, A + "spcAft")
    etree.SubElement(spcAft, A + "spcPts").set("val", "0")
    etree.SubElement(pPr, A + "buNone")
    if text:
        r = etree.SubElement(p, A + "r")
        rPr = etree.SubElement(r, A + "rPr")
        rPr.set("lang", "en-AU")
        rPr.set("dirty", "0")
        t = etree.SubElement(r, A + "t")
        t.text = text
    else:
        # endParaRPr keeps PowerPoint happy on empty paragraphs
        endRPr = etree.SubElement(p, A + "endParaRPr")
        endRPr.set("lang", "en-AU")
    return p


def set_notes(slide, notes_text: str) -> None:
    """Replace the notes text with notes_text. Preserves the notes_slide shape."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    txBody = tf._txBody
    # Remove all existing <a:p> elements
    for p in list(txBody.findall(A + "p")):
        txBody.remove(p)
    # Split notes_text into lines and build paragraphs
    lines = notes_text.split("\n")
    # Strip a single leading and trailing blank line if present
    while lines and lines[0] == "":
        lines.pop(0)
    while lines and lines[-1] == "":
        lines.pop()
    for line in lines:
        txBody.append(make_paragraph(line))


# ============================================================================
# NOTES BY SLIDE
# ============================================================================
# Convention: each entry is a multiline string. Section headers are UPPERCASE
# followed by a colon. Blank lines separate sections. No leading "- " on
# content lines (the deck's notes style does not use bullets).

NOTES = {}

NOTES[1] = """\
TEACHER NOTES:
Unit cover slide for Week 4, Term 2. Open here as students settle. No teaching action needed.
"""

NOTES[2] = """\
TEACHER NOTES:
Title slide announcing this week's focus on 2D shape. Move on quickly to the planning slides.
"""

NOTES[3] = """\
TEACHER NOTES:
Curriculum reference for teacher planning. Year 1 focus: Recognise and classify familiar 2D shapes and three-dimensional objects using obvious features. Skip past for student-facing teaching.
"""

NOTES[4] = """\
SAY:
Read the learning intention with me.
These are the three things we are practising this week.
Ask: what does '2D' mean in your own words? Expected: a flat shape.

DO:
Point to the learning intention first.
Track each success criterion with your finger as you read it.
Keep this slide brief - one or two minutes.

TEACHER NOTES:
This is the week-level LI/SC, not just one lesson. Refer back at the start of each session. The first criterion (explain what a 2D shape is) is the most concrete entry point.

WATCH FOR:
Students who confuse 2D with 3D objects in the room. Be ready to point to a flat picture next to a real cube during the launch.
"""

NOTES[5] = """\
SAY:
These are the words we will hear and use this week.
Some you may know already. Some will be new.
We will keep meeting them in our lessons.

DO:
Read each word aloud with the class.
Point to any word students already know and ask one student to use it in a sentence.
Do not pre-teach every word here. Most are introduced again on their own keyword slide.

TEACHER NOTES:
Use this slide as a quick scan. Spend longer on 'vertex/vertices' and 'polygon' if those are new. The dedicated keyword slides will revisit each term inside the lessons.

WATCH FOR:
Students who say 'diamond' when they mean rhombus or square. That is fine for Year 1; you will firm up names across the week.
"""

NOTES[6] = """\
TEACHER NOTES:
Teacher planning overview. Use to check resources before each session. Sessions 2, 4 and 5 need pre-prepared materials: print mystery bag sheets and make 12 sets before Lesson 2; print the Shape Identifying BLM before Lesson 4; pre-make or borrow bingo cards before Lesson 5. Skip past for student-facing teaching.
"""

NOTES[7] = """\
TEACHER NOTES:
Lesson 1 divider. Note: this slide reads 'Topic: Patterns' but Session 1 in the overview is 'Naming, identifying and describing two-dimensional shapes'. Treat today as a 2D shape lesson and ignore the 'Patterns' label. Do not change the slide unless asked.
"""

NOTES[8] = """\
SAY:
Daily review. Quick warm up before our shape lesson.
Count by 2s with me as I point to each pair of elephants.
Ask: what number comes after 18 when we count by 2s? Expected: 20.

DO:
Point to each pair of elephants as the class counts.
Read the answer numbers aloud as they appear: 2, 4, 6 and so on.
If students notice that the visible numbers jump out of order, pause and rebuild the count by tapping pairs.

TEACHER NOTES:
This is review of skip counting by 2s, not new teaching. The visible sequence on the slide jumps from 18 to 22 to 24 to 20 - if students notice, point to the pairs in order and rebuild together. Keep the pace brisk.

WATCH FOR:
Students saying every number instead of every second number. Tap two elephants together as you say each new number.
"""

NOTES[9] = """\
SAY:
Daily review. Pulling numbers apart into two parts.
Watch how I split 11 into 4 and another number.
Ask: 4 and what makes 11? Expected: 7.
Now your turn. Split 15 into two parts. Show me on your board.

DO:
Model 11 = 4 + 7 first using fingers or counters.
Cue boards up after a short think time.
Accept any correct partition for 15: 10 + 5, 8 + 7, 9 + 6 and so on.

TEACHER NOTES:
The slide says answers may vary. Praise different correct splits. This is review of partitioning, not the focus of today's new lesson - keep it short.

WATCH FOR:
Students writing only one number. Re-cue: I need two numbers that add to 15.
"""

NOTES[10] = """\
SAY:
Now we split into three parts.
Watch this one first. 11 splits into 3, 3 and 5.
Ask: do 3 plus 3 plus 5 make 11? Expected: yes.
Your turn. Split 11 a different way, into three parts. Show me on your board.

DO:
Point to each of 3, 3 and 5 on the worked example.
Give 30 seconds think time.
Take two student examples to share aloud.

TEACHER NOTES:
Many correct answers exist. Accept any three numbers that add to 11.

WATCH FOR:
Students using one big number and two zeros. Prompt: try three numbers that are bigger than zero.
"""

NOTES[11] = """\
SAY:
Quick fluency. Skip counting by 10s.
Forwards first: 10, 20, 30 and so on.
Now backwards from 120: 120, 110, 100 and so on.

DO:
Set a brisk pace with claps or finger taps.
Use a class counting chart on display if you have one.
Keep this to about a minute.

TEACHER NOTES:
Fluency only - no new teaching. The 'Linked Enabler and Extender Prompts' line on the slide is a generic label; use your usual differentiation.

WATCH FOR:
Students who lose the count past 100. Pause at 90, restart at 100, finish to 120 together.
"""

NOTES[12] = """\
SAY:
Today we are learning about 2D shapes.
Read the learning intention with me.
Ask: what is one shape you already know the name of? Expected: any 2D shape name.

DO:
Point to the learning intention.
Track each success criterion with your pointer or finger.
Take one or two student names of shapes before moving on.

TEACHER NOTES:
This LI/SC is for today only. We will name shapes and label features (sides and vertices). Keep this slide to about a minute.

WATCH FOR:
Students naming 3D objects (a ball, a box). Note these and come back to the difference in the launch.
"""

NOTES[13] = """\
TEACHER NOTES:
Orientation slide for the keyword set that follows. The existing teacher note explains how to use keyword slides - revisit it before front-loading or revising. No new teaching action here.
"""

NOTES[14] = """\
SOURCES:
VIC Mathematics glossary - 'An object with length and width is two-dimensional. A polygon is an example of a two-dimensional geometric object.'

SAY:
Say it with me: two-dimensional. 2D for short.
2D means flat. Like a picture on paper.
Ask: is this slide flat or solid? Expected: flat.

DO:
Hold up a piece of paper next to a classroom block or ball.
Show 'flat' with your hand sliding across the page.
Have students say 'two-dimensional' and 'flat' twice.

TEACHER NOTES:
Anchor 'flat' before any other feature. Source line preserved above.

WATCH FOR:
Students who say a real classroom object is 2D. Hold the object next to a paper picture to compare.
"""

NOTES[15] = """\
SOURCES:
The boundary of a 2D shape (preserved from existing notes).

SAY:
Say it with me: sides.
Sides are the lines around the edge of a flat shape.
Ask: how many sides on a triangle? Expected: 3.

DO:
Trace around a triangle in the air.
Tap each side as you count it.
Have students trace a square in the air and count sides aloud.

TEACHER NOTES:
Build the routine of pointing and counting sides - we will use it every lesson this week.

WATCH FOR:
Students counting the same side twice. Mark the start point with your finger so they know when to stop.
"""

NOTES[16] = """\
SOURCES:
The point where two lines meet (preserved from existing notes).

SAY:
Say it with me: vertex. More than one vertex is vertices.
A vertex is a corner. Where two lines meet.
Ask: how many vertices on a square? Expected: 4.

DO:
Point to a corner of a square and pinch your fingers together.
Have students hold up four fingers for square corners.
Have students say 'vertex' then 'vertices' twice.

TEACHER NOTES:
Many Year 1 students will say 'corners' first. Accept 'corners' as a starting point and add 'vertex' beside it. We will use both words this week.

WATCH FOR:
Students pointing inside the shape, not at the corner. Re-show by pinching at a real corner.
"""

NOTES[17] = """\
TEACHER NOTES:
Divider into the activate prior learning sequence. Move straight to the next slide.
"""

NOTES[18] = """\
SAY:
Together. What is this shape?
Show me on your board: the name, the number of sides, the number of vertices.
Ask: how many vertices does a circle have? Expected: zero.

DO:
Give 20 seconds for boards.
Scan all boards before revealing.
Reveal the answers on the slide once boards are up.

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
Ask: what is this shape called? Expected: circle.
Scan for: most boards showing 'circle', '1 side', '0 vertices'.
PROCEED:
80 percent or more correct. Move to the rectangle.
PIVOT:
Most likely: students write 1 vertex because the curve closes.
Reteach: trace the circle with your finger in the air. Your finger never stops or turns. No corners means no vertices.
Re-check: ask the same question with a fresh circle drawn on the board.

TEACHER NOTES:
Circle is the 'no vertices' anchor for this week. If students get this idea now, the polygon checklist in Lesson 2 will be easier.

WATCH FOR:
Students writing 'round' for shape name. Accept and add 'circle' beside it.
Students writing 4 vertices because they see corners on the slide layout, not the shape itself.
"""

NOTES[19] = """\
SAY:
Same routine. Together.
Name. Number of sides. Number of vertices.
Ask: are all sides the same length on a rectangle? Expected: no, two are longer and two are shorter.

DO:
Point to each side as the class counts.
Pinch each corner as you count vertices.
Reveal the answers after boards are up.

TEACHER NOTES:
Use this slide to firm up the side-and-vertex routine. Some students will say 'square' - prompt them to look at the side lengths.

WATCH FOR:
Students who count 4 sides but trace each one twice. Slow your point.
"""

NOTES[20] = """\
SAY:
Last together one. Triangle.
Name. Sides. Vertices.
Ask: do the sides have to be straight? Expected: yes.

DO:
Point and count sides and vertices together.
Reveal the answers.

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
Ask: how many sides AND how many vertices on a triangle? Expected: 3 and 3.
Scan for: matching number of sides and vertices on every board.
PROCEED:
80 percent or more correct. Move to the launch.
PIVOT:
Most likely: students count sides correctly but write a different number for vertices.
Reteach: tap each side, then pinch each corner. The numbers match for triangle.
Re-check: draw a different triangle on the board and ask for sides and vertices.

TEACHER NOTES:
Useful pattern - for many shapes, the number of sides equals the number of vertices. Note this aloud for stronger students.

WATCH FOR:
Students who insist a shape with one curved line is a triangle if it has 3 corners. Reinforce 'straight sides'.
"""

NOTES[21] = """\
SAY:
Watch the screen. Two-dimensional shapes are flat and closed.
Flat means I could draw it on paper.
Closed means there are no gaps. The line meets up.
Ask: is a triangle closed? Expected: yes.

DO:
Point to each of the three shapes on the slide: circle, quadrilateral, triangle.
Trace the outline of one with your finger to show 'closed'.
Pause and ask one student to trace another shape in the air.

TEACHER NOTES:
This launch sets up the two big rules: flat AND closed. We will add 'straight sides' in Lesson 2 to reach the polygon definition. Do not introduce 'polygon' here yet.

WATCH FOR:
Students who think 'closed' means coloured in. Show an unfilled outline that is still closed.
"""

NOTES[22] = """\
SAY:
Two-dimensional shapes also have length and width.
Length is one way. Width is the other way.
I can draw them on flat paper.

DO:
Stretch your hands left and right for length, then up and down for width.
Point to a shape on the slide and show its length and width with two finger taps.

TEACHER NOTES:
For Year 1, 'length' and 'width' are introduced gently. Do not test the words in isolation; the goal is that students see 2D as having two directions, not three.

WATCH FOR:
Students adding 'height' or 'depth'. Gently say: that is for solid shapes, not flat.
"""

NOTES[23] = """\
SAY:
Look at the shapes. Two of them are labelled A and B.
Show me on your fingers: 1 if A is 2D, 2 if B is 2D, both fingers if both are 2D.
Ask: how do you know? Expected: it is flat.

DO:
Give 15 seconds silent thinking.
Cue fingers up together.
Take one student to point to a flat clue and one to point to a solid clue.

CFU CHECKPOINT:
Technique: Finger Vote
Script:
Ask: which shapes are two-dimensional? Expected: depends on the visible images on the slide.
Scan for: most students picking the flat shape and not the solid one.
PROCEED:
80 percent or more correct. Move to the explore.
PIVOT:
Most likely: students choose a 3D object that looks colourful or familiar.
Reteach: hold up a sheet of paper next to a real cube or ball. Flat versus solid.
Re-check: point to one item in the room and ask: 2D or 3D?

TEACHER NOTES:
The expected answer depends on the actual images shown on the slide. Confirm A and B before teaching.
"""

NOTES[24] = """\
SAY:
Look around our room. What flat shapes do you see?
Open your workbook. Draw and label one shape you can see.
Ask: can you find a rectangle in our room? Expected: a window, a book, a door.

DO:
Model on the board: draw a rectangle, write 'rectangle', write '4 sides', write '4 vertices'.
Place wooden shape blocks on the helper table for enabler students.
Write the four basic shape names on the board so enabler students can copy.

TEACHER NOTES:
Existing teacher note preserved: complete in workbooks; model drawing and labelling on the whiteboard; enablers may use wooden shapes to trace around with shape names written on the board to copy. Keep teacher modelling short - one shape - so students have most of the time to draw and label.

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: trace around a wooden shape and copy the name from the board.
EXTENDING PROMPT:
Task: find and label a shape that is NOT a circle, square, rectangle or triangle (for example a hexagon on the floor mat).

WATCH FOR:
Students who only draw one shape and stop. Prompt them to find a second shape in a different part of the room.
"""

NOTES[25] = """\
SAY:
Now draw as many flat shapes as you know.
Label each one with the name, number of sides, number of vertices.
Like the square example on the screen.

DO:
Point to the worked square example.
Set a 5 to 8 minute timer.
Circulate to enabler students first - place wooden shapes near them.

TEACHER NOTES:
Existing teacher note preserved: complete in workbooks; model drawing and labelling on the whiteboard; enablers may use wooden shapes to trace around. The square on the slide is the model - point to it before students start.

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: trace around two wooden shapes and copy names off the board.
EXTENDING PROMPT:
Task: include a hexagon, pentagon or rhombus and label all features.

WATCH FOR:
Students drawing 3D objects (a cube). Redirect to flat shapes only.
Students who draw shapes well but skip the labels. Cue: now write the name and the numbers.
"""

NOTES[26] = """\
TEACHER NOTES:
Teacher menu of close-the-lesson moves. Pick one to fit your class today. Do not run all options. The most useful for Year 1 is usually 'Ask 3 students to share a strategy' or a quick gallery walk around the workbook drawings.
"""

NOTES[27] = """\
TEACHER NOTES:
Lesson 2 divider. Move straight on.
"""

NOTES[28] = """\
SAY:
Daily review. Looking at numbers on our number chart.
I am looking for numbers larger than the ones I show.
Watch how I find a number larger than 74.
Ask: is 79 larger than 74? Expected: yes.

DO:
Point to 74 on the chart, then to 79 and trace the path forwards.
Then do 118 - point and confirm it is past 100.
Reveal the answers visible on the slide.

TEACHER NOTES:
This is the I Do for the larger-numbers routine. Use the chart on the slide as the anchor. Keep it brisk - the focus is the new lesson, not this review.

WATCH FOR:
Students who confuse 'larger' with 'later in counting'. Both are correct here. Use whichever language makes sense to the student.
"""

NOTES[29] = """\
SAY:
Together. Three numbers on screen: 75, 33, 104.
Are they larger than 71? Show me thumbs up or thumbs down for each.

DO:
Cue thumbs after each number is read.
Use the chart to confirm by pointing to each number.

CFU CHECKPOINT:
Technique: Thumbs Up / Down
Script:
Ask: is 33 larger than 71? Expected: no.
Scan for: most thumbs down for 33, thumbs up for 75 and 104.
PROCEED:
80 percent or more correct. Move to the You Do.
PIVOT:
Most likely: students see '3' in 33 and pick by first digit only.
Reteach: count from 71 forwards on the chart. We do not reach 33.
Re-check: ask: is 50 larger than 71?

TEACHER NOTES:
33 is the trap because it has a '3' that may look big. Use the chart to confirm.

WATCH FOR:
Students who answer fast without checking the chart. Cue: show me on the chart first.
"""

NOTES[30] = """\
SAY:
Your turn on your own.
Write any number larger than 75 on your board.
Hold them up when ready.

DO:
Give 30 seconds.
Cue boards up.
Scan from back row to front row. Take one student to point to their number on the chart.

TEACHER NOTES:
Many correct answers. Accept any number greater than 75 up to 120. If a student writes 75, gently ask: is 75 larger than 75, or the same?

WATCH FOR:
Students who write a 1- or 2-digit number smaller than 75. Walk them to the chart and start at 75.
"""

NOTES[31] = """\
SAY:
Now smaller numbers.
Watch how I find a number smaller than 97.
Ask: is 52 smaller than 97? Expected: yes.

DO:
Point to 97 on the chart, then count back to 52 with finger taps.
Repeat the move with 40.

TEACHER NOTES:
I Do for the 'smaller than' version of the routine. Keep modelling brief.

WATCH FOR:
Students who try to call out before you have modelled. Cue: watch first.
"""

NOTES[32] = """\
SAY:
Together. Are 73, 111 and 50 smaller than 93?
Thumbs up or thumbs down for each.

DO:
Cue thumbs after each number.
Confirm on the chart.

CFU CHECKPOINT:
Technique: Thumbs Up / Down
Script:
Ask: is 111 smaller than 93? Expected: no.
Scan for: most thumbs down for 111, thumbs up for 73 and 50.
PROCEED:
80 percent or more correct. Move to the You Do.
PIVOT:
Most likely: students see '1' in 111 and call it small.
Reteach: count past 100 on the chart. 111 sits past the bottom of the 100s grid.
Re-check: is 105 smaller than 93?

TEACHER NOTES:
111 is the three-digit trap. Use the chart's bottom row.

WATCH FOR:
Students confusing 'smaller than' with 'first in counting'. Use the chart, not just memory.
"""

NOTES[33] = """\
SAY:
On your own. Write a number smaller than 104.
Hold up when ready.

DO:
Give 30 seconds.
Scan boards.
Take one student to confirm on the chart.

TEACHER NOTES:
Many correct answers. If a student writes 104, prompt: smaller, not the same.

WATCH FOR:
Students who write 105 or 110. Walk them to the chart.
"""

NOTES[34] = """\
SAY:
Now we put numbers in order, smallest to largest.
Watch me with these three: 103, 25 and 80.
Ask: which is smallest? Expected: 25.
I find them on the chart. The earliest one in counting is the smallest.

DO:
Point to each number on the chart.
Read them in order: 25, 80, 103.

TEACHER NOTES:
I Do model. The phrase 'earliest in counting' connects ordering to counting forwards.

WATCH FOR:
Students who pick the number with the smallest first digit. Use the chart.
"""

NOTES[35] = """\
SAY:
Now together. Largest to smallest.
Three numbers: 103, 119 and 116.
Show me on your boards in order from biggest to smallest.

DO:
Give 45 seconds.
Reveal: 119, 116, 103.
Use the chart to confirm: 103 is at the start of the bottom row, 119 is near the end.

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
Ask: which is largest of 103, 119 and 116? Expected: 119.
Scan for: most boards starting with 119.
PROCEED:
80 percent or more correct. Move to the fluency.
PIVOT:
Most likely: students struggle to order numbers that are all in the 100s.
Reteach: cover the '1' in front of each. Ask: is 03, 19 or 16 the largest? Then put the 1s back.
Re-check: order 105, 102, 109 from largest to smallest.

TEACHER NOTES:
All three numbers are past 100, which is the trickier case. Use the chart.

WATCH FOR:
Students writing largest to smallest in reverse. Re-cue with arrows: big first.
"""

NOTES[36] = """\
SAY:
Quick count. Forwards to 120.
Now backwards from 120.

DO:
Set a brisk pace with claps.
Use a number chart on display if you have one.

TEACHER NOTES:
Fluency, not new teaching. About a minute total.

WATCH FOR:
Students who lose the count past 100. Pause, restart at 100.
"""

NOTES[37] = """\
SAY:
Today we are learning to identify 2D shapes in our environment.
Read the learning intention with me.
Ask: what is one shape you saw on your way to school? Expected: any 2D shape name.

DO:
Point to the LI.
Read each success criterion in order.

TEACHER NOTES:
Four success criteria here is a lot. The first criterion (name shapes in the environment) is the most achievable. The exit point will be the shape hunt later.

WATCH FOR:
Students who name a 3D object. Note for the launch.
"""

NOTES[38] = """\
TEACHER NOTES:
Repeat keyword orientation slide. The existing note about how to use keyword slides is preserved - revisit before front-loading or revising.
"""

NOTES[39] = """\
SAY:
I want to teach you a new word: polygon.
A polygon is a closed shape. With straight sides only.
Watch me sort. This shape is a polygon. This one is not.
Ask: why is the curved one not a polygon? Expected: it has a curved side, not straight.

DO:
Point to each shape in the Polygon column. Trace one outline with your finger.
Point to each shape in the Not a polygon column. Name what makes it different (curved, open, not flat).
Build the rule on the board: closed + straight sides = polygon.

TEACHER NOTES:
This is the new word for today. Keep the rule visible on the board for the rest of the lesson. The shapes on the slide are the source of the contrast - point to them, do not name shapes that are not visible.

WATCH FOR:
Students who guess shape names instead of checking the rule. Bring them back to: closed and straight.
"""

NOTES[40] = """\
SAY:
I am going to test if this shape is a polygon.
Checklist: closed shape - tick. Straight sides - tick. Two-dimensional - tick.
Ask: is a parallelogram a polygon? Expected: yes.

DO:
Point to each checklist item and the matching feature on the parallelogram.
Reveal each tick one at a time as you check.
Conclude: yes, parallelogram is a polygon.

TEACHER NOTES:
I Do for the checklist routine. Move slowly through the three checks. The word 'parallelogram' is new for Year 1; acknowledge it and move on - the routine is the lesson, not the word.

WATCH FOR:
Students racing to 'yes' before the checks. Stay on each tick.
"""

NOTES[41] = """\
SAY:
This time the checklist has a mistake.
One tick should not be there.
Watch me check each one carefully.
Ask: which one is wrong? Expected: depends on the shape shown - check the slide.

DO:
Read each ticked item out loud.
Test it against the actual shape on screen.
Mark the wrong tick.

TEACHER NOTES:
The existing note (identify the wrong tick) is preserved. The expected answer depends on the actual shape on the slide - confirm before teaching which feature is the error.

WATCH FOR:
Students copying the ticks without checking. Cue: look at the shape, not just the boxes.
"""

NOTES[42] = """\
SAY:
Together. Use the checklist for the kite.
Closed? Straight sides? Two-dimensional?
Show me thumbs up or down for each, then for: is it a polygon?
Ask: is a kite a polygon? Expected: yes.

DO:
Read each checklist line and cue thumbs.
Reveal each tick after the class has voted.
Conclude: kite is a polygon.

TEACHER NOTES:
Together version of the checklist routine. Move at the class's pace. 'Kite' is a useful real-life shape name for Year 1.

WATCH FOR:
Students who say 'no' because the kite looks 'wonky'. Re-anchor: closed and straight, that is enough.
"""

NOTES[43] = """\
SAY:
Together. Use the checklist for the cube.
Closed? Yes.
Straight sides? Yes.
Two-dimensional? Hmm. Show me thumbs up or down.
Ask: is a cube two-dimensional? Expected: no, it is solid.

DO:
Pick up a real cube or block from the room and hold it next to a paper square.
Reveal that the third box should not be ticked.
Conclude: cube is not a polygon.

CFU CHECKPOINT:
Technique: Thumbs Up / Down
Script:
Ask: is a cube two-dimensional? Expected: no.
Scan for: most thumbs down.
PROCEED:
80 percent or more correct. Move to the next slide.
PIVOT:
Most likely: students tick all three because the cube has straight edges.
Reteach: hold the cube. It is solid - I can hold it. A polygon must be flat.
Re-check: is a ball two-dimensional?

TEACHER NOTES:
Cube is the 'trap' - it has straight edges and looks polygon-like. The 2D test is what stops it. Hold a real cube if possible.

WATCH FOR:
Students who say a cube is 2D because the picture on screen is flat. Hold a real cube to break the picture trick.
"""

NOTES[44] = """\
SAY:
Last together one. The arrow.
Closed? Straight sides? Two-dimensional? Show me thumbs.
Ask: is an arrow shape a polygon? Expected: yes.

DO:
Trace the outline of the arrow with your finger to show it is closed.
Point to each straight edge.
Reveal: yes, it is a polygon.

TEACHER NOTES:
Some students will say no because an arrow does not 'look like a normal shape'. Reinforce that the rule is the rule - shape names do not change the rule.

WATCH FOR:
Students who count the arrow as 'not closed' because of the point. The line still meets up.
"""

NOTES[45] = """\
SAY:
Today we play a shape mystery game.
Each pod gets a bag with shape cards.
Pick a card. Do not show your partner. Describe the shape - sides, vertices, what it looks like - without saying the name.
Your partner guesses. Then swap.

DO:
Hand out one mystery bag per pod.
Model with one student first - pick a shape, describe its features, see if the class can guess.
Set a 5 minute play timer with at least 4 swaps per pod.

TEACHER NOTES:
Existing teacher note preserved: pre-prepare shape mystery bags before the lesson. Bags should be ready before the lesson - print Sheets 1 and 2 from the resource list and make 12 sets. Keep the modelled describe-don't-name routine clear.

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: pick a card and say only the number of sides; partner guesses from a smaller list.
EXTENDING PROMPT:
Task: describe a shape using two features only - guess in three or fewer guesses.

WATCH FOR:
Students who name the shape straight away. Re-cue: describe first, do not say the name.
"""

NOTES[46] = """\
SAY:
Look at the picture on the screen.
Find the flat shapes inside it.
Count how many you see, and tell me what kind.
Ask: how many of one shape can you spot? Expected: depends on the picture.

DO:
Point to one shape in the picture as a model.
Give 30 to 60 seconds for the class to scan.
Take a few responses and check together.

TEACHER NOTES:
Shape spotting in a picture. The numbers visible on the slide (2, 5, 6) appear to be answers for some categories - confirm against the visible image before sharing answers.

WATCH FOR:
Students naming colours instead of shapes. Re-cue to features.
"""

NOTES[47] = """\
SAY:
True or false: there are two squares in this picture.
Look carefully. Show me one finger for true, two fingers for false.
Ask: how do you know? Expected: I counted, or I cannot see two squares.

DO:
Cue fingers up after a short look.
Take one student to point and count squares aloud.

CFU CHECKPOINT:
Technique: Finger Vote (1 = true, 2 = false)
Script:
Ask: are there two squares in this picture? Expected: depends on the visible picture - confirm before teaching.
Scan for: clear finger vote from most students.
PROCEED:
80 percent or more agree on the correct answer. Move on.
PIVOT:
Most likely: students count rectangles as squares.
Reteach: a square has all four sides the same length. A rectangle has two long, two short.
Re-check: how many squares can you see now?

TEACHER NOTES:
Confirm the actual answer against the visible picture before teaching. The trap is usually rectangles being counted as squares.

WATCH FOR:
Students who count quickly without checking equal sides.
"""

NOTES[48] = """\
SAY:
Together. Draw the flat shapes you see in this picture in your workbook.
Count and write how many of each.
Ask: how many triangles can you see? Expected: depends on the picture - confirm against the slide.

DO:
Point to one shape as a model.
Set a short timer (3 to 5 minutes).
Circulate to enabler students first.

TEACHER NOTES:
The visible numbers on the slide (14, 3) appear to be answers - confirm before sharing. Use the modelled square-rectangle routine if needed.

WATCH FOR:
Students who only count one shape type. Cue: now check for circles.
"""

NOTES[49] = """\
SAY:
Same job. Look at this new picture.
Draw and count the flat shapes.

DO:
Brief modelling - one shape - then release.
Circulate.

TEACHER NOTES:
Visible numbers (10, 5) appear to be answer counts - confirm before sharing.

WATCH FOR:
Students who do not write the count beside the drawing.
"""

NOTES[50] = """\
SAY:
Last together one. Find and count flat shapes.

DO:
Set the timer.
Take a few responses at the end.

TEACHER NOTES:
Visible numbers (1, 5, 13) appear to be answer counts - confirm before sharing. After this slide, students are ready for the shape hunt.

WATCH FOR:
Students who get tired by the third picture. Keep pace brisk.
"""

NOTES[51] = """\
SOURCES:
Shape Hunt BLM (existing reference).

SAY:
Now we hunt for shapes outside.
Take your board, pencil and worksheet.
We look for circles, squares, rectangles and triangles.
When you find one, mark it on your sheet.

DO:
Distribute clipboards (or boards), pencils and worksheets at the door.
Model one example outside before students spread out.
Walk among students - prompt enablers to a specific feature you can see.

TEACHER NOTES:
Ensure the BLM is printed before the lesson. Plan a meeting point and time before leaving the room. Keep the hunt area small enough to scan all students.

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: focus only on circles and rectangles in one named area (the door, the sandpit edge).
EXTENDING PROMPT:
Task: also record rhombus, parallelogram or kite shapes if you can spot them on the school building.

WATCH FOR:
Students who tally the same window twice. Cue: mark each thing once.
"""

NOTES[52] = """\
TEACHER NOTES:
Pick one review move. The fishbowl works well after the shape hunt - watch one pod compare their tally sheets and ask: how did you decide? The existing note prompts you to notice which strategies students used and whether they used them efficiently.
"""

NOTES[53] = """\
TEACHER NOTES:
Lesson 3 divider. Move on.
"""

NOTES[54] = """\
SAY:
Daily review. Looking at numbers on the number line to 20.
Where is 10? Show me with a finger in the air, then together we point on the screen.
Ask: is 10 closer to 0 or to 20? Expected: in the middle, equal.

DO:
Point to 0 first, count up to 10 by ones.
Confirm the position is the middle.

TEACHER NOTES:
Existing teacher note preserved: ask students to locate various numbers on the number line. This is review, not new teaching. Use it to refresh number-line skills before today's quadrilateral focus.

WATCH FOR:
Students who jump to a position by guessing. Cue them to count from 0.
"""

NOTES[55] = """\
SAY:
Now where is 20? Point in the air, then together on screen.
Ask: is 20 at the start or the end? Expected: at the end.

DO:
Sweep your finger from 0 to 20.

TEACHER NOTES:
20 is the easy anchor at the end of this number line.

WATCH FOR:
Students who point to 12 or another number. Re-anchor at 0 and count.
"""

NOTES[56] = """\
SAY:
Where is 13? Point in the air.
Ask: between which two numbers is 13? Expected: between 12 and 14.

DO:
Count from 10 onwards: 10, 11, 12, 13.

TEACHER NOTES:
This builds the 'count from a known landmark' strategy. Use 10 as the anchor for teen numbers.

WATCH FOR:
Students who count from 0 every time. Once they have shown they can, prompt them to start at 10.
"""

NOTES[57] = """\
SAY:
Where is 4? Point in the air.
Ask: how do you know it is 4? Expected: I started at 0 and counted 1, 2, 3, 4.

DO:
Count slowly with the class.

TEACHER NOTES:
4 is close to the start. A useful contrast to 13 and 18.

WATCH FOR:
Students who confuse 4 with 14. Slow the count.
"""

NOTES[58] = """\
SAY:
Last one. Where is 18?
Ask: which numbers does 18 sit between? Expected: between 17 and 19.

DO:
Count from 10 forwards to 18.

TEACHER NOTES:
18 is close to 20 - useful for 'near 20' thinking.

WATCH FOR:
Students who pick 8 instead of 18. Slow the count from 10.
"""

NOTES[59] = """\
SAY:
Skip count by 5s. Forwards first.
Now backwards from 120.

DO:
Brisk pace, claps for each number.
Use a chart on display if available.

TEACHER NOTES:
Fluency only. About a minute.

WATCH FOR:
Students who count by 1s instead. Re-anchor at 5, 10, 15.
"""

NOTES[60] = """\
SAY:
Today we are learning about quadrilaterals.
A quadrilateral is a flat shape with 4 sides.
Read the learning intention with me.

DO:
Point to the LI.
Track each success criterion.
Hold up four fingers and say 'quadrilateral' twice.

TEACHER NOTES:
New big word for Year 1. The first SC (identify quadrilaterals) is the most achievable. Today's exit point is making and drawing quadrilaterals.

WATCH FOR:
Students who confuse quadrilateral with quadruped or other 'quad' words. That is fine - use the four-fingers cue.
"""

NOTES[61] = """\
SAY:
Quick polygon review. Is the circle a polygon?
Closed? Straight sides? Two-dimensional?
Show me thumbs up or down for each box.
Ask: is the circle a polygon? Expected: no.

DO:
Read each checklist line.
Reveal: no straight sides - so not a polygon.
Show the thumbs-down icon on screen as confirmation.

TEACHER NOTES:
Reactivates yesterday's polygon checklist. Circle is the cleanest 'no' because it has no straight sides.

WATCH FOR:
Students who tick all three. Re-anchor at 'straight sides' - the circle has none.
"""

NOTES[62] = """\
SAY:
Look at the four shapes A, B, C, D.
Show me on your fingers which ones are NOT polygons.

DO:
Give 20 seconds.
Cue fingers up together.
Take one student to explain: B is not a polygon because...

CFU CHECKPOINT:
Technique: Show Me Fingers (A=1, B=2, C=3, D=4)
Script:
Ask: which shapes are not polygons? Expected: depends on the visible shapes - confirm before teaching.
Scan for: agreement on the curved or open shapes.
PROCEED:
80 percent or more agree on the correct answer. Move to the keywords.
PIVOT:
Most likely: students miss the open shape because it looks like a normal letter.
Reteach: trace each shape with your finger - if it does not meet up, it is not closed.
Re-check: point to one shape and ask: closed or open?

TEACHER NOTES:
Confirm the visible shapes A, B, C, D before teaching to know which are the not-polygons.

WATCH FOR:
Students who pick all four. Cue them to apply the checklist.
"""

NOTES[63] = """\
TEACHER NOTES:
Repeat keyword orientation slide. Existing note about how to use keyword slides is preserved.
"""

NOTES[64] = """\
SOURCES:
VIC Mathematics glossary - 'An object with width and length is two-dimensional. A polygon is an example of a two-dimensional geometric object.'

SAY:
Quick revision. What does two-dimensional mean?
Ask: 2D means flat or solid? Expected: flat.

DO:
Hold up paper for flat and a block for solid as a contrast.

TEACHER NOTES:
Brief revision - this word was introduced in Lesson 1. Source line preserved above.

WATCH FOR:
Students who hesitate. Repeat the paper-vs-block contrast.
"""

NOTES[65] = """\
SOURCES:
Any enclosed shape made of 3 or more straight lines (preserved from existing notes).

SAY:
New keyword today: polygon.
A polygon is a closed shape with three or more straight sides.
Ask: is a triangle a polygon? Expected: yes.

DO:
Point to each example shape on the slide.
Trace one outline with your finger to show 'closed'.

TEACHER NOTES:
This formalises the polygon idea from Lesson 2.

WATCH FOR:
Students who think only 'fancy' shapes are polygons. A triangle is a polygon too.
"""

NOTES[66] = """\
SOURCES:
The boundary of a 2D shape (preserved from existing notes).

SAY:
Sides are the lines around the edge of a flat shape.
Ask: how many sides on a quadrilateral? Expected: 4.

DO:
Hold up four fingers. Say 'quad means four'.

TEACHER NOTES:
Revisit - already taught in Lesson 1. Add the 'quad means four' link today.

WATCH FOR:
Students who count the same side twice.
"""

NOTES[67] = """\
SOURCES:
The point where two lines meet (preserved from existing notes).

SAY:
Vertices are corners. Where two lines meet.
Ask: how many vertices on a quadrilateral? Expected: 4.

DO:
Pinch four corners on a square shape on the board.

TEACHER NOTES:
Revisit. Today, link vertices to quadrilaterals: 4 vertices.

WATCH FOR:
Students who say 'vertex' when they mean 'vertices' (plural). Gently restate.
"""

NOTES[68] = """\
TEACHER NOTES:
Divider. Move on.
"""

NOTES[69] = """\
SAY:
Look at this sports picture. What flat shapes can you see?
I can see rectangles in the field. I can see circles around the players.
Ask: why are there so many rectangles here? Expected: it is the way the field is marked.

DO:
Point to one rectangle and one circle on the screen.
Ask one student to point to another shape.
Compare to other sports if time (existing note suggestion).

TEACHER NOTES:
Existing note preserved: discuss the choice of shapes and compare them to other sports. I Do for shape spotting in a real picture.

WATCH FOR:
Students who name colours, players or objects instead of shapes. Re-cue to flat shapes.
"""

NOTES[70] = """\
SAY:
Big idea today. A quadrilateral is a polygon with 4 sides and 4 vertices.
Watch the screen. Three different quadrilaterals.
Ask: do they all have 4 sides? Expected: yes.

DO:
Point to each shape and count the sides aloud.
Pinch each vertex aloud: 1, 2, 3, 4.
Repeat the rule: 4 sides, 4 vertices.

TEACHER NOTES:
This is the quadrilateral definition slide. Note: the existing teacher note about mini whiteboards and number lines is from a different lesson and does not apply here - ignore it. Keep the rule visible on the board for the rest of the lesson.

WATCH FOR:
Students who think only squares are quadrilaterals. The slide shows three different ones for a reason.
"""

NOTES[71] = """\
SAY:
Two shapes on screen. Both labelled with 4 sides and 4 vertices.
Ask: are they both quadrilaterals? Expected: yes.

DO:
Point to each shape.
Confirm: both have 4 sides, both have 4 vertices, both are quadrilaterals.

CFU CHECKPOINT:
Technique: Show Me Boards (write Y or N for each)
Script:
Ask: is the parallelogram a quadrilateral? Expected: yes.
Scan for: most boards showing Y for both shapes.
PROCEED:
80 percent or more correct. Move on.
PIVOT:
Most likely: students say only the parallelogram counts because it looks 'neat'.
Reteach: count the irregular polygon's sides aloud. 1, 2, 3, 4. It is a quadrilateral.
Re-check: draw a wonky 4-sided shape on the board. Quadrilateral?

TEACHER NOTES:
This slide breaks the 'quadrilateral must look like a square' idea. Praise students who include the irregular shape.

WATCH FOR:
Students who only choose the parallelogram. Re-anchor at '4 sides, 4 vertices = quadrilateral'.
"""

NOTES[72] = """\
SAY:
Watch how I make a quadrilateral with popsticks.
Four sticks. I join them at the corners.
Ask: how many sticks do I need? Expected: 4.

DO:
Take 4 popsticks. Join with blu-tack or a pipe cleaner at each corner.
Hold up the shape - count sides and vertices aloud.
Show that you can change the angles to make a different quadrilateral.

TEACHER NOTES:
Prepare popsticks and blu-tack or pipe cleaners before the lesson. Existing teacher note preserved: use blu-tack or pipe cleaners to hold corners together. Make at least two different quadrilaterals as the model.

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: pre-stick four popsticks into a square outline; student rearranges only one corner.
EXTENDING PROMPT:
Task: make a quadrilateral that does not look like a square or rectangle.

WATCH FOR:
Students who stick only 3 corners and leave one open. Cue: closed shape - all four corners must join.
"""

NOTES[73] = """\
SAY:
Four shapes on screen. A, B, C, D.
Show me on your fingers which ones are quadrilaterals.

DO:
Give 20 seconds.
Cue fingers.
Take one student to explain.

CFU CHECKPOINT:
Technique: Show Me Fingers
Script:
Ask: which shapes are quadrilaterals? Expected: depends on visible shapes - confirm before teaching.
Scan for: students choosing 4-sided closed shapes only.
PROCEED:
80 percent or more correct. Move to the explore.
PIVOT:
Most likely: students choose triangles or pentagons because they have straight sides.
Reteach: count the sides on each shape aloud. Only 4 sides is a quadrilateral.
Re-check: point to a fresh shape on the board.

TEACHER NOTES:
Confirm the four visible shapes A, B, C, D before teaching to know which are quadrilaterals.

WATCH FOR:
Students who include circles. Re-anchor: closed AND straight AND 4 sides.
"""

NOTES[74] = """\
SAY:
Now your turn with popsticks.
Make as many different quadrilaterals as you can.
Each one must have 4 straight sides and 4 vertices.

DO:
Distribute popsticks and joiners (blu-tack or pipe cleaners).
Set a 5 to 8 minute timer.
Circulate to enabler students first.

TEACHER NOTES:
Existing instructions preserved: use popsticks to make as many different quadrilaterals as you can; draw the quadrilaterals you have made into your workbook under the heading Quadrilaterals (see the next slide).

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: make one square, then move only one corner to make a new quadrilateral.
EXTENDING PROMPT:
Task: make a quadrilateral with all four sides different lengths.

WATCH FOR:
Students who make the same shape twice. Cue: try one with longer top and shorter sides.
"""

NOTES[75] = """\
SAY:
Now draw the quadrilaterals you made into your workbook.
Write the heading Quadrilaterals.
Each drawing must have 4 sides and 4 vertices.

DO:
Show the workbook layout on the board: heading, then 4 to 6 drawings.
Circulate as students draw.

TEACHER NOTES:
Existing instructions preserved: use popsticks to make as many different quadrilaterals as you can; draw the quadrilaterals into your workbook under the heading Quadrilaterals. Drawing is harder than building - some students may need to keep their popstick model beside them as a reference.

WATCH FOR:
Students who draw 3- or 5-sided shapes. Re-cue to count.
"""

NOTES[76] = """\
TEACHER NOTES:
Optional review slides follow if time permits. Skip the next set if the popstick task ran long. Note: the existing teacher note refers to mini whiteboard strategies and does not apply to this divider - ignore it.
"""

NOTES[77] = """\
SAY:
I am going to sort shapes into two groups.
Quadrilaterals on this side. Not quadrilaterals on the other.
Watch me check each one. 4 sides? Yes. 4 vertices? Yes. So it goes here.
Ask: why does this shape go in 'not quadrilaterals'? Expected: it does not have 4 sides.

DO:
Point to each example in the Quadrilaterals column.
Point to each example in the Not quadrilaterals column.
Name the difference - 3 sides, 5 sides, curved.

TEACHER NOTES:
I Do for sorting. Slow modelling. The shapes shown are the source of the contrast - point to them.

WATCH FOR:
Students who race ahead to call answers. Cue: watch first.
"""

NOTES[78] = """\
SAY:
Together. Each shape goes into Quadrilaterals or Other 2D shape.
Show me on your fingers which group as I point to each shape.

DO:
Point to each shape one at a time.
Cue 1 finger for Quadrilateral, 2 fingers for Other.
Confirm together.

TEACHER NOTES:
Together version of the sort. Use a brisk pace - aim for 6 to 8 shapes in 2 minutes.

WATCH FOR:
Students who say 'yes' to every shape. Reset with a clearly non-quadrilateral example.
"""

NOTES[79] = """\
SAY:
Multiple choice. How many sides does a quadrilateral have?
A is 8. B is 0. C is 1. D is 4.
Show me on your fingers.

DO:
Give 15 seconds.
Cue fingers.
Reveal: D is 4.

CFU CHECKPOINT:
Technique: Multiple Choice Fingers (1=A, 2=B, 3=C, 4=D)
Script:
Ask: how many sides does a quadrilateral have? Expected: 4 (option D).
Scan for: most fingers showing 4.
PROCEED:
80 percent or more correct. Move on.
PIVOT:
Most likely: students pick C (1) thinking 'quadrilateral' means one shape.
Reteach: hold up 4 fingers. Quad means four. Count it on a square on the board.
Re-check: how many vertices does a quadrilateral have?

TEACHER NOTES:
Hinge check on the meaning of 'quad'. If most students get this, the unit's key word has landed.

WATCH FOR:
Students who hold up 4 fingers but say 'I do not know' - they have the answer, they just need confidence.
"""

NOTES[80] = """\
SAY:
Together. Sort each shape into Quadrilaterals or Other 2D shapes.
Show me on your fingers which group as I point.

DO:
Point to each shape.
Cue fingers.
Confirm together.

TEACHER NOTES:
Continued sorting practice. Different shape mix from slide 78.

WATCH FOR:
Students who lose focus on the seventh or eighth shape. Take a breath; re-anchor at '4 sides, 4 vertices'.
"""

NOTES[81] = """\
SAY:
Last together one. Same job.
Quadrilaterals on this side. Other 2D shapes on the other.

DO:
Brisk pace - aim for 1 to 2 minutes.

TEACHER NOTES:
Final sorting practice for today. After this, students should be ready to identify quadrilaterals independently in Lesson 4.

WATCH FOR:
Students who can sort but cannot say why. Ask one to explain a choice in their own words.
"""

NOTES[82] = """\
TEACHER NOTES:
Pick one review move. The fishbowl works well after the popstick task - watch one pod's quadrilateral set and ask: how do you know they are all quadrilaterals?
"""

NOTES[83] = """\
TEACHER NOTES:
Lesson 4 divider. Move on.
"""

NOTES[84] = """\
TEACHER NOTES:
Divider into the Daily Review sequence. The existing note suggests modelling each strategy on mini whiteboards - keep that in mind during slides 85 to 89.
"""

NOTES[85] = """\
SAY:
Daily review. Finding one more.
I have 45. I add one more block. Now I have 46.
Watch the picture - one extra block on the right.
Ask: 1 more than 45 is what? Expected: 46.

DO:
Point to the block picture for 45.
Point to the extra block in the new picture.
Read the sentence: 1 more than 45 is 46.

TEACHER NOTES:
This is the I Do for one more. Use the visible base-10 blocks to anchor the idea before moving to the chart.

WATCH FOR:
Students who count all blocks each time. Cue: just add one more.
"""

NOTES[86] = """\
SAY:
Same idea. Pulling it back.
The picture shows blocks for 23 and 24.
Ask: 1 more than 23 is what? Expected: 24.

DO:
Point to the picture.
Read the sentence aloud.

TEACHER NOTES:
This slide shows 23 and 24 with blocks - confirm against the visible image which is the start number and which is '1 more'.

WATCH FOR:
Students who reverse the sentence. Cue: we add one to the first number.
"""

NOTES[87] = """\
SAY:
Now on the hundred chart.
1 more than 66 is what?
Watch me - I find 66, I move one to the right.
Ask: where do I land? Expected: 67.

DO:
Point to 66 on the chart.
Move your finger one square right to 67.
Read the sentence aloud.

TEACHER NOTES:
I Do for the chart-based method. The 'one square to the right' move is the rule - make it visible.

WATCH FOR:
Students who jump down a row instead of right. Re-anchor at 'right means 1 more'.
"""

NOTES[88] = """\
SAY:
Together. 1 more than 34.
Show me on your board.
Ask: what is 1 more than 34? Expected: 35.

DO:
Give 15 seconds.
Cue boards.
Confirm on the chart together.

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
Ask: 1 more than 34 is what? Expected: 35.
Scan for: most boards showing 35.
PROCEED:
80 percent or more correct. Move on.
PIVOT:
Most likely: students write 33 (one less) or 44 (one row down).
Reteach: point to 34 on the chart, move one square right.
Re-check: 1 more than 41?

TEACHER NOTES:
Quick hinge check. If most students struggle, return to base-10 blocks for one more example.

WATCH FOR:
Students who confuse 'more' with 'less'. Use directional language: right is more, left is less.
"""

NOTES[89] = """\
SAY:
On the chart, some numbers are missing.
On your board, write the five missing numbers in order.
Ask: which row is 13 in? Expected: the second row.

DO:
Point to each gap on the chart.
Give 1 to 2 minutes for boards.
Reveal one at a time.

TEACHER NOTES:
This brings together '1 more' thinking applied to the chart. Students who can name the missing numbers are ready for today's main lesson.

WATCH FOR:
Students who fill in any 5 numbers without checking the gaps. Cue: look at where the empty boxes are.
"""

NOTES[90] = """\
SAY:
Quick count. Forwards then backwards to 120.

DO:
Brisk pace.

TEACHER NOTES:
Fluency only.

WATCH FOR:
Students who lose count past 100. Pause and restart at 100.
"""

NOTES[91] = """\
SAY:
Today we are learning to sort 2D shapes by their features.
Read the learning intention with me.
Ask: what is one feature you could use to sort shapes? Expected: number of sides, number of vertices.

DO:
Point to LI and each SC.
Take one or two student suggestions for sorting features.

TEACHER NOTES:
New focus today is sorting by features. Students already know names; today they group by what shapes have in common.

WATCH FOR:
Students who only think of colour as a feature. Re-anchor at sides and vertices.
"""

NOTES[92] = """\
TEACHER NOTES:
Repeat keyword orientation. Existing note preserved.
"""

NOTES[93] = """\
SOURCES:
VIC Mathematics glossary - 'An object with width and length is two-dimensional. A polygon is an example of a two-dimensional geometric object.'

SAY:
Quick revision. 2D means flat.

DO:
Hold up paper for flat.

TEACHER NOTES:
Brief revision. Source line preserved above.

WATCH FOR:
Students who are still uncertain. Hold up a real solid object as a contrast.
"""

NOTES[94] = """\
SOURCES:
Merriam-Webster - 'feature: the structure, form, or appearance.'

SAY:
New keyword: features.
Features are the parts or look of something.
A triangle has features: three sides, and on this slide it is orange.
Ask: what is one feature of a circle? Expected: it is round, or 1 curved side.

DO:
Point to 'triangle', then 'three sides', then 'orange' on the slide.
Take one student to name a feature of a square.

TEACHER NOTES:
'Features' is the lesson's key word. Today we will sort shapes by their features. Source line preserved above.

WATCH FOR:
Students who think colour is the only feature. Add 'sides' and 'vertices' to the list.
"""

NOTES[95] = """\
SOURCES:
The point where two lines meet (preserved from existing notes).

SAY:
Vertices are corners.
Ask: how many vertices on a square? Expected: 4.

DO:
Pinch four corners on a square shape.

TEACHER NOTES:
Quick revision.

WATCH FOR:
Students who still say 'vertex' for plural. Restate.
"""

NOTES[96] = """\
SOURCES:
The boundary of a 2D shape (preserved from existing notes).

SAY:
Sides are the lines around a flat shape.
Ask: how many sides on a triangle? Expected: 3.

DO:
Trace a triangle in the air.

TEACHER NOTES:
Quick revision.

WATCH FOR:
Students who count one side twice.
"""

NOTES[97] = """\
TEACHER NOTES:
Divider into the launch. Note: the existing teacher note about number lines and 100 charts is from a different lesson and does not apply here - ignore it. Move on to the next slides.
"""

NOTES[98] = """\
SAY:
Look at the shape. Tell me three things on your board.
Name. Sides. Vertices.
Ask: how many sides on a rectangle? Expected: 4.

DO:
Give 30 seconds for boards.
Reveal: rectangle, 4 sides, 4 vertices.

TEACHER NOTES:
Quick warm-up describing a known shape. Sets up the sort-by-features focus.

WATCH FOR:
Students who write 'square' - prompt them to look at the side lengths.
"""

NOTES[99] = """\
SAY:
Same job. Name, sides, vertices for the square.

DO:
Give 30 seconds.
Reveal: square, 4 sides, 4 vertices.

TEACHER NOTES:
Quick warm-up. Note: square and rectangle have the same side and vertex count - the different feature is the side lengths.

WATCH FOR:
Students who write 'rectangle' - all sides equal here, so it is a square.
"""

NOTES[100] = """\
SAY:
Last one. Triangle.

DO:
Give 30 seconds.
Reveal: triangle, 3 sides, 3 vertices.

TEACHER NOTES:
Sets up the contrast (3 vs 4 sides) for the explore that follows.

WATCH FOR:
Students who confuse 3 sides with 4. Slow the count.
"""

NOTES[101] = """\
SAY:
I want you to think of features as a checklist.
This shape has 4 vertices. This shape has 4 straight sides. So this shape is a quadrilateral.
Ask: what features make a shape a quadrilateral? Expected: 4 straight sides and 4 vertices.

DO:
Point to '4 vertices' then '4 straight sides' then 'quadrilateral'.
Read the connector sentence: it is a quadrilateral.

TEACHER NOTES:
I Do for using features to name shapes. Two features lead to one name.

WATCH FOR:
Students who say only '4 sides' without saying vertices. Both features are needed today.
"""

NOTES[102] = """\
TEACHER NOTES:
Explore divider. Move on.
"""

NOTES[103] = """\
SAY:
Look at the group of shapes.
What is the same about all of them?
Ask: how many sides does each one have? Expected: 3.

DO:
Point to each shape and count sides aloud.
Reveal: 3 sides.

TEACHER NOTES:
First 'find the common feature' task. The answer is on screen - keep boards down so students focus on the visual reasoning.

WATCH FOR:
Students who name colour or size as the common feature. Cue: look at the lines, not the colour.
"""

NOTES[104] = """\
SAY:
New group. What is the same?
Ask: how many vertices does each one have? Expected: 4.

DO:
Point to each corner.
Reveal: 4 vertices.

TEACHER NOTES:
Same routine, vertex focus. Note this uses vertices, not sides.

WATCH FOR:
Students who go to 'sides' because that was the answer last time. Re-cue them to vertices.
"""

NOTES[105] = """\
SAY:
Look at the group. Does each shape have 4 sides?
Show me thumbs up for yes, thumbs down for no.

DO:
Cue thumbs.
Reveal: No.
Ask one student to point to the shape that does not have 4 sides.

CFU CHECKPOINT:
Technique: Thumbs Up / Down
Script:
Ask: does each shape have 4 sides? Expected: no.
Scan for: thumbs down from most students.
PROCEED:
80 percent or more correct. Move on.
PIVOT:
Most likely: students see one rectangle and say yes for all.
Reteach: count sides on each shape one at a time.
Re-check: point to one shape and ask: 4 sides?

TEACHER NOTES:
This slide shifts the question - it is a 'do all of them' check, not 'is this one'. Slow the count.

WATCH FOR:
Students who skip checking each shape. Cue: check every one.
"""

NOTES[106] = """\
SAY:
Now we sort. Some shapes have 3 sides. Some have 4 sides.
Ask: which group does this shape go in? Expected: depends on the shape.

DO:
Point to each shape and ask: 3 or 4?
Place each in the correct column.

TEACHER NOTES:
I Do for sorting by feature. Use the visible columns on the slide as the headers.

WATCH FOR:
Students who count diagonally. Slow the count.
"""

NOTES[107] = """\
SAY:
Together. Each shape into 3 sides or more than 3 sides.
Show me on your fingers - 1 for 3 sides, 2 for more than 3.

DO:
Point to each shape.
Cue fingers.
Place each in the correct column.

TEACHER NOTES:
Together sort with a different category split (3 vs more). The numbers visible on screen appear to be answer counts - confirm against the visible shapes before sharing.

WATCH FOR:
Students who default to '3 sides' for every shape. Re-cue to count first.
"""

NOTES[108] = """\
SAY:
Same routine, vertex focus.
1 for 4 vertices, 2 for more than 4.

DO:
Point to each shape.
Cue fingers.
Place each in the correct column.

TEACHER NOTES:
Together sort by vertex count. The numbers visible on screen appear to be answer counts - confirm against the visible shapes before sharing.

WATCH FOR:
Students who muddle vertices and sides. Pinch each corner as you count.
"""

NOTES[109] = """\
SOURCES:
2D Shape Matching BLM. Existing link preserved: https://drive.google.com/file/d/1qjxfJI6bYWlUKgbrA_ch9_MMzcOp3fo9/view?usp=drive_link
Page 2 for all to cut and paste. Pages 3 and 4 for early finishers (can be printed back to back).

SAY:
Now your turn on paper.
Cut and paste the shapes onto the matching feature on Page 2.
Early finishers - try Pages 3 and 4.

DO:
Distribute Page 2 of the BLM to all students.
Hand Pages 3 and 4 to early finishers.
Set a 10 to 15 minute timer.
Circulate to enabler students first.

TEACHER NOTES:
Existing teacher notes (Drive link and page-by-page instructions) preserved in SOURCES above. Print and prepare these BLM pages before the lesson.

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: only sort the first row of shapes; refer back to the sorting slides if stuck.
EXTENDING PROMPT:
Task: complete pages 3 and 4 and then write one sentence: 'I sorted these shapes by...'

WATCH FOR:
Students who guess. Cue: count the sides first.
Students who finish all pages quickly. Move them to the extension sentence.
"""

NOTES[110] = """\
TEACHER NOTES:
Pick one review move. After the BLM, a quick gallery walk works well - students see how others sorted.
"""

NOTES[111] = """\
TEACHER NOTES:
Lesson 5 divider. Today is consolidation through bingo. Move on.
"""

NOTES[112] = """\
SOURCES:
www.dailymathsreview.au
www.ochre.org.au

TEACHER NOTES:
Generic Daily Review divider. The existing note links to dailymathsreview.au and ochre.org.au - use those for today's review activity. The slides that follow show 'one less' practice.
"""

NOTES[113] = """\
SAY:
Today: 1 less than a number.
I have 37. I take one block away. Now I have 36.
Watch the picture - one fewer block.
Ask: 1 less than 37 is what? Expected: 36.

DO:
Point to the 37 block picture.
Cover or remove the last block to show 36.
Read the sentence aloud.

TEACHER NOTES:
Mirror of yesterday's '1 more' routine. The blocks make 'less' visible.

WATCH FOR:
Students who add one instead of taking one. Re-anchor at 'less means smaller'.
"""

NOTES[114] = """\
SAY:
Same idea. I have 22. 1 less is what?
Ask: 1 less than 22 is what? Expected: 21.

DO:
Point to the picture.
Read the sentence aloud.

TEACHER NOTES:
Quick recall. Use the picture to confirm.

WATCH FOR:
Students who say 23. Cue: we are taking one away, not adding.
"""

NOTES[115] = """\
SAY:
Now on the chart. 1 less than 66.
Watch - I find 66, I move one square left.
Ask: 1 less than 66 is what? Expected: 65.

DO:
Point to 66.
Move finger one square left to 65.

TEACHER NOTES:
I Do for the chart move - left means less.

WATCH FOR:
Students who go right (more) instead. Re-anchor: left is less.
"""

NOTES[116] = """\
SAY:
Together. 1 less than 34.
Show me on your board.

DO:
Give 15 seconds.
Cue boards.
Reveal: 33.

CFU CHECKPOINT:
Technique: Show Me Boards
Script:
Ask: 1 less than 34 is what? Expected: 33.
Scan for: most boards showing 33.
PROCEED:
80 percent or more correct. Move on.
PIVOT:
Most likely: students write 35 (one more) or 24 (one row up).
Reteach: point to 34 on chart, move one left to 33.
Re-check: 1 less than 50?

TEACHER NOTES:
Quick hinge check. Use the chart on screen.

WATCH FOR:
Students who confuse less with more. Re-cue: left is less.
"""

NOTES[117] = """\
SAY:
On your own. Find 1 less than each missing number on the chart.
Use your board.

DO:
Point to the missing squares on the chart.
Set a 3 to 5 minute timer.
Circulate.

TEACHER NOTES:
The existing internal note 'Need a template' suggests this slide may not have a finished student template. Either project the chart and have students write 1 less for each missing number on their boards, or provide a printed copy.

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: only find 1 less for the first two missing numbers.
EXTENDING PROMPT:
Task: also find 1 more and 10 more for each missing number.

WATCH FOR:
Students who read the missing number wrong. Re-anchor at the chart.
"""

NOTES[118] = """\
SAY:
Skip count by 2s. Forwards first, then backwards from 120.

DO:
Brisk pace, claps.

TEACHER NOTES:
Fluency only. About a minute.

WATCH FOR:
Students who say every number. Slow the pace and tap pairs.
"""

NOTES[119] = """\
SAY:
Today is our last 2D shape lesson for the week.
Read the learning intention with me.
Ask: what is one shape you can name now that you could not name on Monday? Expected: any shape from this week.

DO:
Point to the LI.
Track each success criterion.
Take one or two student responses.

TEACHER NOTES:
Same LI/SC as the week-level slide (slide 4) - the games today should give students a chance to show all three. Use student responses to gauge confidence before bingo.

WATCH FOR:
Students who hesitate. Note them for the bingo group choice.
"""

NOTES[120] = """\
SAY:
We are going to hunt one more time.
Find something in the room with 4 edges or sides.
Now find something with 3 corners or vertices.
Now something with 1 round edge.
Now something that is a different 2D shape - tell me its name.

DO:
Read each prompt one at a time.
Give 20 to 30 seconds between prompts.
Take one student to point to and name each shape they found.

TEACHER NOTES:
Quick whole-class shape hunt to warm up before bingo. Use the room itself - door, window, clock, signs.

WATCH FOR:
Students who pick the same object every time. Cue: try a different object now.
"""

NOTES[121] = """\
SAY:
Today we play 2D shape bingo.
In your group, one person is the caller.
The caller calls out shape names.
Check your board. Place a counter on a shape with that name.
Cover all your squares to win.

DO:
Form groups of 4 or 5.
Hand out a bingo board and counters per group.
Choose or assign the first caller.
Walk between groups to support callers.

TEACHER NOTES:
Bingo cards must be pre-prepared (or borrowed - see overview slide). The existing overview suggests playing as a whole class or in 3 to 4 groups. Set a clear rotation rule for the caller (winner becomes next caller).

ENABLING & EXTENDING:
ENABLING PROMPT:
Task: pair up with a partner who can read the shape names; they help find the shape on the board.
EXTENDING PROMPT:
Task: as caller, give a feature clue ('a shape with 3 sides') instead of the name.

WATCH FOR:
Students who cannot read the shape names. Pair them with a stronger reader.
Students who race ahead and call bingo without checking. Cue a quick check before declaring the winner.
"""

NOTES[122] = """\
TEACHER NOTES:
Final review for the week. A walk-around or fishbowl works well after the bingo - which shapes did students find easy or hard to spot?
"""


# ============================================================================
# MAIN
# ============================================================================

def main():
    if not SRC.exists():
        raise SystemExit(f"Source not found: {SRC}")
    pres = Presentation(str(SRC))
    n_slides = len(pres.slides)
    print(f"Source: {SRC}")
    print(f"Slides: {n_slides}")
    missing = [i for i in range(1, n_slides + 1) if i not in NOTES]
    if missing:
        raise SystemExit(f"Missing notes for slides: {missing}")
    extra = [k for k in NOTES if k > n_slides]
    if extra:
        print(f"Warning: notes defined for non-existent slides: {extra}")
    for idx, slide in enumerate(pres.slides, start=1):
        text = NOTES[idx]
        set_notes(slide, text)
    pres.save(str(DST))
    print(f"Saved: {DST}")


if __name__ == "__main__":
    main()
