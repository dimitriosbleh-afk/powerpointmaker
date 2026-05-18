#!/usr/bin/env python3
"""Combine Grade 1 Numeracy 'Numbers to 120' unit into one PPTX and one
Resources folder.

- Merges G1N_Lesson1 ... G1N_Lesson10 into a single deck (in lesson order).
- Copies every PDF from each per-lesson resources-sessionN/ folder into one
  flat output/G1N_Unit_Numbers_to_120/Resources/ folder.
"""

import copy
import os
import shutil
from pptx import Presentation

LESSONS = [
    "G1N_Lesson1_Reading_Writing_Numerals",
    "G1N_Lesson2_Ordering_On_Number_Line",
    "G1N_Lesson3_Hundreds_Chart",
    "G1N_Lesson4_Extending_To_120",
    "G1N_Lesson5_Pairs_That_Make_10",
    "G1N_Lesson6_Tens_And_Ones",
    "G1N_Lesson7_Different_Ways",
    "G1N_Lesson8_Think_Board_Addition",
    "G1N_Lesson9_Think_Board_Subtraction",
    "G1N_Lesson10_Money_Stories",
]

INPUT_PPTX = [
    f"output/{name}/{name}.pptx" for name in LESSONS
]

OUTPUT_DIR = "output/G1N_Unit_Numbers_to_120"
OUTPUT_PPTX = f"{OUTPUT_DIR}/G1N_Numbers_to_120_Full_Unit.pptx"
RESOURCES_DIR = f"{OUTPUT_DIR}/Resources"


def copy_slide(src_slide, dst_prs):
    """Copy one slide (shapes, background, images, notes) into dst_prs."""
    dst_layout = dst_prs.slide_layouts[6] if len(dst_prs.slide_layouts) > 6 else dst_prs.slide_layouts[0]
    dst_slide = dst_prs.slides.add_slide(dst_layout)

    # Replace shape tree
    src_spTree = src_slide.shapes._spTree
    dst_spTree = dst_slide.shapes._spTree
    dst_spTree.getparent().replace(dst_spTree, copy.deepcopy(src_spTree))

    # Copy background
    src_bg = src_slide.background._element
    dst_bg = dst_slide.background._element
    dst_bg.getparent().replace(dst_bg, copy.deepcopy(src_bg))

    # Remap image rIds
    rid_map = {}
    for rId, rel in list(src_slide.part.rels.items()):
        if "image" in rel.reltype:
            new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rId] = new_rId

    if rid_map:
        for old_rId, new_rId in rid_map.items():
            if old_rId != new_rId:
                for elem in dst_slide.shapes._spTree.iter():
                    for attr_name in (
                        "embed",
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                    ):
                        val = elem.get(attr_name)
                        if val == old_rId:
                            elem.set(attr_name, new_rId)

    # Copy notes
    if src_slide.has_notes_slide:
        notes_text = src_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            dst_slide.notes_slide.notes_text_frame.text = notes_text

    return dst_slide


def merge_decks():
    """Merge all lesson decks into one PPTX in order."""
    base_prs = Presentation(INPUT_PPTX[0])

    for fpath in INPUT_PPTX[1:]:
        src_prs = Presentation(fpath)
        for slide in src_prs.slides:
            copy_slide(slide, base_prs)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    base_prs.save(OUTPUT_PPTX)
    return len(base_prs.slides)


def copy_resources():
    """Copy every PDF from each lesson's session folder into one Resources folder."""
    os.makedirs(RESOURCES_DIR, exist_ok=True)
    copied = []
    for idx, name in enumerate(LESSONS, start=1):
        src_folder = f"output/{name}/resources-session{idx}"
        if not os.path.isdir(src_folder):
            print(f"  WARN: missing {src_folder}")
            continue
        for fname in sorted(os.listdir(src_folder)):
            if not fname.lower().endswith(".pdf"):
                continue
            # Drop the leading "Session N " prefix so files read as "Lesson N ..."
            clean = fname
            prefix = f"Session {idx} "
            if clean.startswith(prefix):
                clean = clean[len(prefix):]
            src = os.path.join(src_folder, fname)
            dst = os.path.join(RESOURCES_DIR, clean)
            shutil.copyfile(src, dst)
            copied.append(clean)
    return copied


def main():
    print("Merging decks...")
    total_slides = merge_decks()
    print(f"  Combined PPTX written: {OUTPUT_PPTX}")
    print(f"  Total slides: {total_slides}")

    print("\nCopying resources...")
    copied = copy_resources()
    print(f"  Resources folder: {RESOURCES_DIR}")
    print(f"  Total PDFs copied: {len(copied)}")

    print("\nFinal structure:")
    print(f"  {OUTPUT_DIR}/")
    print(f"    G1N_Numbers_to_120_Full_Unit.pptx  ({total_slides} slides)")
    print(f"    Resources/  ({len(copied)} PDFs)")


if __name__ == "__main__":
    main()
