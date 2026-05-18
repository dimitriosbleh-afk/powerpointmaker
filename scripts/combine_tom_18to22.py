#!/usr/bin/env python3
"""Combine Tom Appleby sessions 18-22 (+ recap) into a single PPTX and
gather every companion PDF into a single Resources/ folder.

Order:
    18 -> 19 -> 20 -> 21 -> 18-21 recap -> 22
"""

import copy
import shutil
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent

# Source decks in teaching order. Each is a (label, pptx_path, resources_dir) tuple.
SOURCES = [
    ("Session 18", ROOT / "output/Tom_Lesson18_Christmas_Hanging_Sentence_Expansion/Tom_Lesson18.pptx",
     ROOT / "output/Tom_Lesson18_Christmas_Hanging_Sentence_Expansion/resources-session18"),
    ("Session 19", ROOT / "output/Tom_Lesson19_Disease_and_Settlement_SPO/Tom_Lesson19.pptx",
     ROOT / "output/Tom_Lesson19_Disease_and_Settlement_SPO/resources-session19"),
    ("Session 20", ROOT / "output/Tom_Lesson20_Write_Body_Paragraph_Sydney_Cove/Tom_Lesson20.pptx",
     ROOT / "output/Tom_Lesson20_Write_Body_Paragraph_Sydney_Cove/resources-session20"),
    ("Session 21", ROOT / "output/Tom_Lesson21_Eora_Woman_Subordinating_Conjunctions/Tom_Lesson21.pptx",
     ROOT / "output/Tom_Lesson21_Eora_Woman_Subordinating_Conjunctions/resources-session21"),
    ("Sessions 18-21 Recap", ROOT / "output/Tom_Review_Lessons18to21/Tom_Review_Lessons18to21.pptx",
     ROOT / "output/Tom_Review_Lessons18to21/resources-session22"),
    ("Session 22", ROOT / "output/Tom_Lesson22_Food_Shortage_KPAS/Tom_Lesson22.pptx",
     ROOT / "output/Tom_Lesson22_Food_Shortage_KPAS/resources-session22"),
]

OUT_DIR = ROOT / "output/Tom_Sessions18to22_Combined"
OUT_PPTX = OUT_DIR / "Tom_Sessions18to22.pptx"
OUT_RESOURCES = OUT_DIR / "Resources"


def copy_slide(src_prs, src_slide, dst_prs):
    """Copy a slide from src into dst, preserving images and notes."""
    dst_layout = dst_prs.slide_layouts[0]
    dst_slide = dst_prs.slides.add_slide(dst_layout)

    # Replace shape tree
    src_spTree = src_slide.shapes._spTree
    dst_spTree = dst_slide.shapes._spTree
    dst_spTree.getparent().replace(dst_spTree, copy.deepcopy(src_spTree))

    # Copy background
    src_bg = src_slide.background._element
    dst_bg = dst_slide.background._element
    dst_bg.getparent().replace(dst_bg, copy.deepcopy(src_bg))

    # Copy image parts and remap rIds where they differ
    rid_map = {}
    for rId, rel in list(src_slide.part.rels.items()):
        if "image" in rel.reltype:
            new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rId] = new_rId

    if rid_map:
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        for old_rId, new_rId in rid_map.items():
            if old_rId != new_rId:
                for elem in dst_slide.shapes._spTree.iter():
                    for attr_name in ("embed", f"{{{r_ns}}}embed"):
                        val = elem.get(attr_name)
                        if val == old_rId:
                            elem.set(attr_name, new_rId)

    # Copy notes
    if src_slide.has_notes_slide:
        notes_text = src_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            dst_slide.notes_slide.notes_text_frame.text = notes_text

    return dst_slide


def gather_resources():
    """Copy every PDF from each source resources folder into a single
    Resources/ folder. Filenames already include 'Session N' prefixes so
    collisions are unlikely; if any occur, fall back to a numbered suffix."""
    if OUT_RESOURCES.exists():
        shutil.rmtree(OUT_RESOURCES)
    OUT_RESOURCES.mkdir(parents=True)

    copied = []
    for label, _pptx, res_dir in SOURCES:
        if not res_dir.exists():
            print(f"  WARN: {label} resources folder missing: {res_dir}")
            continue
        for pdf in sorted(res_dir.glob("*.pdf")):
            target = OUT_RESOURCES / pdf.name
            if target.exists():
                stem = target.stem
                suffix = target.suffix
                n = 2
                while (OUT_RESOURCES / f"{stem} ({n}){suffix}").exists():
                    n += 1
                target = OUT_RESOURCES / f"{stem} ({n}){suffix}"
            shutil.copy2(pdf, target)
            copied.append((label, target.name))
    return copied


def merge_decks():
    label0, pptx0, _ = SOURCES[0]
    base_prs = Presentation(str(pptx0))
    counts = {label0: len(base_prs.slides)}

    for label, pptx, _ in SOURCES[1:]:
        src_prs = Presentation(str(pptx))
        counts[label] = len(src_prs.slides)
        for slide in src_prs.slides:
            copy_slide(src_prs, slide, base_prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    base_prs.save(str(OUT_PPTX))
    return counts, len(base_prs.slides)


def main():
    print("Combining Tom Appleby sessions 18-22...\n")
    counts, total = merge_decks()
    print(f"PPTX -> {OUT_PPTX}")
    for label, n in counts.items():
        print(f"  {label}: {n} slides")
    print(f"  Combined total: {total} slides\n")

    copied = gather_resources()
    print(f"Resources -> {OUT_RESOURCES}")
    for label, fname in copied:
        print(f"  [{label}]  {fname}")
    print(f"  Total resources: {len(copied)}\n")

    print(f"DONE")
    print(f"Folder: {OUT_DIR}")


if __name__ == "__main__":
    main()
