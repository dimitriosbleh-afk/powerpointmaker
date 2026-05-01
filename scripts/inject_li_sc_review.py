"""Inject LI/SC slide (where missing) and Review-and-Reflect slide (always) into
all 50 decks in 'Challenge Slides'. Save to a parallel 'Challenge Slides Updated'
tree so originals stay untouched.

Rules:
- If a deck already has a specific LI + 3 'I can' SC slide, leave it alone.
- If a deck only has the generic OCHRE 'Learning objectives' slide, insert a
  new 'Learning Intention & Success Criteria' slide IMMEDIATELY AFTER the
  generic slide. The generic slide is preserved.
- Always append a 'Review and Reflect' slide as the final slide. It echoes the
  same three SCs and a thumbs-up self-assessment routine.
- Pull the SC text from the existing LI/SC slide where present; otherwise use
  the lesson-specific content drafted in LESSON_LI_SC below.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
from lxml import etree
import re
import shutil

ROOT = Path(r"C:\Users\09560329\Downloads\Challenge Slides")
OUT = Path(r"C:\Users\09560329\Downloads\Challenge Slides Updated")

# ---------------------------------------------------------------------------
# Lesson-specific LI/SC for the 12 OCHRE-only decks. Drafted from each deck's
# actual teaching content (pages read, sentence-writing focus, paragraph type).
# ---------------------------------------------------------------------------
LESSON_LI_SC = {
    12: {
        "li": "We are learning to read and understand Storm Boy and to use a single paragraph outline (SPO) to summarise the main ideas of a text.",
        "sc": [
            "I can read and discuss pages 49-54 of Storm Boy and answer questions about the text",
            "I can identify the topic sentence, supporting details and concluding sentence in a paragraph",
            "I can complete a single paragraph outline that summarises a section of the text",
        ],
    },
    13: {
        "li": "We are learning to read and understand Storm Boy and to take notes that capture the main ideas of a text.",
        "sc": [
            "I can read and discuss pages 54-57 of Storm Boy and answer questions about the text",
            "I can identify the keywords and phrases that carry the main idea in a sentence",
            "I can convert a sentence into keywords, phrases, abbreviations or symbols to take useful notes",
        ],
    },
    16: {
        "li": "We are learning to read and understand Storm Boy and to combine two simple sentences into a compound sentence using a FANBOYS conjunction.",
        "sc": [
            "I can read and discuss pages 57-62 of Storm Boy and answer questions about the text",
            "I can identify a compound sentence and the FANBOYS conjunction that joins the two ideas",
            "I can combine two simple sentences into a compound sentence using a comma and a FANBOYS conjunction",
        ],
    },
    17: {
        "li": "We are learning to read and understand Storm Boy and to use a single paragraph outline (SPO) to summarise the main ideas of a text.",
        "sc": [
            "I can read and discuss pages 62-66 of Storm Boy and answer questions about the text",
            "I can identify the topic sentence, supporting details and concluding sentence of a paragraph",
            "I can complete a single paragraph outline that summarises a section of the text",
        ],
    },
    18: {
        "li": "We are learning to read and understand Storm Boy and to take notes that capture the main ideas of a text.",
        "sc": [
            "I can read and discuss pages 67-69 of Storm Boy and answer questions about the text",
            "I can identify the keywords and phrases that carry the main idea in a sentence",
            "I can convert a sentence into keywords, phrases, abbreviations or symbols to take useful notes",
        ],
    },
    19: {
        "li": "We are learning to plan a body paragraph for our information report on the Coorong, focused on flora and fauna.",
        "sc": [
            "I can identify the topic sentence, supporting details and concluding sentence in a body paragraph",
            "I can choose facts from the Coorong text that support the topic of flora and fauna",
            "I can complete a single paragraph outline that plans body paragraph three",
        ],
    },
    20: {
        "li": "We are learning to write a body paragraph for our information report by turning our plan into a full paragraph.",
        "sc": [
            "I can write a topic sentence that expresses the main idea of body paragraph three",
            "I can write supporting detail sentences using facts from my single paragraph outline",
            "I can write a concluding sentence that closes the paragraph without repeating the topic sentence",
        ],
    },
    21: {
        "li": "We are learning to read and understand Storm Boy and to punctuate direct speech accurately.",
        "sc": [
            "I can read and discuss pages 69-72 of Storm Boy and answer questions about the text",
            "I can identify direct speech in a text by the speech marks and the speaker's words",
            "I can punctuate a sentence of direct speech using speech marks, a comma and a capital letter",
        ],
    },
    22: {
        "li": "We are learning to read and understand Storm Boy and to use a single paragraph outline (SPO) to summarise the main ideas of a text.",
        "sc": [
            "I can read and discuss pages 72-73 of Storm Boy and answer questions about the text",
            "I can identify the topic sentence, supporting details and concluding sentence of a paragraph",
            "I can complete a single paragraph outline that summarises a section of the text",
        ],
    },
    23: {
        "li": "We are learning to read and understand Storm Boy and to take notes that capture the main ideas of a text.",
        "sc": [
            "I can read and discuss pages 73-76 of Storm Boy and answer questions about the text",
            "I can identify the keywords and phrases that carry the main idea in a sentence",
            "I can convert a sentence into keywords, phrases, abbreviations or symbols to take useful notes",
        ],
    },
    24: {
        "li": "We are learning to plan a concluding paragraph for our information report on the Coorong using the TSG formula.",
        "sc": [
            "I can explain that a concluding paragraph summarises the report and does not introduce new information",
            "I can identify the thesis statement, summary of body paragraphs and general statement that make up a TSG conclusion",
            "I can complete a TSG plan for the concluding paragraph of our information report",
        ],
    },
    25: {
        "li": "We are learning to write a concluding paragraph for our information report by turning our TSG plan into a full paragraph.",
        "sc": [
            "I can write a thesis statement that links back to the introduction",
            "I can write a sentence that summarises the main ideas covered in the body paragraphs",
            "I can write a general statement that closes the report without introducing new information",
        ],
    },
}

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def lesson_no(name: str):
    m = re.match(r"\s*(\d+)\.", name)
    return int(m.group(1)) if m else None


def slide_text(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                if t.strip():
                    out.append(t)
    return "\n".join(out)


def find_li_sc_index(prs):
    """Return 1-based slide index of an existing specific LI/SC slide, or None."""
    for i, s in enumerate(prs.slides, 1):
        t = slide_text(s).lower()
        i_can = [ln for ln in t.splitlines() if ln.strip().startswith("i can ")]
        has_li = ("learning intention" in t) or ("we are learning to" in t)
        if has_li and len(i_can) >= 2:
            return i
    return None


def find_ochre_objectives_index(prs):
    """Return 1-based slide index of the OCHRE generic 'Learning objectives' slide, or None."""
    for i, s in enumerate(prs.slides, 1):
        t = slide_text(s).lower()
        i_can = [ln for ln in t.splitlines() if ln.strip().startswith("i can ")]
        if "learning objective" in t and "we are learning to" in t and len(i_can) == 0:
            return i
    return None


def extract_existing_li_sc(prs):
    """Pull LI text and 3 SC bullets from the deck's existing LI/SC slide.
    Returns (li, [sc1, sc2, sc3]) or None."""
    idx = find_li_sc_index(prs)
    if idx is None:
        return None
    slide = list(prs.slides)[idx - 1]
    text = slide_text(slide)
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    li = None
    sc = []
    # LI: line starting with 'We are learning to'
    for ln in lines:
        low = ln.lower()
        if low.startswith("we are learning to") and li is None:
            li = ln
        elif low.startswith("i can "):
            sc.append(ln)
    if li and len(sc) >= 3:
        return li, sc[:3]
    return None


def get_layout(prs, *names):
    """Return the first slide layout whose name matches any in `names`."""
    for layout in prs.slide_layouts:
        if layout.name in names:
            return layout
    # Fall back to first layout
    return prs.slide_layouts[0]


def find_placeholder(slide, *idxs):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in idxs:
            return ph
    return None


def remove_placeholder(slide, idx):
    ph = find_placeholder(slide, idx)
    if ph is not None:
        sp = ph._element
        sp.getparent().remove(sp)


def set_text_runs(tf, paragraphs):
    """Replace the text frame contents with the given list of paragraph dicts.
    paragraphs: list of dicts with keys: text, bold (default False), size (Pt int),
    spaceBefore (Pt, optional). Default font is Arial-like inherited from layout."""
    # Clear existing paragraphs
    tf.clear()
    # The first paragraph is created by clear()
    first = tf.paragraphs[0]
    for i, p in enumerate(paragraphs):
        para = first if i == 0 else tf.add_paragraph()
        # Wipe pre-existing runs (for the first paragraph after clear, there is one empty run by default — replace its text)
        # Easier: nuke runs and add fresh.
        for r in list(para.runs):
            r._r.getparent().remove(r._r)
        run = para.add_run()
        run.text = p["text"]
        if p.get("bold"):
            run.font.bold = True
        if p.get("size"):
            run.font.size = Pt(p["size"])
        if p.get("color"):
            run.font.color.rgb = p["color"]
        if p.get("align"):
            para.alignment = p["align"]
        # Spacing handled via XML if requested
        if p.get("space_after_pt") is not None:
            para.space_after = Pt(p["space_after_pt"])
        if p.get("space_before_pt") is not None:
            para.space_before = Pt(p["space_before_pt"])


def _add_text_box(slide, left_emu, top_emu, width_emu, height_emu, paragraphs):
    """Add a textbox at fixed coordinates with the given paragraphs."""
    tb = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
    tf = tb.text_frame
    tf.word_wrap = True
    set_text_runs(tf, paragraphs)
    return tb


def _force_white_background(slide):
    """Set the slide-level background fill to plain white so the new slide does
    not inherit a coloured master that would clash with our dark text."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_li_sc_slide(prs, li_text, sc_list, after_index_one_based):
    """Append a new LI/SC slide using BLANK layout with manually positioned text
    boxes, then move it to immediately after the given 1-based index."""
    layout = get_layout(prs, "BLANK", "TITLE_ONLY", "TITLE")
    slide = prs.slides.add_slide(layout)
    # Drop any inherited placeholders so they don't interfere
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 12:  # keep slide-number placeholder if present
            ph._element.getparent().remove(ph._element)
    _force_white_background(slide)

    sw = prs.slide_width
    sh = prs.slide_height

    # Margins
    left = Emu(int(sw * 0.06))
    width = Emu(int(sw * 0.88))

    # Title
    _add_text_box(slide,
        left, Emu(int(sh * 0.06)), width, Emu(int(sh * 0.13)),
        [{"text": "Learning Intention & Success Criteria", "bold": True, "size": 32}],
    )

    # Body
    _add_text_box(slide,
        left, Emu(int(sh * 0.22)), width, Emu(int(sh * 0.72)),
        [
            {"text": "Learning Intention", "bold": True, "size": 22, "space_after_pt": 4},
            {"text": li_text, "size": 20, "space_after_pt": 16},
            {"text": "Success Criteria - I can...", "bold": True, "size": 22, "space_after_pt": 6},
            {"text": sc_list[0], "size": 20, "space_after_pt": 6},
            {"text": sc_list[1], "size": 20, "space_after_pt": 6},
            {"text": sc_list[2], "size": 20, "space_after_pt": 6},
        ],
    )

    # Move new slide (last) to after_index_one_based
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    new_id = sld_ids[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(after_index_one_based, new_id)
    return slide


def add_review_slide(prs, li_text, sc_list):
    """Append a 'Review and Reflect' slide as the final slide using BLANK layout."""
    layout = get_layout(prs, "BLANK", "TITLE_ONLY", "TITLE")
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 12:
            ph._element.getparent().remove(ph._element)
    _force_white_background(slide)

    sw = prs.slide_width
    sh = prs.slide_height
    left = Emu(int(sw * 0.06))
    width = Emu(int(sw * 0.88))

    _add_text_box(slide,
        left, Emu(int(sh * 0.06)), width, Emu(int(sh * 0.13)),
        [{"text": "Review and Reflect", "bold": True, "size": 32}],
    )

    _add_text_box(slide,
        left, Emu(int(sh * 0.22)), width, Emu(int(sh * 0.72)),
        [
            {"text": "Success Criteria - I can...", "bold": True, "size": 22, "space_after_pt": 6},
            {"text": sc_list[0], "size": 20, "space_after_pt": 6},
            {"text": sc_list[1], "size": 20, "space_after_pt": 6},
            {"text": sc_list[2], "size": 20, "space_after_pt": 16},
            {"text": "Show your thinking", "bold": True, "size": 22, "space_after_pt": 6},
            {"text": "Thumbs up - I can do this on my own.", "size": 18, "space_after_pt": 4},
            {"text": "Thumbs sideways - I can do this with some help.", "size": 18, "space_after_pt": 4},
            {"text": "Thumbs down - I need more practice.", "size": 18, "space_after_pt": 4},
        ],
    )
    # Add presenter note for the review slide
    notes = (
        "SAY:\n"
        "- Read each I can statement with me.\n"
        "- Show me on your thumbs - up, sideways or down - for each one.\n"
        "- Pick the I can statement you feel most confident about today.\n\n"
        "DO:\n"
        "- Choral read each SC, scan thumbs each time.\n"
        "- Note which SC has the most thumbs sideways or down to plan tomorrow.\n\n"
        "TEACHER NOTES:\n"
        "Use this data to decide tomorrow's launch and any small group reteach.\n\n"
        "WATCH FOR:\n"
        "- Students avoiding the rating - prompt them to commit to one.\n"
        "- Patterns where SC2 or SC3 is mostly thumbs down - flag for reteach."
    )
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def process_deck(src: Path, dst: Path):
    prs = Presentation(str(src))
    lno = lesson_no(src.name)

    existing = extract_existing_li_sc(prs)
    li_sc_idx = find_li_sc_index(prs)
    ochre_idx = find_ochre_objectives_index(prs)

    if existing:
        li_text, sc_list = existing
    else:
        if lno not in LESSON_LI_SC:
            raise RuntimeError(f"No LI/SC source for lesson {lno} in {src.name}")
        data = LESSON_LI_SC[lno]
        li_text = data["li"]
        sc_list = data["sc"]

    # 1. Insert new LI/SC slide where missing
    if li_sc_idx is None and ochre_idx is not None:
        add_li_sc_slide(prs, li_text, sc_list, ochre_idx)  # inserts at position ochre_idx (so it sits AFTER the ochre slide)

    # 2. Append Review slide
    add_review_slide(prs, li_text, sc_list)

    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))


def main():
    if OUT.exists():
        shutil.rmtree(OUT)
    OUT.mkdir(parents=True)

    decks = sorted(ROOT.rglob("*.pptx"))
    for src in decks:
        rel = src.relative_to(ROOT)
        dst = OUT / rel
        try:
            process_deck(src, dst)
            print(f"OK   {rel}")
        except Exception as e:
            print(f"FAIL {rel}: {e}")
            raise


if __name__ == "__main__":
    main()
