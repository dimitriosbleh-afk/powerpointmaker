"""Dump teaching content from each lesson that lacks specific LI/SC.

Writes one .txt per lesson into scripts/_ochre_dump/ so we can read each one
without Windows console encoding issues."""
from pathlib import Path
from pptx import Presentation
import re

ROOT = Path(r"C:\Users\09560329\Downloads\Challenge Slides")
OUT = Path(r"C:\Users\09560329\Documents\Scripts\Claude Powerpoint\scripts\_ochre_dump")
OUT.mkdir(exist_ok=True)

NEEDS = {12, 13, 16, 17, 18, 19, 20, 21, 22, 23, 24, 25}

def lesson_no(name):
    m = re.match(r"\s*(\d+)\.", name)
    return int(m.group(1)) if m else None

def slide_text(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                if t.strip():
                    out.append(t)
    return "\n".join(out)

# Use 'with notes' versions since they have richer notes content
files = sorted(ROOT.rglob("*with notes*.pptx"))
for f in files:
    lno = lesson_no(f.name)
    if lno not in NEEDS:
        continue
    p = Presentation(str(f))
    lines = [f"# Lesson {lno}: {f.name}", ""]
    for i, s in enumerate(p.slides, 1):
        t = slide_text(s)
        notes = ""
        if s.has_notes_slide:
            ntf = s.notes_slide.notes_text_frame
            if ntf:
                notes = ntf.text or ""
        lines.append(f"--- Slide {i} ---")
        lines.append(t)
        if notes.strip():
            lines.append("[NOTES]")
            lines.append(notes)
        lines.append("")
    out_path = OUT / f"lesson_{lno:02d}.txt"
    out_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {out_path}")

print(f"Done. Files in {OUT}")
